"""
Generate editable .pptx file with figures for EJA submission.
One figure per slide with title and caption.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
fig_dir = os.path.join(SCRIPT_DIR, 'figures')
out_dir = os.path.join(SCRIPT_DIR, 'papers')
os.makedirs(out_dir, exist_ok=True)

# Widescreen dimensions
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def create_eja_pptx():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    figures = [
        ('fig1_price_timeseries.png',
         'Fig. 1',
         'Time series of eBay completed sale prices for desflurane (red), sevoflurane (blue) '
         'and isoflurane (green) vaporisers over three years (March 2023 to March 2026). '
         'Vertical dashed lines indicate key EU regulatory milestones. '
         'Curved lines represent LOWESS trend estimates (fraction = 0.3).'),
        ('fig2_boxplot_comparison.png',
         'Fig. 2',
         'Box plot comparison of vaporiser prices before and after the EU desflurane ban '
         '(1 January 2026). Individual data points are shown as jittered dots.'),
        ('fig3_monthly_median.png',
         'Fig. 3',
         'Monthly median prices of anaesthetic vaporisers on eBay. '
         'Annotations indicate the number of transactions per month (n).'),
        ('fig4_histograms.png',
         'Fig. 4',
         'Price distribution histograms for each vaporiser type, comparing pre-ban '
         '(solid fill) and post-ban (hatched) periods.'),
        ('fig5_regulatory_timeline.png',
         'Fig. 5',
         'Anaesthetic vaporiser prices mapped against the EU regulatory timeline. '
         'Shaded regions indicate regulatory phases.'),
        ('fig6_quarterly_trends.png',
         'Fig. 6',
         'Quarterly median price trends (upper panel) and sales volume (lower panel).'),
    ]

    for filename, title, caption in figures:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Title at top
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title
        run.font.size = Pt(20)
        run.font.bold = True

        # Figure image (centered)
        fig_path = os.path.join(fig_dir, filename)
        if os.path.exists(fig_path):
            img_width = Inches(10)
            img_height = Inches(5.5)
            left = (SLIDE_WIDTH - img_width) // 2
            top = Inches(0.9)
            slide.shapes.add_picture(fig_path, left, top, img_width, img_height)

        # Caption at bottom
        caption_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.8))
        tf = caption_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f'{title}. {caption}'
        run.font.size = Pt(11)
        run.font.italic = True

    path = os.path.join(out_dir, 'eja_figures_english.pptx')
    prs.save(path)
    print(f"EJA figures PPTX saved: {path}")
    return path


if __name__ == '__main__':
    create_eja_pptx()
