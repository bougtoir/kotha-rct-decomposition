#!/usr/bin/env python3
"""Build an editable PPTX for JBS supplementary figures (S1a/b, S2-S3)."""
import argparse
import os
import re
import shutil
import zipfile
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def clean_math(latex):
    """Convert a small subset of LaTeX math to plain ASCII text for PPTX captions."""
    text = latex
    replacements = [
        ('\\text{', ''), ('}', ''),
        ('\\sim', '~'), ('\\cdot', '*'),
        ('\\sqrt', 'sqrt'), ('\\sum', 'sum'),
        ('\\prod', 'prod'), ('\\alpha', 'alpha'),
        ('\\beta', 'beta'), ('\\mu', 'mu'),
        ('\\tau', 'tau'), ('\\theta', 'theta'),
        ('\\sigma', 'sigma'), ('\\delta', 'delta'),
        ('\\rho', 'rho'), ('\\Phi', 'Phi'),
        ('\\phi', 'phi'), ('\\in', ' in '),
        ('\\leq', '<='), ('\\geq', '>='),
        ('\\neq', '!='), ('\\times', 'x'),
        ('\\to', '->'), ('\\approx', '~'),
        ('\\frac', ''), ('\\left', ''),
        ('\\right', ''), ('\\quad', '  '),
        ('\\hat', ''), ('\\log', 'log'),
        ('\\exp', 'exp'), ('\\mid', '|'),
        ('^2', '^2'),
    ]
    for old, new in replacements:
        text = text.replace(old, new)
    text = re.sub(r'_\{([^}]+)\}', r'_\1', text)
    text = re.sub(r'\^\{([^}]+)\}', r'^\1', text)
    text = text.replace('\\', '')
    return text


def _sanitize_pptx(path):
    """Replace typographic and multibyte characters in PPTX XML with ASCII."""
    tmp = path + '.tmp'
    try:
        with zipfile.ZipFile(path, 'r') as zin, zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as zout:
            for info in zin.infolist():
                data = zin.read(info.filename)
                if info.filename.endswith('.xml'):
                    text = data.decode('utf-8')
                    text = text.replace('\u2039', '').replace('\u203a', '')
                    text = text.replace('\u2018', "'").replace('\u2019', "'")
                    text = text.replace('\u201c', '"').replace('\u201d', '"')
                    text = text.replace('\u2013', '-').replace('\u2014', '-')
                    text = text.replace('\u00a0', ' ')
                    data = text.encode('utf-8')
                zout.writestr(info, data)
        shutil.move(tmp, path)
    except Exception:
        if os.path.exists(tmp):
            os.remove(tmp)
        raise


# Trace plots are referenced in the Methods but are not given full figure blocks in the manuscript.
TRACE_CAPTIONS = {
    'figS1a_trace_mcmc_magnesium.png': (
        'MCMC trace plots for selected parameters in the Bayesian power-prior model for '
        'magnesium in acute myocardial infarction (discounting factor alpha = 0.3). Chains show no obvious non-stationarity.'
    ),
    'figS1b_trace_mcmc_statins.png': (
        'MCMC trace plots for selected parameters in the Bayesian power-prior model for '
        'statins in heart failure (discounting factor alpha = 0.3). Chains show no obvious non-stationarity.'
    ),
}


def _parse_cctc_figures(cctc_md_path):
    """Parse CCTC figure blocks for Figs 5/7 and map them to JBS supplementary labels."""
    md = Path(cctc_md_path).read_text()
    mapping = {}
    label_map = {
        '5': 'S2',
        '7': 'S3',
    }
    pattern = r'\*\*Fig\.\s*(5|7)\*\*\s*(.*?)\n\s*!\[.*?\]\((.*?)\)'
    for m in re.finditer(pattern, md, re.S):
        fig_num = m.group(1)
        caption = m.group(2).strip()
        img_path = m.group(3).strip()
        caption = re.sub(r'\$([^$]+)\$', lambda x: clean_math(x.group(1)), caption)
        mapping[label_map[fig_num]] = {
            'caption': caption,
            'img': Path(cctc_md_path).parent / img_path,
        }
    return mapping


def build(cctc_md_path, out_path, base_dir=None):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    base_dir = Path(base_dir) if base_dir else Path(cctc_md_path).parent

    # S1a and S1b trace plots
    figures = []
    for img_name, caption in [
        ('figS1a_trace_mcmc_magnesium.png', 'Magnesium'),
        ('figS1b_trace_mcmc_statins.png', 'Statins'),
    ]:
        src = base_dir / 'validation' / 'figures' / img_name
        figures.append((f'S1{chr(ord("a") + len(figures))}', src, TRACE_CAPTIONS.get(img_name, '')))

    # S2-S3 from CCTC manuscript
    cctc_figs = _parse_cctc_figures(cctc_md_path)
    for label, info in [('S2', '5'), ('S3', '7')]:
        if label in cctc_figs:
            figures.append((label, cctc_figs[label]['img'], cctc_figs[label]['caption']))

    for label, img_path, caption in figures:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f'Supplementary Figure {label}'
        p.font.size = Pt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        if img_path.exists():
            slide.shapes.add_picture(str(img_path), Inches(1.0), Inches(1.1), width=Inches(11.333))
        else:
            print(f'  WARNING: supplementary figure image not found: {img_path}')

        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(1.0))
        tf2 = cap_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = f'Supplementary Figure {label}. {caption}'
        p2.font.size = Pt(14)
        p2.word_wrap = True
        p2.alignment = PP_ALIGN.LEFT

    prs.save(out_path)
    _sanitize_pptx(out_path)
    print(f'Saved {out_path}')


if __name__ == '__main__':
    parser = argparse.ArgumentParser()
    parser.add_argument('--cctc-md', default='05_paper_cctc.md')
    parser.add_argument('--out', default='KOTHA_Framework_JBS_supplementary_figures.pptx')
    parser.add_argument('--base-dir', default=None)
    args = parser.parse_args()
    build(args.cctc_md, args.out, args.base_dir)
