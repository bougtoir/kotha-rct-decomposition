"""
Generate A2 manuscript docx for JAMIA/JBI submission.
Ensemble Causal Discovery with Feedback Quantification.
"""

import os
import re
from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
FIG_DIR = os.path.join(BASE_DIR, 'figures')


GREEK = {
    r'\alpha': '\u03b1', r'\beta': '\u03b2', r'\gamma': '\u03b3',
    r'\delta': '\u03b4', r'\epsilon': '\u03b5', r'\varepsilon': '\u03b5',
    r'\zeta': '\u03b6', r'\eta': '\u03b7', r'\theta': '\u03b8',
    r'\iota': '\u03b9', r'\kappa': '\u03ba', r'\lambda': '\u03bb',
    r'\mu': '\u03bc', r'\nu': '\u03bd', r'\xi': '\u03be',
    r'\pi': '\u03c0', r'\rho': '\u03c1', r'\sigma': '\u03c3',
    r'\tau': '\u03c4', r'\phi': '\u03c6', r'\chi': '\u03c7',
    r'\psi': '\u03c8', r'\omega': '\u03c9', r'\Phi': '\u03a6',
    r'\Psi': '\u03a8', r'\Omega': '\u03a9', r'\Delta': '\u0394',
}

SYMBOLS = {
    r'\infty': '\u221e', r'\partial': '\u2202', r'\times': '\u00d7',
    r'\in': '\u2208', r'\to': '\u2192', r'\rightarrow': '\u2192',
    r'\leftarrow': '\u2190', r'\leftrightarrow': '\u2194',
    r'\leq': '\u2264', r'\geq': '\u2265', r'\neq': '\u2260',
    r'\approx': '\u2248', r'\equiv': '\u2261', r'\sim': '\u223c',
    r'\cdot': '\u00b7', r'\pm': '\u00b1', r'\perp': '\u22a5',
}

SUB_MAP = {
    '0': '\u2080', '1': '\u2081', '2': '\u2082', '3': '\u2083',
    '4': '\u2084', '5': '\u2085', '6': '\u2086', '7': '\u2087',
    '8': '\u2088', '9': '\u2089', 'i': '\u1d62', 'j': '\u2c7c',
    'k': '\u2096', 'n': '\u2099',
}

SUP_MAP = {
    '0': '\u2070', '1': '\u00b9', '2': '\u00b2', '3': '\u00b3',
    '4': '\u2074', '5': '\u2075', '6': '\u2076', '7': '\u2077',
    '8': '\u2078', '9': '\u2079', '*': '*', 'n': '\u207f',
}


def convert_math(text):
    """Convert inline $...$ to Unicode."""
    def process(m):
        s = m.group(1)
        for k, v in GREEK.items():
            s = s.replace(k, v)
        for k, v in SYMBOLS.items():
            s = s.replace(k, v)
        s = re.sub(r'\\text\{([^}]+)\}', r'\1', s)
        s = re.sub(r'\\mathrm\{([^}]+)\}', r'\1', s)
        s = re.sub(r'\\mathcal\{([^}])\}', lambda m2: m2.group(1), s)
        s = re.sub(r'\\hat\{([^}]+)\}', lambda m2: m2.group(1) + '\u0302', s)
        s = re.sub(r'\\bar\{([^}]+)\}', lambda m2: m2.group(1) + '\u0304', s)
        def sub_repl(m2):
            return ''.join(SUB_MAP.get(c, c) for c in m2.group(1))
        s = re.sub(r'_\{([^}]+)\}', sub_repl, s)
        s = re.sub(r'_([a-z0-9])', lambda m2: SUB_MAP.get(m2.group(1), '_' + m2.group(1)), s)
        def sup_repl(m2):
            return ''.join(SUP_MAP.get(c, c) for c in m2.group(1))
        s = re.sub(r'\^\{([^}]+)\}', sup_repl, s)
        s = re.sub(r'\^([a-z0-9*])', lambda m2: SUP_MAP.get(m2.group(1), '^' + m2.group(1)), s)
        s = s.replace('\\{', '\x00LB\x00').replace('\\}', '\x00RB\x00')
        s = s.replace('\\', '').replace('{', '').replace('}', '')
        s = s.replace('\x00LB\x00', '{').replace('\x00RB\x00', '}')
        return s
    # Inline math
    text = re.sub(r'(?<!\$)\$(?!\$)(.+?)(?<!\$)\$(?!\$)', process, text)
    return text


def add_figure(doc, image_path, caption, width=Inches(5.5)):
    """Add figure with caption."""
    if not os.path.exists(image_path):
        p = doc.add_paragraph(f'[Figure missing: {os.path.basename(image_path)}]')
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        return
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(12)
    run = p.add_run()
    run.add_picture(image_path, width=width)
    cap_p = doc.add_paragraph()
    cap_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    cap_p.paragraph_format.space_before = Pt(4)
    cap_p.paragraph_format.space_after = Pt(12)
    cap_run = cap_p.add_run(convert_math(caption))
    cap_run.font.size = Pt(9)
    cap_run.italic = True


def add_rich_paragraph(doc, text, style=None):
    """Add paragraph with bold/italic formatting."""
    text = convert_math(text)
    p = doc.add_paragraph(style=style)
    parts = re.split(r'\*\*(.+?)\*\*', text)
    for i, part in enumerate(parts):
        if i % 2 == 1:
            run = p.add_run(part)
            run.bold = True
        else:
            iparts = re.split(r'\*(.+?)\*', part)
            for j, ipart in enumerate(iparts):
                run = p.add_run(ipart)
                if j % 2 == 1:
                    run.italic = True
    return p


def generate_manuscript(md_file='manuscript_a2_ecd.md', out_file='manuscript_a2_ecd.docx', lang='en'):
    """Generate the A2 manuscript docx from markdown."""
    md_path = os.path.join(BASE_DIR, md_file)
    with open(md_path, 'r', encoding='utf-8') as f:
        lines = f.readlines()

    doc = Document()
    style = doc.styles['Normal']
    style.font.name = 'Times New Roman'
    style.font.size = Pt(11)
    style.paragraph_format.space_after = Pt(6)
    style.paragraph_format.line_spacing = 1.15

    # Figure insertion points (after which section heading's content)
    if lang == 'ja':
        FIGURES = [
            ('### 3.1', 'fig5_hill_radar.png',
             '図1. Hillの9基準カバレッジ。HIKONE（紫）が単一手法より広範なカバレッジを達成。'),
            ('### 3.5', 'fig_ecd_pipeline.png',
             '図2. HIKONE実運用パイプラインの4ステップ。LiNGAMブートストラップ→DPIスペクトル解析→Hodge分解→介入可能性スコアリング。'),
            ('### 4.2', 'fig_interventionability.png',
             '\u56f33. \u56e0\u679c\u30dd\u30c6\u30f3\u30b7\u30e3\u30eb\u03c6 vs \u81e8\u5e8a\u7684\u4ecb\u5165\u53ef\u80fd\u6027\u03b9\u3002\u03c6\u3068\u03b9\u306e\u9006\u76f8\u95a2\u304c\u6570\u7406\u91cf\u3068\u81e8\u5e8a\u7684\u884c\u52d5\u53ef\u80fd\u6027\u306e\u5bfe\u5fdc\u3092\u5b9f\u8a3c\u3002'),
            ('## 5.', 'fig6_causal_dag.png',
             '\u56f34. DirectLiNGAM\u306b\u3088\u308b\u63a8\u5b9a\u56e0\u679cDAG\u3002\u77e2\u5370\u306e\u592a\u3055\u306f\u63a8\u5b9a\u56e0\u679c\u52b9\u679c|B_ij|\u3092\u793a\u3059\u3002'),
            ('### 6.1', 'fig4_direction_comparison.png',
             '\u56f35. 3\u624b\u6cd5\u306b\u3088\u308b\u56e0\u679c\u65b9\u5411\u6bd4\u8f03\u3002\u5168\u8fba\u30da\u30a2\u3067\u306e\u4e00\u81f4\u30fb\u4e0d\u4e00\u81f4\u3002'),
            ('### 6.2', 'fig7_lingam_vs_spectral.png',
             '\u56f36. LiNGAM DAG\uff08\u5de6\uff09vs \u30b9\u30da\u30af\u30c8\u30eb\u56e0\u679c\u6027DCG\uff08\u53f3\uff09\u3002\u7834\u7dda\uff1a\u30d5\u30a3\u30fc\u30c9\u30d0\u30c3\u30af\u8fba\u3002'),
            ('### 8.2', 'fig9_ecd_pruning_analysis.png',
             '図7. HIKONEアンサンブル・プルーニング解析。辺レベルフィードバック率。'),
            ('### 8.3', 'fig_feedback_network.png',
             '\u56f38. \u81e8\u5e8a\u30d5\u30a3\u30fc\u30c9\u30d0\u30c3\u30af\u30cd\u30c3\u30c8\u30ef\u30fc\u30af\u3002\u8fba\u8272: \u7dd1=\u4e00\u65b9\u5411\u3001\u30aa\u30ec\u30f3\u30b8=\u5f31\u30d5\u30a3\u30fc\u30c9\u30d0\u30c3\u30af\u3001\u8d64=\u5f37\u30d5\u30a3\u30fc\u30c9\u30d0\u30c3\u30af\u3002'),
        ]
    else:
        FIGURES = [
            ('### 3.1', 'fig5_hill_radar.png',
             "Figure 1. Hill's nine criteria coverage. HIKONE (purple) achieves broader coverage than any single method."),
            ('### 3.5', 'fig_ecd_pipeline.png',
             'Figure 2. Four-step HIKONE deployment pipeline: LiNGAM bootstrap → DPI spectral analysis → Hodge decomposition → Interventionability scoring.'),
            ('### 4.2', 'fig_interventionability.png',
             'Figure 3. Causal potential \u03c6 vs. clinical interventionability \u03b9. The inverse relationship demonstrates that mathematical quantities correspond to clinical actionability.'),
            ('## 5.', 'fig6_causal_dag.png',
             'Figure 4. Estimated causal DAG from DirectLiNGAM. Arrow weight indicates estimated causal effect |B_ij|.'),
            ('### 6.1', 'fig4_direction_comparison.png',
             'Figure 5. Three-method causal direction comparison across all edge pairs.'),
            ('### 6.2', 'fig7_lingam_vs_spectral.png',
             'Figure 6. LiNGAM DAG (left) vs. Spectral Causality DCG (right). Dashed lines: feedback edges.'),
            ('### 8.2', 'fig9_ecd_pruning_analysis.png',
             'Figure 7. HIKONE ensemble and pruning analysis with edge-level feedback rates.'),
            ('### 8.3', 'fig_feedback_network.png',
             'Figure 8. Clinical feedback network. Edge colors: green = unidirectional, orange = weak feedback, red = strong feedback.'),
        ]

    figure_triggers = {}  # section_key -> (fig_file, caption)
    for trigger, fig_file, caption in FIGURES:
        figure_triggers[trigger] = (fig_file, caption)

    i = 0
    current_heading = ''
    in_display_math = False
    math_buffer = []
    table_buffer = []
    in_table = False
    figures_inserted = set()

    while i < len(lines):
        line = lines[i].rstrip('\n')
        i += 1

        # Display math block
        if line.strip() == '$$':
            if in_display_math:
                # End of display math
                math_text = ' '.join(math_buffer)
                p = doc.add_paragraph()
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                p.paragraph_format.space_before = Pt(6)
                p.paragraph_format.space_after = Pt(6)
                run = p.add_run(convert_math('$' + math_text + '$'))
                run.font.name = 'Cambria Math'
                in_display_math = False
                math_buffer = []
            else:
                in_display_math = True
                math_buffer = []
            continue

        if in_display_math:
            math_buffer.append(line.strip())
            continue

        # Table handling
        if line.strip().startswith('|') and '|' in line.strip()[1:]:
            if not in_table:
                in_table = True
                table_buffer = []
            table_buffer.append(line)
            continue
        elif in_table:
            # Flush table
            _flush_table(doc, table_buffer)
            in_table = False
            table_buffer = []

        # Skip image links
        if re.match(r'^!\[.*\]\(.*\)', line.strip()):
            continue

        # Horizontal rule
        if line.strip() == '---':
            continue

        # Headings
        if line.startswith('# ') and not line.startswith('## '):
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            p.paragraph_format.space_before = Pt(24)
            run = p.add_run(line[2:].strip())
            run.bold = True
            run.font.size = Pt(14)
            continue

        if line.startswith('## '):
            heading = line[3:].strip()
            doc.add_heading(heading, level=1)
            current_heading = line.strip()
            continue

        if line.startswith('### '):
            heading = line[4:].strip()
            doc.add_heading(heading, level=2)
            current_heading = line.strip()
            continue

        if line.startswith('#### '):
            heading = line[5:].strip()
            doc.add_heading(heading, level=3)
            continue

        # Blockquote
        if line.startswith('> '):
            text = line[2:].strip()
            p = add_rich_paragraph(doc, text)
            p.paragraph_format.left_indent = Inches(0.5)
            continue

        # Bullet list
        if line.strip().startswith('- ') or line.strip().startswith('* '):
            text = line.strip()[2:]
            add_rich_paragraph(doc, text, style='List Bullet')
            continue

        # Numbered list
        num_match = re.match(r'^(\d+)\.\s+(.+)', line.strip())
        if num_match:
            text = f'{num_match.group(1)}. {num_match.group(2)}'
            add_rich_paragraph(doc, text)
            continue

        # Inline display math ($$...$$on single line)
        if line.strip().startswith('$$') and line.strip().endswith('$$') and len(line.strip()) > 4:
            math_text = line.strip()[2:-2]
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = p.add_run(convert_math('$' + math_text + '$'))
            run.font.name = 'Cambria Math'
            continue

        # Empty line
        if not line.strip():
            # Check if we should insert a figure
            for trigger, (fig_file, caption) in figure_triggers.items():
                if trigger in current_heading and fig_file not in figures_inserted:
                    fig_path = os.path.join(FIG_DIR, fig_file)
                    add_figure(doc, fig_path, caption)
                    figures_inserted.add(fig_file)
                    break
            continue

        # Regular paragraph
        add_rich_paragraph(doc, line.strip())

    # Flush any remaining table
    if in_table and table_buffer:
        _flush_table(doc, table_buffer)

    out_path = os.path.join(BASE_DIR, out_file)
    doc.save(out_path)
    print(f'Saved: {out_path}')
    return out_path


def _flush_table(doc, table_lines):
    """Parse and add a markdown table."""
    data_lines = [l for l in table_lines if not re.match(r'^\s*\|[-:\s|]+\|\s*$', l)]
    if not data_lines:
        return
    rows = []
    for line in data_lines:
        cells = [c.strip() for c in line.strip().strip('|').split('|')]
        rows.append(cells)
    if len(rows) < 1:
        return
    ncols = max(len(r) for r in rows)
    table = doc.add_table(rows=len(rows), cols=ncols)
    table.style = 'Table Grid'
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for ri, row_data in enumerate(rows):
        for ci, cell_text in enumerate(row_data):
            if ci < ncols:
                cell = table.rows[ri].cells[ci]
                cell.text = ''
                p = cell.paragraphs[0]
                run = p.add_run(convert_math(cell_text))
                run.font.size = Pt(9)
                if ri == 0:
                    run.bold = True
    doc.add_paragraph('')


def generate_cover_letter():
    """Generate cover letter for JBI."""
    doc = Document()
    style = doc.styles['Normal']
    style.font.name = 'Times New Roman'
    style.font.size = Pt(11)
    style.paragraph_format.space_after = Pt(6)
    style.paragraph_format.line_spacing = 1.15

    doc.add_paragraph('[Date]')
    doc.add_paragraph('')

    p = doc.add_paragraph()
    p.add_run('Editors-in-Chief').bold = True
    doc.add_paragraph('Journal of Biomedical Informatics')
    doc.add_paragraph('Elsevier')
    doc.add_paragraph('')

    p = doc.add_paragraph()
    p.add_run('Re: Submission of manuscript entitled ')
    run = p.add_run('"HIKONE: Hodge-Integrated Knowledge-Optional Network Estimation '
        'with Feedback Quantification for Clinical Causal Inference"')
    run.italic = True
    doc.add_paragraph('')

    doc.add_paragraph('Dear Editors,')
    doc.add_paragraph('')

    # Paragraph 1: What and why
    doc.add_paragraph(
        'We are pleased to submit the above-titled manuscript (~5,000 words, 8 figures, '
        '10 tables) for consideration as an Original Research Article in the Journal of '
        'Biomedical Informatics. This manuscript presents HIKONE (Hodge-Integrated '
        'Knowledge-Optional Network Estimation), '
        'a novel pipeline that integrates LiNGAM\u2019s identifiability guarantees with '
        'spectral causality\u2019s Hodge decomposition to enable clinically meaningful causal '
        'inference that accommodates feedback loops\u2014a ubiquitous feature of pathophysiology '
        'that existing DAG-based methods systematically exclude.'
    )

    # Paragraph 2: Scope fit
    doc.add_paragraph(
        'We believe JBI is the ideal venue for this work for three reasons. First, it directly '
        'addresses the gap between causal discovery methods and clinical applicability\u2014a '
        'topic of growing interest in biomedical informatics, as highlighted by the recent '
        'scoping review by Liu et al. (BMC MRM, 2024). Second, the HIKONE pipeline produces '
        'interventionability scores that map graph-theoretic quantities (Hodge causal potential) '
        'to clinical actionability (treatment target prioritization), providing a principled '
        'bridge between mathematical methods and clinical decision support. Third, the manuscript '
        'introduces the Directional Predictability Index (DPI), an ensemble of three asymmetric '
        'statistics that enables causal direction estimation from cross-sectional data alone, '
        'without requiring domain knowledge\u2014expanding the practical applicability of causal '
        'discovery to settings where expert annotation is unavailable.'
    )

    # Paragraph 3: Novelty / key contributions
    p = doc.add_paragraph()
    p.add_run('Key contributions').bold = True
    p.add_run(' of this manuscript include:')
    items = [
        'A Directional Predictability Index (DPI) that combines regression asymmetry, '
        'additive noise model residual independence, and conditional entropy reduction '
        'to estimate causal directions from data alone (\u03b1 = 0: 9 directed edges detected, '
        '67% LiNGAM agreement; \u03b1 = 0.6 with domain knowledge: r_gradient = 0.859).',
        'Edge-level feedback quantification via Hodge decomposition, revealing clinically '
        'meaningful cycles (MaxHR \u2194 STDepression: 73% feedback rate) that the DAG assumption masks.',
        'An interventionability score linking Hodge causal potential \u03c6 to clinical actionability, '
        'validated against established treatment guidelines (statins for cholesterol, '
        'antihypertensives for blood pressure).',
        'Broader coverage of Hill\u2019s nine epidemiological criteria than any single method alone, '
        'with spectral causality uniquely contributing H6 (biological plausibility), H7 (coherence), '
        'and H9 (analogy).',
        'A DAG transition analysis demonstrating that knowledge quality (p*_flip \u2248 0.15) '
        'rather than quantity is the critical threshold, with clinical implications for '
        'how much domain expertise is needed to produce valid causal structures.',
    ]
    for item in items:
        doc.add_paragraph(item, style='List Bullet')

    # Paragraph 4: Relationship to companion paper
    doc.add_paragraph('')
    p = doc.add_paragraph()
    p.add_run('Relationship to companion paper. ').bold = True
    p.add_run(
        'A companion paper presenting the full mathematical foundations, formal proofs, '
        'and identifiability theory of spectral causality has been submitted separately '
        'to a theoretical venue (Reference [7] in the manuscript). The present manuscript '
        'is fully self-contained and focuses exclusively on the ensemble framework, '
        'clinical interpretation, and practical deployment pipeline that are not covered '
        'in the companion paper. The two manuscripts share no overlapping text; '
        'theoretical results from the companion paper are cited but not reproduced.'
    )

    # Paragraph 5: Ethics / COI / simultaneous submission
    doc.add_paragraph('')
    doc.add_paragraph(
        'We confirm that this manuscript has not been previously published and is not under '
        'simultaneous consideration elsewhere (except as noted above for the companion paper '
        'targeting a different scope and audience). The study uses publicly available data '
        '(UCI Heart Disease Dataset, Cleveland subset) and does not involve human subjects '
        'research requiring ethical approval. All authors have approved the manuscript and '
        'agree with its submission. The authors declare no conflicts of interest.'
    )

    # Suggested reviewers
    doc.add_paragraph('')
    p = doc.add_paragraph()
    p.add_run('Suggested reviewers:').bold = True
    reviewers = [
        'Prof. Markus P\u00fcschel (ETH Z\u00fcrich, Switzerland) \u2014 Pioneer of '
        'graph signal processing on DAGs and causal Fourier analysis; his group\u2019s '
        'work on spectral methods for directed graphs (ICASSP 2024) is directly relevant.',
        'Prof. Kun Zhang (Carnegie Mellon University, USA) \u2014 Leading expert on '
        'causal discovery, non-linear causal models, and identifiability theory; '
        'developer of key additive noise model results underlying our DPI.',
        'Prof. Lucila Ohno-Machado (Yale University, USA) \u2014 Distinguished '
        'researcher in biomedical informatics and clinical causal inference; '
        'her perspective on clinical applicability is ideal for evaluating our pipeline.',
    ]
    for r in reviewers:
        doc.add_paragraph(r, style='List Bullet')

    doc.add_paragraph('')
    doc.add_paragraph(
        'We thank you for your consideration and look forward to your response.'
    )
    doc.add_paragraph('')
    doc.add_paragraph('Sincerely,')
    doc.add_paragraph('')
    p = doc.add_paragraph()
    p.add_run('[Author Name]').bold = True
    doc.add_paragraph('[Affiliation]')
    doc.add_paragraph('[Email]')

    out_path = os.path.join(BASE_DIR, 'cover_letter_a2_jbi.docx')
    doc.save(out_path)
    print(f'Saved: {out_path}')


def generate_pptx():
    """Generate editable figures pptx."""
    from pptx import Presentation
    from pptx.util import Inches as PInches, Pt as PPt

    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)

    FIGURES = [
        ('fig5_hill_radar.png', "Figure 1: Hill's Nine Criteria Coverage",
         'ECD (purple) achieves broader coverage than any single method.'),
        ('fig_ecd_pipeline.png', 'Figure 2: ECD Deployment Pipeline',
         'Four-step pipeline: LiNGAM bootstrap \u2192 DPI spectral analysis \u2192 Hodge decomposition \u2192 Interventionability scoring.'),
        ('fig_interventionability.png', 'Figure 3: Causal Potential vs. Interventionability',
         'Hodge potential \u03c6 vs. clinical interventionability \u03b9. Inverse relationship demonstrates mathematical-clinical correspondence.'),
        ('fig6_causal_dag.png', 'Figure 4: Estimated Causal DAG (DirectLiNGAM)',
         'Arrow weight indicates estimated causal effect size |B_ij|.'),
        ('fig4_direction_comparison.png', 'Figure 5: Three-Method Direction Comparison',
         'Green = agreement, Red = disagreement between methods.'),
        ('fig7_lingam_vs_spectral.png', 'Figure 6: LiNGAM DAG vs. Spectral Causality DCG',
         'Dashed lines indicate feedback edges absent in the DAG.'),
        ('fig9_ecd_pruning_analysis.png', 'Figure 7: ECD Ensemble & Pruning Analysis',
         'Edge-level feedback rates and bootstrap confidence intervals.'),
        ('fig_feedback_network.png', 'Figure 8: Clinical Feedback Network',
         'Edge colors: green = unidirectional (<10%), orange = weak feedback (10\u201350%), red = strong feedback (>50%). MaxHR \u2194 STDep: 73%.'),
    ]

    for fname, title, caption in FIGURES:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        # Title
        txBox = slide.shapes.add_textbox(PInches(0.5), PInches(0.2), PInches(12.3), PInches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = PPt(20)
        p.font.bold = True
        p.alignment = 1  # center

        # Image
        img_path = os.path.join(FIG_DIR, fname)
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, PInches(2.0), PInches(1.2), PInches(9.3), PInches(5.2))

        # Caption
        txBox2 = slide.shapes.add_textbox(PInches(1.0), PInches(6.6), PInches(11.3), PInches(0.7))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = caption
        p2.font.size = PPt(12)
        p2.font.italic = True
        p2.alignment = 1

    out_path = os.path.join(BASE_DIR, 'figures_a2_ecd.pptx')
    prs.save(out_path)
    print(f'Saved: {out_path}')


if __name__ == '__main__':
    import sys
    if '--ja' in sys.argv:
        generate_manuscript('manuscript_a2_ecd_ja.md', 'manuscript_a2_ecd_ja.docx', lang='ja')
    else:
        generate_manuscript()
        generate_manuscript('manuscript_a2_ecd_ja.md', 'manuscript_a2_ecd_ja.docx', lang='ja')
        generate_cover_letter()
        generate_pptx()
