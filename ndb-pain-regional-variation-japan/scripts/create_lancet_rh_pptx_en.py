#!/usr/bin/env python3
"""Create editable PPTX with figures for Lancet Regional Health – WP (English).

Widescreen slides (13.333 x 7.5 inches), one figure per slide with title and caption.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
BASE_DIR = os.path.dirname(SCRIPT_DIR)
OUTPUT_DIR = os.path.join(BASE_DIR, 'output')
LRH_DIR = os.path.join(OUTPUT_DIR, 'lancet_rh_wp')
os.makedirs(LRH_DIR, exist_ok=True)
FIG_DIR = OUTPUT_DIR

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

figures = [
    {
        'path': os.path.join(FIG_DIR, 'fig1_neuropathic_unadjusted_en.png'),
        'title': 'Figure 1',
        'caption': (
            'Outpatient neuropathic pain drug prescribing per surgery by prefecture (unadjusted). '
            'Tohoku prefectures (red) cluster at the high end. Dashed line = national mean.'
        ),
    },
    {
        'path': os.path.join(FIG_DIR, 'fig2_confounder_correlations_en.png'),
        'title': 'Figure 2',
        'caption': (
            'Correlation between neuropathic pain prescribing and confounder disease proxies. '
            'Each dot represents one prefecture. Tohoku prefectures have red borders. '
            'Diabetes drugs show the strongest correlation (r=0\u00b787).'
        ),
    },
    {
        'path': os.path.join(FIG_DIR, 'fig4_region_unadj_vs_adj_en.png'),
        'title': 'Figure 3',
        'caption': (
            'Regional comparison of neuropathic pain prescribing: (a) unadjusted and '
            '(b) after confounder adjustment. Tohoku (red) shifts from the highest region '
            'to mid-range after adjustment. Error bars = SD.'
        ),
    },
]

for fig in figures:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title at top
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = fig['title']
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 0, 0)

    # Image centered
    if os.path.exists(fig['path']):
        from PIL import Image
        im = Image.open(fig['path'])
        img_w, img_h = im.size
        im.close()

        max_w = Inches(11.5)
        max_h = Inches(5.5)
        scale_w = max_w / Emu(int(img_w * 914400 / 96))
        scale_h = max_h / Emu(int(img_h * 914400 / 96))
        scale = min(scale_w, scale_h, 1.0)

        final_w = int(img_w * 914400 / 96 * scale)
        final_h = int(img_h * 914400 / 96 * scale)
        left = (prs.slide_width - final_w) // 2
        top = Inches(1.0)

        slide.shapes.add_picture(fig['path'], left, top, final_w, final_h)

    # Caption at bottom
    cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.8))
    ctf = cap_box.text_frame
    ctf.word_wrap = True
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    crun = cp.add_run()
    crun.text = fig['caption']
    crun.font.size = Pt(12)
    crun.font.color.rgb = RGBColor(80, 80, 80)

outpath = os.path.join(LRH_DIR, 'LRH_figures_EN.pptx')
prs.save(outpath)
print(f'Saved: {outpath}')
