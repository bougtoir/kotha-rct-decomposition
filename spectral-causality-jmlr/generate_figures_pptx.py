#!/usr/bin/env python3
"""Generate editable .pptx file with one figure per slide for JMLR submission."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

FIGURES = [
    {
        "file": "figures/fig1_three_approaches.png",
        "title": "Figure 1: Three Approaches to Causal Inference",
        "caption": (
            "Three complementary approaches to causal inference: (1) structural equation "
            "models and DAG-based methods (e.g., LiNGAM), (2) spectral graph methods via "
            "the magnetic Laplacian, and (3) Hodge-theoretic flow decomposition. "
            "Spectral causality unifies (2) and (3)."
        ),
    },
    {
        "file": "figures/fig_dpi_architecture.png",
        "title": "Figure 2: DPI Architecture",
        "caption": (
            "Architecture of the Directional Predictability Index (DPI). Three independent "
            "asymmetric statistics (regression coefficient asymmetry, ANM residual independence, "
            "conditional entropy reduction) are normalized and averaged. The composite asymmetric "
            "score modulates the symmetric correlation into the directional DPI."
        ),
    },
    {
        "file": "figures/fig_cci_complex_plane.png",
        "title": "Figure 3: Complex Causal Index in the Complex Plane",
        "caption": (
            "CCI(i,j) for all variable pairs plotted in the complex plane. "
            "Real axis = SCC (coupling strength); imaginary axis = SCD (causal direction). "
            "Red = forward direction (i -> j); blue = reverse. Most pairs cluster in the "
            "lower half-plane, reflecting Age's dominance as the upstream root node."
        ),
    },
    {
        "file": "figures/fig2_magnetic_laplacian_q.png",
        "title": "Figure 4: Magnetic Laplacian Eigenvectors in Complex Plane",
        "caption": (
            "Fiedler eigenvector of the magnetic Laplacian plotted in the complex plane "
            "for q=0, 0.1, and 0.25. At q=0, all points lie on the real axis. As q increases, "
            "variables spread into the complex plane, with phase angle ordering reflecting "
            "causal flow direction."
        ),
    },
    {
        "file": "figures/fig3_hodge_decomposition.png",
        "title": "Figure 5: Hodge Decomposition Results",
        "caption": (
            "Hodge decomposition results. (A) 85.9% of flow energy is in the gradient (DAG) "
            "component; 14.1% is in the curl (feedback) component. (B) Causal potential phi: "
            "Age is most upstream; ST Depression is most downstream."
        ),
    },
    {
        "file": "figures/fig_pflip_ucurve.png",
        "title": "Figure 6: Knowledge Quality Phase Transition",
        "caption": (
            "U-shaped quality curve: r_gradient vs p_flip. Empirical results (200 trials, blue) "
            "closely match theoretical prediction (1-2p)^2 * r* (red dashed). "
            "Critical threshold p*_flip ~ 0.15: at least 85% correct directions needed. "
            "Partial misinformation is worse than complete ignorance."
        ),
    },
    {
        "file": "figures/fig8_alpha_sweep.png",
        "title": "Figure 7: Alpha-Sweep Phase Transition Analysis",
        "caption": (
            "Alpha-sweep analysis with DPI. (A) r_gradient increases smoothly from 0.581 "
            "(alpha=0) to 0.859 (alpha=1). (B) Number of detected edges and LiNGAM agreement rate. "
            "(C) Asymmetric norm. (D) Phase diagram."
        ),
    },
]


def create_pptx(output_path="spectral_causality_figures.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank

    for fig in FIGURES:
        slide = prs.slides.add_slide(blank_layout)

        # Title
        txBox = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.6)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = fig["title"]
        p.font.size = Pt(20)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Image
        img_path = fig["file"]
        if os.path.exists(img_path):
            pic = slide.shapes.add_picture(
                img_path, Inches(1.5), Inches(1.0), Inches(10.333), Inches(5.0)
            )

        # Caption
        txBox2 = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.2), Inches(12.333), Inches(1.0)
        )
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = fig["caption"]
        p2.font.size = Pt(12)
        p2.font.italic = True
        p2.alignment = PP_ALIGN.CENTER

    prs.save(output_path)
    print(f"Saved {output_path} with {len(FIGURES)} slides")


if __name__ == "__main__":
    create_pptx()
