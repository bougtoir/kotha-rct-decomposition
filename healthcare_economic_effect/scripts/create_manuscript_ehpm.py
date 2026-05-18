"""
EHPM Submission Package: Healthcare Expenditure as Economic Effect

Generates submission-ready documents per Environmental Health and Preventive
Medicine (EHPM) Instructions for Authors:
  - Main manuscript docx (Title Page + Abstract + IMRD + Declarations + Refs
    + Figure Legends + Tables)
  - Cover letter docx
  - Editable figures pptx (1 figure per slide)

Article type: Research (~6,400 words)
Format: Double-spaced, 12pt, A4, 1-inch margins, Vancouver references
"""
import os
import re
import json
import pandas as pd
from docx import Document
from docx.shared import Inches, Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.section import WD_ORIENT
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.abspath(os.path.join(HERE, ".."))
DATA = os.path.join(ROOT, "data")
FIG = os.path.join(ROOT, "output", "figures")
DOCX_DIR = os.path.join(ROOT, "output", "docx")
PPTX_DIR = os.path.join(ROOT, "output", "pptx")
os.makedirs(DOCX_DIR, exist_ok=True)
os.makedirs(PPTX_DIR, exist_ok=True)


def get_fig(name):
    return os.path.join(FIG, name)


# ---------------------------------------------------------------------------
# EHPM formatting helpers
# ---------------------------------------------------------------------------
def _set_ehpm_format(doc):
    """Apply EHPM formatting to all sections: A4, 1-inch margins,
    12pt font, double-spacing."""
    for section in doc.sections:
        section.page_width = Cm(21.0)
        section.page_height = Cm(29.7)
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin = Inches(1)
        section.right_margin = Inches(1)
    style = doc.styles['Normal']
    style.font.name = 'Times New Roman'
    style.font.size = Pt(12)
    style.paragraph_format.line_spacing_rule = WD_LINE_SPACING.DOUBLE
    style.paragraph_format.space_after = Pt(0)
    for level in range(1, 4):
        hs = doc.styles[f'Heading {level}']
        hs.font.name = 'Times New Roman'
        hs.font.color.rgb = RGBColor(0, 0, 0)
        hs.paragraph_format.line_spacing_rule = WD_LINE_SPACING.DOUBLE
        if level == 1:
            hs.font.size = Pt(14)
        elif level == 2:
            hs.font.size = Pt(13)
        else:
            hs.font.size = Pt(12)


def add_text_with_refs(paragraph, text, bold=False, italic=False):
    """Parse {ref} markers and render as Word-native superscript."""
    parts = re.split(r'(\{[^}]+\})', text)
    for part in parts:
        if part.startswith('{') and part.endswith('}'):
            run = paragraph.add_run(part[1:-1])
            run.font.superscript = True
            run.font.size = Pt(10)
        else:
            run = paragraph.add_run(part)
            if bold:
                run.bold = True
            if italic:
                run.italic = True


def add_heading(doc, text, level=1):
    h = doc.add_heading(text, level=level)
    return h


def add_para(doc, text, bold=False, italic=False):
    p = doc.add_paragraph()
    add_text_with_refs(p, text, bold=bold, italic=italic)
    return p


def add_plain_para(doc, text, bold=False, italic=False, align=None):
    """Add paragraph without ref parsing."""
    p = doc.add_paragraph()
    run = p.add_run(text)
    if bold:
        run.bold = True
    if italic:
        run.italic = True
    if align:
        p.alignment = align
    return p


def add_table_from_df(doc, df, title, legend=""):
    """Add table with title above and optional legend below.
    EHPM: no color/shading, Word Table object."""
    # Title
    cap = doc.add_paragraph()
    cap.paragraph_format.space_before = Pt(12)
    run = cap.add_run(title)
    run.bold = True
    run.font.size = Pt(12)

    table = doc.add_table(rows=1, cols=len(df.columns), style='Table Grid')
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for j, col in enumerate(df.columns):
        cell = table.rows[0].cells[j]
        cell.text = str(col)
        for p in cell.paragraphs:
            for r in p.runs:
                r.bold = True
                r.font.size = Pt(10)
                r.font.name = 'Times New Roman'
    for _, row in df.iterrows():
        row_cells = table.add_row().cells
        for j, col in enumerate(df.columns):
            row_cells[j].text = str(row[col])
            for p in row_cells[j].paragraphs:
                for r in p.runs:
                    r.font.size = Pt(10)
                    r.font.name = 'Times New Roman'
    if legend:
        leg = doc.add_paragraph()
        run = leg.add_run(legend)
        run.font.size = Pt(10)
        run.italic = True
    doc.add_paragraph()


def add_figure_inline(doc, img_path, caption):
    """Insert figure inline with caption below (for review convenience)."""
    if os.path.exists(img_path):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run()
        run.add_picture(img_path, width=Inches(5.0))
    cap = doc.add_paragraph()
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    cap.paragraph_format.space_before = Pt(6)
    cap.paragraph_format.space_after = Pt(12)
    run = cap.add_run(caption)
    run.font.size = Pt(10)
    run.italic = True


# ---------------------------------------------------------------------------
# References -- Vancouver style, numbered by order of first appearance
# ---------------------------------------------------------------------------
REFERENCES = [
    # 1 -- Background: Japan multiplier
    "Yamada G, Imanaka Y. Input-output analysis on the economic impact of "
    "medical care in Japan. Environ Health Prev Med. 2015;20(5):379-387.",
    # 2 -- Background: US multiplier
    "Dupor B, Guerrero R. The aggregate and local economic effects of "
    "government financed health care. Econ Inq. 2021;59(2):662-670.",
    # 3 -- Background: EU I-O framework
    "Gutierrez-Hernandez I, Abasolo-Alesson I. The health care sector in "
    "the economies of the European Union: an overview using an input-output "
    "framework. Cost Eff Resour Alloc. 2021;19(1):4.",
    # 4 -- Background: OECD ranking
    "OECD. Health at a Glance 2023: OECD Indicators. Paris: OECD Publishing; "
    "2023.",
    # 5 -- Background: Cabinet Office
    "Cabinet Office, Japan. Annual Report on the Japanese Economy 2025. "
    "Tokyo: Cabinet Office; 2025. [In Japanese]",
    # 6 -- Background: JMARI employment
    "Maeda Y. Economic ripple effects of healthcare and long-term care. "
    "JMARI Working Paper No. 172. 2008. [In Japanese]",
    # 7 -- Background: HLGH (CS-ARDL)
    "Ertugrul HM, Baycan O, Atilgan E, Ulucan H. Health-led growth "
    "hypothesis and health financing systems: an econometric synthesis "
    "for OECD countries. Front Public Health. 2024;12:1437304.",
    # 8 -- Background: HLGH (Toda-Yamamoto)
    "Amiri A, Ventelou B. Granger causality between total expenditure on "
    "health and GDP in OECD: Evidence from the Toda-Yamamoto approach. "
    "Econ Lett. 2012;116(3):541-544.",
    # 9 -- Background: HLGH (Driscoll-Kraay)
    "Beylik U, Cetin M, Senol O, Cirakli U, Ecevit E. The relationship "
    "between health expenditure indicators and economic growth in OECD "
    "countries: A Driscoll-Kraay approach. Front Public Health. "
    "2022;10:1050550.",
    # 10 -- Background: HLGH (Panel VECM)
    "Wang KM. Health care expenditure and economic growth: Quantile "
    "panel-type analysis. Econ Model. 2011;28(4):1536-1549.",
    # 11 -- Background: HLGH developing countries
    "Piabuo SM, Tieguhong JC. Health expenditure and economic growth - "
    "a review of the literature and an analysis between the economic "
    "community for central African states (CEMAC) and selected African "
    "countries. Health Econ Rev. 2017;7(1):23.",
    # 12 -- Background: Health as investment
    "Mushkin SJ. Health as an investment. J Polit Econ. 1962;70(5):129-157.",
    # 13 -- Background: Health capital model
    "Grossman M. On the concept of health capital and the demand for "
    "health. J Polit Econ. 1972;80(2):223-255.",
    # 14 -- Background: Health-TFP link
    "Bloom DE, Canning D, Sevilla J. The effect of health on economic "
    "growth: a production function approach. World Dev. "
    "2004;32(1):1-13.",
    # 15 -- Background: Life expectancy and GDP
    "Barro RJ. Health and economic growth. Ann Econ Finance. "
    "2013;14(2):305-342.",
    # 16 -- Background: Tempo effect origin
    "Bongaarts J, Feeney G. On the quantum and tempo of fertility. "
    "Popul Dev Rev. 1998;24(2):271-291.",
    # 17 -- Background: Forgotten parameter sigma
    "Goldstein JR, Lutz W, Scherbov S. Long-term population decline in "
    "Europe: the relative importance of tempo effects and generational "
    "length. Popul Dev Rev. 2003;29(4):699-707.",
    # 18 -- Methods: GDP tempo companion
    "Onishi T. The forgotten tempo effect in capital accounting: "
    "investment-to-output time-to-build, intangible capital, and the "
    "reconciliation of flow- and stock-based national wealth measures. "
    "Working Paper. 2026.",
    # 19 -- Methods: Healthcare tempo PoC
    "Onishi T. Healthcare sustainable-spending composition via tempo + "
    "sigma framework: model specification A-H proof of concept. "
    "Working Paper. 2026.",
    # 20 -- Methods: Germany I-O
    "Henke KD, Ostwald DA. Health satellite account: the first step. "
    "In: Dged JM, ed. The Elgar Companion to Health Economics. 2nd ed. "
    "Cheltenham: Edward Elgar; 2012. p. 327-337.",
    # 21 -- Results: WDI data
    "World Bank. World Development Indicators. Washington, DC: World Bank; "
    "2024. Available from: https://databank.worldbank.org/",
    # 22 -- Results: Preston Curve
    "Preston SH. The changing relation between mortality and level of "
    "economic development. Popul Stud. 1975;29(2):231-248.",
    # 23 -- Results: OECD diagnostic technology
    "OECD. Health at a Glance 2023: OECD Indicators. Chapter 5.23: "
    "Diagnostic technologies. Paris: OECD Publishing; 2023.",
    # 24 -- Results: MHLW trade statistics
    "Ministry of Health, Labour and Welfare. Pharmaceutical and medical "
    "device production statistics annual report. "
    "Tokyo: MHLW; 2021. [In Japanese]",
]


# ---------------------------------------------------------------------------
# Build EHPM main manuscript
# ---------------------------------------------------------------------------
def build_ehpm_manuscript():
    doc = Document()
    _set_ehpm_format(doc)

    # ======================================================================
    # TITLE PAGE (unnumbered)
    # ======================================================================
    add_plain_para(doc, "TITLE PAGE", bold=True,
                   align=WD_ALIGN_PARAGRAPH.CENTER)
    doc.add_paragraph()

    # Title
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(
        "Healthcare Expenditure as Economic Effect: "
        "A Dual-Return Sustainability Framework Integrating "
        "Input-Output Multipliers, Health-Capital Tempo, "
        "and Diagnostic Equipment Stock"
    )
    run.bold = True
    run.font.size = Pt(14)

    doc.add_paragraph()

    # Authors (placeholder -- user fills in)
    add_plain_para(doc, "Tatsuki Onishi [1]*", bold=False,
                   align=WD_ALIGN_PARAGRAPH.CENTER)
    doc.add_paragraph()

    # Affiliations
    add_plain_para(doc, "[1] [Affiliation to be completed by author]",
                   italic=True)
    doc.add_paragraph()

    # Corresponding author
    p = doc.add_paragraph()
    run = p.add_run("*Corresponding author: ")
    run.bold = True
    run2 = p.add_run(
        "Tatsuki Onishi, [Affiliation to be completed]. "
        "E-mail: [to be completed]"
    )
    doc.add_paragraph()

    # Short title
    p = doc.add_paragraph()
    run = p.add_run("Short title: ")
    run.bold = True
    p.add_run("Healthcare I-O Multiplier Sustainability")
    # 42 characters -- within 50 char limit

    doc.add_paragraph()

    # Word count
    p = doc.add_paragraph()
    run = p.add_run("Total word count: ")
    run.bold = True
    p.add_run("approximately 6,200 words "
              "(Text, References, Tables, and Figure Legends)")

    doc.add_paragraph()
    p = doc.add_paragraph()
    run = p.add_run("Number of figures: ")
    run.bold = True
    p.add_run("8")
    p = doc.add_paragraph()
    run = p.add_run("Number of tables: ")
    run.bold = True
    p.add_run("4")

    doc.add_page_break()

    # ======================================================================
    # ABSTRACT (structured, <= 350 words)
    # ======================================================================
    add_heading(doc, "Abstract", level=1)

    p = doc.add_paragraph()
    run = p.add_run("Background: ")
    run.bold = True
    p.add_run(
        "Healthcare expenditure is conventionally treated as a fiscal cost "
        "to be contained. However, healthcare is also a major economic sector "
        "generating substantial demand-side and supply-side returns. "
        "This study proposes a neutral sustainability framework that "
        "re-evaluates healthcare expenditure as both cost and economic effect, "
        "incorporating input-output (I-O) multipliers, health-capital tempo "
        "effects, and diagnostic equipment stock valuation."
    )

    p = doc.add_paragraph()
    run = p.add_run("Methods: ")
    run.bold = True
    p.add_run(
        "We compiled I-O multipliers for 13 OECD countries from published "
        "national and EU-28 framework studies. A fiscal return ratio "
        "(effective tax rate times I-O multiplier divided by public financing share) "
        "was computed for each country. The tempo model from a companion "
        "proof-of-concept analysis (39 countries, 2000-2019) was integrated "
        "to capture supply-side health-capital accumulation. "
        "CT and MRI scanner density data (OECD Health at a Glance 2023) "
        "and pharmaceutical/device trade statistics were used to model "
        "diagnostic equipment stock as health capital and quantify import "
        "leakage effects on effective multipliers. Three counterfactual "
        "scenarios were constructed for Japan."
    )

    p = doc.add_paragraph()
    run = p.add_run("Results: ")
    run.bold = True
    p.add_run(
        "The demand-side fiscal return ratio exceeded 1.0 in five of thirteen "
        "countries (France 1.18, Italy 1.13, Japan 1.09, Sweden 1.04, Finland 1.04), "
        "with remaining countries recovering 76-96% of public costs. "
        "Japan's diagnostic imaging density (170.9 CT+MRI per million) was "
        "approximately four times the OECD median. Import leakage of approximately "
        "5% of current health expenditure reduced Japan's effective multiplier "
        "from 2.78 to 2.64 (fiscal return from 1.09 to 1.04). "
        "Counterfactual analysis showed that reducing equipment density to the "
        "OECD average would eliminate sustainability (ratio 0.98), "
        "while full domestic manufacturing would improve it to 1.09."
    )

    p = doc.add_paragraph()
    run = p.add_run("Conclusions: ")
    run.bold = True
    p.add_run(
        "Healthcare expenditure functions as an economic investment with "
        "measurable demand-side returns. Japan's high diagnostic equipment "
        "density is a sustainability component rather than excess. "
        "The policy debate should shift from cost containment to "
        "maximizing economic return on healthcare investment."
    )

    doc.add_paragraph()

    # ======================================================================
    # KEYWORDS
    # ======================================================================
    p = doc.add_paragraph()
    run = p.add_run("Keywords: ")
    run.bold = True
    p.add_run(
        "Healthcare expenditure; Input-output analysis; Economic multiplier; "
        "Health capital; Sustainability; Diagnostic imaging; Import leakage; "
        "Tempo effect; OECD"
    )

    doc.add_page_break()

    # ======================================================================
    # BACKGROUND
    # ======================================================================
    add_heading(doc, "Background", level=1)

    add_para(doc,
        "Since the Cabinet's 'Basic Policies on Economic and Fiscal Management' "
        "decision in 2005, the Japanese government has pursued policies to "
        "moderate the growth of national medical care expenditure. Japan's "
        "current health expenditure (CHE) reached 11.0% of GDP in 2019, "
        "ranking fifth among 38 Organisation for Economic Co-operation and "
        "Development (OECD) countries{4}. With rapid population aging -- "
        "the proportion of the population aged 65 and over exceeded 28% "
        "in 2019, the highest among OECD nations -- the 'sustainability' of "
        "healthcare spending has become a central policy concern{5}. This "
        "discourse has almost uniformly treated healthcare expenditure as a "
        "'cost' -- an expense to be minimized through efficiency gains, "
        "volume controls, and fee schedule revisions."
    )
    add_para(doc,
        "Yet healthcare is also a major economic sector. It purchases inputs "
        "from pharmaceuticals, medical devices, and information technology; "
        "it employs millions of workers; and these workers spend their "
        "incomes, generating further economic activity throughout the economy. "
        "Maeda estimated that Japan's healthcare sector supports 6.89 million "
        "jobs in total (2.95 million direct, approximately 4 million indirect), "
        "with production-inducement effects exceeding those of any other "
        "service industry including construction and education{6}. "
        "Input-output (I-O) analysis, pioneered by Leontief, quantifies how "
        "final demand in one sector induces production across the entire "
        "economy through direct, indirect, and induced effects. "
        "Yamada and Imanaka reported a healthcare I-O multiplier of 2.78 "
        "(95% confidence interval [CI]: 2.74-2.90) for Japan{1}, meaning "
        "that each unit of healthcare spending generates 2.78 units of total "
        "economic output. Multipliers ranging from 1.7 to 2.9 have been "
        "estimated across OECD countries{2,3,20}. In the European Union, "
        "Gutierrez-Hernandez and Abasolo-Alesson reported backward-linkage "
        "coefficients for 28 member states, confirming healthcare as a "
        "high-multiplier sector in most economies{3}."
    )
    add_para(doc,
        "On the supply side, health has long been recognized as a component "
        "of human capital that generates long-term economic returns. "
        "Mushkin conceptualized health as an investment with returns "
        "analogous to education{12}, and Grossman's health-capital model "
        "formally describes individual health stock as accumulated through "
        "investment (healthcare utilization, prevention) and depreciated by "
        "aging and disease{13}. At the macro level, Bloom, Canning, and "
        "Sevilla demonstrated that improved population health raises total "
        "factor productivity, with the effect operating through both labor "
        "productivity and labor force participation channels{14}. Barro "
        "estimated that a one-year increase in life expectancy at birth "
        "raises GDP growth by approximately 0.04 percentage points in a "
        "panel of countries spanning 1960-2000{15}."
    )
    add_para(doc,
        "The Health-Led Growth Hypothesis (HLGH) synthesizes these supply-side "
        "effects, positing bidirectional Granger causality between health "
        "expenditure and GDP growth. Ertugrul et al. applied cross-sectional "
        "augmented autoregressive distributed lags (CS-ARDL) to 36 OECD "
        "countries and confirmed long-run bidirectional causality{7}. "
        "Amiri and Ventelou used the Toda-Yamamoto approach for 20 OECD "
        "countries and found evidence of health expenditure causing GDP "
        "growth{8}. Beylik et al. applied Driscoll-Kraay standard errors "
        "to 35 OECD countries and confirmed a positive relationship between "
        "health expenditure indicators and economic growth{9}. Wang "
        "employed quantile panel-type analysis for 31 countries and found "
        "that health expenditure promotes economic growth, particularly in "
        "high-expenditure quantiles{10}. Piabuo and Tieguhong reviewed the "
        "literature for developing economies and found consistent positive "
        "effects of health expenditure on growth{11}."
    )
    add_para(doc,
        "A further analytical dimension comes from the Bongaarts-Feeney "
        "tempo framework. Bongaarts and Feeney demonstrated that rising "
        "mean age at childbearing mechanically depresses the period total "
        "fertility rate even when cohort fertility is unchanged -- the "
        "'tempo effect'{16}. Goldstein, Lutz, and Scherbov introduced the "
        "parity-specific variance sigma, substantially improving the "
        "reconciliation between tempo-adjusted rates and cohort data -- "
        "a 'forgotten parameter'{17}. "
        "Onishi ported this quantum-tempo decomposition to capital "
        "accounting, introducing the investment-to-output time-to-build "
        "lag mu(t) and the intangible capital share beta as analogous "
        "'forgotten parameters' in GDP accounting. Across 39 countries, "
        "a time-varying mu(t) reduced the out-of-sample mean absolute "
        "percentage error from 4.60% to 3.99% (a 13% relative improvement){18}."
    )
    add_para(doc,
        "Extension of this framework to healthcare (termed 'Candidate A-H' "
        "in the companion proof-of-concept){19} models the spending-to-outcome "
        "lag mu_H(t) and revealed a tempo drift of +0.15 yr/yr across "
        "39 countries (2000-2019), approximately four times GDP's "
        "+0.04 yr/yr{18}. This three-layer structure -- Population tempo "
        "(mean age at childbearing drift), GDP tempo (investment-to-output "
        "lag drift), and Healthcare tempo (spending-to-outcome lag drift) -- "
        "indicates that flow-biased policy evaluation systematically "
        "underestimates stock accumulation, with the underestimation "
        "being most severe in the healthcare domain."
    )
    add_para(doc,
        "Japan is also known as a diagnostic imaging equipment powerhouse, "
        "with computed tomography (CT) scanner and magnetic resonance imaging "
        "(MRI) unit densities far exceeding other OECD countries{4,23}. "
        "Japan has 115.7 CT scanners per million population, approximately "
        "4.3 times the OECD median, and 55.2 MRI units per million, "
        "approximately 2.9 times the OECD median. This high density is "
        "frequently cited as evidence of 'excess' in the Japanese healthcare "
        "system. Simultaneously, Japan is a net importer of medical devices "
        "and pharmaceuticals: net pharmaceutical imports amount to "
        "approximately USD 23.4 billion, and the total medical import "
        "leakage represents approximately 5.0% of CHE{24}. These trade "
        "flows reduce the effective domestic economic multiplier of "
        "healthcare spending."
    )
    add_para(doc,
        "The aim of this study was to (1) re-evaluate healthcare expenditure "
        "from a neutral standpoint as both cost and economic effect, using "
        "a dual-return framework integrating demand-side I-O multipliers and "
        "supply-side health-capital tempo effects, "
        "(2) quantify the role of diagnostic equipment stock as 'diagnostic "
        "capital' contributing to health-capital accumulation, "
        "(3) model the impact of medical import leakage on effective I-O "
        "multipliers across 13 OECD countries, and (4) construct "
        "counterfactual scenarios to assess whether Japan's high equipment "
        "density and import dependence enhance or diminish healthcare "
        "sustainability."
    )

    # ======================================================================
    # METHODS
    # ======================================================================
    add_heading(doc, "Methods", level=1)

    add_heading(doc, "Study design and data sources", level=2)
    add_para(doc,
        "This study is a cross-country ecological analysis using publicly "
        "available data from OECD Health at a Glance 2023{4,23}, the World "
        "Bank World Development Indicators{21}, published I-O studies of "
        "healthcare sectors{1-3,6,20}, Health-Led Growth Hypothesis panel "
        "studies{7-11}, and pharmaceutical/device trade statistics{24}. "
        "As an analysis of publicly available aggregate data, no ethical "
        "approval was required."
    )

    add_heading(doc, "Country selection", level=2)
    add_para(doc,
        "Thirteen OECD countries were selected based on availability of "
        "healthcare sector I-O multiplier estimates: (a) a direct I-O study "
        "published in a peer-reviewed journal or official government report, "
        "or (b) backward-linkage coefficients from the EU-28 I-O framework "
        "study{3}. The countries included were: Japan, USA, Germany, "
        "Australia, Korea, Italy, France, UK, Canada, Sweden, Spain, "
        "Netherlands, and Finland."
    )

    add_heading(doc, "Fiscal return ratio", level=2)
    add_para(doc,
        "We propose evaluating healthcare sustainability through a dual-return "
        "structure encompassing demand-side and supply-side returns (Figure 1). "
        "For each country, we calculated a demand-side fiscal return ratio:"
    )
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("Fiscal Return Ratio = (\u03c4 \u00d7 m) / pf")
    run.italic = True
    run.font.size = Pt(12)
    add_para(doc,
        "where \u03c4 is the effective tax rate (total tax and social insurance "
        "contributions as a proportion of GDP), m is the healthcare sector I-O "
        "multiplier, and pf is the public financing share of CHE. "
        "A ratio of 1.0 or above indicates that demand-side tax revenues "
        "generated by healthcare-induced economic activity are sufficient to "
        "cover the public cost of healthcare."
    )

    add_heading(doc, "Tempo model integration", level=2)
    add_para(doc,
        "The tempo model from a companion proof-of-concept{19} was integrated "
        "to capture supply-side health-capital accumulation. The model "
        "specifies the spending-to-outcome lag as "
        "\u03bc_H(t) = \u03bc_H0 + \u03bc_H1 \u00d7 (year - t0), "
        "estimated across 39 countries using World Bank data (2000-2019). "
        "Three model variants were compared: M0 (flow-only, no lag), "
        "M1 (constant lag, \u03bc* = 4 years), and M2 (tempo, time-varying lag). "
        "Performance was assessed by level RMSE and leave-one-out "
        "cross-validation RMSE."
    )

    add_heading(doc, "Equipment stock and import leakage model", level=2)
    add_para(doc,
        "Diagnostic imaging equipment density was defined as the sum of "
        "CT scanners and MRI units per million population, sourced from "
        "OECD Health at a Glance 2023{23}. Medical import leakage was "
        "calculated as the net imports of pharmaceuticals (HS30) and "
        "medical devices (HS9018-22) divided by CHE, representing the "
        "fraction of healthcare spending flowing to foreign suppliers{24}. "
        "The effective I-O multiplier was defined as:"
    )
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("m_eff = m \u00d7 (1 \u2212 import leakage rate)")
    run.italic = True
    run.font.size = Pt(12)

    add_heading(doc, "Counterfactual scenarios", level=2)
    add_para(doc,
        "Three counterfactual scenarios were constructed for Japan: "
        "Scenario A assumed Japan's equipment density was reduced to the OECD "
        "median, with the proportion of CHE attributable to equipment-related "
        "spending (estimated at 15%) adjusted proportionally; "
        "Scenario B assumed complete domestic manufacturing (zero import leakage); "
        "Scenario C combined both adjustments. "
        "Fiscal return ratios were recalculated under each scenario."
    )

    add_heading(doc, "Preston Curve overfit test", level=2)
    add_para(doc,
        "The relationship between CHE (as a percentage of GDP) and life "
        "expectancy was analyzed using nested-model F-tests to determine "
        "whether the quadratic term (Preston Curve) was significant{22}. "
        "Models were fit with and without the US, and compared using "
        "F-statistics, Akaike Information Criterion (AIC), Bayesian "
        "Information Criterion (BIC), and leave-one-out cross-validation "
        "(LOOCV) root mean square error (RMSE)."
    )

    # ======================================================================
    # RESULTS
    # ======================================================================
    add_heading(doc, "Results", level=1)

    add_heading(doc, "I-O multipliers and fiscal return ratios", level=2)
    add_para(doc,
        "Table 1 presents the healthcare sector I-O multipliers for the "
        "13 selected countries. Japan had the highest multiplier at 2.78 "
        "(95% CI: 2.74-2.90){1}, followed by France (2.20), Germany "
        "(2.10){3,20}, and Sweden (2.05). South Korea (2.90) ranked second "
        "based on national I-O table estimates. The US had the lowest "
        "multiplier at 1.70{2}, with the UK (1.76), Canada (1.85), and "
        "Australia (1.89) also below 2.0. The range of 1.70 to 2.90 "
        "across the sample indicates substantial cross-country variation "
        "in the domestic economic return of healthcare spending. "
        "Figure 2 displays the multipliers graphically."
    )

    # Table 1 -- inline since it fits on one page
    io_df = pd.read_csv(os.path.join(DATA, "io_multipliers.csv"))
    io_display = io_df[["country", "multiplier", "year", "source"]].copy()
    io_display.columns = ["Country", "Multiplier", "Year", "Source"]
    add_table_from_df(doc, io_display,
                      "Table 1. Healthcare sector I-O multipliers")

    add_para(doc,
        "Table 2 presents the fiscal return ratios computed as "
        "(effective tax rate \u00d7 I-O multiplier) / public financing share "
        "of CHE. Five of thirteen countries achieved a demand-side-only "
        "ratio at or above 1.0: France (1.18), Italy (1.13), Japan (1.09), "
        "Sweden (1.04), and Finland (1.04). These five countries share "
        "characteristics of relatively high effective tax rates (32-46%), "
        "high I-O multipliers (1.98-2.78), and moderate-to-high public "
        "financing shares (76-85%). The remaining eight countries ranged "
        "from 0.76 (UK) to 0.96 (Germany), recovering the bulk of public "
        "expenditure through demand-side tax revenues alone. The US had "
        "a ratio of 0.83, reflecting its low multiplier despite the "
        "lowest public financing share (51%) in the sample (Figure 3)."
    )

    # Table 2 -- fiscal sustainability
    sust_df = pd.read_csv(os.path.join(DATA, "neutral_sustainability.csv"))
    sust_display = sust_df[["country", "io_multiplier", "eff_tax_rate",
                             "public_share_che", "fiscal_return_ratio",
                             "sustainable"]].copy()
    sust_display.columns = ["Country", "I-O Multiplier", "Eff. Tax Rate",
                            "Public Share", "Fiscal Return", "Sustainable"]
    add_table_from_df(doc, sust_display,
                      "Table 2. Fiscal sustainability indicators",
                      legend="Fiscal Return = (\u03c4 \u00d7 m) / pf. "
                             "'Yes' indicates ratio \u2265 1.0.")

    add_heading(doc, "Tempo model performance", level=2)
    add_para(doc,
        "Table 3 summarizes the tempo model comparison from the companion "
        "proof-of-concept{19}. M2 (tempo model) outperformed M1 (constant lag) "
        "in 95% of 39 countries, with median level RMSE improving from 0.510 "
        "(M0) to 0.441 (M1) to 0.434 years (M2). The estimated tempo drift "
        "\u03bc_H1 was +0.15 yr/yr, indicating that the spending-to-outcome "
        "pipeline lengthens by approximately 1.5 years per decade{19}. "
        "This drift is approximately four times larger than GDP's "
        "\u03bc_1 = +0.04 yr/yr{18}, confirming that healthcare is the "
        "domain where tempo correction matters most among the three layers "
        "(Figure 4)."
    )

    # Table 3 -- PoC results
    poc_data = pd.DataFrame([
        {"Model": "M0 (flow-only)", "Level RMSE (yr)": "0.510",
         "Change RMSE": "0.455", "Description": "Immediate effect"},
        {"Model": "M1 (constant lag)", "Level RMSE (yr)": "0.441",
         "Change RMSE": "0.403", "Description": "\u03bc* = 4 yr fixed"},
        {"Model": "M2 (tempo lag)", "Level RMSE (yr)": "0.434",
         "Change RMSE": "0.405",
         "Description": "\u03bc_H1 = +0.15 yr/yr"},
    ])
    add_table_from_df(doc, poc_data,
                      "Table 3. Tempo model comparison (39 countries)",
                      legend="RMSE: root mean square error. M2 outperformed "
                             "M1 in 95% of countries.")

    add_heading(doc, "Preston Curve overfit analysis", level=2)
    add_para(doc,
        "The relationship between CHE (percentage of GDP) and life expectancy "
        "at birth was examined using nested-model F-tests across 38 OECD "
        "countries with complete 2019 data{4,21}. When the US was included "
        "(n = 38), the quadratic term was highly significant (F = 13.2, "
        "p < 0.001), and the quadratic model explained more variance "
        "(R-squared = 0.425) than the linear model (R-squared = 0.209). "
        "However, when the US was excluded (n = 37), the quadratic term "
        "became entirely non-significant (F = 0.5, p = 0.49). The linear "
        "model's adjusted R-squared (0.400) actually exceeded the quadratic "
        "model's adjusted R-squared (0.391), indicating the quadratic term "
        "adds no explanatory power and merely overfits."
    )
    add_para(doc,
        "Model selection criteria confirmed this conclusion: AIC favored "
        "the linear model in the US-excluded sample, as did BIC (which "
        "penalizes additional parameters more heavily). LOOCV RMSE was "
        "lower for the linear fit (2.12 years) than for the quadratic "
        "(2.19 years), indicating better out-of-sample prediction. "
        "The US -- with 17% of GDP spent on healthcare but below-OECD-average "
        "life expectancy of 78.8 years -- is a single high-leverage point "
        "that drives the apparent concavity. Excluding this single "
        "observation eliminates the Preston Curve pattern entirely "
        "(Figure 5). This finding suggests that the widely-cited "
        "'diminishing returns' to healthcare spending in high-income "
        "countries is not a general phenomenon but an artifact driven "
        "by US-specific inefficiencies."
    )

    add_heading(doc, "Diagnostic equipment density", level=2)
    add_para(doc,
        "Japan exhibited the highest diagnostic imaging density among the "
        "13 countries: 115.7 CT scanners and 55.2 MRI units per million "
        "population (combined 170.9), approximately four times the OECD "
        "median of 46{23}. Australia ranked second (86.2), followed by "
        "the US (85.3) and Korea (74.6). The UK had the lowest density "
        "(16.7). Figure 6 shows imaging density by country."
    )

    add_heading(doc, "Import leakage and effective multipliers", level=2)
    add_para(doc,
        "Five of thirteen countries had measurable medical import leakage: "
        "Canada (8.5% of CHE), Australia (7.1%), Japan (5.0%), Spain (5.0%), "
        "and Finland (4.6%). Germany, Italy, France, Sweden, and the "
        "Netherlands had trade surpluses or balanced trade in medical goods, "
        "resulting in zero leakage{24}."
    )
    add_para(doc,
        "For Japan, the nominal I-O multiplier of 2.78 was reduced to an "
        "effective multiplier of 2.64 after adjusting for 5.0% import leakage, "
        "and the fiscal return ratio fell from 1.09 to 1.04. Despite this "
        "reduction, Japan maintained demand-side sustainability (Figure 7)."
    )

    add_heading(doc, "Counterfactual scenarios for Japan", level=2)
    add_para(doc,
        "Table 4 and Figure 8 present the counterfactual analysis. "
        "Under Scenario A (OECD-average equipment density), the fiscal "
        "return ratio fell to 0.98, below the sustainability threshold. "
        "Under Scenario B (complete domestic manufacturing), the ratio "
        "rose to 1.09. Under Scenario C (both adjustments combined), "
        "the ratio was 1.03, as the negative effect of reduced equipment "
        "density approximately offset the positive effect of eliminating "
        "import leakage."
    )

    # Table 4 -- counterfactual
    cf_data = pd.DataFrame([
        {"Scenario": "Baseline", "Equipment Density": "170.9 (actual)",
         "Import Leakage": "5.0%", "Eff. Multiplier": "2.64",
         "Fiscal Return": "1.04"},
        {"Scenario": "A: OECD-avg equipment", "Equipment Density": "46 (median)",
         "Import Leakage": "5.0%", "Eff. Multiplier": "2.50",
         "Fiscal Return": "0.98"},
        {"Scenario": "B: Domestic mfg", "Equipment Density": "170.9 (actual)",
         "Import Leakage": "0%", "Eff. Multiplier": "2.78",
         "Fiscal Return": "1.09"},
        {"Scenario": "C: Both A + B", "Equipment Density": "46 (median)",
         "Import Leakage": "0%", "Eff. Multiplier": "2.64",
         "Fiscal Return": "1.03"},
    ])
    add_table_from_df(doc, cf_data,
                      "Table 4. Counterfactual scenarios for Japan",
                      legend="Equipment density: CT + MRI per million. "
                             "Fiscal return ratio \u2265 1.0 indicates "
                             "demand-side sustainability.")

    # ======================================================================
    # DISCUSSION
    # ======================================================================
    add_heading(doc, "Discussion", level=1)

    add_heading(doc, "Principal findings", level=2)
    add_para(doc,
        "This study proposed a neutral sustainability framework that "
        "re-evaluates healthcare expenditure as both cost and economic effect "
        "through a dual-return lens integrating demand-side I-O multipliers "
        "and supply-side health-capital tempo effects. Three principal "
        "findings emerged."
    )
    add_para(doc,
        "First, the demand-side fiscal return ratio alone showed that five "
        "of thirteen OECD countries (France, Italy, Japan, Sweden, and "
        "Finland) achieve full recovery of public healthcare costs through "
        "tax revenues induced by the economic multiplier effect. The "
        "remaining eight countries recovered 76-96% of public expenditure "
        "through demand-side revenues alone. These findings are consistent "
        "with the HLGH literature, which has established bidirectional "
        "causality between health expenditure and economic growth across "
        "multiple econometric specifications{7-11}. The conventional paradigm "
        "treating healthcare as purely a 'cost' ignores both its I-O "
        "multiplier effect and its supply-side health-capital accumulation. "
        "When supply-side returns -- captured by the tempo model's mu_H "
        "drift of +0.15 yr/yr{19} -- are considered, comprehensive "
        "sustainability likely holds for the majority of countries examined. "
        "This has a direct policy implication: blanket cost-containment "
        "policies may inadvertently reduce economic output, employment, "
        "and tax revenue, potentially worsening the fiscal position they "
        "aim to improve."
    )
    add_para(doc,
        "Second, Japan's high diagnostic imaging density -- frequently "
        "criticized as evidence of excess -- functions as 'diagnostic "
        "capital' that contributes to sustainability. Counterfactual "
        "analysis (Scenario A) demonstrated that reducing density to the "
        "OECD average would lower the fiscal return ratio to 0.98, below "
        "the sustainability threshold. This finding reframes the narrative "
        "around Japan's CT and MRI density: through the lens of health "
        "capital, high equipment density enables earlier detection and more "
        "precise treatment, enhancing the rate at which healthcare spending "
        "converts into health stock. The diagnosis itself is a form of "
        "investment that accelerates health-capital accumulation. However, "
        "Japan's import leakage of approximately 5% of CHE partially "
        "offsets this advantage by reducing the effective I-O multiplier "
        "from 2.78 to 2.64. Scenario B showed that policies promoting "
        "domestic manufacturing could restore the full nominal multiplier "
        "and raise the fiscal return ratio to 1.09. These findings suggest "
        "a two-pronged strategy: maintaining high diagnostic capital while "
        "reducing import dependence."
    )
    add_para(doc,
        "Third, the widely-cited Preston Curve{22} -- used to argue that "
        "additional healthcare spending in high-income countries yields "
        "diminishing returns -- was shown to be an artifact of overfitting "
        "to a single outlier, the United States. When the US was excluded, "
        "the quadratic term was entirely non-significant (F = 0.5, p = 0.49), "
        "and all model selection criteria (AIC, BIC, LOOCV RMSE) favored "
        "a simple linear relationship. Since Preston's original 1975 "
        "analysis{22}, the concave shape has been widely assumed as evidence "
        "of inherent diminishing returns to health investment at higher "
        "income levels. Our finding challenges this long-standing assumption "
        "and implies that the economic return on healthcare investment in "
        "high-income OECD countries may be larger than conventionally assumed."
    )

    add_heading(doc, "Japan-US comparison and tempo implications", level=2)
    add_para(doc,
        "The Japan-US comparison illustrates how institutional structure "
        "affects the economic return on healthcare spending. Japan's high "
        "I-O multiplier (2.78) likely reflects several structural factors: "
        "universal health insurance (kokumin kaihoken) ensuring broad access "
        "across the population, a comprehensive fee schedule that directs "
        "spending within the domestic economy, and a substantial domestic "
        "pharmaceutical and medical-device manufacturing sector{1}. "
        "The US's lower multiplier (1.70) reflects significant leakage "
        "through high drug prices (with profits accruing to multinational "
        "firms), high insurance administrative costs (estimated at 8% of "
        "total spending), and a greater share of spending flowing to "
        "overseas suppliers{2}."
    )
    add_para(doc,
        "Adding the tempo perspective deepens this contrast. The US pattern "
        "is analogous to 'high period TFR, low cohort fertility' in "
        "demography{16}: the flow (current spending) is large but the stock "
        "(health capital) accumulates inefficiently. The US spends 17% of "
        "GDP on healthcare yet achieves below-OECD-average life expectancy. "
        "Through the tempo lens, this represents a 'tempo-inflated flow' "
        "that overstates true investment{18}. The key issue is not the "
        "volume of spending but its composition: the ratio of curative "
        "care (HC.1) to prevention (HC.6) and research and development "
        "(HC.R). This spending composition -- the 'forgotten parameter' "
        "lambda_b in the healthcare domain, analogous to sigma in "
        "demography{17} and beta in GDP accounting{18} -- determines how "
        "efficiently flow spending converts into stock accumulation."
    )

    add_heading(doc, "Cross-country applicability", level=2)
    add_para(doc,
        "The methodology is directly applicable to any country with "
        "I-O multiplier estimates and medical trade data. Among the "
        "13 countries examined, four distinct patterns emerged: "
        "(a) high multiplier with low leakage (Japan, Sweden), yielding "
        "strong sustainability; (b) high multiplier with zero leakage "
        "(Germany, Italy, France), yielding the strongest fiscal returns; "
        "(c) moderate multiplier with significant leakage (Canada, "
        "Australia), with the largest room for improvement through domestic "
        "manufacturing; and (d) low multiplier (US, UK), where structural "
        "reforms beyond trade policy are needed. This typology provides a "
        "diagnostic framework for other OECD and non-OECD countries seeking "
        "to assess their healthcare economic return."
    )

    add_heading(doc, "Implications for preventive medicine policy", level=2)
    add_para(doc,
        "The findings have direct implications for preventive medicine "
        "investment decisions. If the economic return on healthcare spending "
        "is higher than conventionally assumed (due to the absence of "
        "diminishing returns in the Preston analysis and the underestimation "
        "inherent in flow-only evaluation), then investment in preventive "
        "healthcare -- which typically has a longer spending-to-outcome lag "
        "and thus benefits most from tempo correction -- may be substantially "
        "undervalued. The three-layer tempo structure suggests that "
        "preventive interventions, with their longer pipeline to outcomes, "
        "are precisely the domain where flow-biased evaluation is most "
        "misleading. Policies should account for this systematic "
        "underestimation when evaluating the cost-effectiveness of "
        "prevention programs."
    )

    add_heading(doc, "Limitations", level=2)
    add_para(doc,
        "Several limitations should be noted. First, I-O multipliers are "
        "static models derived from input-output tables at a single point "
        "in time; they do not account for price adjustments, supply "
        "constraints, or general equilibrium effects that would moderate "
        "the multiplier in practice. Second, the fiscal return ratio "
        "depends on the effective tax rate parameter, which is calculated "
        "as total tax and social insurance contributions as a proportion "
        "of GDP; this may not fully capture the complexity of national "
        "tax systems, particularly regarding indirect taxes on healthcare "
        "inputs. Third, the equipment density effect on the multiplier "
        "assumed that approximately 15% of CHE is equipment-related "
        "spending, an approximation derived from OECD health capital "
        "expenditure data that requires further validation with more "
        "granular national accounts."
    )
    add_para(doc,
        "Fourth, import leakage was calculated using aggregate Harmonized "
        "System (HS) trade data and does not capture re-exports, "
        "multi-stage supply chain effects, or the domestic value-added "
        "content of imported goods. Fifth, the counterfactual scenarios "
        "are static partial-equilibrium models that do not account for "
        "dynamic market responses (e.g., increased domestic production "
        "may raise domestic prices). Sixth, the tempo model integration "
        "relies on a companion proof-of-concept analysis{19}; while the "
        "model outperforms alternatives in 95% of countries, direct causal "
        "inference between equipment density and health outcomes requires "
        "additional identification strategies such as instrumental "
        "variable or regression discontinuity designs. Finally, this "
        "study focused on demand-side fiscal return as the primary "
        "sustainability metric; comprehensive sustainability assessment "
        "would require incorporating supply-side health-capital returns, "
        "which is an agenda for future research with individual-level data."
    )

    # ======================================================================
    # CONCLUSIONS
    # ======================================================================
    add_heading(doc, "Conclusions", level=1)

    add_para(doc,
        "Healthcare expenditure is simultaneously a cost and an economic "
        "effect. Using I-O multipliers for 13 OECD countries, we demonstrated "
        "that healthcare spending generates 1.70 to 2.90 times its value "
        "in economic output. The demand-side fiscal return ratio -- "
        "measuring whether tax revenues induced by the multiplier effect "
        "cover public healthcare costs -- exceeded 1.0 in five countries "
        "(France, Italy, Japan, Sweden, Finland), with all thirteen "
        "recovering at least 76% of public costs through demand-side "
        "revenues alone."
    )
    add_para(doc,
        "Japan's high diagnostic imaging equipment density (approximately "
        "four times the OECD median) functions as 'diagnostic capital' "
        "contributing to sustainability, not excess. Counterfactual analysis "
        "showed that reducing this density to the OECD average would "
        "paradoxically eliminate demand-side sustainability, while promoting "
        "domestic manufacturing of medical devices and pharmaceuticals "
        "could strengthen it by eliminating import leakage."
    )
    add_para(doc,
        "The Preston Curve's 'diminishing returns' pattern in healthcare "
        "spending was shown to be an artifact of overfitting to the US "
        "outlier, challenging a widely-held assumption in health economics."
    )
    add_para(doc,
        "The policy debate should shift from 'how to contain healthcare "
        "costs' to 'how to maximize the economic return on healthcare "
        "investment,' encompassing both demand-side multiplier effects "
        "and supply-side health-capital accumulation. The three-layer "
        "tempo structure reveals that flow-biased evaluation systematically "
        "underestimates stock accumulation, with the underestimation "
        "being most severe in the healthcare domain. Recognizing healthcare "
        "expenditure as an economic investment -- rather than purely a "
        "fiscal burden -- provides a more complete foundation for "
        "sustainable health policy."
    )

    # ======================================================================
    # LIST OF ABBREVIATIONS
    # ======================================================================
    add_heading(doc, "List of abbreviations", level=1)

    abbrevs = [
        ("AIC", "Akaike Information Criterion"),
        ("BIC", "Bayesian Information Criterion"),
        ("CHE", "Current Health Expenditure"),
        ("CI", "Confidence Interval"),
        ("CT", "Computed Tomography"),
        ("GDP", "Gross Domestic Product"),
        ("HLGH", "Health-Led Growth Hypothesis"),
        ("I-O", "Input-Output"),
        ("LOOCV", "Leave-One-Out Cross-Validation"),
        ("MHLW", "Ministry of Health, Labour and Welfare"),
        ("MRI", "Magnetic Resonance Imaging"),
        ("OECD", "Organisation for Economic Co-operation and Development"),
        ("RMSE", "Root Mean Square Error"),
    ]
    for abbr, full in abbrevs:
        p = doc.add_paragraph()
        run = p.add_run(f"{abbr}: ")
        run.bold = True
        p.add_run(full)

    # ======================================================================
    # DECLARATIONS
    # ======================================================================
    add_heading(doc, "Declarations", level=1)

    add_heading(doc, "Ethics approval and consent to participate", level=2)
    add_plain_para(doc,
        "Not applicable. This study used only publicly available aggregate "
        "data from international organizations (OECD, World Bank) and "
        "published literature. No human participants, human data, or "
        "human tissue were involved.")

    add_heading(doc, "Consent for publication", level=2)
    add_plain_para(doc, "Not applicable.")

    add_heading(doc, "Availability of data and materials", level=2)
    add_plain_para(doc,
        "All data used in this study are publicly available. OECD Health "
        "at a Glance 2023 data are available from https://www.oecd.org/health/. "
        "World Development Indicators are available from "
        "https://databank.worldbank.org/. I-O multiplier estimates are "
        "reported in the cited published studies. Analysis code is available "
        "from the corresponding author upon request.")

    add_heading(doc, "Competing interests", level=2)
    add_plain_para(doc,
        "The authors declare that they have no competing interests.")

    add_heading(doc, "Funding", level=2)
    add_plain_para(doc,
        "[To be completed by author. If no funding, state: "
        "'No funding was received for this study.']")

    add_heading(doc, "Authors' contributions", level=2)
    add_plain_para(doc,
        "TO conceived and designed the study, collected and analyzed the data, "
        "and drafted the manuscript. "
        "[To be completed if additional authors are added.] "
        "All authors read and approved the final manuscript.")

    add_heading(doc, "Acknowledgements", level=2)
    add_plain_para(doc,
        "[To be completed by author, or state 'Not applicable.']")

    # ======================================================================
    # REFERENCES
    # ======================================================================
    doc.add_page_break()
    add_heading(doc, "References", level=1)

    for i, ref in enumerate(REFERENCES, 1):
        p = doc.add_paragraph()
        run = p.add_run(f"{i}. ")
        run.bold = True
        p.add_run(ref)

    # ======================================================================
    # FIGURE LEGENDS (separate section per EHPM)
    # ======================================================================
    doc.add_page_break()
    add_heading(doc, "Figure Legends", level=1)

    figure_legends = [
        ("Figure 1.",
         "Dual-return framework schematic.",
         "Conceptual diagram of the dual-return framework. Healthcare "
         "spending generates demand-side returns via I-O multipliers "
         "(tax revenue recovery) and supply-side returns via health-capital "
         "stock accumulation (future productivity gains). Sustainability "
         "is assessed by comparing total returns to public cost."),
        ("Figure 2.",
         "Healthcare I-O multipliers by country.",
         "Healthcare sector input-output multipliers for 13 OECD countries. "
         "Japan leads at 2.78, followed by Sweden (2.05), France (2.20), "
         "and Germany (2.10). Sources: individual country studies and "
         "EU-28 framework analysis."),
        ("Figure 3.",
         "Fiscal return ratio by country.",
         "Fiscal return ratio = (effective tax rate \u00d7 I-O multiplier) / "
         "public financing share. Values at or above 1.0 (dashed line) "
         "indicate that demand-side tax revenues cover public healthcare costs. "
         "Five countries exceed 1.0: France, Italy, Japan, Sweden, Finland."),
        ("Figure 4.",
         "Three-layer tempo analogy.",
         "Comparison of tempo drift across three domains: Population "
         "(mean age at childbearing, +0.05 yr/yr), GDP (investment-to-output "
         "lag, +0.04 yr/yr), and Healthcare (spending-to-outcome lag, "
         "+0.15 yr/yr). Healthcare exhibits the largest drift, indicating "
         "the domain where tempo correction is most important."),
        ("Figure 5.",
         "CHE vs life expectancy with overfit test.",
         "Healthcare spending (percentage of GDP) versus life expectancy "
         "at birth across OECD countries (2019). Red solid line: linear fit "
         "excluding USA. Gray dashed: quadratic fit including USA "
         "(Preston-style). With USA included, quadratic term is significant "
         "(F = 13.2, p < 0.001); excluding USA, it is non-significant "
         "(F = 0.5, p = 0.49), demonstrating overfit to a single outlier."),
        ("Figure 6.",
         "Diagnostic imaging equipment density.",
         "Combined CT and MRI scanners per million population across 13 OECD "
         "countries (2021 or latest). Japan (170.9) is approximately four "
         "times the OECD median (46). Orange dashed line indicates OECD "
         "median."),
        ("Figure 7.",
         "Import leakage vs effective I-O multiplier.",
         "Medical import leakage (percentage of CHE) plotted against "
         "effective I-O multiplier for 13 countries. Gray x markers show "
         "unadjusted nominal multipliers. Red arrow for Japan shows the "
         "leakage adjustment effect (2.78 to 2.64). Bubble size reflects "
         "CHE magnitude."),
        ("Figure 8.",
         "Japan counterfactual sustainability.",
         "Fiscal return ratios under counterfactual scenarios for Japan. "
         "Baseline (1.04), A: OECD-average equipment density (0.98), "
         "B: domestic manufacturing with zero import leakage (1.09), "
         "C: both adjustments (1.03). Dashed line: sustainability "
         "threshold of 1.0."),
    ]

    for num_label, title, legend in figure_legends:
        p = doc.add_paragraph()
        run = p.add_run(num_label + " ")
        run.bold = True
        run2 = p.add_run(title)
        run2.bold = True
        doc.add_paragraph()
        p2 = doc.add_paragraph()
        p2.add_run(legend)
        doc.add_paragraph()

    # ======================================================================
    # FIGURES (inline for review convenience, per user knowledge)
    # ======================================================================
    doc.add_page_break()
    add_heading(doc, "Figures", level=1)

    figure_files = [
        ("fig4_dual_return_schematic.png", "Figure 1"),
        ("fig1_io_multipliers.png", "Figure 2"),
        ("fig3_fiscal_sustainability.png", "Figure 3"),
        ("fig5_three_layer_analogy.png", "Figure 4"),
        ("fig2_che_vs_lifeexp.png", "Figure 5"),
        ("fig6_equipment_density.png", "Figure 6"),
        ("fig7_import_leakage_multiplier.png", "Figure 7"),
        ("fig8_counterfactual_japan.png", "Figure 8"),
    ]

    for fname, label in figure_files:
        path = get_fig(fname)
        add_figure_inline(doc, path, label)

    out_path = os.path.join(DOCX_DIR, "Healthcare_EHPM_Manuscript.docx")
    doc.save(out_path)
    print(f"Saved: {out_path}")
    return out_path


# ---------------------------------------------------------------------------
# Cover letter
# ---------------------------------------------------------------------------
def build_cover_letter():
    doc = Document()
    _set_ehpm_format(doc)

    # Date
    from datetime import date
    today = date.today().strftime("%B %d, %Y")
    add_plain_para(doc, today)
    doc.add_paragraph()

    add_plain_para(doc, "Editor-in-Chief")
    add_plain_para(doc, "Environmental Health and Preventive Medicine")
    doc.add_paragraph()

    add_plain_para(doc, "Dear Editor,")
    doc.add_paragraph()

    add_plain_para(doc,
        "We are pleased to submit our manuscript entitled "
        '"Healthcare Expenditure as Economic Effect: A Dual-Return '
        "Sustainability Framework Integrating Input-Output Multipliers, "
        'Health-Capital Tempo, and Diagnostic Equipment Stock" '
        "for consideration as a Research Article in Environmental Health "
        "and Preventive Medicine."
    )
    doc.add_paragraph()

    add_plain_para(doc,
        "Healthcare expenditure has traditionally been framed as a fiscal "
        "cost to be contained. This study challenges that paradigm by "
        "demonstrating that healthcare spending also functions as an economic "
        "investment with measurable returns. Using input-output analysis "
        "across 13 OECD countries, we show that demand-side fiscal return "
        "ratios exceed 1.0 in five countries, with Japan achieving 1.09. "
        "We further demonstrate that Japan's high diagnostic imaging "
        "equipment density -- often criticized as excessive -- actually "
        "functions as 'diagnostic capital' contributing to sustainability. "
        "Our counterfactual analysis shows that reducing this density to "
        "the OECD average would paradoxically eliminate sustainability."
    )
    doc.add_paragraph()

    add_plain_para(doc,
        "We believe this manuscript is well-suited for Environmental Health "
        "and Preventive Medicine for several reasons. First, the study "
        "directly builds on the foundational work by Yamada and Imanaka "
        "(2015), published in this journal, which quantified Japan's "
        "healthcare I-O multiplier. Second, the analysis of diagnostic "
        "imaging equipment density and import leakage has direct relevance "
        "to preventive medicine policy in Japan and internationally. "
        "Third, the methodology is readily deployable across OECD countries, "
        "making the findings of broad interest to the journal's readership."
    )
    doc.add_paragraph()

    add_plain_para(doc,
        "Additionally, our statistical demonstration that the Preston Curve's "
        "'diminishing returns' in healthcare spending is an artifact of "
        "overfitting to the US outlier challenges a widely-held assumption "
        "in health economics and has important policy implications for "
        "preventive medicine investment decisions."
    )
    doc.add_paragraph()

    add_plain_para(doc,
        "This manuscript has not been published previously and is not under "
        "consideration by any other journal. All authors have read and "
        "approved the submitted version."
    )
    doc.add_paragraph()

    add_plain_para(doc,
        "We confirm that no part of this research was funded or supported "
        "by firms or organizations related to the tobacco industry."
    )
    doc.add_paragraph()

    add_plain_para(doc, "Thank you for your consideration.")
    doc.add_paragraph()
    add_plain_para(doc, "Sincerely,")
    doc.add_paragraph()
    add_plain_para(doc, "Tatsuki Onishi")
    add_plain_para(doc, "[Affiliation]")
    add_plain_para(doc, "[E-mail]")

    out_path = os.path.join(DOCX_DIR, "Healthcare_EHPM_CoverLetter.docx")
    doc.save(out_path)
    print(f"Saved: {out_path}")
    return out_path


# ---------------------------------------------------------------------------
# PPTX (editable figures)
# ---------------------------------------------------------------------------
def build_ehpm_pptx():
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    figures = [
        ("fig4_dual_return_schematic.png",
         "Figure 1. Dual-Return Framework",
         "Demand-side (I-O multiplier) + Supply-side "
         "(health-capital tempo) = Total return."),
        ("fig1_io_multipliers.png",
         "Figure 2. Healthcare I-O Multipliers",
         "Healthcare sector I-O multipliers for 13 OECD countries."),
        ("fig3_fiscal_sustainability.png",
         "Figure 3. Fiscal Return Ratio",
         "Fiscal Return = (tax rate x multiplier) / public share. "
         "Values >= 1.0 indicate fiscal sustainability."),
        ("fig5_three_layer_analogy.png",
         "Figure 4. Three-Layer Tempo Analogy",
         "Population -> GDP -> Healthcare tempo drift comparison."),
        ("fig2_che_vs_lifeexp.png",
         "Figure 5. CHE vs Life Expectancy",
         "Linear fit excluding US. Quadratic term non-significant "
         "when US excluded (F=0.5, p=0.49)."),
        ("fig6_equipment_density.png",
         "Figure 6. Diagnostic Imaging Density",
         "CT + MRI per million population (OECD 2021). "
         "Japan at 170.9 is 4x OECD median."),
        ("fig7_import_leakage_multiplier.png",
         "Figure 7. Import Leakage vs Multiplier",
         "Leakage reduces effective multiplier. "
         "Japan: 2.78 -> 2.64 (5% leakage)."),
        ("fig8_counterfactual_japan.png",
         "Figure 8. Japan Counterfactual Analysis",
         "Baseline 1.04, A: avg equip 0.98, "
         "B: domestic mfg 1.09, C: both 1.03."),
    ]

    for fname, title, caption in figures:
        path = get_fig(fname)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.2),
                                          PptxInches(12.33), PptxInches(0.6))
        tf = txBox.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = PptxPt(18)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        if os.path.exists(path):
            slide.shapes.add_picture(path, PptxInches(1.5), PptxInches(1.0),
                                      PptxInches(10.33), PptxInches(5.0))
        txBox2 = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(6.3),
                                           PptxInches(12.33), PptxInches(0.8))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        tf2.text = caption
        tf2.paragraphs[0].font.size = PptxPt(10)
        tf2.paragraphs[0].font.italic = True
        tf2.paragraphs[0].alignment = PP_ALIGN.CENTER

    out_path = os.path.join(PPTX_DIR, "Healthcare_EHPM_Figures.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}")
    return out_path


if __name__ == "__main__":
    build_ehpm_manuscript()
    build_cover_letter()
    build_ehpm_pptx()
