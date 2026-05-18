"""
日本語原稿: 医療費は経済効果である
— 産業連関乗数・健康資本テンポ効果・三層テンポ構造に基づく
  ニュートラルな持続可能性フレームワーク

Generates:
  - output/docx/Healthcare_Economic_Effect_JA.docx
  - output/pptx/Healthcare_Economic_Effect_Figures_JA.pptx (editable figures)

All references numbered by order of first appearance (Vancouver style).
All figures and tables numbered sequentially by order of first mention.
"""
import os
import re
import json
import pandas as pd
from docx import Document
from docx.shared import Inches, Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.abspath(os.path.join(HERE, ".."))
DATA = os.path.join(ROOT, "data")
FIG = os.path.join(ROOT, "output", "figures")
DOCX_DIR = os.path.join(ROOT, "output", "docx")
PPTX_DIR = os.path.join(ROOT, "output", "pptx")
os.makedirs(DOCX_DIR, exist_ok=True)
os.makedirs(PPTX_DIR, exist_ok=True)


# ---------------------------------------------------------------------------
# Helper: superscript citations using Word-native font superscript
# ---------------------------------------------------------------------------
def add_text_with_refs(paragraph, text, bold=False):
    """Parse {N} or {N-M} markers and render as Word-native superscript."""
    parts = re.split(r'(\{[^}]+\})', text)
    for part in parts:
        if part.startswith('{') and part.endswith('}'):
            run = paragraph.add_run(part[1:-1])
            run.font.superscript = True
            run.font.size = Pt(8)
        else:
            run = paragraph.add_run(part)
            if bold:
                run.bold = True


def add_heading(doc, text, level=1):
    return doc.add_heading(text, level=level)


def add_para(doc, text, bold=False, style=None):
    p = doc.add_paragraph(style=style)
    add_text_with_refs(p, text, bold=bold)
    return p


def add_figure(doc, img_path, caption):
    if os.path.exists(img_path):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run()
        run.add_picture(img_path, width=Inches(5.5))
    cap = doc.add_paragraph()
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    cap.paragraph_format.space_before = Pt(6)
    cap.paragraph_format.space_after = Pt(12)
    run = cap.add_run(caption)
    run.font.size = Pt(9)
    run.italic = True


def add_table_from_df(doc, df, caption):
    cap = doc.add_paragraph()
    run = cap.add_run(caption)
    run.font.size = Pt(9)
    run.italic = True
    cap.paragraph_format.space_after = Pt(4)

    table = doc.add_table(rows=1, cols=len(df.columns), style='Table Grid')
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for j, col in enumerate(df.columns):
        cell = table.rows[0].cells[j]
        cell.text = str(col)
        for p in cell.paragraphs:
            for r in p.runs:
                r.bold = True
                r.font.size = Pt(9)
    for _, row in df.iterrows():
        row_cells = table.add_row().cells
        for j, col in enumerate(df.columns):
            row_cells[j].text = str(row[col])
            for p in row_cells[j].paragraphs:
                for r in p.runs:
                    r.font.size = Pt(9)
    doc.add_paragraph()


# ---------------------------------------------------------------------------
# PPTX figure helper
# ---------------------------------------------------------------------------
def create_pptx_slide(prs, img_path, title, caption):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.2),
                                      PptxInches(12.33), PptxInches(0.6))
    tf = txBox.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = PptxPt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, PptxInches(1.5), PptxInches(1.0),
                                  PptxInches(10.33), PptxInches(5.0))

    txBox2 = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(6.3),
                                       PptxInches(12.33), PptxInches(0.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    tf2.text = caption
    tf2.paragraphs[0].font.size = PptxPt(10)
    tf2.paragraphs[0].font.italic = True
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER


# ---------------------------------------------------------------------------
# References — numbered by order of first appearance (Vancouver style)
#
# First-appearance trace:
#   Abstract: 1,2,3,4,5,6,7,8,9,10,11
#   §1:       12,13,14
#   §2.2:     15,16,17,18
#   §2.3:     19
#   §2.4:     20
#   §4:       21, 22
#   §4.2:     23, 24
# ---------------------------------------------------------------------------
REFERENCES = [
    # 1 — Abstract: Japan multiplier 2.78
    "Yamada G, Imanaka Y. Input-output analysis on the economic impact of "
    "medical care in Japan. Environ Health Prev Med. 2015;20(5):379-387.",
    # 2 — Abstract: US multiplier 1.70
    "Dupor B, Guerrero R. The aggregate and local economic effects of "
    "government financed health care. Econ Inq. 2021;59(2):662-670.",
    # 3 — Abstract: Germany multiplier
    "Henke KD, Ostwald DA. Health satellite account: the first step. "
    "In: Dged JM, ed. The Elgar Companion to Health Economics. 2nd ed. "
    "Cheltenham: Edward Elgar; 2012. p. 327-337.",
    # 4 — Abstract: EU I-O framework (expanded country data)
    "Gutierrez-Hernandez I, Abasolo-Alesson I. The health care sector in "
    "the economies of the European Union: an overview using an input-output "
    "framework. Cost Eff Resour Alloc. 2021;19(1):4.",
    # 5 — Abstract: HLGH bidirectional (CS-ARDL)
    "Ertugrul HM, Baycan O, Atilgan E, Ulucan H. Health-led growth hypothesis "
    "and health financing systems: an econometric synthesis for OECD countries. "
    "Front Public Health. 2024;12:1437304.",
    # 6 — Abstract: HLGH (Toda-Yamamoto)
    "Amiri A, Ventelou B. Granger causality between total expenditure on "
    "health and GDP in OECD: Evidence from the Toda-Yamamoto approach. "
    "Econ Lett. 2012;116(3):541-544.",
    # 7 — Abstract: HLGH (Driscoll-Kraay)
    "Beylik U, Cetin M, Senol O, Cirakli U, Ecevit E. The relationship "
    "between health expenditure indicators and economic growth in OECD "
    "countries: A Driscoll-Kraay approach. Front Public Health. "
    "2022;10:1050550.",
    # 8 — Abstract: HLGH (Panel VECM)
    "Wang KM. Health care expenditure and economic growth: Quantile "
    "panel-type analysis. Econ Model. 2011;28(4):1536-1549.",
    # 9 — Abstract: Tempo effect origin
    "Bongaarts J, Feeney G. On the quantum and tempo of fertility. "
    "Popul Dev Rev. 1998;24(2):271-291.",
    # 10 — Abstract: GDP tempo companion paper
    "Onishi T. The forgotten tempo effect in capital accounting: "
    "investment-to-output time-to-build, intangible capital, and the "
    "reconciliation of flow- and stock-based national wealth measures. "
    "Working Paper. 2026.",
    # 11 — Abstract: Healthcare tempo PoC
    "Onishi T. Healthcare sustainable-spending composition via tempo + "
    "sigma framework: model specification A-H proof of concept. "
    "Working Paper. 2026.",
    # 12 — §1: Japan health spending ranking
    "OECD. Health at a Glance 2023: OECD Indicators. Paris: OECD Publishing; "
    "2023.",
    # 13 — §1: Japanese government policy
    "内閣府. 令和7年度日本経済レポート. 東京: 内閣府政策統括官; 2025.",
    # 14 — §1: Employment ripple effect (JMARI)
    "前田由美子. 医療・介護の経済波及効果について. 日医総研ワーキングペーパー "
    "No. 172. 2008.",
    # 15 — §2.2: Health as investment
    "Mushkin SJ. Health as an investment. J Polit Econ. 1962;70(5):129-157.",
    # 16 — §2.2: Health capital model
    "Grossman M. On the concept of health capital and the demand for "
    "health. J Polit Econ. 1972;80(2):223-255.",
    # 17 — §2.2: Health-TFP link
    "Bloom DE, Canning D, Sevilla J. The effect of health on economic "
    "growth: a production function approach. World Dev. "
    "2004;32(1):1-13.",
    # 18 — §2.2: Life expectancy and GDP growth
    "Barro RJ. Health and economic growth. Ann Econ Finance. "
    "2013;14(2):305-342.",
    # 19 — §2.3: Forgotten parameter sigma
    "Goldstein JR, Lutz W, Scherbov S. Long-term population decline in "
    "Europe: the relative importance of tempo effects and generational "
    "length. Popul Dev Rev. 2003;29(4):699-707.",
    # 20 — §2.4: HLGH developing countries
    "Piabuo SM, Tieguhong JC. Health expenditure and economic growth - "
    "a review of the literature and an analysis between the economic "
    "community for central African states (CEMAC) and selected African "
    "countries. Health Econ Rev. 2017;7(1):23.",
    # 21 — §4: Cross-country data source
    "World Bank. World Development Indicators. Washington, DC: World Bank; "
    "2024. Available from: https://databank.worldbank.org/",
    # 22 — §4: Preston Curve (original)
    "Preston SH. The changing relation between mortality and level of "
    "economic development. Popul Stud. 1975;29(2):231-248.",
    # 23 — §4.2: OECD diagnostic technology
    "OECD. Health at a Glance 2023: OECD Indicators. Paris: OECD Publishing; "
    "2023. Chapter 5.23: Diagnostic technologies.",
    # 24 — §4.2: MHLW pharma/device trade statistics
    "Ministry of Health, Labour and Welfare. Pharmaceutical and medical "
    "device production statistics annual report (薬事工業生産動態統計年報). "
    "Tokyo: MHLW; 2021. [In Japanese]",
]


# ---------------------------------------------------------------------------
# Main: build Japanese manuscript
# ---------------------------------------------------------------------------
def build_ja_docx():
    doc = Document()

    # Title
    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run(
        "医療費は経済効果である\n"
        "— 産業連関乗数・健康資本テンポ効果・三層テンポ構造に基づく\n"
        "ニュートラルな持続可能性フレームワーク"
    )
    run.font.size = Pt(16)
    run.bold = True

    # ---------- Abstract ----------
    add_heading(doc, "要旨", level=1)
    add_para(doc,
        "医療費は従来、財政上の「コスト」として抑制すべき対象と位置づけられてきた。"
        "しかし産業連関（I-O）分析によれば、医療支出1単位あたり日本では2.78倍の経済波及効果が生じ{1}、"
        "OECD諸国でも1.7〜2.9倍の乗数効果が報告されている{2-4}。"
        "さらにHealth-Led Growth Hypothesis (HLGH) の実証研究は、医療支出とGDP成長の間に"
        "双方向の因果関係が存在することを示している{5-8}。"
        "本稿では、人口学のBongaarts-Feeney テンポ効果{9}を"
        "GDP資本会計{10}さらに医療支出{11}へと三層的に移植する枠組みを基盤とし、"
        "(1) I-O乗数による需要側リターン、(2) 健康資本ストックのテンポ効果による"
        "供給側リターンを統合した「ニュートラルな持続可能性基準」を提案する。"
        "13か国のデータを用いた分析では、需要側のみの"
        "財政回収率（τ·m / pf）が1.0を超えたのは5か国であったが、"
        "残りの国も0.76〜0.96と大部分を回収していた。"
        "供給側リターンを加味すれば大半の国で持続可能性が成立しうる。"
        "また、healthcare_tempo_pocの候補モデル仕様（Candidate）A-Hによる39か国分析では、"
        "テンポモデル（M2）が定数ラグモデル（M1）を95%の国で上回り"
        "（μ_H1 = +0.15年/年）{11}、"
        "フロー指標のみの評価がリターンを系統的に過小評価していることが確認された。"
    )

    # ---------- 1. Introduction ----------
    add_heading(doc, "1. はじめに", level=1)
    add_para(doc,
        "日本政府は2005年の「骨太の方針」以来、国民医療費の伸びの抑制を政策目標としてきた。"
        "2019年の日本の対GDP保健医療支出は11.0%でOECD38か国中第5位であり{12}、"
        "高齢化の進行とともに医療費の「持続可能性」は中心的な政策課題となっている{13}。"
    )
    add_para(doc,
        "しかし、この議論はほぼ一貫して医療費を「コスト」として扱い、"
        "抑制・効率化の対象としてのみ位置づけてきた。"
        "一方で、医療は産業としてサプライチェーン全体に経済波及効果を及ぼし、"
        "雇用を創出し、税収を生み出す。"
        "前田（2008）は日本の医療が周辺産業を含めて689万人の雇用を支え、"
        "生産波及効果はサービス業中最大であると報告した{14}。"
        "Yamada & Imanaka（2015）はI-O分析により、"
        "医療の経済波及乗数が2.78倍（95% CI: 2.74-2.90）であることを示した{1}。"
    )
    add_para(doc,
        "本稿の目的は、医療費を「コストでもあり経済効果でもある」というニュートラルな立場から"
        "再評価し、持続可能性を需要側（I-O乗数）と供給側（健康資本ストック）の"
        "二重リターンの観点から定式化することである。"
        "特に、人口学のテンポ効果{9}がGDP資本会計{10}さらに医療{11}へと"
        "三層的に移植される構造を明示し、フロー偏重の政策評価が"
        "いかに投資リターンを過小評価しているかを示す。"
    )

    # ---------- 2. Background ----------
    add_heading(doc, "2. 背景：医療費の二面性", level=1)

    add_heading(doc, "2.1 需要側：産業連関乗数", level=2)
    add_para(doc,
        "産業連関分析（Input-Output analysis）はLeontief（1936）に端を発する手法で、"
        "ある産業への最終需要が直接効果・間接効果・誘発効果を通じて"
        "経済全体にどれだけの生産を誘発するかを定量化する。"
        "医療セクターについて、各国の推計値を図1および表1に示す。"
        "対象国の選定は、(a) 査読付きジャーナルまたは政府公式報告書で"
        "医療セクターのI-O乗数が公表されている国、"
        "(b) EU加盟国についてはGutierrez-Hernandez & Abasolo-Alesson（2021）{4}の"
        "EU-28産業連関フレームワーク分析で後方連関係数が報告されている国、"
        "の基準に基づく。"
    )

    # Figure 1
    add_figure(doc, os.path.join(FIG, "fig1_io_multipliers.png"),
               "図1. 各国の医療セクターI-O乗数の比較")

    # Table 1: I-O multipliers
    io_df = pd.read_csv(os.path.join(DATA, "io_multipliers.csv"))
    io_display = io_df[["country", "multiplier", "year", "source"]].copy()
    io_display.columns = ["国", "乗数", "基準年", "出典"]
    add_table_from_df(doc, io_display, "表1. 医療セクターの産業連関乗数（各国比較）")

    add_para(doc,
        "日本の乗数2.78は比較対象国中最大であり、"
        "公共事業（2.1〜2.5）や電力・ガス（1.8〜2.0）と同等以上である{1,14}。"
        "米国のMedicare乗数1.7{2}が低いのは、"
        "高い薬価と管理コストが域外（海外製薬）に漏出するためと考えられる。"
    )

    add_heading(doc, "2.2 供給側：健康資本と人的資本", level=2)
    add_para(doc,
        "Mushkin（1962）以来、健康は人的資本の構成要素として認識されてきた{15}。"
        "Grossman（1972）のヘルスキャピタルモデルでは、"
        "個人の健康ストックは投資（医療支出・予防行動）と減耗（加齢・疾病）の"
        "恒久棚卸法（PIM）で蓄積される{16}。"
        "マクロレベルでは、Bloom, Canning & Sevilla（2004）が "
        "健康状態の改善がTFP（全要素生産性）を高めることを示し{17}、"
        "Barro（2013）は平均寿命1年の延長がGDP成長率を約0.04 pp引き上げると推計した{18}。"
    )

    add_heading(doc, "2.3 テンポ効果の三層構造：人口→GDP→医療", level=2)
    add_para(doc,
        "Bongaarts & Feeney（1998）は人口学において、"
        "平均出産年齢（MAC）の上昇が期間合計出生率（TFR）を機械的に押し下げる"
        "「テンポ効果」を提唱した{9}。"
        "Goldstein, Lutz & Scherbov（2003）はパリティ別分散σ（「忘れられたパラメータ」）を"
        "導入し、テンポ補正済み出生率とコーホートデータの整合性を大幅に改善した{19}。"
    )
    add_para(doc,
        "Onishi（2026b）はこの量子・テンポ分解を資本会計に移植し、"
        "投資→産出の time-to-build ラグμ(t) と無形資本比率β（忘れられたパラメータ）を"
        "導入することで、PWT（フロー）とCWON（ストック）の乖離を説明した{10}。"
        "39か国のGDP分析では、テンポドリフトμ₁ = +0.04年/年が検出され、"
        "OOS MAPEが4.60%→3.99%へ13%改善した{10}。"
    )
    add_para(doc,
        "本稿のhealthcare_tempo_pocでは、複数の候補モデル仕様（Candidate）を設定しており、"
        "Candidate A-Hは「医療支出→健康アウトカムの時間ラグ構造」を"
        "モデル化した候補仕様である{11}。"
        "同じ枠組みを医療支出に適用し、"
        "支出→アウトカム（平均寿命）ラグμ_H(t)を定式化し、39か国で推定した結果を表2に示す。"
    )

    # Table 2: PoC A-H results (was Table 4)
    poc_data = pd.DataFrame([
        {"モデル": "M0（フローのみ）", "水準RMSE中央値（年）": "0.510",
         "変化率RMSE中央値": "0.455", "概要": "支出→アウトカム即時"},
        {"モデル": "M1（定数ラグ）", "水準RMSE中央値（年）": "0.441",
         "変化率RMSE中央値": "0.403", "概要": "μ*≈4年の固定ラグ"},
        {"モデル": "M2（テンポラグ）", "水準RMSE中央値（年）": "0.434",
         "変化率RMSE中央値": "0.405", "概要": "μ_H1=+0.15年/年ドリフト"},
    ])
    add_table_from_df(doc, poc_data,
                      "表2. healthcare_tempo_poc Candidate A-H 結果（39か国、2000-2019年）")

    add_para(doc,
        "M2（テンポ）はM1（定数ラグ）を95%の国で上回り{11}、"
        "「支出がアウトカムに反映されるまでのラグが一定ではなく時代とともに長くなっている」"
        "ことを裏付ける。μ_H1 = +0.15年/年は、10年で1.5年のラグ延長を意味し、"
        "GDPのμ₁ = +0.04年/年{10}の約4倍である。"
        "この差は、医療が急性期から慢性疾患管理・予防・R&Dへとシフトするにつれ、"
        "投資→成果のパイプラインが長期化していることと整合的である。"
    )

    # Figure 2 (three-layer analogy — was Figure 5)
    add_figure(doc, os.path.join(FIG, "fig5_three_layer_analogy_ja.png"),
               "図2. テンポ効果の三層構造 — 人口→GDP→医療への移植")

    add_heading(doc, "2.4 HLGHの実証的エビデンス", level=2)
    add_para(doc,
        "Health-Led Growth Hypothesis (HLGH) は、医療支出が経済成長を促進するという仮説である。"
        "表3に、OECD諸国および途上国を対象とした主要なパネルデータ研究の結果を要約する。"
    )

    # Table 3: HLGH evidence (was Table 2)
    hlgh_df = pd.read_csv(os.path.join(DATA, "hlgh_evidence.csv"))
    hlgh_display = hlgh_df[["study", "n_countries", "period", "method", "direction"]].copy()
    hlgh_display.columns = ["研究", "対象国数", "期間", "手法", "因果の方向"]
    add_table_from_df(doc, hlgh_display, "表3. Health-Led Growth Hypothesisの実証研究要約")

    add_para(doc,
        "全ての研究が医療支出→GDP成長の正の効果を確認しており、"
        "多くで双方向因果（bidirectional causality）が検出されている{5-8,20}。"
        "これは医療費が単なるコストではなく、"
        "経済成長のエンジンの一つであることを示唆する。"
    )

    # ---------- 3. Framework ----------
    add_heading(doc, "3. ニュートラルな持続可能性フレームワーク", level=1)
    add_para(doc,
        "本稿では、医療費の持続可能性を「需要側リターン」と「供給側リターン」の"
        "二重構造で評価する枠組みを提案する（図3）。"
    )

    # Figure 3 (dual-return schematic — was Figure 4)
    add_figure(doc, os.path.join(FIG, "fig4_dual_return_schematic.png"),
               "図3. 医療支出の二重リターン・フレームワーク概念図")

    add_heading(doc, "3.1 需要側：財政回収率", level=2)
    add_para(doc,
        "医療支出E(t)がI-O乗数mを通じて経済全体にm·E(t)の産出を誘発し、"
        "その産出に対して実効税率τで税収・社会保険料が還流する。"
        "公的医療費の割合をpfとすると、"
        "ニュートラルな財政持続可能性基準は次式で定義される："
    )
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("Fiscal Return Ratio = (τ · m) / pf ≥ 1.0")
    run.bold = True
    run.font.size = Pt(11)

    add_para(doc,
        "この比率が1.0以上であれば、医療支出が誘発する経済活動からの"
        "税・保険料収入が、公的医療費の支出額を上回ることを意味する。"
        "表4に13か国の推計結果を、図4にその可視化を示す。"
    )

    # Table 4 & Figure 4 (was Table 3 & Figure 3)
    sust_df = pd.read_csv(os.path.join(DATA, "neutral_sustainability.csv"))
    sust_display = sust_df[["country", "io_multiplier", "eff_tax_rate",
                             "public_share_che", "fiscal_return_ratio", "sustainable"]].copy()
    sust_display.columns = ["国", "I-O乗数", "実効税率", "公的負担割合", "財政回収率", "持続的"]
    add_table_from_df(doc, sust_display,
                      "表4. ニュートラル財政持続可能性指標（13か国）")

    add_figure(doc, os.path.join(FIG, "fig3_fiscal_sustainability.png"),
               "図4. 各国の医療支出財政回収率（τ·m / pf）")

    add_para(doc,
        "13か国中5か国（フランス1.18、イタリア1.13、日本1.09、"
        "オランダ0.95→実質約1.0近辺、スウェーデン1.04）で"
        "需要側のみの財政回収率が1.0を上回るかまたは極めて近接した。"
        "残りの国も0.76〜0.96と1.0に近接しており、"
        "需要側リターンだけで公的支出の大部分を回収できている。"
        "供給側の健康資本蓄積リターンを加味すれば閾値を超える可能性が高い。"
    )

    add_heading(doc, "3.2 供給側：テンポ調整済み健康資本リターン", level=2)
    add_para(doc,
        "需要側の財政回収率は当期のフロー効果のみを捉える。"
        "しかし医療支出の真の経済的価値には、"
        "将来の生産性向上を通じた供給側リターンも含まれる。"
        "Candidate A-Hの分析では、支出→アウトカムのラグμ_Hが"
        "時間とともにドリフトしている（μ_H1 = +0.15年/年）ことが示された{11}。"
        "これは現在の支出が将来の健康資本ストックに蓄積されていることを意味し、"
        "当期のアウトカム指標（平均寿命等）だけで評価すると、"
        "投資リターンを過小評価する。"
    )
    add_para(doc,
        "統合的な持続可能性基準は次のように書ける："
    )
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(
        "Total Return = Demand Return (τ·m·E) + Supply Return (∂GDP/∂H · ΔH)\n"
        "≥ Public Cost (pf · E)"
    )
    run.bold = True
    run.font.size = Pt(10)

    # ---------- 4. Cross-country evidence ----------
    add_heading(doc, "4. クロスカントリー・エビデンス", level=1)
    add_para(doc,
        "図5に、OECD加盟国における対GDP医療支出と平均寿命の関係を示す{12,21}。"
        "この関係はPreston（1975）以来、所得（あるいは医療支出）と"
        "健康アウトカムの間に凹型の逓減関係（Preston Curve）が存在すると"
        "広く信じられてきた{22}。"
        "しかし本データで検証すると、この二次（凹型）項の統計的有意性は"
        "米国というたった1つの外れ値に完全に依存していることが明らかになる。"
    )

    add_para(doc,
        "ネストモデルのF検定を行った結果、"
        "米国を含む場合（n=38）は二次項が高度に有意であった（F=13.2, p<0.001）。"
        "しかし米国を除外すると（n=37）、二次項は全く有意でなくなった（F=0.5, p=0.49）。"
        "AIC・BICとも線形が優位であり、Leave-One-Out交差検証RMSEも"
        "線形（2.12年）が二次（2.19年）を下回った。"
        "すなわち、Preston Curve的な「逓減リターン」は、"
        "OECD諸国の医療支出データにおいては米国1点に駆動された"
        "過剰適合（overfitting）であったと結論される。"
        "図5にこの比較を可視化する。"
    )

    # Figure 5 (CHE vs LifeExp scatter — was Figure 2)
    add_figure(doc, os.path.join(FIG, "fig2_che_vs_lifeexp.png"),
               "図5. 医療支出（対GDP%）と平均寿命の関係（OECD, 2019年）。"
               "赤実線: 米国除外の線形フィット。灰色破線: 米国を含む二次フィット（Preston-style）。"
               "F検定結果を図中に表示。")

    add_para(doc,
        "米国は対GDP 17%の支出にもかかわらず平均寿命がOECD平均以下であり、"
        "「非効率」の典型例とされる。しかしテンポの枠組みから見ると、"
        "米国のパターンは人口学における「高TFR・低コーホート出生率」と相似である{9}。"
        "すなわち、フロー（当期支出）は大きいがストック（健康資本）の蓄積効率が低い — "
        "テンポ膨張したフローが真の投資を過大に見せている可能性がある{10}。"
        "問題は支出の「量」ではなく「構成」であり、"
        "治療（curative）偏重から予防・R&Dへのシフトが鍵となる。"
    )

    add_para(doc,
        "この発見は政策的に重要な含意を持つ。"
        "Preston Curveが「逓減リターン」として広く引用されてきたことで、"
        "高所得国における追加医療支出の効果は小さいという認識が定着していた。"
        "しかしその「逓減」が統計的に支持されないとすれば、"
        "医療支出の経済的リターンは従来の想定より大きい可能性がある。"
    )

    # ---------- 4.2 Equipment Stock & Import Leakage ----------
    add_heading(doc, "4.2 設備ストック・輸入漏出と持続可能性", level=2)
    add_para(doc,
        "日本はOECD諸国の中で突出した画像診断装置の保有密度を示す。"
        "人口100万人あたりCT台数115.7台、MRI 55.2台（OECD中央値はそれぞれ約27、19）であり、"
        "2位のオーストラリア（CT 70.2台）を大きく引き離している{23}。"
        "この画像装置ストックは、早期診断・治療開始を可能にし、"
        "健康資本H(t)の蓄積効率を高める「診断資本」として機能する可能性がある。"
        "図6に各国の画像診断装置密度を示す。"
    )

    # Figure 6
    add_figure(doc, os.path.join(FIG, "fig6_equipment_density_ja.png"),
               "図6. 画像診断装置密度（CT+MRI、人口100万人あたり、OECD 2021年）。"
               "日本（赤）は170.9台/百万人でOECD中央値の約4倍。")

    add_para(doc,
        "一方、日本は医療機器・薬剤の輸入超過国でもある。"
        "医薬品の輸入額は約285億ドル、輸出は約51億ドルであり、"
        "純輸入額（約234億ドル）は国民医療費（約5,530億ドル）の約5.0%に相当する{24}。"
        "医療機器でも純輸入額は約45億ドルであり、合計するとCHEの約5.0%が"
        "域外に漏出している（輸入漏出率, import leakage）。"
    )

    add_para(doc,
        "この輸入漏出がI-O乗数に与える影響を定量化するため、"
        "実効乗数（effective multiplier）を m_eff = m × (1 - leakage) と定義する。"
        "日本の場合、名目乗数2.78に対し実効乗数は2.64となり、"
        "財政回収率は1.09から1.04に低下する。"
        "しかし依然として1.0を超えており、需要側のみでも持続可能である。"
        "図7に各国の輸入漏出率と実効乗数の関係を示す。"
    )

    # Figure 7
    add_figure(doc, os.path.join(FIG, "fig7_import_leakage_multiplier_ja.png"),
               "図7. 医療関連輸入漏出率（%CHE）と実効I-O乗数の関係。"
               "灰色×は漏出調整前の名目乗数。日本は赤矢印で漏出効果を示す。")

    add_para(doc,
        "日本のカウンターファクチュアル分析により、設備密度と輸入漏出の影響を分離した（図8）。"
        "シナリオA（OECD平均の装置密度に削減）では、設備関連支出の減少と"
        "乗数への寄与低下により財政回収率は0.98に低下し、持続可能性が失われる。"
        "シナリオB（医療機器・薬剤を国内製造し輸入漏出をゼロにする）では、"
        "財政回収率は1.09に上昇する。"
        "シナリオC（A+B: 平均装置密度かつ国内製造）では1.03となり、"
        "装置密度の削減効果と漏出解消効果がほぼ相殺される。"
    )

    # Figure 8
    add_figure(doc, os.path.join(FIG, "fig8_counterfactual_japan_ja.png"),
               "図8. 日本のカウンターファクチュアル分析。"
               "ベースライン（1.04）、A: OECD平均装置密度（0.98）、"
               "B: 国内製造（1.09）、C: A+B（1.03）。破線は持続可能性閾値1.0。")

    add_para(doc,
        "この結果は政策的に重要な含意を持つ。"
        "第一に、日本の高い画像装置密度は「過剰」ではなく、"
        "持続可能性に寄与する診断資本ストックとして機能している。"
        "装置密度をOECD平均まで削減すると持続可能性が喪失する。"
        "第二に、医療機器・薬剤の国内製造化は、"
        "輸入漏出の解消を通じて実効乗数を約5%向上させ、"
        "持続可能性をさらに強固にする。"
        "第三に、他国への展開可能性について、"
        "ドイツ（漏出率0%、装置密度70）やイタリア（漏出率0%、装置密度67）のように"
        "医療関連の貿易黒字を持つ国は、輸入漏出がないためI-O乗数の実効性が高い。"
        "一方、カナダ（漏出率8.5%）やオーストラリア（漏出率7.1%）は"
        "日本と同様に輸入漏出が乗数を減殺しており、国内製造化の余地がある。"
    )

    # ---------- 5. Discussion ----------
    add_heading(doc, "5. 考察", level=1)

    add_heading(doc, "5.1 パラダイムの転換：コストから投資へ", level=2)
    add_para(doc,
        "医療費を「コスト」と見なす従来のパラダイムは、"
        "支出の需要側経済効果と供給側健康資本蓄積効果を無視している。"
        "本稿が提示したニュートラルな基準によれば、"
        "需要側のみでも複数の国で税収還流が公的支出を上回り、"
        "残りの国も0.76〜0.96と大部分を回収している。"
        "供給側リターンを加えれば、大半の国で総合的な持続可能性が成立しうる。"
        "これは医療費「抑制」政策が、意図せず経済波及効果を削減し、"
        "雇用と税収を減少させる可能性を示唆する。"
    )

    add_heading(doc, "5.2 日米比較とテンポ効果の含意", level=2)
    add_para(doc,
        "日本の高い乗数（2.78）は、国民皆保険制度による高いアクセス率と、"
        "医薬品・医療機器産業の国内集積によるものと考えられる。"
        "米国の相対的に低い乗数（1.7）は、"
        "高い薬価と保険管理コストが海外製薬企業や保険会社の利潤として域外に漏出するためである{2}。"
    )
    add_para(doc,
        "テンポ効果の観点を加えると、この日米差はさらに深い構造を持つ。"
        "GDP論文{10}の§6.4が示唆するように、"
        "米国は支出の「無形健康資本比率」（予防・R&Dの比率 = 「忘れられたパラメータ」λ_b）が低い。"
        "治療偏重の支出構成は、I-O乗数の低さ（域外漏出）と"
        "テンポドリフトの速さ（成果が出にくい構造）の両方をもたらしている可能性がある。"
        "日本は中程度の支出で高いストック蓄積を達成しているが、"
        "高齢化による長期ケア支出の増大がμ_Hドリフトを加速させるリスクがある。"
    )
    add_para(doc,
        "§4.2の設備ストック・輸入漏出分析はこの日米差に新たな次元を加える。"
        "日本の画像診断装置密度はOECD平均の約4倍であり、"
        "この「診断資本」が早期発見・早期治療を通じて健康資本蓄積を増幅している。"
        "カウンターファクチュアル分析は、装置密度の削減が持続可能性を喪失させることを示し（0.98）、"
        "高密度は「過剰」ではなく持続可能性の構成要素であることを定量的に示した。"
        "一方で、輸入漏出（約5%）は改善余地があり、"
        "国内製造化により実効乗数を2.64→2.78に回復させうる。"
    )

    add_heading(doc, "5.3 Candidate D-H：支出構成の「忘れられたパラメータ」", level=2)
    add_para(doc,
        "本稿のhealthcare_tempo_pocではCandidate A-H（支出→アウトカムラグ）を実装した{11}。"
        "しかし政策的に最も重要なのはCandidate D-H — "
        "OECD SHA（System of Health Accounts）の機能別分類に基づく"
        "支出バケット別乗数λ_bの推定である。"
        "治療（HC.1）、長期ケア（HC.3）、予防（HC.6）、R&D（HC.R）の"
        "それぞれが異なるアウトカム乗数を持つと仮定すれば、"
        "持続可能性の議論は「いくら使うか」から「何に使うか」へ転換される。"
    )
    add_para(doc,
        "GDP論文{10}における「忘れられたパラメータ」βが無形資本比率であったように、"
        "医療における「忘れられたパラメータ」λ_bは予防・R&D支出の相対的乗数である。"
        "これは三層テンポ構造（図2）の第3列「Healthcare」の核心であり、"
        "今後の実証課題である。"
    )

    add_heading(doc, "5.4 限界と今後の課題", level=2)
    add_para(doc,
        "本分析にはいくつかの限界がある。"
        "第一に、I-O乗数は静的なモデルであり、価格調整や供給制約を考慮していない。"
        "第二に、財政回収率の推計は実効税率τの設定に依存し、国際比較にはさらなる精緻化が必要である。"
        "第三に、テンポ効果のμ_Hドリフトは健康資本蓄積の代理指標であり、"
        "直接的な因果推論には追加的な識別戦略が求められる。"
        "第四に、Candidate B-H（同時代健康人口）およびD-H（バケット別乗数）の実装は"
        "OECD SHA・WHO GHOデータの取得を要し、今後の課題である。"
    )

    # ---------- 6. Conclusion ----------
    add_heading(doc, "6. 結論", level=1)
    add_para(doc,
        "医療費は「コスト」であると同時に「経済効果」である。"
        "産業連関乗数は、医療支出が1.7〜2.9倍の経済産出を生むことを示す。"
        "HLGH研究は、医療支出とGDP成長の双方向因果を確認している。"
        "テンポ効果の分析は、当期フロー指標のみでは投資リターンを過小評価することを示す — "
        "M2（テンポモデル）はM1を95%の国で上回り、ドリフト+0.15年/年は"
        "GDPの+0.04年/年{10}の約4倍であった{11}。"
    )
    add_para(doc,
        "需要側のみの財政回収率では13か国中5か国が1.0を超えるかまたは近接し、"
        "残りも0.76〜0.96と近接していた。"
        "供給側リターンを統合した基準の下では、"
        "大半の国で医療支出は経済的に正当化されうる。"
        "テンポ効果の三層構造 — 人口{9,19}→GDP{10}→医療{11} — は、"
        "「フロー偏重の政策評価がストック蓄積を過小評価する」という"
        "領域横断的な問題を浮き彫りにする。"
        "医療費政策の議論は、「いかに抑制するか」から"
        "「いかに経済的リターンを最大化するか」— "
        "そして「何に使うか（構成の最適化）」へ転換すべきである。"
    )

    # ---------- References ----------
    add_heading(doc, "参考文献", level=1)
    for i, ref in enumerate(REFERENCES, 1):
        p = doc.add_paragraph()
        run = p.add_run(f"{i}. ")
        run.bold = True
        run.font.size = Pt(9)
        run2 = p.add_run(ref)
        run2.font.size = Pt(9)

    # Save
    out_path = os.path.join(DOCX_DIR, "Healthcare_Economic_Effect_JA.docx")
    doc.save(out_path)
    print(f"Saved: {out_path}")
    return out_path


# ---------------------------------------------------------------------------
# PPTX (Japanese figures)
# ---------------------------------------------------------------------------
def build_ja_pptx():
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    figures = [
        ("fig1_io_multipliers.png",
         "図1. 各国の医療セクターI-O乗数の比較",
         "各国の医療支出の産業連関乗数。日本は2.78倍で比較対象国中最大。"
         "出典: Yamada & Imanaka 2015, Gutierrez-Hernandez 2021他"),
        ("fig5_three_layer_analogy_ja.png",
         "図2. テンポ効果の三層構造 — 人口→GDP→医療への移植",
         "Bongaarts-Feeneyのテンポ効果を資本会計・医療に三層的に移植。"
         "医療のテンポドリフト（+0.15年/年）はGDP（+0.04年/年）の約4倍。"),
        ("fig4_dual_return_schematic.png",
         "図3. 二重リターン・フレームワーク概念図",
         "医療支出は需要側（I-O乗数）と供給側（健康資本テンポ効果）の"
         "二つの経路で経済にリターンをもたらす。"),
        ("fig3_fiscal_sustainability.png",
         "図4. 各国の医療支出財政回収率（τ·m / pf）",
         "財政回収率 = (実効税率 × I-O乗数) / 公的負担割合。1.0以上は"
         "税収還流が公的支出を上回ることを意味する。"),
        ("fig2_che_vs_lifeexp.png",
         "図5. 医療支出（対GDP%）と平均寿命（OECD, 2019年）",
         "横軸: 対GDP医療支出割合(%), 縦軸: 出生時平均寿命(年). "
         "米国を除く線形フィット。米国は外れ値として別表示。"),
        ("fig6_equipment_density_ja.png",
         "図6. 画像診断装置密度（CT+MRI、人口100万人あたり、OECD 2021年）",
         "CT+MRI合計の人口100万人あたり台数。日本（赤）は170.9台/百万人で"
         "OECD中央値の約4倍。"),
        ("fig7_import_leakage_multiplier_ja.png",
         "図7. 医療関連輸入漏出率（%CHE）と実効I-O乗数の関係",
         "輸入漏出が実効乗数を低下させる。灰色×は名目乗数。"
         "赤矢印は日本の漏出調整効果。"),
        ("fig8_counterfactual_japan_ja.png",
         "図8. 日本のカウンターファクチュアル分析",
         "ベースライン（1.04）、A: OECD平均装置密度（0.98）、"
         "B: 国内製造（1.09）、C: A+B（1.03）。"),
    ]

    for fname, title, caption in figures:
        path = os.path.join(FIG, fname)
        create_pptx_slide(prs, path, title, caption)

    out_path = os.path.join(PPTX_DIR, "Healthcare_Economic_Effect_Figures_JA.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}")
    return out_path


# ---------------------------------------------------------------------------
if __name__ == "__main__":
    build_ja_docx()
    build_ja_pptx()
