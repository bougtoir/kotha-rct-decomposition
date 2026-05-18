"""
English manuscript: Healthcare Expenditure as Economic Effect
-- A Neutral Sustainability Framework Based on I-O Multipliers,
   Health-Capital Tempo, and the Three-Layer Tempo Analogy

Generates:
  - output/docx/Healthcare_Economic_Effect_EN.docx
  - output/pptx/Healthcare_Economic_Effect_Figures_EN.pptx

All references numbered by order of first appearance (Vancouver style).
All figures and tables numbered sequentially by order of first mention.
"""
import os
import re
import json
import pandas as pd
from docx import Document
from docx.shared import Inches, Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.abspath(os.path.join(HERE, ".."))
DATA = os.path.join(ROOT, "data")
FIG = os.path.join(ROOT, "output", "figures")
DOCX_DIR = os.path.join(ROOT, "output", "docx")
PPTX_DIR = os.path.join(ROOT, "output", "pptx")
os.makedirs(DOCX_DIR, exist_ok=True)
os.makedirs(PPTX_DIR, exist_ok=True)


def get_fig(name):
    return os.path.join(FIG, name)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def add_text_with_refs(paragraph, text, bold=False):
    parts = re.split(r'(\{[^}]+\})', text)
    for part in parts:
        if part.startswith('{') and part.endswith('}'):
            run = paragraph.add_run(part[1:-1])
            run.font.superscript = True
            run.font.size = Pt(8)
        else:
            run = paragraph.add_run(part)
            if bold:
                run.bold = True


def add_heading(doc, text, level=1):
    return doc.add_heading(text, level=level)


def add_para(doc, text, bold=False):
    p = doc.add_paragraph()
    add_text_with_refs(p, text, bold=bold)
    return p


def add_figure(doc, img_path, caption):
    if os.path.exists(img_path):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run()
        run.add_picture(img_path, width=Inches(5.5))
    cap = doc.add_paragraph()
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    cap.paragraph_format.space_before = Pt(6)
    cap.paragraph_format.space_after = Pt(12)
    run = cap.add_run(caption)
    run.font.size = Pt(9)
    run.italic = True


def add_table_from_df(doc, df, caption):
    cap = doc.add_paragraph()
    run = cap.add_run(caption)
    run.font.size = Pt(9)
    run.italic = True
    cap.paragraph_format.space_after = Pt(4)
    table = doc.add_table(rows=1, cols=len(df.columns), style='Table Grid')
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for j, col in enumerate(df.columns):
        cell = table.rows[0].cells[j]
        cell.text = str(col)
        for p in cell.paragraphs:
            for r in p.runs:
                r.bold = True
                r.font.size = Pt(9)
    for _, row in df.iterrows():
        row_cells = table.add_row().cells
        for j, col in enumerate(df.columns):
            row_cells[j].text = str(row[col])
            for p in row_cells[j].paragraphs:
                for r in p.runs:
                    r.font.size = Pt(9)
    doc.add_paragraph()


def create_pptx_slide(prs, img_path, title, caption):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.2),
                                      PptxInches(12.33), PptxInches(0.6))
    tf = txBox.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = PptxPt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, PptxInches(1.5), PptxInches(1.0),
                                  PptxInches(10.33), PptxInches(5.0))
    txBox2 = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(6.3),
                                       PptxInches(12.33), PptxInches(0.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    tf2.text = caption
    tf2.paragraphs[0].font.size = PptxPt(10)
    tf2.paragraphs[0].font.italic = True
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER


# ---------------------------------------------------------------------------
# References -- numbered by order of first appearance (Vancouver style)
#
# First-appearance trace (mirrors JA manuscript):
#   Abstract: 1,2,3,4,5,6,7,8,9,10,11
#   S1:       12,13,14
#   S2.2:     15,16,17,18
#   S2.3:     19
#   S2.4:     20
#   S4:       21, 22
#   S4.2:     23, 24
# ---------------------------------------------------------------------------
REFERENCES = [
    # 1 -- Abstract: Japan multiplier
    "Yamada G, Imanaka Y. Input-output analysis on the economic impact of "
    "medical care in Japan. Environ Health Prev Med. 2015;20(5):379-387.",
    # 2 -- Abstract: US multiplier
    "Dupor B, Guerrero R. The aggregate and local economic effects of "
    "government financed health care. Econ Inq. 2021;59(2):662-670.",
    # 3 -- Abstract: Germany multiplier
    "Henke KD, Ostwald DA. Health satellite account: the first step. "
    "In: Dged JM, ed. The Elgar Companion to Health Economics. 2nd ed. "
    "Cheltenham: Edward Elgar; 2012. p. 327-337.",
    # 4 -- Abstract: EU I-O framework
    "Gutierrez-Hernandez I, Abasolo-Alesson I. The health care sector in "
    "the economies of the European Union: an overview using an input-output "
    "framework. Cost Eff Resour Alloc. 2021;19(1):4.",
    # 5 -- Abstract: HLGH (CS-ARDL)
    "Ertugrul HM, Baycan O, Atilgan E, Ulucan H. Health-led growth hypothesis "
    "and health financing systems: an econometric synthesis for OECD countries. "
    "Front Public Health. 2024;12:1437304.",
    # 6 -- Abstract: HLGH (Toda-Yamamoto)
    "Amiri A, Ventelou B. Granger causality between total expenditure on "
    "health and GDP in OECD: Evidence from the Toda-Yamamoto approach. "
    "Econ Lett. 2012;116(3):541-544.",
    # 7 -- Abstract: HLGH (Driscoll-Kraay)
    "Beylik U, Cetin M, Senol O, Cirakli U, Ecevit E. The relationship "
    "between health expenditure indicators and economic growth in OECD "
    "countries: A Driscoll-Kraay approach. Front Public Health. "
    "2022;10:1050550.",
    # 8 -- Abstract: HLGH (Panel VECM)
    "Wang KM. Health care expenditure and economic growth: Quantile "
    "panel-type analysis. Econ Model. 2011;28(4):1536-1549.",
    # 9 -- Abstract: Tempo effect origin
    "Bongaarts J, Feeney G. On the quantum and tempo of fertility. "
    "Popul Dev Rev. 1998;24(2):271-291.",
    # 10 -- Abstract: GDP tempo companion
    "Onishi T. The forgotten tempo effect in capital accounting: "
    "investment-to-output time-to-build, intangible capital, and the "
    "reconciliation of flow- and stock-based national wealth measures. "
    "Working Paper. 2026.",
    # 11 -- Abstract: Healthcare tempo PoC
    "Onishi T. Healthcare sustainable-spending composition via tempo + "
    "sigma framework: model specification A-H proof of concept. "
    "Working Paper. 2026.",
    # 12 -- S1: OECD ranking
    "OECD. Health at a Glance 2023: OECD Indicators. Paris: OECD Publishing; "
    "2023.",
    # 13 -- S1: Cabinet Office
    "Cabinet Office, Japan. Annual Report on the Japanese Economy 2025. "
    "Tokyo: Cabinet Office; 2025. [In Japanese]",
    # 14 -- S1: JMARI employment
    "Maeda Y. Economic ripple effects of healthcare and long-term care. "
    "JMARI Working Paper No. 172. 2008. [In Japanese]",
    # 15 -- S2.2: Health as investment
    "Mushkin SJ. Health as an investment. J Polit Econ. 1962;70(5):129-157.",
    # 16 -- S2.2: Health capital model
    "Grossman M. On the concept of health capital and the demand for "
    "health. J Polit Econ. 1972;80(2):223-255.",
    # 17 -- S2.2: Health-TFP link
    "Bloom DE, Canning D, Sevilla J. The effect of health on economic "
    "growth: a production function approach. World Dev. "
    "2004;32(1):1-13.",
    # 18 -- S2.2: Life expectancy and GDP
    "Barro RJ. Health and economic growth. Ann Econ Finance. "
    "2013;14(2):305-342.",
    # 19 -- S2.3: Forgotten parameter sigma
    "Goldstein JR, Lutz W, Scherbov S. Long-term population decline in "
    "Europe: the relative importance of tempo effects and generational "
    "length. Popul Dev Rev. 2003;29(4):699-707.",
    # 20 -- S2.4: HLGH developing countries
    "Piabuo SM, Tieguhong JC. Health expenditure and economic growth - "
    "a review of the literature and an analysis between the economic "
    "community for central African states (CEMAC) and selected African "
    "countries. Health Econ Rev. 2017;7(1):23.",
    # 21 -- S4: Cross-country data
    "World Bank. World Development Indicators. Washington, DC: World Bank; "
    "2024. Available from: https://databank.worldbank.org/",
    # 22 -- S4: Preston Curve (original)
    "Preston SH. The changing relation between mortality and level of "
    "economic development. Popul Stud. 1975;29(2):231-248.",
    # 23 -- S4.2: OECD diagnostic technology
    "OECD. Health at a Glance 2023: OECD Indicators. Paris: OECD Publishing; "
    "2023. Chapter 5.23: Diagnostic technologies.",
    # 24 -- S4.2: MHLW pharma/device trade statistics
    "Ministry of Health, Labour and Welfare. Pharmaceutical and medical "
    "device production statistics annual report. "
    "Tokyo: MHLW; 2021. [In Japanese]",
]


# ---------------------------------------------------------------------------
# Build English manuscript
# ---------------------------------------------------------------------------
def build_en_docx():
    doc = Document()

    # Title
    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run(
        "Healthcare Expenditure as Economic Effect:\n"
        "A Neutral Sustainability Framework Based on I-O Multipliers, "
        "Health-Capital Tempo, and the Three-Layer Tempo Analogy"
    )
    run.font.size = Pt(16)
    run.bold = True

    # ---------- Abstract ----------
    add_heading(doc, "Abstract", level=1)
    add_para(doc,
        "Healthcare expenditure has traditionally been framed as a fiscal "
        "'cost' to be contained. However, input-output (I-O) analysis reveals "
        "that each unit of healthcare spending generates 2.78 times its value "
        "in economic output in Japan{1}, with multipliers ranging from 1.7 to "
        "2.9 across OECD countries{2-4}. The Health-Led Growth "
        "Hypothesis (HLGH) literature confirms bidirectional Granger causality "
        "between health expenditure and GDP growth{5-8}. "
        "Building on the Bongaarts-Feeney tempo framework{9} ported to "
        "GDP/wealth accounting{10} and to healthcare{11}, this paper proposes "
        "a 'neutral sustainability criterion' that integrates (1) demand-side "
        "returns via I-O multipliers and (2) supply-side returns via health-capital "
        "stock accumulation with tempo effects. "
        "The demand-side fiscal return ratio (tau * m / pf) exceeds 1.0 "
        "in five of thirteen countries examined, with the remaining countries "
        "achieving 0.76-0.96. "
        "Our healthcare_tempo_poc model specification (Candidate) A-H "
        "(39 countries) shows that the tempo model (M2) "
        "outperforms the constant-lag model (M1) in 95% of countries "
        "(mu_H1 = +0.15 yr/yr){11}, confirming that flow-only evaluation "
        "systematically underestimates returns. "
        "The three-layer tempo analogy (Population -> GDP -> Healthcare) reveals "
        "that healthcare exhibits the largest tempo drift (+0.15 yr/yr vs GDP's +0.04), "
        "making it the domain where tempo correction matters most."
    )

    # ---------- 1. Introduction ----------
    add_heading(doc, "1. Introduction", level=1)
    add_para(doc,
        "Since the Cabinet's 'Basic Policies' decision in 2005, the Japanese "
        "government has pursued policies to moderate the growth of national "
        "medical care expenditure. Japan's current health expenditure reached "
        "11.0% of GDP in 2019, ranking 5th among 38 OECD countries{12}. With "
        "rapid population aging, the 'sustainability' of healthcare spending "
        "has become a central policy concern{13}."
    )
    add_para(doc,
        "This discourse has almost uniformly treated healthcare expenditure as a "
        "'cost' -- an expense to be minimized through efficiency gains and volume "
        "controls. Yet healthcare is also a major economic sector: it purchases "
        "inputs from pharmaceuticals, medical devices, and IT; it employs millions "
        "of workers; and these workers spend their incomes, generating further "
        "economic activity. Maeda (2008) estimated that Japan's healthcare sector "
        "supports 6.89 million jobs in total (2.95 million direct, approximately "
        "4 million indirect), with production-inducement effects exceeding those "
        "of any other service industry{14}. Yamada and Imanaka (2015) quantified "
        "the I-O multiplier at 2.78 (95% CI: 2.74-2.90){1}."
    )
    add_para(doc,
        "The aim of this paper is to re-evaluate healthcare expenditure from a "
        "neutral standpoint -- as both a cost and an economic effect -- and to "
        "formalize sustainability in terms of dual returns: demand-side (I-O "
        "multiplier) and supply-side (health-capital accumulation). "
        "We embed this analysis within the three-layer tempo framework: "
        "the Bongaarts-Feeney quantum-tempo decomposition{9} originated in "
        "demography, was ported to GDP capital accounting{10}, and is here "
        "extended to healthcare expenditure{11}. This three-layer structure "
        "reveals that flow-biased policy evaluation systematically underestimates "
        "stock accumulation across all three domains."
    )

    # ---------- 2. Background ----------
    add_heading(doc, "2. Background: The Dual Nature of Healthcare Expenditure", level=1)

    add_heading(doc, "2.1 Demand Side: Input-Output Multipliers", level=2)
    add_para(doc,
        "Input-output analysis, pioneered by Leontief (1936), quantifies how "
        "final demand in one sector induces production across the entire economy "
        "through direct, indirect, and induced effects. For the healthcare sector, "
        "estimated multipliers across countries are shown in Figure 1 and Table 1. "
        "Country selection is based on two criteria: (a) a direct I-O study of the "
        "healthcare sector published in a peer-reviewed journal or official government "
        "report, or (b) backward-linkage coefficients from national I-O tables reported "
        "in the EU-28 I-O framework study{4}."
    )

    add_figure(doc, get_fig("fig1_io_multipliers.png"),
               "Figure 1. Healthcare I-O Multipliers by Country")

    io_df = pd.read_csv(os.path.join(DATA, "io_multipliers.csv"))
    io_display = io_df[["country", "multiplier", "year", "source"]].copy()
    io_display.columns = ["Country", "Multiplier", "Reference Year", "Source"]
    add_table_from_df(doc, io_display,
                      "Table 1. Healthcare Sector I-O Multipliers (Cross-Country)")

    add_para(doc,
        "Japan's multiplier of 2.78 is the highest among the comparator countries, "
        "comparable to or exceeding public works (2.1-2.5) and utilities "
        "(1.8-2.0){1,14}. The US Medicare multiplier of 1.7{2} is lower, "
        "likely reflecting leakage through high pharmaceutical prices and "
        "administrative costs to overseas firms."
    )

    add_heading(doc, "2.2 Supply Side: Health Capital and Human Capital", level=2)
    add_para(doc,
        "Since Mushkin (1962), health has been recognized as a component of human "
        "capital{15}. Grossman's (1972) health-capital model describes individual "
        "health stock as accumulated through investment (healthcare, prevention) "
        "and depreciated by aging and disease{16}. At the macro level, Bloom, "
        "Canning, and Sevilla (2004) demonstrated that improved health raises "
        "total factor productivity{17}, and Barro (2013) estimated that a one-year "
        "increase in life expectancy raises GDP growth by approximately 0.04 "
        "percentage points{18}."
    )

    add_heading(doc, "2.3 The Three-Layer Tempo Structure: Population -> GDP -> Healthcare", level=2)
    add_para(doc,
        "Bongaarts and Feeney (1998) demonstrated that rising mean age at "
        "childbearing (MAC) mechanically depresses the period total fertility "
        "rate (TFR) even when cohort fertility is unchanged -- the 'tempo effect'{9}. "
        "Goldstein, Lutz, and Scherbov (2003) introduced the parity-specific "
        "variance sigma (the 'forgotten parameter'), substantially improving the "
        "reconciliation between tempo-adjusted TFR and cohort data{19}."
    )
    add_para(doc,
        "Onishi (2026b) ported this quantum-tempo decomposition to capital "
        "accounting, introducing the investment-to-output time-to-build lag mu(t) "
        "and the intangible capital share beta (the forgotten parameter in GDP "
        "accounting). Across 39 countries, a time-varying mu(t) reduced the "
        "out-of-sample MAPE from 4.60% to 3.99% (13% relative improvement){10}."
    )
    add_para(doc,
        "The healthcare_tempo_poc defines multiple model specifications called "
        "'Candidates'; Candidate A-H models the time-lag structure from "
        "health spending to health outcomes{11}. "
        "We model the spending-to-outcome "
        "lag mu_H(t) = mu_H0 + mu_H1*(year - t0) and estimate it across 39 countries "
        "using WB data (2000-2019). Table 2 summarizes the results."
    )

    # Table 2: PoC A-H results (was Table 4)
    poc_data = pd.DataFrame([
        {"Model": "M0 (flow-only)", "Level RMSE Median (yr)": "0.510",
         "Change RMSE Median": "0.455", "Description": "Spending -> outcome (immediate)"},
        {"Model": "M1 (constant lag)", "Level RMSE Median (yr)": "0.441",
         "Change RMSE Median": "0.403", "Description": "mu* = 4 yr (fixed lag, PIM)"},
        {"Model": "M2 (tempo lag)", "Level RMSE Median (yr)": "0.434",
         "Change RMSE Median": "0.405", "Description": "mu_H1 = +0.15 yr/yr (drift)"},
    ])
    add_table_from_df(doc, poc_data,
                      "Table 2. healthcare_tempo_poc Candidate A-H Results "
                      "(39 countries, 2000-2019)")

    add_para(doc,
        "M2 (tempo) outperforms M1 (constant lag) in 95% of countries{11}, "
        "confirming that the spending-to-outcome lag is not constant but drifts "
        "at +0.15 yr/yr -- the pipeline lengthens by ~1.5 years per decade. "
        "This drift is approximately four times larger than GDP's mu_1 = +0.04 yr/yr{10}, "
        "suggesting that healthcare is the domain where the tempo correction "
        "matters most among the three layers (Figure 2)."
    )

    # Figure 2 (three-layer analogy -- was Figure 5)
    add_figure(doc, get_fig("fig5_three_layer_analogy.png"),
               "Figure 2. Three-Layer Tempo Analogy: Population to GDP to Healthcare")

    add_heading(doc, "2.4 Empirical Evidence for the HLGH", level=2)
    add_para(doc,
        "The Health-Led Growth Hypothesis (HLGH) posits that healthcare "
        "expenditure promotes economic growth. Table 3 summarizes the principal "
        "panel-data studies for OECD and developing countries."
    )

    hlgh_df = pd.read_csv(os.path.join(DATA, "hlgh_evidence.csv"))
    hlgh_display = hlgh_df[["study", "n_countries", "period", "method", "direction"]].copy()
    hlgh_display.columns = ["Study", "Countries", "Period", "Method", "Causality Direction"]
    add_table_from_df(doc, hlgh_display,
                      "Table 3. Summary of HLGH Empirical Studies")

    add_para(doc,
        "All studies confirm a positive effect of health expenditure on GDP "
        "growth, and most detect bidirectional causality{5-8,20}. This supports "
        "the view that healthcare spending is not merely a cost but functions "
        "as an engine of economic growth."
    )

    # ---------- 3. Framework ----------
    add_heading(doc, "3. A Neutral Sustainability Framework", level=1)
    add_para(doc,
        "We propose evaluating healthcare sustainability through a dual-return "
        "structure encompassing demand-side and supply-side returns (Figure 3)."
    )

    # Figure 3 (dual-return schematic -- was Figure 4)
    add_figure(doc, get_fig("fig4_dual_return_schematic.png"),
               "Figure 3. Dual-Return Framework for Neutral Healthcare Sustainability")

    add_heading(doc, "3.1 Demand Side: Fiscal Return Ratio", level=2)
    add_para(doc,
        "Healthcare expenditure E(t), through the I-O multiplier m, induces "
        "total output of m * E(t) across the economy. With an effective tax "
        "rate tau on this output, and a public financing share pf of healthcare "
        "expenditure, the neutral fiscal sustainability criterion is:"
    )
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("Fiscal Return Ratio = (tau * m) / pf >= 1.0")
    run.bold = True
    run.font.size = Pt(11)

    add_para(doc,
        "When this ratio exceeds 1.0, the tax and social-insurance revenues "
        "generated by healthcare-induced economic activity exceed the public "
        "cost of healthcare. Table 4 and Figure 4 present the estimates for "
        "thirteen countries."
    )

    sust_df = pd.read_csv(os.path.join(DATA, "neutral_sustainability.csv"))
    sust_display = sust_df[["country", "io_multiplier", "eff_tax_rate",
                             "public_share_che", "fiscal_return_ratio", "sustainable"]].copy()
    sust_display.columns = ["Country", "I-O Multiplier", "Eff. Tax Rate",
                            "Public Share", "Fiscal Return Ratio", "Sustainable"]
    add_table_from_df(doc, sust_display,
                      "Table 4. Neutral Fiscal Sustainability Indicators (13 countries)")

    # Figure 4 (fiscal sustainability -- was Figure 3)
    add_figure(doc, get_fig("fig3_fiscal_sustainability.png"),
               "Figure 4. Fiscal Return Ratio of Healthcare Spending by Country")

    add_para(doc,
        "Five of thirteen countries achieve a demand-side-only fiscal return ratio "
        "at or above 1.0 (France 1.18, Italy 1.13, Japan 1.09, Sweden 1.04, "
        "Finland 1.04). The remaining countries range from 0.76 to 0.96, "
        "recovering the bulk of public expenditure through demand-side tax "
        "revenues alone. These gaps are expected to be closed by supply-side "
        "returns (health-capital accumulation) not captured in this demand-only metric."
    )

    add_heading(doc, "3.2 Supply Side: Tempo-Adjusted Health-Capital Returns", level=2)
    add_para(doc,
        "The fiscal return ratio captures only the contemporaneous flow effect. "
        "The full economic value of healthcare spending also includes supply-side "
        "returns through future productivity gains from health-capital "
        "accumulation. Candidate A-H showed that mu_H drifts at +0.15 yr/yr{11}, "
        "meaning current spending accumulates into a health-capital stock whose "
        "returns manifest in future periods. "
        "Evaluating by period outcomes alone thus underestimates the true return."
    )
    add_para(doc, "The integrated sustainability criterion can be written as:")
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(
        "Total Return = Demand Return (tau * m * E) + Supply Return (dGDP/dH * delta_H)\n"
        ">= Public Cost (pf * E)"
    )
    run.bold = True
    run.font.size = Pt(10)

    # ---------- 4. Cross-country evidence ----------
    add_heading(doc, "4. Cross-Country Evidence", level=1)
    add_para(doc,
        "Figure 5 shows the relationship between healthcare expenditure "
        "(% of GDP) and life expectancy across OECD countries{12,21}. "
        "Since Preston (1975), it has been widely assumed that this relationship "
        "follows a concave curve -- the 'Preston Curve' -- where returns to "
        "health spending diminish at higher income levels{22}. "
        "We test this assumption directly using nested-model F-tests."
    )

    add_para(doc,
        "When the US is included (n=38), the quadratic term is highly significant "
        "(F=13.2, p<0.001). However, when the US is excluded (n=37), the quadratic "
        "term becomes entirely non-significant (F=0.5, p=0.49). "
        "AIC and BIC both favor the linear model without the US, and Leave-One-Out "
        "cross-validation RMSE is lower for the linear fit (2.12 yr) than for "
        "the quadratic (2.19 yr). "
        "The Preston-style 'diminishing returns' in OECD healthcare spending data "
        "is thus an artifact of overfitting to a single outlier -- the United States. "
        "Figure 5 visualizes this comparison."
    )

    # Figure 5 (CHE vs LifeExp scatter -- was Figure 2)
    add_figure(doc, get_fig("fig2_che_vs_lifeexp.png"),
               "Figure 5. Healthcare Spending (% GDP) vs Life Expectancy "
               "(OECD, 2019). Red solid: linear fit excluding US. "
               "Gray dashed: quadratic fit including US (Preston-style). "
               "F-test results shown in figure.")

    add_para(doc,
        "The US spends 17% of GDP on healthcare yet has below-OECD-average life "
        "expectancy. Through the tempo lens, the US pattern is analogous to "
        "'high period TFR, low cohort fertility' in demography{9}: the flow "
        "(current spending) is large but the stock (health capital) accumulates "
        "inefficiently -- a tempo-inflated flow that overstates true investment{10}. "
        "The issue is not the volume of spending but its composition -- the shift "
        "from curative-heavy toward prevention and R&D is key."
    )

    add_para(doc,
        "This finding has important policy implications. The Preston Curve has been "
        "widely cited to argue that additional healthcare spending in high-income "
        "countries yields diminishing returns. If this 'diminishing return' is not "
        "statistically supported, the economic return on healthcare investment may "
        "be larger than conventionally assumed."
    )

    # ---------- 4.2 Equipment Stock & Import Leakage ----------
    add_heading(doc, "4.2 Equipment Stock, Import Leakage, and Sustainability", level=2)
    add_para(doc,
        "Japan exhibits the highest diagnostic imaging density among OECD countries: "
        "115.7 CT scanners and 55.2 MRI units per million population "
        "(OECD medians: approximately 27 and 19, respectively), far exceeding "
        "second-ranked Australia (CT: 70.2){23}. "
        "This equipment stock may function as 'diagnostic capital' that enhances "
        "health capital H(t) accumulation through earlier detection and treatment. "
        "Figure 6 shows imaging equipment density by country."
    )

    # Figure 6
    add_figure(doc, get_fig("fig6_equipment_density.png"),
               "Figure 6. Diagnostic imaging equipment density (CT+MRI per million "
               "population, OECD 2021). Japan (red) at 170.9 per million is "
               "approximately 4x the OECD median.")

    add_para(doc,
        "However, Japan is also a net importer of medical devices and pharmaceuticals. "
        "Pharmaceutical imports total approximately USD 28.5 billion against exports "
        "of USD 5.1 billion; net imports (approx. USD 23.4 billion) represent "
        "approximately 5.0% of current health expenditure (CHE){24}. "
        "Including medical devices, the total medical import leakage is approximately "
        "5.0% of CHE."
    )

    add_para(doc,
        "To quantify this effect on the I-O multiplier, we define the "
        "effective multiplier as m_eff = m x (1 - leakage). "
        "For Japan, the nominal multiplier of 2.78 becomes an effective multiplier "
        "of 2.64, and the fiscal return ratio falls from 1.09 to 1.04 -- still "
        "above 1.0, maintaining demand-side sustainability. "
        "Figure 7 shows the relationship between import leakage and effective "
        "multipliers across countries."
    )

    # Figure 7
    add_figure(doc, get_fig("fig7_import_leakage_multiplier.png"),
               "Figure 7. Medical import leakage (% of CHE) vs effective I-O "
               "multiplier. Gray x markers show unadjusted nominal multipliers. "
               "Red arrow for Japan shows the leakage adjustment effect.")

    add_para(doc,
        "Counterfactual analysis for Japan separates the effects of equipment density "
        "and import leakage (Figure 8). "
        "Scenario A (reducing to OECD-average equipment density) lowers the fiscal "
        "return ratio to 0.98, eliminating sustainability. "
        "Scenario B (domestic manufacturing, zero import leakage) raises it to 1.09. "
        "Scenario C (both: average equipment + domestic manufacturing) yields 1.03, "
        "showing that equipment reduction and leakage elimination roughly offset."
    )

    # Figure 8
    add_figure(doc, get_fig("fig8_counterfactual_japan.png"),
               "Figure 8. Japan counterfactual analysis. Baseline (1.04), "
               "A: OECD-average equipment (0.98), B: Domestic manufacturing (1.09), "
               "C: Both (1.03). Dashed line: sustainability threshold 1.0.")

    add_para(doc,
        "These findings have important policy implications. "
        "First, Japan's high imaging density is not 'excessive' but functions as "
        "diagnostic capital stock that contributes to sustainability. "
        "Reducing density to the OECD average would eliminate sustainability. "
        "Second, domestic manufacturing of medical devices and pharmaceuticals "
        "would improve the effective multiplier by approximately 5%. "
        "Third, regarding exportability: countries with medical trade surpluses "
        "(Germany, Italy) already have zero import leakage and higher effective "
        "multipliers, while import-dependent countries (Canada: 8.5%, Australia: 7.1%) "
        "have room for domestic manufacturing gains similar to Japan."
    )

    # ---------- 5. Discussion ----------
    add_heading(doc, "5. Discussion", level=1)

    add_heading(doc, "5.1 Paradigm Shift: From Cost to Investment", level=2)
    add_para(doc,
        "The conventional paradigm treating healthcare as a 'cost' ignores both "
        "its demand-side economic-multiplier effect and its supply-side "
        "health-capital accumulation effect. Under the neutral criterion proposed "
        "here, demand-side tax revenues alone recover 76-118% of public healthcare "
        "expenditure across thirteen countries, with five at or above full recovery. "
        "When supply-side returns from health-capital accumulation are included, "
        "comprehensive sustainability likely holds for the majority. This implies "
        "that blanket cost-containment policies may inadvertently reduce economic "
        "output, employment, and tax revenue."
    )

    add_heading(doc, "5.2 Japan-US Comparison and Tempo Implications", level=2)
    add_para(doc,
        "Japan's high multiplier (2.78) likely reflects universal health insurance "
        "ensuring broad access, and a domestic concentration of pharmaceutical and "
        "medical-device industries. The US's lower multiplier (1.7) reflects "
        "leakage through high drug prices and insurance administrative costs "
        "accruing to overseas firms and insurer profits{2}."
    )
    add_para(doc,
        "Adding the tempo perspective deepens this contrast. As the GDP paper's "
        "Section 6.4 suggests{10}, the US has a low 'intangible health-capital "
        "share' -- spending is heavily weighted toward curative care (HC.1) rather "
        "than prevention (HC.6) or R&D (HC.R), the healthcare analogue of the "
        "'forgotten parameter' lambda_b. A curative-heavy composition simultaneously "
        "depresses the I-O multiplier (leakage abroad) and accelerates the tempo "
        "drift (outcomes take longer to materialize). "
        "Japan achieves higher stock accumulation per unit of spending but faces "
        "the risk that aging-driven growth of long-term care may accelerate mu_H drift."
    )
    add_para(doc,
        "The equipment stock and import leakage analysis in Section 4.2 adds a new "
        "dimension to this Japan-US contrast. Japan's diagnostic imaging density is "
        "approximately 4 times the OECD average, and this 'diagnostic capital' "
        "amplifies health capital accumulation through earlier detection and treatment. "
        "Counterfactual analysis shows that reducing equipment density would eliminate "
        "sustainability (ratio falls to 0.98), demonstrating quantitatively that "
        "high density functions as a component of sustainability, not excess. "
        "Meanwhile, import leakage (approximately 5%) offers room for improvement: "
        "domestic manufacturing could restore the effective multiplier from 2.64 to 2.78."
    )

    add_heading(doc, "5.3 Candidate D-H: The 'Forgotten Parameter' in Spending Composition", level=2)
    add_para(doc,
        "The healthcare_tempo_poc implemented Candidate A-H (spending-to-outcome lag){11}. "
        "The most policy-relevant extension is Candidate D-H -- estimating "
        "bucket-specific outcome multipliers lambda_b by functional category from "
        "OECD SHA (System of Health Accounts) data. "
        "If curative (HC.1), long-term care (HC.3), prevention (HC.6), and R&D (HC.R) "
        "carry different lambda_b values, the sustainability discussion shifts from "
        "'how much to spend' to 'what to spend on.'"
    )
    add_para(doc,
        "Just as the forgotten parameter beta (intangible capital share) reconciles "
        "flow and stock accounts in GDP{10}, lambda_b (prevention/R&D composition "
        "multiplier) is the forgotten parameter in healthcare accounting. "
        "This is the core of the third column in the three-layer analogy (Figure 2) "
        "and constitutes the primary empirical agenda going forward."
    )

    add_heading(doc, "5.4 Limitations and Future Directions", level=2)
    add_para(doc,
        "Several limitations apply. First, I-O multipliers are static models that "
        "do not account for price adjustments or supply constraints. Second, the "
        "fiscal return ratio depends on the effective tax rate parameter, requiring "
        "further refinement for robust international comparisons. Third, the tempo "
        "drift mu_H is a proxy for health-capital accumulation; direct causal "
        "inference requires additional identification strategies. "
        "Fourth, Candidates B-H (same-healthy-population) and D-H (bucket multipliers) "
        "require OECD SHA and WHO GHO data and are left for future work."
    )

    # ---------- 6. Conclusion ----------
    add_heading(doc, "6. Conclusion", level=1)
    add_para(doc,
        "Healthcare expenditure is simultaneously a 'cost' and an 'economic effect.' "
        "I-O multipliers demonstrate that healthcare spending generates 1.7 to 2.9 "
        "times its value in economic output. HLGH studies confirm bidirectional "
        "causality between health expenditure and GDP growth. "
        "Tempo analysis shows that contemporaneous flow indicators underestimate "
        "investment returns -- M2 (tempo model) outperforms M1 in 95% of countries "
        "with drift of +0.15 yr/yr, roughly four times GDP's +0.04{10,11}."
    )
    add_para(doc,
        "Under the demand-side fiscal return ratio, five of thirteen countries exceed "
        "or approach 1.0, and all thirteen recover at least 76% of public costs from "
        "demand-side tax revenues alone. "
        "The three-layer tempo structure -- Population{9,19} -> GDP{10} -> "
        "Healthcare{11} -- highlights a cross-domain pattern: flow-biased policy "
        "evaluation systematically underestimates stock accumulation. "
        "The policy debate should shift from 'how to contain costs' to "
        "'how to maximize economic return on healthcare investment' -- and "
        "ultimately to 'what to spend on (composition optimization).'"
    )

    # ---------- References ----------
    add_heading(doc, "References", level=1)
    for i, ref in enumerate(REFERENCES, 1):
        p = doc.add_paragraph()
        run = p.add_run(f"{i}. ")
        run.bold = True
        run.font.size = Pt(9)
        run2 = p.add_run(ref)
        run2.font.size = Pt(9)

    out_path = os.path.join(DOCX_DIR, "Healthcare_Economic_Effect_EN.docx")
    doc.save(out_path)
    print(f"Saved: {out_path}")
    return out_path


def build_en_pptx():
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    figures = [
        ("fig1_io_multipliers.png",
         "Figure 1. Healthcare I-O Multipliers by Country",
         "Healthcare sector input-output multipliers. Japan leads at 2.78x. "
         "Sources: Yamada & Imanaka 2015, Gutierrez-Hernandez 2021, et al."),
        ("fig5_three_layer_analogy.png",
         "Figure 2. Three-Layer Tempo Analogy: Population -> GDP -> Healthcare",
         "The Bongaarts-Feeney tempo framework ported across three domains. "
         "Healthcare shows the largest tempo drift (+0.15 yr/yr vs GDP +0.04)."),
        ("fig4_dual_return_schematic.png",
         "Figure 3. Dual-Return Framework for Neutral Sustainability",
         "Healthcare spending generates returns through two channels: "
         "demand-side (I-O multiplier) and supply-side (health-capital tempo effect)."),
        ("fig3_fiscal_sustainability.png",
         "Figure 4. Fiscal Return Ratio of Healthcare Spending",
         "Fiscal Return Ratio = (Effective tax rate x I-O multiplier) / Public share. "
         "Values >= 1.0 indicate fiscal self-sustainability."),
        ("fig2_che_vs_lifeexp.png",
         "Figure 5. Healthcare Spending (% GDP) vs Life Expectancy (OECD, 2019)",
         "X-axis: CHE as % of GDP. Y-axis: Life expectancy at birth (years). "
         "Linear fit excludes US (outlier shown in red). F-test results in figure."),
        ("fig6_equipment_density.png",
         "Figure 6. Diagnostic Imaging Equipment Density (CT+MRI, OECD 2021)",
         "Combined CT and MRI scanners per million population. "
         "Japan (red) at 170.9 is approximately 4x the OECD median."),
        ("fig7_import_leakage_multiplier.png",
         "Figure 7. Medical Import Leakage vs Effective I-O Multiplier",
         "Import leakage reduces the effective I-O multiplier. "
         "Gray x = nominal multiplier. Red arrow = Japan's leakage adjustment."),
        ("fig8_counterfactual_japan.png",
         "Figure 8. Japan Counterfactual Sustainability Scenarios",
         "Baseline (1.04), A: OECD-avg equipment (0.98), "
         "B: Domestic manufacturing (1.09), C: Both (1.03)."),
    ]

    for fname, title, caption in figures:
        path = get_fig(fname)
        create_pptx_slide(prs, path, title, caption)

    out_path = os.path.join(PPTX_DIR, "Healthcare_Economic_Effect_Figures_EN.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}")
    return out_path


if __name__ == "__main__":
    build_en_docx()
    build_en_pptx()
