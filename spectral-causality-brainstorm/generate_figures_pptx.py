"""
Generate editable pptx with all figures (1 figure per slide).
Widescreen 13.333 x 7.5 inches.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pathlib import Path

FIGURES_DIR = Path('/home/ubuntu/repos/wip/spectral-causality-brainstorm/figures')
OUTPUT = Path('/home/ubuntu/repos/wip/spectral-causality-brainstorm/spectral_causality_figures.pptx')

FIGURES = [
    {
        'file': 'fig1_three_approaches.png',
        'title': 'Figure 1: Three Approaches to Causal Analysis',
        'caption': '(A) Conventional variable-level analysis, (B) LiNGAM causal DAG, (C) Spectral causality theme-level causal flow. UCI Heart Disease Data (n=297).'
    },
    {
        'file': 'fig2_magnetic_laplacian_q.png',
        'title': 'Figure 2: Magnetic Laplacian Eigenvectors',
        'caption': 'Second eigenvector of the magnetic Laplacian plotted in the complex plane for q=0, 0.1, 0.25. Phase angle order encodes causal flow direction.'
    },
    {
        'file': 'fig3_hodge_decomposition.png',
        'title': 'Figure 3: Hodge Decomposition',
        'caption': '(A) Gradient (85.9%) vs curl (14.1%) energy decomposition. (B) Causal potential phi for each variable. Age=most upstream, STDep=most downstream.'
    },
    {
        'file': 'fig4_direction_comparison.png',
        'title': 'Figure 4: Causal Direction Comparison (3 Methods)',
        'caption': 'Comparison of causal directions across all 10 variable pairs. LiNGAM (red), SCD (blue), Hodge potential (green). Green background = all methods agree.'
    },
    {
        'file': 'fig5_hill_radar.png',
        'title': "Figure 5: Hill's 9 Criteria Coverage by Method",
        'caption': "Radar chart showing Hill's 9 criteria coverage. LiNGAM covers H1/H3, Utility Causality covers H6/H7/H9. ECD ensemble covers nearly all criteria."
    },
    {
        'file': 'fig6_causal_dag.png',
        'title': 'Figure 6: DirectLiNGAM Causal DAG',
        'caption': 'Estimated causal DAG by DirectLiNGAM. UCI Heart Disease (Cleveland, n=297). Blue=positive effect, red=negative effect. Edge width proportional to effect size.'
    },
    {
        'file': 'fig7_lingam_vs_spectral.png',
        'title': 'Figure 7: LiNGAM DAG vs Spectral Causality DCG',
        'caption': '(A) LiNGAM DAG, (B) Spectral (alpha=0.6, domain knowledge), (C) Spectral (alpha=0, pure data-driven: no edges detected).'
    },
    {
        'file': 'fig8_alpha_sweep.png',
        'title': 'Figure 8: Alpha Sweep — DAG Transition Analysis',
        'caption': 'Gradient ratio r_gradient as a function of alpha. Discontinuous jump at alpha=0+ demonstrates that knowledge amount is irrelevant; only structure matters.'
    },
    {
        'file': 'fig9_ecd_pruning_analysis.png',
        'title': 'Figure 9: ECD Ensemble + Pruning Analysis',
        'caption': '(A) Hodge gradient/curl decomposition, (B) Causal potential vs interventionability, (C) Per-edge feedback fraction, (D) DAG-ness vs knowledge quality (U-shaped curve, p_flip*≈0.15).'
    },
]


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for fig_info in FIGURES:
        fig_path = FIGURES_DIR / fig_info['file']
        if not fig_path.exists():
            print(f"Warning: {fig_path} not found, skipping")
            continue

        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)

        # Title at top
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = fig_info['title']
        p.font.size = Pt(22)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Image centered
        from PIL import Image as PILImage
        img = PILImage.open(fig_path)
        img_w, img_h = img.size
        aspect = img_w / img_h

        max_w = Inches(11)
        max_h = Inches(5.5)

        if aspect > (11 / 5.5):
            width = max_w
            height = int(width / aspect)
        else:
            height = max_h
            width = int(height * aspect)

        left = int((prs.slide_width - width) / 2)
        top = Inches(1.0)
        slide.shapes.add_picture(str(fig_path), left, top, width, height)

        # Caption at bottom
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.7))
        tf2 = cap_box.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = fig_info['caption']
        p2.font.size = Pt(12)
        p2.font.italic = True
        p2.alignment = PP_ALIGN.CENTER

    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")


if __name__ == '__main__':
    main()
