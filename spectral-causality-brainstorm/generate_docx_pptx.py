"""
Generate docx (inline figures) and pptx (editable, 1 figure per slide)
for the Spectral Causality Explainer document.
"""

import os
import re
from docx import Document
from docx.shared import Inches, Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGBColor

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
FIGURES_DIR = os.path.join(BASE_DIR, "figures")
OUTPUT_DIR = BASE_DIR

# ============================================================
# Figure metadata
# ============================================================
FIGURES = [
    {
        "file": "fig1_three_approaches.png",
        "num": 1,
        "title": "Three Approaches to Causal Inference",
        "caption": (
            "Figure 1: (A) Conventional approach — isolated variables. "
            "(B) LiNGAM — variable-level causal DAG. "
            "(C) Spectral Causality — theme-level causal flow. "
            "Vertical position represents causal potential: higher = more upstream (cause), "
            "lower = more downstream (effect)."
        ),
    },
    {
        "file": "fig2_magnetic_laplacian_q.png",
        "num": 2,
        "title": "Magnetic Laplacian Eigenvectors in the Complex Plane",
        "caption": (
            "Figure 2: Second eigenvector of the magnetic Laplacian plotted on the complex plane. "
            "At q=0, all points lie on the real axis (no directionality). "
            "At q=0.1 and q=0.25, variables spread into the complex plane, "
            "with phase angle ordering encoding causal flow direction."
        ),
    },
    {
        "file": "fig3_hodge_decomposition.png",
        "num": 3,
        "title": "Hodge Decomposition of Information Flow",
        "caption": (
            "Figure 3: (A) Hodge decomposition of information flow: "
            "85.9% gradient (causal, one-directional) vs 14.1% curl (feedback loops). "
            "(B) Causal potential for each variable. "
            "Age is most upstream; ST Depression is most downstream."
        ),
    },
    {
        "file": "fig4_direction_comparison.png",
        "num": 4,
        "title": "Causal Direction Comparison Across Three Methods",
        "caption": (
            "Figure 4: Causal direction for all 10 variable pairs comparing "
            "LiNGAM (red), Spectral Causal Direction SCD (blue), and "
            "Hodge potential (green). +1 = first variable causes second; "
            "-1 = reverse. Green background = all three methods agree."
        ),
    },
    {
        "file": "fig5_hill_radar.png",
        "num": 5,
        "title": "Hill's 9 Criteria Coverage by Method",
        "caption": (
            "Figure 5: Radar plots showing coverage of Hill's 9 causal criteria. "
            "LiNGAM excels at H1 (strength) and H3 (specificity) but lacks H6/H7/H9. "
            "Utility Causality covers H6 (plausibility), H7 (coherence), H9 (analogy). "
            "Ensemble (ECD) combines both, covering nearly all criteria."
        ),
    },
]


# ============================================================
# DOCX Generation
# ============================================================
def generate_docx():
    doc = Document()

    # -- Styles --
    style = doc.styles['Normal']
    style.font.name = 'Times New Roman'
    style.font.size = Pt(11)
    style.paragraph_format.space_after = Pt(6)

    # Title
    title_p = doc.add_paragraph()
    title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_run = title_p.add_run(
        "スペクトル因果性 — 「音の科学」で因果関係を見つける"
    )
    title_run.bold = True
    title_run.font.size = Pt(16)

    subtitle_p = doc.add_paragraph()
    subtitle_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub_run = subtitle_p.add_run(
        "Spectral Causality: Finding Causal Relations Through the Science of Sound"
    )
    sub_run.italic = True
    sub_run.font.size = Pt(12)
    sub_run.font.color.rgb = RGBColor(100, 100, 100)

    # Read markdown source
    md_path = os.path.join(BASE_DIR, "05_spectral_causality_explainer.md")
    with open(md_path, "r", encoding="utf-8") as f:
        md_text = f.read()

    # Split into sections by ## headings
    sections = re.split(r'\n(## .+)', md_text)

    # Track which figures have been inserted
    inserted_figures = set()

    def insert_figure_after_mention(paragraph_text, doc):
        """Check if paragraph mentions a figure and insert it."""
        for fig in FIGURES:
            fig_num = fig["num"]
            patterns = [
                f"図{fig_num}",
                f"Figure {fig_num}",
                f"fig{fig_num}",
                f"Fig {fig_num}",
                f"（図{fig_num}",
                f"(Fig {fig_num}",
                f"（図{fig_num}）",
            ]
            if fig_num not in inserted_figures:
                for pat in patterns:
                    if pat.lower() in paragraph_text.lower():
                        # Insert figure
                        fig_path = os.path.join(FIGURES_DIR, fig["file"])
                        if os.path.exists(fig_path):
                            # Add spacing
                            spacer = doc.add_paragraph()
                            spacer.paragraph_format.space_before = Pt(12)

                            # Add image
                            p = doc.add_paragraph()
                            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                            run = p.add_run()
                            run.add_picture(fig_path, width=Inches(5.5))

                            # Add caption
                            cap_p = doc.add_paragraph()
                            cap_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                            cap_p.paragraph_format.space_before = Pt(6)
                            cap_p.paragraph_format.space_after = Pt(12)
                            cap_run = cap_p.add_run(fig["caption"])
                            cap_run.font.size = Pt(9)
                            cap_run.italic = True

                            inserted_figures.add(fig_num)
                        return

    def add_text_block(text, doc):
        """Add a block of markdown text as docx paragraphs."""
        lines = text.strip().split('\n')
        i = 0
        while i < len(lines):
            line = lines[i].strip()

            # Skip image links (we handle figures separately)
            if line.startswith('![') or line.startswith('*図') or line.startswith('*Figure'):
                i += 1
                continue

            # Skip empty lines
            if not line:
                i += 1
                continue

            # Code blocks
            if line.startswith('```'):
                code_lines = []
                i += 1
                while i < len(lines) and not lines[i].strip().startswith('```'):
                    code_lines.append(lines[i])
                    i += 1
                i += 1  # skip closing ```
                code_p = doc.add_paragraph()
                code_p.paragraph_format.left_indent = Cm(1)
                code_run = code_p.add_run('\n'.join(code_lines))
                code_run.font.name = 'Courier New'
                code_run.font.size = Pt(9)
                continue

            # Table rows
            if line.startswith('|'):
                # Collect all table lines
                table_lines = []
                while i < len(lines) and lines[i].strip().startswith('|'):
                    tl = lines[i].strip()
                    # Skip separator rows
                    if not re.match(r'^\|[\s\-|:]+\|$', tl):
                        table_lines.append(tl)
                    i += 1

                if len(table_lines) >= 1:
                    # Parse table
                    rows_data = []
                    for tl in table_lines:
                        cells = [c.strip() for c in tl.split('|')[1:-1]]
                        rows_data.append(cells)

                    if rows_data:
                        n_cols = max(len(r) for r in rows_data)
                        table = doc.add_table(rows=len(rows_data), cols=n_cols)
                        table.style = 'Table Grid'
                        for ri, row_data in enumerate(rows_data):
                            for ci, cell_text in enumerate(row_data):
                                if ci < n_cols:
                                    cell = table.cell(ri, ci)
                                    cell.text = cell_text.replace('**', '')
                                    for paragraph in cell.paragraphs:
                                        for run in paragraph.runs:
                                            run.font.size = Pt(9)
                                    # Bold header row
                                    if ri == 0:
                                        for paragraph in cell.paragraphs:
                                            for run in paragraph.runs:
                                                run.bold = True

                        doc.add_paragraph()  # spacing after table
                continue

            # Heading ### (H3)
            if line.startswith('### '):
                h_text = line[4:].replace('**', '')
                p = doc.add_heading(h_text, level=3)
                i += 1
                continue

            # Bulleted list
            if line.startswith('- ') or line.startswith('* '):
                text_content = line[2:]
                # Clean markdown
                text_content = re.sub(r'\*\*(.+?)\*\*', r'\1', text_content)
                p = doc.add_paragraph(text_content, style='List Bullet')
                insert_figure_after_mention(text_content, doc)
                i += 1
                continue

            # Numbered list
            if re.match(r'^\d+\. ', line):
                text_content = re.sub(r'^\d+\. ', '', line)
                text_content = re.sub(r'\*\*(.+?)\*\*', r'\1', text_content)
                p = doc.add_paragraph(text_content, style='List Number')
                i += 1
                continue

            # Blockquote
            if line.startswith('> '):
                text_content = line[2:]
                text_content = re.sub(r'\*\*(.+?)\*\*', r'\1', text_content)
                p = doc.add_paragraph()
                p.paragraph_format.left_indent = Cm(1.5)
                run = p.add_run(text_content)
                run.italic = True
                i += 1
                continue

            # Regular paragraph
            text_content = line
            # Clean bold markdown
            text_content = re.sub(r'\*\*(.+?)\*\*', r'\1', text_content)
            if text_content:
                p = doc.add_paragraph(text_content)
                insert_figure_after_mention(text_content, doc)

            i += 1

    # Process sections
    for idx, section in enumerate(sections):
        section = section.strip()
        if not section:
            continue

        # Section heading
        if section.startswith('## '):
            heading_text = section[3:].strip()
            # Skip the title/purpose section marker
            if heading_text.startswith('この文書の目的') or '---' in heading_text:
                continue
            doc.add_heading(heading_text, level=2)
            continue

        # Section body
        add_text_block(section, doc)

    # Save
    output_path = os.path.join(OUTPUT_DIR, "05_spectral_causality_explainer.docx")
    doc.save(output_path)
    print(f"Saved: {output_path}")
    return output_path


# ============================================================
# PPTX Generation (English, editable, 1 figure per slide)
# ============================================================
def generate_pptx():
    prs = Presentation()
    # Widescreen 16:9
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    # Title slide
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    from pptx.util import Emu

    # Title text box
    txBox = slide.shapes.add_textbox(
        PptxInches(1), PptxInches(2), PptxInches(11.333), PptxInches(2)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Spectral Causality"
    p.font.size = PptxPt(40)
    p.font.bold = True
    p.font.color.rgb = PptxRGBColor(0, 51, 102)
    p.alignment = 1  # center

    p2 = tf.add_paragraph()
    p2.text = "Finding Causal Relations Through the Science of Sound"
    p2.font.size = PptxPt(24)
    p2.font.color.rgb = PptxRGBColor(100, 100, 100)
    p2.alignment = 1

    p3 = tf.add_paragraph()
    p3.text = "UCI Heart Disease Data (n=297) — Real-World Analysis"
    p3.font.size = PptxPt(18)
    p3.font.color.rgb = PptxRGBColor(120, 120, 120)
    p3.alignment = 1

    # Figure slides
    for fig in FIGURES:
        fig_path = os.path.join(FIGURES_DIR, fig["file"])
        if not os.path.exists(fig_path):
            print(f"  WARNING: {fig_path} not found, skipping")
            continue

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

        # Title
        title_box = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(0.3), PptxInches(12.333), PptxInches(0.8)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"Figure {fig['num']}: {fig['title']}"
        p.font.size = PptxPt(24)
        p.font.bold = True
        p.font.color.rgb = PptxRGBColor(0, 51, 102)

        # Image (centered, scaled to fit)
        from PIL import Image
        img = Image.open(fig_path)
        img_w, img_h = img.size
        aspect = img_w / img_h

        max_w = PptxInches(11.5)
        max_h = PptxInches(4.8)

        if aspect > (11.5 / 4.8):
            width = max_w
            height = int(width / aspect)
        else:
            height = max_h
            width = int(height * aspect)

        left = (prs.slide_width - width) // 2
        top = PptxInches(1.3)

        slide.shapes.add_picture(fig_path, left, top, width, height)

        # Caption
        cap_box = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(6.3), PptxInches(12.333), PptxInches(1.0)
        )
        tf = cap_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = fig["caption"]
        p.font.size = PptxPt(12)
        p.font.color.rgb = PptxRGBColor(80, 80, 80)
        p.font.italic = True

    # Save
    output_path = os.path.join(OUTPUT_DIR, "05_spectral_causality_figures.pptx")
    prs.save(output_path)
    print(f"Saved: {output_path}")
    return output_path


if __name__ == "__main__":
    print("Generating docx...")
    generate_docx()
    print("\nGenerating pptx...")
    generate_pptx()
    print("\nDone!")
