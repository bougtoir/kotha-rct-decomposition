"""Generate English editable pptx with all 8 figures (one per slide)."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
FIG_DIR = os.path.join(BASE_DIR, 'figures')

FIGURES = [
    {
        'file': 'fig6_causal_dag.png',
        'title': 'Figure 1: Estimated Causal DAG (DirectLiNGAM)',
        'caption': 'Causal DAG estimated by DirectLiNGAM on the UCI Heart Disease dataset '
                   '(Cleveland subset, n = 297). Five clinical variables. Blue = positive effect; '
                   'red = negative effect. Causal order: Age > MaxHR > STDep > RestBP > Chol.',
    },
    {
        'file': 'fig5_hill_radar.png',
        'title': "Figure 2: Hill's 9 Criteria Coverage",
        'caption': "Radar chart of each method's coverage of Hill's nine criteria. "
                   'LiNGAM: H1/H3. Utility Causality: H6/H7/H9. ECD Ensemble: nearly all criteria.',
    },
    {
        'file': 'fig2_magnetic_laplacian_q.png',
        'title': 'Figure 3: Magnetic Laplacian Eigenvectors',
        'caption': 'Second eigenvector of the magnetic Laplacian on the complex plane. '
                   'q=0: all points on real axis. q=0.25: phase angle ordering encodes causal flow.',
    },
    {
        'file': 'fig3_hodge_decomposition.png',
        'title': 'Figure 4: Hodge Decomposition of Causal Flows',
        'caption': 'Orthogonal decomposition into gradient (DAG-like, 85.9%), '
                   'curl (feedback, 14.1%), and harmonic components.',
    },
    {
        'file': 'fig4_direction_comparison.png',
        'title': 'Figure 5: Three-Method Direction Comparison',
        'caption': 'Causal directions: LiNGAM vs. spectral causality (alpha=0.6) vs. DPI (alpha=0). '
                   '5 variables, 10 possible edges. 67% agreement between DPI and LiNGAM.',
    },
    {
        'file': 'fig7_lingam_vs_spectral.png',
        'title': 'Figure 6: LiNGAM DAG vs. Spectral Causality DCG',
        'caption': 'LiNGAM enforces DAG (acyclic). Spectral causality permits feedback loops. '
                   'Edge feedback rates quantified by Hodge decomposition.',
    },
    {
        'file': 'fig9_ecd_pruning_analysis.png',
        'title': 'Figure 7: ECD Ensemble & Pruning Analysis',
        'caption': '(A) Hodge decomposition. (B) Causal potential vs. interventionability. '
                   '(C) Edge-level feedback rates. (D) p_flip U-shaped curve.',
    },
    {
        'file': 'fig8_alpha_sweep.png',
        'title': 'Figure 8: Alpha-Sweep DAG Transition Analysis',
        'caption': '(A) r_gradient: 0.581 (alpha=0) to 0.859 (alpha=1.0), smooth transition with DPI. '
                   '(B) Edge count & LiNGAM agreement. (C) Asymmetric norm. (D) Phase diagram.',
    },
]

def generate_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank

    for fig in FIGURES:
        slide = prs.slides.add_slide(blank_layout)
        img_path = os.path.join(FIG_DIR, fig['file'])

        # Title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.6))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = fig['title']
        p.font.size = Pt(20)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Image (centered)
        if os.path.exists(img_path):
            from PIL import Image
            with Image.open(img_path) as img:
                w, h = img.size
            max_w = Inches(11)
            max_h = Inches(5.2)
            scale = min(max_w / Emu(int(w * 914400 / 96)),
                       max_h / Emu(int(h * 914400 / 96)))
            img_w = Emu(int(w * 914400 / 96 * scale))
            img_h = Emu(int(h * 914400 / 96 * scale))
            left = (prs.slide_width - img_w) // 2
            top = Inches(1.0)
            slide.shapes.add_picture(img_path, left, top, img_w, img_h)
        else:
            txBox2 = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(1))
            txBox2.text_frame.paragraphs[0].text = f'[Image not found: {fig["file"]}]'

        # Caption
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.8))
        tf2 = cap_box.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = fig['caption']
        p2.font.size = Pt(12)
        p2.font.italic = True
        p2.font.color.rgb = RGBColor(80, 80, 80)
        p2.alignment = PP_ALIGN.CENTER

    out_path = os.path.join(BASE_DIR, 'spectral_causality_figures_en.pptx')
    prs.save(out_path)
    print(f'Saved: {out_path}')
    return out_path


if __name__ == '__main__':
    generate_pptx()
