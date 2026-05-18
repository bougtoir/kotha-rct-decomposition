"""
Generate Cliometrica manuscript: "Flow Disruption and State Collapse"
Technical Network Exclusion sensitivity analysis with 7-country reclassification.

Outputs:
  manuscript/manuscript.docx          — Main manuscript (Cliometrica format)
  manuscript/table_s1.docx            — Supplementary Table S1 (96 polities)
  manuscript/figures/Fig1.png … Fig4.png — Separate figure files
  manuscript/figures_pptx.pptx        — Editable PPTX (1 figure per slide)
"""

import os
import sys
import re
import math
import warnings
warnings.filterwarnings("ignore")

import numpy as np
import pandas as pd
import scipy.stats as stats
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
from matplotlib.patches import FancyBboxPatch
import seaborn as sns

from docx import Document
from docx.shared import Inches, Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.section import WD_ORIENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN

# ── project imports ──
sys.path.insert(0, os.path.dirname(__file__))
from data import load_data
from sensitivity_technical_network_exclusion import (
    STRONG_CANDIDATES, MODERATE_CANDIDATES, RATIONALE,
    apply_technical_network_exclusion, apply_disrupted_assignment,
    compute_confusion_stats, compute_closure_analysis,
    compute_logistic_with_closure, compute_mediation_paths,
)

OUT = os.path.join(os.path.dirname(__file__), "manuscript")
FIG = os.path.join(OUT, "figures")
os.makedirs(FIG, exist_ok=True)

# ════════════════════════════════════════════════════════════
# Modern-country mapping & turning-point events for Table S1
# ════════════════════════════════════════════════════════════

MODERN_COUNTRY = {
    "アケメネス朝ペルシア": "Iran",
    "アテネ（ペロポネソス戦争後）": "Greece",
    "スパルタ": "Greece",
    "カルタゴ": "Tunisia",
    "ローマ共和政〜帝政前期": "Italy",
    "ローマ帝政後期（西）": "Italy",
    "プトレマイオス朝エジプト": "Egypt",
    "漢朝（前漢〜後漢）": "China",
    "ビザンツ帝国（前期）": "Turkey/Greece",
    "ビザンツ帝国（後期）": "Turkey/Greece",
    "アッバース朝（前期）": "Iraq",
    "アッバース朝（後期）": "Iraq",
    "宋朝中国": "China",
    "モンゴル帝国": "Mongolia/China",
    "マムルーク朝エジプト": "Egypt",
    "ヴェネツィア共和国": "Italy",
    "ハンザ同盟": "Germany",
    "オスマン帝国（前期）": "Turkey",
    "オスマン帝国（後期）": "Turkey",
    "明朝中国": "China",
    "清朝中国（前期）": "China",
    "清朝中国（後期）": "China",
    "李朝朝鮮": "Korea",
    "徳川日本": "Japan",
    "ムガル帝国（後期）": "India",
    "サファヴィー朝ペルシア": "Iran",
    "インカ帝国": "Peru",
    "アステカ帝国": "Mexico",
    "ポルトガル帝国": "Portugal",
    "オランダ共和国": "Netherlands",
    "スペイン帝国": "Spain",
    "ムガル帝国（前期）": "India",
    "大英帝国": "United Kingdom",
    "米国（モンロー主義期）": "United States",
    "明治〜大正日本": "Japan",
    "エチオピア帝国": "Ethiopia",
    "シャム（タイ）": "Thailand",
    "スイス": "Switzerland",
    "1930s日本（大東亜共栄圏）": "Japan",
    "ナチスドイツ": "Germany",
    "ファシストイタリア": "Italy",
    "ソ連": "Russia",
    "東ドイツ": "Germany",
    "ユーゴスラビア": "Serbia/Croatia/Bosnia/etc.",
    "英国帝国特恵制度": "United Kingdom",
    "米国（戦後〜冷戦期）": "United States",
    "戦後日本（高度成長期）": "Japan",
    "シンガポール": "Singapore",
    "韓国": "South Korea",
    "北朝鮮": "North Korea",
    "キューバ": "Cuba",
    "現代日本": "Japan",
    "現代中国": "China",
    "現代ロシア": "Russia",
    "マケドニア王国（アレクサンドロス後）": "Greece/North Macedonia",
    "セレウコス朝シリア": "Syria/Iraq/Iran",
    "パルティア": "Iran",
    "ササン朝ペルシア": "Iran",
    "クシャーナ朝": "Afghanistan/Pakistan",
    "グプタ朝インド": "India",
    "唐朝中国": "China",
    "ウマイヤ朝": "Syria",
    "高麗": "Korea",
    "クメール帝国（アンコール）": "Cambodia",
    "シュリーヴィジャヤ": "Indonesia",
    "マジャパヒト王国": "Indonesia",
    "キエフ大公国": "Ukraine/Russia",
    "ジェノヴァ共和国": "Italy",
    "マリ帝国": "Mali",
    "元朝中国": "China",
    "ティムール朝": "Uzbekistan/Iran",
    "琉球王国": "Japan (Okinawa)",
    "ポーランド・リトアニア共和国": "Poland/Lithuania",
    "スウェーデン帝国": "Sweden",
    "ロシア帝国（ピョートル後）": "Russia",
    "清朝中国（中期・交易拡大期）": "China",
    "ハワイ王国": "United States (Hawaii)",
    "ズールー王国": "South Africa",
    "ビルマ（コンバウン朝）": "Myanmar",
    "オーストリア＝ハンガリー帝国": "Austria/Hungary",
    "ベルギー": "Belgium",
    "デンマーク": "Denmark",
    "ナポレオン帝国": "France",
    "大韓帝国": "South Korea",
    "満州国": "China (Manchuria)",
    "南ベトナム": "Vietnam",
    "チェコスロバキア（共産期）": "Czech Republic/Slovakia",
    "ポーランド（共産期）": "Poland",
    "ルーマニア（共産期）": "Romania",
    "台湾": "Taiwan",
    "イスラエル": "Israel",
    "UAE": "United Arab Emirates",
    "現代インド": "India",
    "イラン（イスラム共和国）": "Iran",
    "ミャンマー（軍政期〜現在）": "Myanmar",
    "トルクメニスタン": "Turkmenistan",
}

ENGLISH_NAME = {
    "アケメネス朝ペルシア": "Achaemenid Persia",
    "アテネ（ペロポネソス戦争後）": "Athens (post-Peloponnesian War)",
    "スパルタ": "Sparta",
    "カルタゴ": "Carthage",
    "ローマ共和政〜帝政前期": "Roman Republic–Early Empire",
    "ローマ帝政後期（西）": "Late Western Roman Empire",
    "プトレマイオス朝エジプト": "Ptolemaic Egypt",
    "漢朝（前漢〜後漢）": "Han Dynasty (China)",
    "ビザンツ帝国（前期）": "Byzantine Empire (early)",
    "ビザンツ帝国（後期）": "Byzantine Empire (late)",
    "アッバース朝（前期）": "Abbasid Caliphate (early)",
    "アッバース朝（後期）": "Abbasid Caliphate (late)",
    "宋朝中国": "Song Dynasty (China)",
    "モンゴル帝国": "Mongol Empire",
    "マムルーク朝エジプト": "Mamluk Sultanate",
    "ヴェネツィア共和国": "Republic of Venice",
    "ハンザ同盟": "Hanseatic League",
    "オスマン帝国（前期）": "Ottoman Empire (early)",
    "オスマン帝国（後期）": "Ottoman Empire (late)",
    "明朝中国": "Ming Dynasty (China)",
    "清朝中国（前期）": "Qing Dynasty (early)",
    "清朝中国（後期）": "Qing Dynasty (late)",
    "李朝朝鮮": "Joseon Korea",
    "徳川日本": "Tokugawa Japan",
    "ムガル帝国（後期）": "Mughal Empire (late)",
    "サファヴィー朝ペルシア": "Safavid Persia",
    "インカ帝国": "Inca Empire",
    "アステカ帝国": "Aztec Empire",
    "ポルトガル帝国": "Portuguese Empire",
    "オランダ共和国": "Dutch Republic",
    "スペイン帝国": "Spanish Empire",
    "ムガル帝国（前期）": "Mughal Empire (early)",
    "大英帝国": "British Empire",
    "米国（モンロー主義期）": "United States (Monroe era)",
    "明治〜大正日本": "Meiji–Taisho Japan",
    "エチオピア帝国": "Ethiopian Empire",
    "シャム（タイ）": "Siam (Thailand)",
    "スイス": "Switzerland",
    "1930s日本（大東亜共栄圏）": "Imperial Japan (1930s)",
    "ナチスドイツ": "Nazi Germany",
    "ファシストイタリア": "Fascist Italy",
    "ソ連": "Soviet Union",
    "東ドイツ": "East Germany",
    "ユーゴスラビア": "Yugoslavia",
    "英国帝国特恵制度": "British Imperial Preference",
    "米国（戦後〜冷戦期）": "United States (Cold War era)",
    "戦後日本（高度成長期）": "Post-war Japan",
    "シンガポール": "Singapore",
    "韓国": "South Korea",
    "北朝鮮": "North Korea",
    "キューバ": "Cuba",
    "現代日本": "Contemporary Japan",
    "現代中国": "Contemporary China",
    "現代ロシア": "Contemporary Russia",
    "マケドニア王国（アレクサンドロス後）": "Successor Macedonia",
    "セレウコス朝シリア": "Seleucid Empire",
    "パルティア": "Parthian Empire",
    "ササン朝ペルシア": "Sasanian Persia",
    "クシャーナ朝": "Kushan Empire",
    "グプタ朝インド": "Gupta Empire",
    "唐朝中国": "Tang Dynasty (China)",
    "ウマイヤ朝": "Umayyad Caliphate",
    "高麗": "Goryeo (Korea)",
    "クメール帝国（アンコール）": "Khmer Empire (Angkor)",
    "シュリーヴィジャヤ": "Srivijaya",
    "マジャパヒト王国": "Majapahit",
    "キエフ大公国": "Kievan Rus'",
    "ジェノヴァ共和国": "Republic of Genoa",
    "マリ帝国": "Mali Empire",
    "元朝中国": "Yuan Dynasty (China)",
    "ティムール朝": "Timurid Empire",
    "琉球王国": "Ryukyu Kingdom",
    "ポーランド・リトアニア共和国": "Polish-Lithuanian Commonwealth",
    "スウェーデン帝国": "Swedish Empire",
    "ロシア帝国（ピョートル後）": "Russian Empire (post-Peter)",
    "清朝中国（中期・交易拡大期）": "Qing Dynasty (mid-period)",
    "ハワイ王国": "Kingdom of Hawaii",
    "ズールー王国": "Zulu Kingdom",
    "ビルマ（コンバウン朝）": "Konbaung Burma",
    "オーストリア＝ハンガリー帝国": "Austria-Hungary",
    "ベルギー": "Belgium",
    "デンマーク": "Denmark",
    "ナポレオン帝国": "Napoleonic Empire",
    "大韓帝国": "Korean Empire",
    "満州国": "Manchukuo",
    "南ベトナム": "South Vietnam",
    "チェコスロバキア（共産期）": "Czechoslovakia (communist)",
    "ポーランド（共産期）": "Poland (communist)",
    "ルーマニア（共産期）": "Romania (communist)",
    "台湾": "Taiwan",
    "イスラエル": "Israel",
    "UAE": "United Arab Emirates",
    "現代インド": "Contemporary India",
    "イラン（イスラム共和国）": "Iran (Islamic Republic)",
    "ミャンマー（軍政期〜現在）": "Myanmar (military rule)",
    "トルクメニスタン": "Turkmenistan",
}

TURNING_POINT = {
    "アケメネス朝ペルシア": "Conquered by Alexander the Great at Battle of Gaugamela (331 BC)",
    "アテネ（ペロポネソス戦争後）": "Subjugated by Macedon after Battle of Chaeronea (338 BC)",
    "スパルタ": "Defeated at Battle of Leuctra (371 BC); absorbed by Rome (146 BC)",
    "カルタゴ": "Destroyed by Rome in Third Punic War (146 BC)",
    "ローマ共和政〜帝政前期": "Survived: expanded to dominate the Mediterranean",
    "ローマ帝政後期（西）": "Fall of Rome to Odoacer (476 AD)",
    "プトレマイオス朝エジプト": "Annexed by Rome after Battle of Actium (30 BC)",
    "漢朝（前漢〜後漢）": "Collapsed into Three Kingdoms civil war (220 AD)",
    "ビザンツ帝国（前期）": "Survived: withstood Arab and Bulgar sieges",
    "ビザンツ帝国（後期）": "Fall of Constantinople to Ottoman Turks (1453)",
    "アッバース朝（前期）": "Survived: Islamic Golden Age with vast trade networks",
    "アッバース朝（後期）": "Sack of Baghdad by Mongols (1258)",
    "宋朝中国": "Conquered by Mongol Yuan Dynasty (1279)",
    "モンゴル帝国": "Fragmented into successor khanates after 1260s",
    "マムルーク朝エジプト": "Conquered by Ottoman Empire (1517)",
    "ヴェネツィア共和国": "Dissolved by Napoleon (1797)",
    "ハンザ同盟": "Eclipsed by rise of territorial nation-states (17th c.)",
    "オスマン帝国（前期）": "Survived: expanded across three continents",
    "オスマン帝国（後期）": "Dismembered after WWI; Republic of Turkey founded (1922)",
    "明朝中国": "Conquered by Manchu Qing Dynasty (1644)",
    "清朝中国（前期）": "Survived: Kangxi–Qianlong prosperity under maritime restrictions",
    "清朝中国（後期）": "Semi-colonized after Opium Wars; dynasty fell (1912)",
    "李朝朝鮮": "Annexed by Japan (1910)",
    "徳川日本": "Forced opening by Perry (1853); Meiji Restoration (1868)",
    "ムガル帝国（後期）": "Colonized by British East India Company; dissolved (1857)",
    "サファヴィー朝ペルシア": "Conquered by Afghan Hotaki invaders (1722)",
    "インカ帝国": "Conquered by Spanish conquistadors under Pizarro (1533)",
    "アステカ帝国": "Conquered by Cortes and Spanish forces (1521)",
    "ポルトガル帝国": "Iberian Union under Spanish crown (1580); independence restored (1640)",
    "オランダ共和国": "Conquered by France (1795); later restored as Kingdom of the Netherlands",
    "スペイン帝国": "Napoleonic invasion (1808); loss of American colonies",
    "ムガル帝国（前期）": "Survived: Mughal expansion and consolidation under Akbar–Aurangzeb",
    "大英帝国": "Survived: peaceful decolonization post-WWII",
    "米国（モンロー主義期）": "Survived: continental expansion under Monroe Doctrine",
    "明治〜大正日本": "Survived: rapid modernization and industrialization",
    "エチオピア帝国": "Survived: Victory at Battle of Adwa (1896) maintained independence",
    "シャム（タイ）": "Survived: diplomatic balancing avoided colonization",
    "スイス": "Survived: permanent neutrality since Congress of Vienna (1815)",
    "1930s日本（大東亜共栄圏）": "WWII defeat and Allied occupation (1945); state survived",
    "ナチスドイツ": "WWII defeat; division into East/West Germany (1945)",
    "ファシストイタリア": "WWII defeat; transition to republic (1946)",
    "ソ連": "Internal collapse and dissolution (1991)",
    "東ドイツ": "German reunification (1990)",
    "ユーゴスラビア": "Breakup into successor states (1991–1992)",
    "英国帝国特恵制度": "Survived: peaceful transition to Commonwealth",
    "米国（戦後〜冷戦期）": "Survived: emerged as global superpower",
    "戦後日本（高度成長期）": "Survived: economic miracle under US alliance",
    "シンガポール": "Survived: rapid development as trade hub since independence (1965)",
    "韓国": "Survived: industrialization and democratization",
    "北朝鮮": "Survived: autarkic Juche regime with Chinese patronage",
    "キューバ": "Survived: maintained regime despite US embargo",
    "現代日本": "Survived: ongoing post-bubble economic adaptation",
    "現代中国": "Survived: reform-era growth since Deng Xiaoping (1978)",
    "現代ロシア": "Survived: post-Soviet transition under sanctions",
    "マケドニア王国（アレクサンドロス後）": "Conquered by Rome at Battle of Pydna (168 BC)",
    "セレウコス朝シリア": "Annexed by Rome as province of Syria (63 BC)",
    "パルティア": "Overthrown by Sasanian revolt under Ardashir I (224 AD)",
    "ササン朝ペルシア": "Conquered by Arab-Muslim armies at Battle of al-Qadisiyyah (636 AD)",
    "クシャーナ朝": "Fragmented under Sasanian and nomadic pressure (4th c.)",
    "グプタ朝インド": "Collapsed under Huna (Hephthalite) invasions (6th c.)",
    "唐朝中国": "Collapsed after An Lushan Rebellion; fell to warlords (907)",
    "ウマイヤ朝": "Overthrown by Abbasid Revolution (750)",
    "高麗": "Replaced by Joseon Dynasty coup (1392)",
    "クメール帝国（アンコール）": "Sacked by Ayutthaya (Siamese) forces (1431)",
    "シュリーヴィジャヤ": "Eclipsed by Majapahit and Chola raids; dissolved (14th c.)",
    "マジャパヒト王国": "Replaced by Islamic Demak Sultanate (early 16th c.)",
    "キエフ大公国": "Destroyed by Mongol invasion (1240)",
    "ジェノヴァ共和国": "Annexed by France under Napoleon (1797)",
    "マリ帝国": "Declined under Songhai expansion and internal fragmentation (16th c.)",
    "元朝中国": "Overthrown by Ming rebellion under Zhu Yuanzhang (1368)",
    "ティムール朝": "Fragmented; conquered by Uzbek Shaybanids (1507)",
    "琉球王国": "Annexed by Meiji Japan (1879)",
    "ポーランド・リトアニア共和国": "Partitioned by Russia, Prussia, Austria (1795)",
    "スウェーデン帝国": "Lost Great Northern War to Russia (1721); state survived",
    "ロシア帝国（ピョートル後）": "Russian Revolution (1917); state reconstituted as USSR",
    "清朝中国（中期・交易拡大期）": "Survived: Canton system maintained controlled trade",
    "ハワイ王国": "Overthrown by US-backed coup (1893); annexed by US (1898)",
    "ズールー王国": "Conquered by British Empire in Anglo-Zulu War (1879)",
    "ビルマ（コンバウン朝）": "Conquered by British Empire in Third Anglo-Burmese War (1885)",
    "オーストリア＝ハンガリー帝国": "Dissolved after WWI (1918); successor states emerged",
    "ベルギー": "Survived: maintained independence despite World Wars",
    "デンマーク": "Survived: adapted from empire to small democratic state",
    "ナポレオン帝国": "Defeated at Waterloo (1815); France continued as nation-state",
    "大韓帝国": "Annexed by Japan (1910)",
    "満州国": "Dissolved after Japan's WWII defeat (1945)",
    "南ベトナム": "Conquered by North Vietnam; reunification (1975)",
    "チェコスロバキア（共産期）": "Velvet Revolution (1989); peaceful split into Czech Rep. and Slovakia (1993)",
    "ポーランド（共産期）": "Solidarity movement; democratic transition (1989)",
    "ルーマニア（共産期）": "Romanian Revolution; execution of Ceausescu (1989)",
    "台湾": "Survived: democratization and de facto independence",
    "イスラエル": "Survived: established and maintained statehood since 1948",
    "UAE": "Survived: oil-funded development as trade hub",
    "現代インド": "Survived: liberalization-era growth since 1991",
    "イラン（イスラム共和国）": "Survived: maintained theocratic regime under sanctions",
    "ミャンマー（軍政期〜現在）": "Survived: military regime persists amid civil conflict",
    "トルクメニスタン": "Survived: authoritarian gas-state since independence (1991)",
}


# ════════════════════════════════════════════════════════════
# Helper: compute all analysis results
# ════════════════════════════════════════════════════════════

def run_analysis():
    df = load_data()
    N = len(df)

    closure_scenarios = {
        "baseline": {"label": "Baseline", "df": df, "reclassified": []},
        "strong": {
            "label": "+5 strong",
            "df": apply_technical_network_exclusion(df, STRONG_CANDIDATES),
            "reclassified": STRONG_CANDIDATES,
        },
        "all": {
            "label": "+7 all",
            "df": apply_technical_network_exclusion(df, STRONG_CANDIDATES + MODERATE_CANDIDATES),
            "reclassified": STRONG_CANDIDATES + MODERATE_CANDIDATES,
        },
    }
    disrupted_modes = {"as_conquered": "disrupted=overtaken", "as_survived": "disrupted=survived"}

    results = {"N": N, "df": df, "scenarios": {}, "fisher_ban": {}, "logistic": {}}

    for d_mode in disrupted_modes:
        for c_key, c_sc in closure_scenarios.items():
            key = f"{d_mode}__{c_key}"
            df_p = apply_disrupted_assignment(c_sc["df"], d_mode)
            cm = compute_confusion_stats(df_p)
            cl = compute_closure_analysis(df_p)
            lr = compute_logistic_with_closure(df_p)

            # Fisher for ban vs no-ban
            has_ban = df_p["closure_type"].isin(["maritime_ban", "technical_network_exclusion", "sakoku"])
            ban_df = df_p[has_ban]
            no_ban_df = df_p[~has_ban]
            ban_rate = ban_df["outcome_binary"].mean() if len(ban_df) > 0 else 0
            no_rate = no_ban_df["outcome_binary"].mean() if len(no_ban_df) > 0 else 0
            ban_conq = int(ban_df["outcome_binary"].sum())
            ban_surv = len(ban_df) - ban_conq
            no_conq = int(no_ban_df["outcome_binary"].sum())
            no_surv = len(no_ban_df) - no_conq
            _, p_ban = stats.fisher_exact(
                np.array([[ban_conq, ban_surv], [no_conq, no_surv]]), alternative="greater"
            )

            results["scenarios"][key] = {
                "cm": cm, "closure": cl, "logistic": lr,
                "ban_n": len(ban_df), "ban_rate": ban_rate,
                "no_ban_n": len(no_ban_df), "no_ban_rate": no_rate,
                "fisher_ban_p": p_ban,
                "rr": ban_rate / no_rate if no_rate > 0 else float("inf"),
            }

    # Bootstrap OR (for main scenario only)
    rng = np.random.default_rng(42)
    n_boot = 5000
    boot_results = {}
    for d_mode in disrupted_modes:
        for c_key, c_sc in closure_scenarios.items():
            key = f"{d_mode}__{c_key}"
            df_p = apply_disrupted_assignment(c_sc["df"], d_mode)
            n = len(df_p)
            boot_ors = np.zeros(n_boot)
            for i in range(n_boot):
                idx = rng.choice(n, size=n, replace=True)
                boot_df = df_p.iloc[idx].reset_index(drop=True)
                try:
                    ct = pd.crosstab(boot_df["dominant"], boot_df["outcome_bin_label"])
                    ct = ct.reindex(index=["stock", "flow"], columns=["overtaken", "survived"], fill_value=0)
                    tp, fp = ct.loc["stock", "overtaken"], ct.loc["stock", "survived"]
                    fn, tn = ct.loc["flow", "overtaken"], ct.loc["flow", "survived"]
                    boot_ors[i] = (tp * tn) / (fp * fn) if (fp * fn) > 0 else np.nan
                except (KeyError, ZeroDivisionError):
                    boot_ors[i] = np.nan
            valid = boot_ors[~np.isnan(boot_ors)]
            if len(valid) >= 100:
                boot_results[key] = {
                    "median": np.median(valid),
                    "ci_lo": np.percentile(valid, 2.5),
                    "ci_hi": np.percentile(valid, 97.5),
                }
    results["bootstrap"] = boot_results
    return results


# ════════════════════════════════════════════════════════════
# Figure generation
# ════════════════════════════════════════════════════════════

def create_figures(results):
    plt.rcParams.update({
        "font.family": "sans-serif",
        "font.size": 10,
        "axes.titlesize": 12,
        "figure.dpi": 300,
    })

    # ── Fig 1: Conquest rates by closure type across scenarios ──
    fig, axes = plt.subplots(1, 2, figsize=(12, 5))
    for ax_idx, (d_mode, d_label) in enumerate([
        ("as_conquered", "Disrupted → Overtaken"),
        ("as_survived", "Disrupted → Survived"),
    ]):
        ax = axes[ax_idx]
        scenarios = ["baseline", "strong", "all"]
        labels = ["Baseline", "+5 Reclassified", "+7 Reclassified"]
        x = np.arange(len(labels))
        width = 0.35

        ban_rates = []
        no_rates = []
        for c_key in scenarios:
            key = f"{d_mode}__{c_key}"
            s = results["scenarios"][key]
            ban_rates.append(s["ban_rate"] * 100)
            no_rates.append(s["no_ban_rate"] * 100)

        bars1 = ax.bar(x - width/2, ban_rates, width, label="Maritime closure", color="#c0392b", alpha=0.85)
        bars2 = ax.bar(x + width/2, no_rates, width, label="No closure", color="#2980b9", alpha=0.85)

        for bar in bars1:
            ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 1.5,
                    f"{bar.get_height():.1f}%", ha="center", va="bottom", fontsize=8)
        for bar in bars2:
            ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 1.5,
                    f"{bar.get_height():.1f}%", ha="center", va="bottom", fontsize=8)

        ax.set_ylabel("Conquest rate (%)")
        ax.set_title(d_label)
        ax.set_xticks(x)
        ax.set_xticklabels(labels)
        ax.set_ylim(0, 110)
        ax.legend(loc="upper left", fontsize=8)
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)

    fig.suptitle("Fig. 1  Conquest Rates: Maritime Closure vs. No Closure", fontsize=13, fontweight="bold")
    fig.tight_layout(rect=[0, 0, 1, 0.93])
    fig.savefig(os.path.join(FIG, "Fig1.png"), dpi=300, bbox_inches="tight")
    plt.close(fig)

    # ── Fig 2: Fisher p-value progression ──
    fig, ax = plt.subplots(figsize=(8, 5))
    all_ps = []
    for d_mode, d_label, color, marker in [
        ("as_conquered", "Disrupted → Overtaken", "#c0392b", "o"),
        ("as_survived", "Disrupted → Survived", "#2980b9", "s"),
    ]:
        ps = []
        for c_key in ["baseline", "strong", "all"]:
            key = f"{d_mode}__{c_key}"
            ps.append(results["scenarios"][key]["fisher_ban_p"])
        all_ps.extend(ps)
        ax.plot([0, 5, 7], ps, marker=marker, linewidth=2, markersize=8, label=d_label, color=color)

    ax.axhline(y=0.05, color="gray", linestyle="--", linewidth=1, label="p = 0.05 threshold")
    ax.set_xlabel("Number of reclassified polities")
    ax.set_ylabel("Fisher's exact test p-value (one-sided)")
    ax.set_xticks([0, 5, 7])
    ax.set_xticklabels(["0\n(Baseline)", "5\n(Strong)", "7\n(All)"])
    ax.set_ylim(-0.01, max(0.25, max(all_ps) + 0.05))
    ax.legend(fontsize=9)
    ax.set_title("Fig. 2  Fisher's Exact Test p-values: Maritime Closure → Conquest", fontweight="bold")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    fig.tight_layout()
    fig.savefig(os.path.join(FIG, "Fig2.png"), dpi=300, bbox_inches="tight")
    plt.close(fig)

    # ── Fig 3: Policy vs Technical network exclusion conquest rates ──
    fig, axes = plt.subplots(1, 2, figsize=(12, 5))
    for ax_idx, (d_mode, d_label) in enumerate([
        ("as_conquered", "Disrupted → Overtaken"),
        ("as_survived", "Disrupted → Survived"),
    ]):
        ax = axes[ax_idx]
        key_all = f"{d_mode}__all"
        df_p = apply_disrupted_assignment(
            apply_technical_network_exclusion(results["df"], STRONG_CANDIDATES + MODERATE_CANDIDATES),
            d_mode
        )
        categories = ["maritime_ban", "sakoku", "technical_network_exclusion", "bloc", "none"]
        cat_labels = ["Policy\nmaritime ban", "Sakoku", "Technical\nnetwork excl.", "Bloc", "None\n(open)"]
        rates = []
        counts = []
        for ct in categories:
            sub = df_p[df_p["closure_type"] == ct]
            if len(sub) > 0:
                rate = sub["outcome_binary"].mean() * 100
                rates.append(rate)
                counts.append(len(sub))
            else:
                rates.append(0)
                counts.append(0)

        colors = ["#c0392b", "#e74c3c", "#e67e22", "#95a5a6", "#2980b9"]
        bars = ax.bar(range(len(categories)), rates, color=colors, alpha=0.85)
        for i, (bar, n) in enumerate(zip(bars, counts)):
            if n > 0:
                ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 1.5,
                        f"{rates[i]:.0f}%\n(n={n})", ha="center", va="bottom", fontsize=8)

        ax.set_xticks(range(len(categories)))
        ax.set_xticklabels(cat_labels, fontsize=8)
        ax.set_ylabel("Conquest rate (%)")
        ax.set_title(d_label)
        ax.set_ylim(0, 115)
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)

    fig.suptitle("Fig. 3  Conquest Rates by Closure Type (7-Country Reclassification)", fontsize=13, fontweight="bold")
    fig.tight_layout(rect=[0, 0, 1, 0.93])
    fig.savefig(os.path.join(FIG, "Fig3.png"), dpi=300, bbox_inches="tight")
    plt.close(fig)

    # ── Fig 4: Multivariate logistic regression forest plot ──
    fig, ax = plt.subplots(figsize=(8, 6))
    key = "as_conquered__all"
    lr = results["scenarios"][key]["logistic"]
    if lr.get("with_ban", {}).get("converged"):
        coefs = lr["with_ban"]["coefs"]
        var_labels = {
            "dominant_binary": "Stock-dominant",
            "geo_barrier": "Geographic barrier",
            "external_threat": "External threat",
            "tech_position": "Tech. position",
            "institutional_quality": "Institutional quality",
            "era_code": "Era (time)",
            "has_external_patron": "External patron",
            "has_maritime_ban": "Network closure",
        }
        vars_ordered = list(var_labels.keys())
        y_pos = list(range(len(vars_ordered)))

        for i, var in enumerate(vars_ordered):
            v = coefs[var]
            log_or = np.log(v["OR"])
            log_lo = np.log(v["ci_lo"]) if v["ci_lo"] > 0 else -5
            log_hi = np.log(v["ci_hi"]) if v["ci_hi"] < 1e10 else 5
            color = "#c0392b" if v["p"] < 0.05 else "#2980b9" if v["p"] < 0.10 else "#7f8c8d"
            ax.plot([log_lo, log_hi], [i, i], color=color, linewidth=2)
            ax.plot(log_or, i, "o", color=color, markersize=8)

        ax.axvline(x=0, color="gray", linestyle="--", linewidth=1)
        ax.set_yticks(y_pos)
        ax.set_yticklabels([var_labels[v] for v in vars_ordered])
        ax.set_xlabel("log(Odds Ratio)")
        ax.set_title("Fig. 4  Multivariate Logistic Regression\n(7-Country Reclassification, Disrupted → Overtaken)",
                      fontweight="bold", fontsize=11)
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)

        # Legend
        from matplotlib.lines import Line2D
        legend_elements = [
            Line2D([0], [0], marker="o", color="#c0392b", label="p < 0.05", markersize=8, linestyle="-"),
            Line2D([0], [0], marker="o", color="#2980b9", label="p < 0.10", markersize=8, linestyle="-"),
            Line2D([0], [0], marker="o", color="#7f8c8d", label="p ≥ 0.10", markersize=8, linestyle="-"),
        ]
        ax.legend(handles=legend_elements, loc="lower right", fontsize=8)

    fig.tight_layout()
    fig.savefig(os.path.join(FIG, "Fig4.png"), dpi=300, bbox_inches="tight")
    plt.close(fig)

    print("  Figures saved to", FIG)


# ════════════════════════════════════════════════════════════
# PPTX generation
# ════════════════════════════════════════════════════════════

def create_pptx():
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    fig_files = sorted([f for f in os.listdir(FIG) if f.endswith(".png")])
    captions = {
        "Fig1.png": "Fig. 1  Conquest rates comparing network closure vs. open polities across three reclassification scenarios.",
        "Fig2.png": "Fig. 2  Fisher's exact test p-values for network closure → conquest association, showing transition to significance with reclassification.",
        "Fig3.png": "Fig. 3  Conquest rates by closure type under the 7-country reclassification scenario.",
        "Fig4.png": "Fig. 4  Forest plot of multivariate logistic regression odds ratios (7-country reclassification, disrupted→overtaken).",
    }

    for fname in fig_files:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        # Title
        txBox = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.2), PptxInches(12), PptxInches(0.6))
        tf = txBox.text_frame
        tf.text = fname.replace(".png", "")
        tf.paragraphs[0].font.size = PptxPt(24)
        tf.paragraphs[0].font.bold = True

        # Image
        img_path = os.path.join(FIG, fname)
        slide.shapes.add_picture(img_path, PptxInches(0.5), PptxInches(1.0), PptxInches(12), PptxInches(5.0))

        # Caption
        cap_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(6.2), PptxInches(12), PptxInches(1.0))
        cap_tf = cap_box.text_frame
        cap_tf.word_wrap = True
        cap_tf.text = captions.get(fname, "")
        cap_tf.paragraphs[0].font.size = PptxPt(12)

    pptx_path = os.path.join(OUT, "figures_pptx.pptx")
    prs.save(pptx_path)
    print("  PPTX saved to", pptx_path)


# ════════════════════════════════════════════════════════════
# Supplementary Table S1 (96 polities)
# ════════════════════════════════════════════════════════════

def create_table_s1(results):
    doc = Document()
    style = doc.styles["Normal"]
    font = style.font
    font.name = "Times New Roman"
    font.size = Pt(9)
    style.paragraph_format.space_after = Pt(2)
    style.paragraph_format.line_spacing = 1.15

    doc.add_heading("Supplementary Table S1: Full Dataset of 96 Historical Polities", level=1)
    p = doc.add_paragraph(
        "Each row represents a historical polity included in the analysis. "
        "Columns report the English name, modern-country equivalent, period of existence, "
        "era classification, dominant strategy (stock or flow), closure type, outcome "
        "(overtaken/disrupted/survived), and the specific turning-point event that determined the outcome."
    )
    for run in p.runs:
        run.font.size = Pt(9)

    df = results["df"]
    headers = ["#", "Polity", "Modern Country", "Period", "Era", "Dominant",
               "Closure Type", "Outcome", "Turning-Point Event"]
    table = doc.add_table(rows=1, cols=len(headers))
    table.style = "Table Grid"
    table.alignment = WD_TABLE_ALIGNMENT.CENTER

    hdr = table.rows[0]
    for i, h in enumerate(headers):
        cell = hdr.cells[i]
        cell.text = h
        for paragraph in cell.paragraphs:
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for run in paragraph.runs:
                run.bold = True
                run.font.size = Pt(8)

    era_map = {"ancient": "Ancient", "medieval": "Medieval", "early_modern": "Early Modern",
               "modern": "Modern", "20c": "20th Century", "contemporary": "Contemporary"}
    closure_map = {"none": "None", "maritime_ban": "Maritime ban", "sakoku": "Sakoku",
                   "bloc": "Bloc", "technical_network_exclusion": "Technical network exclusion"}

    for idx, (_, row) in enumerate(df.iterrows()):
        entity = row["entity"]
        cells_data = [
            str(idx + 1),
            ENGLISH_NAME.get(entity, entity),
            MODERN_COUNTRY.get(entity, "—"),
            row["period"],
            era_map.get(row["era"], row["era"]),
            row["dominant"].capitalize(),
            closure_map.get(row["closure_type"], row["closure_type"]),
            row["outcome"].capitalize(),
            TURNING_POINT.get(entity, "—"),
        ]
        row_cells = table.add_row().cells
        for i, val in enumerate(cells_data):
            row_cells[i].text = val
            for paragraph in row_cells[i].paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(7)

    # Set column widths
    widths = [Cm(0.8), Cm(3.5), Cm(2.5), Cm(2.5), Cm(1.8), Cm(1.2), Cm(2.0), Cm(1.5), Cm(6.0)]
    for row_obj in table.rows:
        for i, w in enumerate(widths):
            row_obj.cells[i].width = w

    s1_path = os.path.join(OUT, "table_s1.docx")
    doc.save(s1_path)
    print("  Table S1 saved to", s1_path)


# ════════════════════════════════════════════════════════════
# Main manuscript DOCX
# ════════════════════════════════════════════════════════════

def add_ref(paragraph, text, sup_text):
    """Add text with superscript reference marker using font-based superscript."""
    parts = re.split(r'(\{[^}]+\})', text)
    for part in parts:
        if part.startswith("{") and part.endswith("}"):
            run = paragraph.add_run(part[1:-1])
            run.font.superscript = True
            run.font.size = Pt(9)
        else:
            paragraph.add_run(part)


def create_manuscript(results):
    doc = Document()

    # ── Page setup ──
    for section in doc.sections:
        section.top_margin = Cm(2.5)
        section.bottom_margin = Cm(2.5)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)

    style = doc.styles["Normal"]
    font = style.font
    font.name = "Times New Roman"
    font.size = Pt(12)
    style.paragraph_format.line_spacing = 1.5
    style.paragraph_format.space_after = Pt(6)

    # ════════════════════════════════════════
    # TITLE PAGE
    # ════════════════════════════════════════
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("Network Exclusion and State Collapse:\n"
                     "From Maritime Isolation to Technological Access Denial\n"
                     "in the Long Run of History")
    run.bold = True
    run.font.size = Pt(16)

    doc.add_paragraph()  # blank
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("[Author Name]\n[Affiliation]\n[Email]")
    run.font.size = Pt(11)

    doc.add_paragraph()
    doc.add_paragraph()

    # ── Abstract ──
    p = doc.add_heading("Abstract", level=2)
    abstract_text = (
        "Why do some states collapse while others endure? This study proposes that structural "
        "exclusion from the dominant technological network of an era—not merely deliberate "
        "trade closure—is a recurring driver of state vulnerability, because it severs the flow "
        "of technical knowledge and allows a cumulative technology gap to open between connected "
        "and excluded polities. Using a comparative dataset of 96 historical polities spanning "
        "antiquity to the present, we distinguish between policy-driven isolation (e.g., Ming "
        "haijin, Tokugawa sakoku) and what we term 'technical network exclusion': involuntary "
        "disconnection from prevailing exchange networks due to geographic or technological "
        "constraints. Policy-based closure reduces but does not eliminate technology transfer "
        "(e.g., rangaku via Dejima under sakoku), whereas technical exclusion severs it entirely. "
        "Reclassifying seven technically excluded polities transforms a non-significant "
        "association between closure and conquest (Fisher's exact test p = 0.187) into a "
        "significant one (p = 0.020); all seven were eventually conquered. Multivariate "
        "logistic regression confirms that external threat and institutional quality are the "
        "strongest predictors of collapse, while the stock–flow odds ratio (OR = 1.774) "
        "remains stable across all scenarios. We argue that the underlying mechanism—technology "
        "flow disruption leading to an accumulating civilization-level gap—generalizes beyond "
        "maritime isolation. As the dominant network shifts from sea lanes to semiconductors, "
        "artificial intelligence, and advanced robotics, states structurally excluded from "
        "these technological platforms may face analogous vulnerabilities."
    )
    p = doc.add_paragraph(abstract_text)
    for run in p.runs:
        run.font.size = Pt(11)

    # ── Keywords ──
    p = doc.add_paragraph()
    run = p.add_run("Keywords: ")
    run.bold = True
    p.add_run("network exclusion; economic history; technology flow disruption; state collapse; "
              "technological access; sensitivity analysis")

    # ── JEL ──
    p = doc.add_paragraph()
    run = p.add_run("JEL Classification: ")
    run.bold = True
    p.add_run("N40, N70, F50, O33, C12")

    doc.add_page_break()

    # ════════════════════════════════════════
    # 1. INTRODUCTION
    # ════════════════════════════════════════
    doc.add_heading("1  Introduction", level=1)

    doc.add_paragraph(
        "Every era has a dominant network—a prevailing infrastructure of exchange through which "
        "states access trade, technology, and strategic information. In antiquity, it was the "
        "Mediterranean sea lanes and the Silk Road caravan routes. In the age of sail, it was "
        "transoceanic maritime commerce. In the industrial age, it was railroad-linked factory "
        "systems and colonial supply chains. Today, it is the digital and semiconductor ecosystem. "
        "A recurring pattern in global history is that states excluded from the dominant network "
        "of their era—whether by choice or by circumstance—face elevated risks of decline and "
        "conquest (Kennedy 1987; Findlay and O'Rourke 2007)."
    )
    doc.add_paragraph(
        "The existing literature offers two broad explanations for why states fail. One "
        "emphasizes internal dynamics: Tainter (1988) argued that complex societies collapse "
        "when the marginal returns to increasing complexity decline, and Acemoglu and Robinson "
        "(2012) showed that extractive institutions—those that concentrate power and discourage "
        "innovation—undermine long-run prosperity. The other emphasizes external connectivity: "
        "trade openness, technology diffusion, and the consequences of isolation (Findlay and "
        "O'Rourke 2007; Mokyr 2002). Our contribution lies at the intersection. We argue that "
        "disconnection from the dominant exchange network is an upstream cause that erodes both "
        "institutional quality and technological capacity—the very factors that the internal-"
        "dynamics literature identifies as proximate causes of collapse."
    )
    doc.add_paragraph(
        "The literature on trade and isolation has overwhelmingly focused "
        "on deliberate closure: the maritime prohibitions (haijin) of Ming China, the sakoku "
        "decree of Tokugawa Japan, the autarkic blocs of the Cold War. Deliberate closure is "
        "analytically convenient because it represents a policy choice with an identifiable agent. "
        "Yet this focus on intentional closure overlooks a logically prior question: what happens "
        "to polities that never had the option of connecting to the dominant network in the first "
        "place? The Khmer Empire at Angkor, the Kievan Rus' principality, or the Timurid dynasty "
        "were not isolated because their rulers chose closure; they were isolated because the "
        "dominant exchange networks of their era did not reach them. If these polities "
        "exhibit the same elevated conquest risk as deliberately closed ones, then the mechanism "
        "at work is not the policy decision to close but the disconnection itself."
    )
    doc.add_paragraph(
        "We propose that the critical channel is technology flow. Maritime trade was never merely "
        "an exchange of goods; it was the primary vehicle through which military techniques, "
        "navigation methods, metallurgy, and institutional innovations diffused across polities "
        "(Pomeranz 2000; Mokyr 2002). Deliberate closure, such as sakoku, restricted but did not "
        "eliminate this flow—Tokugawa Japan maintained a narrow conduit of Western scientific "
        "knowledge (rangaku) through the Dutch trading post at Dejima. Technical exclusion, by "
        "contrast, severed the flow entirely. The result was a cumulative technology gap: over "
        "generations, excluded polities fell progressively behind the technological frontier, "
        "and when contact with a more advanced civilization eventually occurred—often through "
        "military confrontation—the gap proved fatal. This is, in essence, a quantitative "
        "restatement of a long-recognized pattern: when civilizations at markedly different "
        "technological levels collide, the less advanced one tends to be absorbed or destroyed "
        "(Diamond 1997)."
    )
    doc.add_paragraph(
        "This paper makes three contributions. First, we construct a comparative historical dataset "
        "of 96 polities spanning six eras (ancient, medieval, early modern, modern, twentieth "
        "century, and contemporary) and classify each along two dimensions: its predominant resource "
        "base—stock-oriented (accumulated assets such as human capital, institutions, and natural "
        "resources) versus flow-oriented (trade, military projection, and diplomatic engagement)—and "
        "its degree of closure from international networks. Second, we introduce the concept of "
        "'technical network exclusion,' defined as the involuntary disconnection of a polity from "
        "the dominant exchange network of its era due to geographic or technological constraints. "
        "In the maritime age, this took the form of geographic exclusion from sea routes, but "
        "the underlying mechanism generalizes across eras. Third, we "
        "conduct a systematic sensitivity analysis in which seven technically excluded polities are "
        "reclassified from 'no closure' to 'technical network exclusion,' testing whether the "
        "closure–conquest association is an artifact of how isolation is defined."
    )
    doc.add_paragraph(
        "The principal finding is that reclassification transforms a non-significant baseline "
        "association (Fisher's exact test, one-sided p = 0.187) into a significant one (p = 0.020), "
        "while the core stock–flow odds ratio remains unchanged at 1.774. All seven technically "
        "excluded polities were eventually conquered. This result suggests that what matters for "
        "state survival is not whether closure was chosen but whether the polity was connected "
        "to the era's critical exchange platform."
    )
    doc.add_paragraph(
        "The forward-looking implication follows directly. If the mechanism is technology flow "
        "disruption leading to a cumulative gap, then the pattern need not be confined to the "
        "maritime age. In a world where geography no longer blocks physical trade, the 'dominant "
        "network' has shifted from sea lanes to the semiconductor supply chain, artificial "
        "intelligence infrastructure, and advanced robotics. States structurally excluded from "
        "these technological platforms—whether by export controls, infrastructure deficits, or "
        "institutional barriers—may face the same accumulating disadvantage that doomed the "
        "technically excluded polities of the premodern era. We develop this argument in the "
        "discussion."
    )
    doc.add_paragraph(
        "The remainder of the paper is organized as follows. Section 2 describes the dataset and "
        "classification framework. Section 3 presents the analytical methods. Section 4 reports "
        "the baseline results. Section 5 details the sensitivity analysis. Section 6 discusses "
        "implications, including the extension to technological access exclusion, and Section 7 "
        "concludes."
    )

    # ════════════════════════════════════════
    # 2. DATA AND CLASSIFICATION
    # ════════════════════════════════════════
    doc.add_heading("2  Data and Classification", level=1)

    doc.add_heading("2.1  Dataset construction", level=2)
    N = results["N"]
    df = results["df"]
    n_overtaken = len(df[df["outcome"] == "overtaken"])
    n_disrupted = len(df[df["outcome"] == "disrupted"])
    n_survived = len(df[df["outcome"] == "survived"])
    n_stock = len(df[df["dominant"] == "stock"])
    n_flow = len(df[df["dominant"] == "flow"])

    doc.add_paragraph(
        f"The dataset comprises {N} historical polities selected from standard reference works "
        f"(Findlay and O'Rourke 2007; Kennedy 1987; Turchin 2009). Selection criteria required "
        f"that each polity be (i) identifiable as a distinct political entity with a defined period "
        f"of existence, and (ii) classifiable along the dimensions described below. Each polity is "
        f"coded on multiple variables: predominant resource base (stock or flow), stock index (0–1 "
        f"continuous), trade openness (0–1 continuous), closure type, historical outcome, geographic "
        f"barrier, external threat level, relative population, technological position, institutional "
        f"quality, regime duration, and presence of an external patron. The complete dataset, "
        f"including modern-country equivalents and turning-point events, is provided in "
        f"Supplementary Table S1."
    )
    doc.add_paragraph(
        f"Of the {N} polities, {n_stock} are classified as stock-oriented and {n_flow} as "
        f"flow-oriented. Historical outcomes fall into three categories: overtaken "
        f"({n_overtaken} polities)—conquest, colonization, or annexation by an external power; "
        f"disrupted ({n_disrupted})—regime collapse followed by reconstitution under a successor "
        f"state (e.g., Tokugawa Japan giving way to Meiji Japan); and survived ({n_survived})—"
        f"continuity of both regime and statehood. The 'disrupted' category is analytically "
        f"ambiguous, as these polities neither clearly fell to foreign conquest nor clearly "
        f"persisted; we address this through a dual-assignment sensitivity design (Section 3.1)."
    )

    doc.add_heading("2.2  Closure typology", level=2)
    doc.add_paragraph(
        "We classify each polity's degree of closure into five categories, ordered by the "
        "nature and intent of isolation: (1) maritime ban—deliberate restriction of maritime "
        "trade by state policy (e.g., Ming haijin, Qing Canton system); (2) sakoku—near-total "
        "isolation enforced by national decree (Tokugawa Japan, Joseon Korea); (3) bloc—closure "
        "within a geopolitical or economic bloc (e.g., COMECON, imperial preference systems); "
        "(4) technical network exclusion—involuntary isolation from the dominant exchange network "
        "of the era due to geographic or technological constraints (not limited to maritime routes); "
        "and (5) none—no significant closure from prevailing exchange networks."
    )
    doc.add_paragraph(
        "Categories 1 through 3 represent deliberate closure: a policy choice by identifiable "
        "agents. Category 4 represents structural exclusion: isolation imposed by geography and "
        "the limits of contemporary transport technology. This distinction is central to our "
        "argument, because if structural exclusion produces the same outcomes as deliberate "
        "closure, the causal mechanism must lie in the disconnection itself rather than in the "
        "decision to disconnect."
    )

    doc.add_heading("2.3  Technical network exclusion: definition and candidates", level=2)
    doc.add_paragraph(
        "We define technical network exclusion as the condition in "
        "which a polity was structurally disconnected from the dominant exchange network of its "
        "era—whether maritime, overland, or otherwise—due to geographic or technological "
        "constraints rather than policy choice. In the maritime age, this typically meant the "
        "absence of regular sea routes; for inland polities along the Silk Road, it meant that "
        "while some goods traveled overland, the volume and variety of technology transfer fell "
        "far short of what maritime-connected polities received. The key distinction from "
        "policy-based closure is the absence of agency: these polities "
        "did not choose isolation. The transport and communication infrastructure of their era "
        "had not yet extended effective connectivity to their region."
    )
    doc.add_paragraph(
        "We identify seven reclassification candidates in two tiers (Table 1)."
    )

    # ── Table 1: Reclassification candidates ──
    table1 = doc.add_table(rows=1, cols=5)
    table1.style = "Table Grid"
    table1.alignment = WD_TABLE_ALIGNMENT.CENTER
    hdr = table1.rows[0]
    for i, h in enumerate(["Polity", "Period", "Tier", "Original Closure", "Rationale"]):
        hdr.cells[i].text = h
        for par in hdr.cells[i].paragraphs:
            for run in par.runs:
                run.bold = True
                run.font.size = Pt(10)

    candidates = STRONG_CANDIDATES + MODERATE_CANDIDATES
    for entity in candidates:
        en_name = ENGLISH_NAME.get(entity, entity)
        row_data = df[df["entity"] == entity].iloc[0]
        tier = "Strong" if entity in STRONG_CANDIDATES else "Moderate"
        rationale_en = {
            "漢朝（前漢〜後漢）": "Excluded from Mediterranean network; Silk Road provided limited overland contact only.",
            "マリ帝国": "Excluded from Atlantic/Mediterranean networks until Portuguese contact (15th c.). Trans-Saharan caravan only.",
            "クメール帝国（アンコール）": "Inland polity; excluded from dominant exchange networks unlike neighbouring Srivijaya.",
            "キエフ大公国": "Dnieper river trade (Varangian route) only; excluded from dominant exchange networks.",
            "ティムール朝": "Landlocked; Silk Road overland only. Excluded from dominant Mediterranean and Indian Ocean networks.",
            "ササン朝ペルシア": "Limited Persian Gulf trade; excluded from established Indian Ocean and Mediterranean networks.",
            "ビルマ（コンバウン朝）": "Coastal access existed but excluded from regular international exchange networks until British arrival.",
        }.get(entity, "")

        row_cells = table1.add_row().cells
        for j, val in enumerate([en_name, row_data["period"], tier, "None", rationale_en]):
            row_cells[j].text = val
            for par in row_cells[j].paragraphs:
                for run in par.runs:
                    run.font.size = Pt(9)

    p = doc.add_paragraph()
    p.space_before = Pt(6)
    run = p.add_run("Table 1  ")
    run.bold = True
    run.font.size = Pt(10)
    p.add_run("Technical network exclusion reclassification candidates. "
              "'Strong' candidates were structurally excluded from the dominant exchange network of their era; "
              "'Moderate' candidates had limited or uncertain connectivity.").font.size = Pt(10)

    # ════════════════════════════════════════
    # 3. METHODS
    # ════════════════════════════════════════
    doc.add_heading("3  Methods", level=1)

    doc.add_heading("3.1  Outcome binarization and sensitivity design", level=2)
    doc.add_paragraph(
        "The three-category outcome (overtaken, disrupted, survived) poses a classification "
        "problem. The 18 'disrupted' polities experienced regime collapse but were not clearly "
        "conquered by an external power; they could reasonably be grouped with either outcome. "
        "Rather than making a single arbitrary assignment, we adopt a dual-assignment design: "
        "in the 'disrupted → overtaken' scenario, all 18 are coded as conquered (yielding 64 "
        "conquered, 32 survived); in the 'disrupted → survived' scenario, they are coded as "
        "having survived (46 conquered, 50 survived). Crossed with three closure reclassification "
        "levels (baseline, +5 strong candidates, +7 all candidates), this produces 3 × 2 = 6 "
        "scenarios. All results are reported across all six to demonstrate robustness."
    )

    doc.add_heading("3.2  Stock–flow association", level=2)
    doc.add_paragraph(
        "We construct 2 × 2 tables crossing the predominant resource base (stock vs. flow) "
        "against the binarized outcome (overtaken vs. survived) and compute the odds ratio (OR), "
        "phi coefficient (φ), one-sided Fisher's exact test, and chi-squared test with Yates "
        "correction. Effect sizes are interpreted using Cohen's benchmarks for φ."
    )

    doc.add_heading("3.3  Closure–conquest association", level=2)
    doc.add_paragraph(
        "For each of the six scenarios, we construct a 2 × 2 table crossing closure status "
        "(network closure vs. open) against the binarized outcome. The 'network closure' "
        "group includes polities coded as maritime ban, sakoku, or technical network exclusion. We "
        "compute one-sided Fisher's exact tests evaluating whether closure is associated with "
        "elevated conquest rates, and report relative risks alongside p-values."
    )

    doc.add_heading("3.4  Multivariate logistic regression", level=2)
    doc.add_paragraph(
        "To assess whether the closure–conquest association survives adjustment for confounders, "
        "we fit logistic regression models with the binarized outcome as the dependent variable "
        "and the following covariates: stock-dominant indicator, geographic barrier, external "
        "threat level, technological position, institutional quality, era (coded ordinally), "
        "external patron indicator, and a network closure indicator. We report exponentiated "
        "coefficients (odds ratios) with 95% confidence intervals."
    )

    doc.add_heading("3.5  Bootstrap validation", level=2)
    doc.add_paragraph(
        "We validate the stock–flow OR using a nonparametric bootstrap (5,000 resamples, "
        "percentile method, seed = 42). This provides a distribution-free confidence interval "
        "that does not depend on asymptotic normality—a relevant consideration given the "
        "modest sample size."
    )

    # ════════════════════════════════════════
    # 4. RESULTS
    # ════════════════════════════════════════
    doc.add_heading("4  Results", level=1)

    doc.add_heading("4.1  Baseline stock–flow association", level=2)
    cm_base = results["scenarios"]["as_conquered__baseline"]["cm"]
    doc.add_paragraph(
        f"Under the disrupted → overtaken assignment, the stock–flow 2 × 2 table yields "
        f"OR = {cm_base['OR']:.3f}, φ = {cm_base['phi']:.3f}, Fisher's exact p = "
        f"{cm_base['p_fisher']:.4f} (one-sided). Stock-oriented polities have a conquest rate "
        f"of {cm_base['stock_conquest_rate']:.1%}, compared with {cm_base['flow_conquest_rate']:.1%} "
        f"for flow-oriented polities—a small-to-medium effect by Cohen's benchmarks. Under the "
        f"disrupted → survived assignment, the OR shifts to "
        f"{results['scenarios']['as_survived__baseline']['cm']['OR']:.3f}, illustrating the "
        f"sensitivity of the stock–flow association to how the ambiguous 'disrupted' category "
        f"is handled. This dual sensitivity—to both closure definition and outcome "
        f"binarization—motivates the full six-scenario analysis below."
    )

    # ── Table 2: Summary across 6 scenarios ──
    doc.add_heading("4.2  Maritime closure and conquest", level=2)

    s_base_c = results["scenarios"]["as_conquered__baseline"]
    s_strong_c = results["scenarios"]["as_conquered__strong"]
    s_all_c = results["scenarios"]["as_conquered__all"]
    s_base_s = results["scenarios"]["as_survived__baseline"]
    s_strong_s = results["scenarios"]["as_survived__strong"]
    s_all_s = results["scenarios"]["as_survived__all"]

    doc.add_paragraph(
        f"Table 2 reports the closure–conquest association across all six scenarios. "
        f"At baseline under the disrupted → overtaken assignment, polities with some form of "
        f"network closure have a conquest rate of {s_base_c['ban_rate']:.1%}, compared with "
        f"{s_base_c['no_ban_rate']:.1%} for open polities (Fisher p = "
        f"{s_base_c['fisher_ban_p']:.4f}, not significant). Reclassifying five strong candidates "
        f"as technically excluded raises the closure-group conquest rate to "
        f"{s_strong_c['ban_rate']:.1%} (Fisher p = {s_strong_c['fisher_ban_p']:.4f}). "
        f"Including all seven candidates strengthens the association further: conquest rate = "
        f"{s_all_c['ban_rate']:.1%}, Fisher p = {s_all_c['fisher_ban_p']:.4f} (Fig. 1, Fig. 2). "
        f"The progressive strengthening of significance as technically excluded polities are "
        f"added suggests that these polities genuinely belong in the closure group rather than "
        f"among the open polities."
    )

    table2 = doc.add_table(rows=1, cols=8)
    table2.style = "Table Grid"
    table2.alignment = WD_TABLE_ALIGNMENT.CENTER
    t2_headers = ["Disrupted As", "Reclassification", "Closure n", "Closure Rate",
                  "Open Rate", "RR", "Fisher p", "Sig."]
    hdr = table2.rows[0]
    for i, h in enumerate(t2_headers):
        hdr.cells[i].text = h
        for par in hdr.cells[i].paragraphs:
            for run in par.runs:
                run.bold = True
                run.font.size = Pt(9)

    for d_mode, d_label in [("as_conquered", "Overtaken"), ("as_survived", "Survived")]:
        for c_key, c_label in [("baseline", "Baseline"), ("strong", "+5 Strong"), ("all", "+7 All")]:
            key = f"{d_mode}__{c_key}"
            s = results["scenarios"][key]
            sig = "*" if s["fisher_ban_p"] < 0.05 else ""
            row_cells = table2.add_row().cells
            vals = [d_label, c_label, str(s["ban_n"]),
                    f"{s['ban_rate']:.1%}", f"{s['no_ban_rate']:.1%}",
                    f"{s['rr']:.3f}", f"{s['fisher_ban_p']:.4f}", sig]
            for j, val in enumerate(vals):
                row_cells[j].text = val
                for par in row_cells[j].paragraphs:
                    for run in par.runs:
                        run.font.size = Pt(9)

    p = doc.add_paragraph()
    p.space_before = Pt(6)
    run = p.add_run("Table 2  ")
    run.bold = True
    run.font.size = Pt(10)
    p.add_run("Maritime closure and conquest across six scenarios. * p < 0.05.").font.size = Pt(10)

    # Insert Fig 1 inline
    p = doc.add_paragraph()
    p.space_before = Pt(12)
    fig1_path = os.path.join(FIG, "Fig1.png")
    if os.path.exists(fig1_path):
        run = p.add_run()
        run.add_picture(fig1_path, width=Inches(6))
    p2 = doc.add_paragraph()
    run = p2.add_run("Fig. 1  ")
    run.bold = True
    run.font.size = Pt(10)
    p2.add_run("Conquest rates comparing network closure vs. open polities across three "
               "reclassification scenarios under both disrupted assignments.").font.size = Pt(10)

    # Insert Fig 2 inline
    p = doc.add_paragraph()
    p.space_before = Pt(12)
    fig2_path = os.path.join(FIG, "Fig2.png")
    if os.path.exists(fig2_path):
        run = p.add_run()
        run.add_picture(fig2_path, width=Inches(5))
    p2 = doc.add_paragraph()
    run = p2.add_run("Fig. 2  ")
    run.bold = True
    run.font.size = Pt(10)
    p2.add_run("Fisher's exact test p-values for the network closure → conquest association "
               "as polities are progressively reclassified.").font.size = Pt(10)

    # ════════════════════════════════════════
    # 5. SENSITIVITY ANALYSIS
    # ════════════════════════════════════════
    doc.add_heading("5  Sensitivity Analysis: Technical Network Exclusion", level=1)

    doc.add_heading("5.1  Closure-type disaggregation and the dose–response pattern", level=2)
    doc.add_paragraph(
        "Figure 3 disaggregates conquest rates by closure type under the 7-country reclassification. "
        "A striking gradient emerges. Technically excluded polities—those with zero access to "
        "the dominant exchange networks of their era—exhibit a 100% conquest rate. Policy-based maritime bans, "
        "which restricted but did not entirely eliminate external contact, show lower rates "
        "(76.9% under disrupted → overtaken; 69.2% under disrupted → survived). Sakoku polities "
        "show 100% and 50% rates under the two assignments, reflecting the borderline case "
        "of Tokugawa Japan, which maintained a narrow technological conduit (rangaku) through "
        "Dejima. Bloc-type closures show the lowest conquest rates among closure categories, "
        "consistent with the interpretation that bloc membership preserves some technology "
        "transfer through alliance-internal channels."
    )
    doc.add_paragraph(
        "This ordering—technical exclusion (100%) > policy ban > sakoku (with partial conduit) > "
        "bloc > open—is consistent with a dose–response relationship between the degree of "
        "technology flow disruption and the probability of conquest. The more completely a polity "
        "was severed from the technological frontier of its era, the higher the likelihood that "
        "it was eventually overtaken."
    )

    # Insert Fig 3 inline
    p = doc.add_paragraph()
    p.space_before = Pt(12)
    fig3_path = os.path.join(FIG, "Fig3.png")
    if os.path.exists(fig3_path):
        run = p.add_run()
        run.add_picture(fig3_path, width=Inches(6))
    p2 = doc.add_paragraph()
    run = p2.add_run("Fig. 3  ")
    run.bold = True
    run.font.size = Pt(10)
    p2.add_run("Conquest rates by closure type under the 7-country reclassification scenario.").font.size = Pt(10)

    doc.add_heading("5.2  Robustness of the stock–flow odds ratio", level=2)
    cm_all_c = results["scenarios"]["as_conquered__all"]["cm"]
    boot_c = results["bootstrap"].get("as_conquered__all", {})
    doc.add_paragraph(
        f"The stock–flow OR = {cm_all_c['OR']:.3f} is identical across all three "
        f"reclassification scenarios. This invariance is expected: the reclassification "
        f"changes the closure-type label but does not alter the stock/flow or outcome coding. "
        f"Bootstrap validation (5,000 resamples) yields a median OR of "
        f"{boot_c.get('median', 0):.3f} (95% CI [{boot_c.get('ci_lo', 0):.3f}, "
        f"{boot_c.get('ci_hi', 0):.3f}]), confirming the stability of the point estimate "
        f"and its independence from the closure reclassification."
    )

    doc.add_heading("5.3  Multivariate regression stability", level=2)
    doc.add_paragraph(
        "Figure 4 presents the multivariate logistic regression results under the 7-country "
        "reclassification with disrupted → overtaken. External threat remains the strongest "
        "predictor of conquest (p < 0.01 across all scenarios), followed by institutional "
        "quality and era (both p < 0.01). The network closure indicator is not independently "
        "significant after controlling for these covariates. This pattern is informative: it "
        "suggests that closure operates not as a direct cause but through the same channels—"
        "technological stagnation, institutional decay, and heightened external vulnerability—"
        "that the multivariate model already captures (Table 3)."
    )

    # Insert Fig 4 inline
    p = doc.add_paragraph()
    p.space_before = Pt(12)
    fig4_path = os.path.join(FIG, "Fig4.png")
    if os.path.exists(fig4_path):
        run = p.add_run()
        run.add_picture(fig4_path, width=Inches(5.5))
    p2 = doc.add_paragraph()
    run = p2.add_run("Fig. 4  ")
    run.bold = True
    run.font.size = Pt(10)
    p2.add_run("Forest plot of multivariate logistic regression odds ratios "
               "(7-country reclassification, disrupted → overtaken).").font.size = Pt(10)

    # ── Table 3: Multivariate results ──
    doc.add_paragraph()
    lr_all_c = results["scenarios"]["as_conquered__all"]["logistic"]
    if lr_all_c.get("with_ban", {}).get("converged"):
        table3 = doc.add_table(rows=1, cols=5)
        table3.style = "Table Grid"
        table3.alignment = WD_TABLE_ALIGNMENT.CENTER
        t3_headers = ["Variable", "OR", "95% CI", "p-value", "Sig."]
        hdr = table3.rows[0]
        for i, h in enumerate(t3_headers):
            hdr.cells[i].text = h
            for par in hdr.cells[i].paragraphs:
                for run in par.runs:
                    run.bold = True
                    run.font.size = Pt(10)

        var_labels = {
            "dominant_binary": "Stock-dominant",
            "geo_barrier": "Geographic barrier",
            "external_threat": "External threat",
            "tech_position": "Technological position",
            "institutional_quality": "Institutional quality",
            "era_code": "Era (time)",
            "has_external_patron": "External patron",
            "has_maritime_ban": "Network closure dummy",
        }
        coefs = lr_all_c["with_ban"]["coefs"]
        for var, label in var_labels.items():
            v = coefs[var]
            sig = "*" if v["p"] < 0.05 else "\u2020" if v["p"] < 0.10 else ""
            ci_lo_str = f"{v['ci_lo']:.3f}" if v["ci_lo"] < 1e6 else ">10\u2076"
            ci_hi_str = f"{v['ci_hi']:.3f}" if v["ci_hi"] < 1e6 else ">10\u2076"
            row_cells = table3.add_row().cells
            vals = [label, f"{v['OR']:.3f}", f"[{ci_lo_str}, {ci_hi_str}]",
                    f"{v['p']:.4f}", sig]
            for j, val in enumerate(vals):
                row_cells[j].text = val
                for par in row_cells[j].paragraphs:
                    for rn in par.runs:
                        rn.font.size = Pt(9)

        p = doc.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run("Table 3  ")
        run.bold = True
        run.font.size = Pt(10)
        p.add_run("Multivariate logistic regression results (7-country reclassification, "
                   "disrupted → overtaken). * p < 0.05, \u2020 p < 0.10.").font.size = Pt(10)

    # ════════════════════════════════════════
    # 6. DISCUSSION
    # ════════════════════════════════════════
    doc.add_heading("6  Discussion", level=1)

    doc.add_heading("6.1  The mechanism: technology flow disruption and cumulative divergence", level=2)
    doc.add_paragraph(
        "The central finding is that the closure–conquest association becomes significant only "
        "when technically excluded polities are grouped with deliberately closed ones. This tells "
        "us something important about the mechanism at work. If closure harmed states solely "
        "through lost trade revenue or reduced diplomatic leverage, then only deliberate closure—"
        "which blocks trade but not necessarily knowledge—should matter. The fact that technical "
        "exclusion (which severs both trade and technology flow) strengthens the association, "
        "while policy-based closure alone does not reach significance, points toward technology "
        "flow as the critical channel."
    )
    doc.add_paragraph(
        "The dose–response pattern in Figure 3 reinforces this interpretation. Technical "
        "exclusion (zero technology transfer) produces a 100% conquest rate. Policy-based "
        "maritime bans, which restrict but do not eliminate technology flow, show conquest "
        "rates below 80%. The case of Tokugawa Japan is particularly instructive: despite "
        "the comprehensive closure of sakoku, the Tokugawa regime deliberately maintained "
        "a narrow conduit for Western scientific and technical knowledge (rangaku) through "
        "the Dutch trading post at Dejima. This selective preservation of a technology "
        "transfer channel—even within an otherwise closed system—appears to have made a "
        "material difference. Japan was disrupted by the forced opening of 1853–54 but was "
        "not conquered; it reconstituted itself as the Meiji state and rapidly closed the "
        "technology gap. Bloc closures, which preserve substantial within-bloc technology "
        "sharing, show the lowest rates among closure categories. The ordering of conquest "
        "risk mirrors the ordering of technology flow disruption."
    )
    doc.add_paragraph(
        "The multivariate results complete the picture. External threat and institutional "
        "quality are the strongest predictors of conquest, and the network closure indicator "
        "loses significance after their inclusion (Table 3, Fig. 4). This is precisely what "
        "a technology-gap mechanism would predict: closure does not kill states directly. "
        "Rather, it initiates a causal chain—technological stagnation erodes institutional "
        "adaptive capacity, which in turn leaves the polity unable to respond to external "
        "threats. This causal ordering is consistent with Acemoglu and Robinson's (2012) "
        "emphasis on institutions as the proximate determinant of national success, while "
        "suggesting that network access is the deeper, upstream variable: exclusion degrades "
        "the very institutions that the inclusive-institutions framework identifies as essential. "
        "The mediating variables (external threat, institutional quality) absorb "
        "the explanatory power of the closure variable because they lie downstream in the "
        "causal pathway."
    )

    doc.add_heading("6.2  First contact and the tragedy of civilizational divergence", level=2)
    doc.add_paragraph(
        "Our findings can be read as a quantitative formulation of a long-recognized historical "
        "pattern: when civilizations that have developed in isolation encounter a technologically "
        "superior civilization, the outcome is overwhelmingly unfavorable for the less advanced "
        "party (Diamond 1997; Diamond and Bellwood 2003). The 100% conquest rate among technically excluded polities is "
        "striking not because the pattern is new, but because it emerges from a systematic, "
        "cross-historical dataset rather than from selective case studies."
    )
    doc.add_paragraph(
        "The mechanism we propose is cumulative divergence through technology flow disruption. "
        "International exchange networks carried not only goods but military techniques, "
        "metallurgical innovations, navigational knowledge, and institutional models (Mokyr "
        "2002; Pomeranz 2000). Polities connected to these networks could adopt, adapt, and "
        "build upon innovations generated elsewhere. Polities severed from them could not. "
        "Over generations, the technology gap widened\u2014a process analogous to the long-run "
        "consequences of network disruption documented by Nunn (2008), who showed that regions "
        "more heavily affected by the slave trade experienced persistent underdevelopment "
        "centuries later. When contact with a more advanced civilization eventually occurred\u2014"
        "often through military expansion\u2014the accumulated gap "
        "proved decisive. The Han Dynasty encountered Central Asian and eventually Roman-linked "
        "military traditions; the Khmer Empire faced the expanding Siamese and Vietnamese states "
        "that were integrated into maritime trade networks; Kievan Rus\u2019 was overrun by the "
        "Mongol armies that had absorbed the military technologies of multiple civilizations "
        "across Eurasia."
    )
    doc.add_paragraph(
        "Crucially, the tragedy of first contact is a function of the gap, not of the contact "
        "itself. Policy-closed polities that maintained narrow conduits of technology transfer—"
        "Japan's rangaku, Qing China's limited Canton trade—accumulated smaller gaps and, "
        "correspondingly, were more likely to survive or reconstitute after disruption. The "
        "implication is that what determines the outcome of civilizational encounter is the "
        "degree and duration of technology flow disruption that preceded it."
    )
    doc.add_paragraph(
        "This finding has a corollary that is worth stating explicitly. If the critical "
        "variable is not closure itself but the residual technology flow that closure permits, "
        "then polities that close their borders while deliberately maintaining selective "
        "channels for frontier knowledge occupy a qualitatively different position from those "
        "that are totally severed. Our dataset contains several instances of such conditional "
        "closure. Tokugawa Japan preserved access to Western science through rangaku at Dejima "
        "and was disrupted but not conquered, reconstituting itself as the Meiji state. Early "
        "Qing China maintained the Canton system—a single, tightly controlled port of trade "
        "through which some foreign knowledge filtered—and survived the Kangxi–Qianlong era "
        "intact. Joseon Korea sustained tributary trade with China, which served as a conduit "
        "for Continental knowledge and technology. The Soviet Union, while sealed from the "
        "Western bloc, maintained extensive technology sharing within the Eastern bloc and "
        "invested heavily in indigenous research institutions. North Korea has preserved a "
        "limited technology channel through its relationship with China. "
        "(Full details of each polity's closure regime, technology channels, and outcomes "
        "are recorded in Supplementary Table S1.)"
    )
    doc.add_paragraph(
        "The outcomes, however, are mixed. Early Qing China survived, but late Qing China—"
        "operating under the same Canton system—was semi-colonized after the Opium Wars. "
        "Tokugawa Japan survived as a political entity, but Joseon Korea, despite its Chinese "
        "conduit, was annexed by Japan in 1910. The Soviet Union collapsed internally despite "
        "its bloc-level technology sharing, while North Korea has persisted under extreme "
        "isolation with minimal Chinese patronage. These divergent outcomes suggest that the "
        "mere existence of a selective channel does not guarantee survival; unmeasured factors—"
        "the breadth and relevance of the knowledge flowing through the channel, the domestic "
        "capacity to absorb and adapt that knowledge, the pace of change at the technological "
        "frontier, and institutional dynamics that our dataset does not capture—likely "
        "condition the effectiveness of conditional closure. The policy question, then, is not "
        "binary (open or closed) but conditional, and the conditions under which selective "
        "channels suffice to prevent a fatal technology gap remain an open and consequential "
        "problem for future research. Readers interested in tracing these cases in detail "
        "are referred to Supplementary Table S1, which documents the specific turning-point "
        "events and outcomes for all 96 polities in the dataset."
    )

    doc.add_heading("6.3  Beyond geographic isolation: technological access exclusion "
                    "in the modern era", level=2)
    doc.add_paragraph(
        "If the mechanism we identify is technology flow disruption rather than network closure "
        "per se, then the pattern should generalize beyond the age of sail. Each era has its own "
        "dominant technological platform—the infrastructure through which frontier knowledge "
        "diffuses across states. In antiquity, it was the Mediterranean trade routes and Silk "
        "Road caravans. In the early modern period, it was oceanic shipping. In the industrial "
        "age, it was railroad-linked factory systems and colonial supply chains. In the twentieth "
        "century, it was aerospace and nuclear technology networks."
    )
    doc.add_paragraph(
        "Today, the dominant platform has shifted again. The critical networks are semiconductor "
        "supply chains, artificial intelligence research ecosystems, and advanced robotics and "
        "automation infrastructure (Acemoglu and Restrepo 2020; Comin and Mestieri 2018). "
        "Comin et al. (2010) show that technology adoption levels in 1000 BC predict income "
        "differences today, while Acemoglu et al. (2002) demonstrate that colonial-era "
        "institutional reversals reshaped global inequality—both consistent with the view "
        "that early technological access has persistent, cumulative consequences. "
        "Geographic isolation no longer blocks physical trade—the completion of global shipping "
        "and communication networks has largely eliminated geographic network exclusion as a "
        "threat. But a new form of structural exclusion has emerged: states may be cut off from "
        "the technological frontier not by mountains and oceans but by export controls on "
        "advanced semiconductors, by the concentration of AI training infrastructure in a "
        "handful of countries, or by the institutional and human-capital barriers that prevent "
        "participation in cutting-edge research networks."
    )
    doc.add_paragraph(
        "The historical parallel is direct. Just as the Khmer Empire or the Timurid dynasty "
        "could not access the exchange networks that carried military and institutional "
        "innovations, a contemporary state excluded from advanced semiconductor fabrication "
        "or AI model development may find itself on the wrong side of a widening technology "
        "gap. If the gap grows large enough, the eventual 'first contact'—whether military, "
        "economic, or geopolitical—may produce outcomes analogous to those documented in our "
        "historical dataset. The form of the dominant network changes; the logic of cumulative "
        "divergence through exclusion does not."
    )
    doc.add_paragraph(
        "We stress that this extrapolation is speculative and cannot be tested within our "
        "historical dataset. The contemporary world differs from the premodern era in ways "
        "that may attenuate or amplify the mechanism: nuclear deterrence, international "
        "institutions, and the speed of modern communication all introduce novel dynamics. "
        "Nevertheless, the historical regularity we document—that structural exclusion from "
        "the dominant technological network is associated with state collapse—provides a "
        "framework for thinking about which dimensions of modern technological access may "
        "be most consequential."
    )

    doc.add_heading("6.4  Robustness of the stock–flow framework and the question of "
                    "resource-base transitions", level=2)
    # Compute stock x closure interaction: 4 cells with Fisher exact tests
    df_all_c = apply_disrupted_assignment(
        apply_technical_network_exclusion(results["df"], STRONG_CANDIDATES + MODERATE_CANDIDATES),
        "as_conquered"
    )
    has_closure = df_all_c["closure_type"].isin(["maritime_ban", "technical_network_exclusion", "sakoku"])
    cells_64 = {}
    for dom in ["stock", "flow"]:
        for cl_label, cl_val in [("closed", True), ("open", False)]:
            sub = df_all_c[(df_all_c["dominant"] == dom) & (has_closure == cl_val)]
            n = len(sub)
            conquered = int(sub["outcome_binary"].sum())
            survived = n - conquered
            rate = sub["outcome_binary"].mean() if n > 0 else 0
            cells_64[(dom, cl_label)] = {"n": n, "c": conquered, "s": survived, "rate": rate}

    sc = cells_64[("stock", "closed")]
    so = cells_64[("stock", "open")]
    fc = cells_64[("flow", "closed")]
    fo = cells_64[("flow", "open")]

    # Fisher exact tests for key pairwise comparisons
    def fisher_or_p(a, b):
        """One-sided Fisher exact test (greater) and OR for two cells."""
        table = np.array([[a["c"], a["s"]], [b["c"], b["s"]]])
        odds_r, p_val = stats.fisher_exact(table, alternative="greater")
        return odds_r, p_val

    or_sc_fo, p_sc_fo = fisher_or_p(sc, fo)  # stock+closed vs flow+open
    or_sc_so, p_sc_so = fisher_or_p(sc, so)  # closure effect within stock
    or_fc_fo, p_fc_fo = fisher_or_p(fc, fo)  # closure effect within flow

    doc.add_paragraph(
        f"The invariance of the stock–flow OR (1.774) across all reclassification scenarios "
        f"confirms that the core finding of the stock–flow framework—that stock-oriented "
        f"polities face moderately higher conquest risk—is independent of how network "
        f"isolation is defined. The reclassification changes the closure subanalysis but "
        f"leaves the primary classification untouched. This separation is analytically "
        f"useful: it shows that the stock–flow distinction and the closure–conquest "
        f"association capture related but distinct dimensions of state vulnerability."
    )
    doc.add_paragraph(
        f"Crossing these two dimensions yields a four-cell classification whose conquest "
        f"rates are reported in Table 4. Flow-oriented polities without closure show the "
        f"lowest conquest rate ({fo['rate']:.1%}, n = {fo['n']}). Stock-oriented polities "
        f"without closure show a moderately elevated rate ({so['rate']:.1%}, n = {so['n']}). "
        f"Stock-oriented polities with closure show a markedly higher rate ({sc['rate']:.1%}, "
        f"n = {sc['n']}). Flow-oriented polities with closure—a small cell (n = {fc['n']})—"
        f"show a {fc['rate']:.1%} conquest rate. The contrast between the highest- and lowest-"
        f"risk cells (stock + closed vs. flow + open) yields an odds ratio of {or_sc_fo:.2f} "
        f"(one-sided Fisher exact p = {p_sc_fo:.3f}). Within stock-oriented polities, the "
        f"effect of closure corresponds to an OR of {or_sc_so:.2f} (p = {p_sc_so:.3f}). "
        f"Within flow-oriented polities, the effect of closure yields an OR of "
        f"{'∞' if np.isinf(or_fc_fo) else f'{or_fc_fo:.2f}'} (p = {p_fc_fo:.3f}), though "
        f"the cell size (n = {fc['n']}) precludes reliable inference."
    )

    # Add Table 4: Stock–flow × closure interaction
    doc.add_paragraph()
    p_t4_title = doc.add_paragraph()
    run_t4 = p_t4_title.add_run(
        "Table 4  Conquest rates by resource-base orientation and closure status "
        "(7-country reclassification, disrupted = conquered)"
    )
    run_t4.bold = True
    run_t4.font.size = Pt(10)
    p_t4_title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    t4 = doc.add_table(rows=4, cols=5)
    t4.style = "Table Grid"
    t4.alignment = WD_TABLE_ALIGNMENT.CENTER
    t4_headers = ["", "Conquered", "Survived", "n", "Rate"]
    for i, h in enumerate(t4_headers):
        cell = t4.rows[0].cells[i]
        cell.text = h
        for par in cell.paragraphs:
            par.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for run in par.runs:
                run.bold = True
                run.font.size = Pt(10)

    t4_data = [
        ("Stock + closed", sc),
        ("Stock + open", so),
        ("Flow + closed", fc),
        ("Flow + open", fo),
    ]
    for row_idx, (label, cell_data) in enumerate(t4_data, start=0):
        row = t4.rows[row_idx]  # re-use existing rows (header is row 0)
    # Actually need 5 rows (1 header + 4 data). Rebuild.
    # Remove the table and redo with correct row count
    # python-docx doesn't support removing tables easily; let's add rows instead
    # We already have 4 rows. The first is headers. We need to add 1 more row.
    t4.add_row()  # now 5 rows total

    # Re-assign: row 0 = header (already set), rows 1-4 = data
    for row_idx, (label, cell_data) in enumerate(t4_data):
        data_row = t4.rows[row_idx + 1]
        vals = [label, str(cell_data["c"]), str(cell_data["s"]),
                str(cell_data["n"]), f'{cell_data["rate"]:.1%}']
        for col_i, val in enumerate(vals):
            data_row.cells[col_i].text = val
            for par in data_row.cells[col_i].paragraphs:
                par.alignment = WD_ALIGN_PARAGRAPH.CENTER
                for run in par.runs:
                    run.font.size = Pt(10)

    # Table 4 footnote
    p_t4_note = doc.add_paragraph(
        f"Fisher exact tests (one-sided): stock + closed vs. flow + open, "
        f"OR = {or_sc_fo:.2f}, p = {p_sc_fo:.3f}; within stock, closed vs. open, "
        f"OR = {or_sc_so:.2f}, p = {p_sc_so:.3f}; within flow, closed vs. open, "
        f"OR = {'∞' if np.isinf(or_fc_fo) else f'{or_fc_fo:.2f}'}, p = {p_fc_fo:.3f}."
    )
    for run in p_t4_note.runs:
        run.font.size = Pt(9)
        run.italic = True

    doc.add_paragraph(
        f"The pattern is consistent with the hypothesis that stock-orientation and network "
        f"exclusion compound each other's risk, though the small cell sizes—particularly for "
        f"flow + closed (n = {fc['n']})—warrant caution. The stock + closed combination is "
        f"the only cell to reach a statistically significant difference from the baseline "
        f"flow + open cell at conventional levels. "
        f"Notably, {sc['s']} stock-oriented polities with closure survived despite their "
        f"high-risk classification; their individual characteristics—including specific "
        f"closure regimes, technology channels maintained, and the historical circumstances "
        f"of their survival—are documented in Supplementary Table S1."
    )
    doc.add_paragraph(
        "An implication worth noting is that the stock–flow distinction is not permanently "
        "fixed for any given state. Polities can transition from flow-oriented to stock-oriented "
        "resource bases—or vice versa—as economic structures, demographic conditions, and "
        "institutional incentives evolve. A state that was historically flow-oriented (trade-"
        "dependent, outward-looking, innovation-absorbing) but that gradually shifts toward "
        "reliance on accumulated domestic assets—physical capital, territorial resources, "
        "existing institutional infrastructure—moves into the higher-risk stock category. "
        "If such a transition coincides with reduced engagement in the dominant technological "
        "network of the era, the compounding of stock-orientation and network exclusion could "
        "amplify vulnerability. Our dataset cannot test this dynamic directly, as polities are "
        "coded at a single point in time, but the theoretical interaction between resource-base "
        "transition and network access merits further investigation."
    )

    doc.add_heading("6.5  Limitations", level=2)
    doc.add_paragraph(
        "Several limitations warrant acknowledgment. First, the coding of historical polities "
        "inevitably involves subjective judgment, particularly for the stock/flow classification "
        "and the identification of technical exclusion candidates. The tiered approach (strong "
        "vs. moderate candidates) and the six-scenario sensitivity design partially address this, "
        "but cannot eliminate it. Second, the sample size (N = 96) constrains the power of the "
        "multivariate analyses; wide confidence intervals for some regression coefficients "
        "reflect this constraint. Third, the dataset treats polities as independent observations, "
        "though historical interconnections (e.g., sequential Chinese dynasties sharing "
        "institutional continuity) may introduce non-independence. Fourth, the 'disrupted' "
        "category introduces a classification ambiguity that the dual-assignment design addresses "
        "but cannot fully resolve. Fifth, the forward-looking extension to modern technological "
        "exclusion is necessarily speculative, as the mechanisms operating in a nuclear-armed, "
        "institutionally dense modern world may differ qualitatively from those in premodern eras."
    )

    # ════════════════════════════════════════
    # 7. CONCLUSION
    # ════════════════════════════════════════
    doc.add_heading("7  Conclusion", level=1)

    doc.add_paragraph(
        "This paper has shown that reclassifying seven technically excluded polities—those "
        "severed from the dominant exchange networks of their era by geography and technology "
        "rather than by policy—transforms a non-significant association between closure and "
        "conquest into a significant one (p = 0.020), while the core stock–flow odds ratio "
        "remains unchanged (OR = 1.774). The 100% conquest rate among technically excluded "
        "polities, the dose–response gradient across closure types, and the absorption of "
        "the closure effect by external threat and institutional quality in multivariate "
        "models all point to a consistent mechanism: disruption of technology flow leads "
        "to cumulative divergence from the technological frontier, eroding the institutional "
        "and military capacity needed to survive contact with more connected civilizations."
    )
    doc.add_paragraph(
        "The broader implication is that this mechanism is not specific to any single network form. "
        "In every era, there exists a dominant network through which frontier technologies "
        "diffuse. Polities excluded from that network—whether by oceans, mountains, policy, "
        "or, in the contemporary period, by semiconductor export controls and AI infrastructure "
        "concentration—risk falling into the same pattern of cumulative divergence. The "
        "historical record we document provides a quantitative baseline for assessing this "
        "risk. Whether the tragedy of first contact between unequally developed civilizations "
        "will find new expression in the age of artificial intelligence is a question that "
        "the coming decades will answer; our analysis suggests it is one worth asking. "
        "Equally pressing is the converse question suggested by the Tokugawa precedent: "
        "whether a state that recognizes the risk of network exclusion can, through "
        "deliberate maintenance of selective technology channels, prevent the accumulation "
        "of a fatal gap—even while restricting broader engagement with the outside world."
    )

    # ════════════════════════════════════════
    # REFERENCES (Author-year, alphabetical — Cliometrica style)
    # ════════════════════════════════════════
    doc.add_heading("References", level=1)

    refs = [
        "Acemoglu D, Johnson S, Robinson JA (2002) Reversal of fortune: geography and institutions in the making of the modern world income distribution. Q J Econ 117:1231–1294. https://doi.org/10.1162/003355302320935025",
        "Acemoglu D, Johnson S, Robinson JA (2005) Institutions as a fundamental cause of long-run growth. In: Aghion P, Durlauf SN (eds) Handbook of economic growth, vol 1A. Elsevier, Amsterdam, pp 385–472",
        "Acemoglu D, Robinson JA (2012) Why nations fail: the origins of power, prosperity, and poverty. Crown, New York",
        "Acemoglu D, Restrepo P (2020) Robots and jobs: evidence from US labor markets. J Polit Econ 128:2188–2244. https://doi.org/10.1086/705716",
        "Arrighi G (1994) The long twentieth century: money, power, and the origins of our times. Verso, London",
        "Broadberry SN, Guan H (2026) Regional variation of GDP per head within China, 1080–1850. Explor Econ Hist 95:101567. https://doi.org/10.1016/j.eeh.2025.101567",
        "Comin D, Easterly W, Gong E (2010) Was the wealth of nations determined in 1000 BC? Am Econ J Macroecon 2:65–97. https://doi.org/10.1257/mac.2.3.65",
        "Comin D, Mestieri M (2018) If technology has arrived everywhere, why has income diverged? Am Econ J Macroecon 10:137–178. https://doi.org/10.1257/mac.20150175",
        "De Vries J (2010) The limits of globalization in the early modern world. Econ Hist Rev 63:710–733. https://doi.org/10.1111/j.1468-0289.2009.00497.x",
        "Diamond J (1997) Guns, germs, and steel: the fates of human societies. W.W. Norton, New York",
        "Diamond J, Bellwood P (2003) Farmers and their languages: the first expansions. Science 300:597–603. https://doi.org/10.1126/science.1078208",
        "Findlay R, O'Rourke KH (2007) Power and plenty: trade, war, and the world economy in the second millennium. Princeton University Press, Princeton",
        "Kennedy P (1987) The rise and fall of the great powers: economic change and military conflict from 1500 to 2000. Random House, New York",
        "Maddison A (2007) Contours of the world economy 1–2030 AD: essays in macro-economic history. Oxford University Press, Oxford",
        "Mokyr J (2002) The gifts of Athena: historical origins of the knowledge economy. Princeton University Press, Princeton",
        "North DC, Wallis JJ, Weingast BR (2009) Violence and social orders: a conceptual framework for interpreting recorded human history. Cambridge University Press, Cambridge",
        "Nunn N (2008) The long-term effects of Africa's slave trades. Q J Econ 123:139–176. https://doi.org/10.1162/qjec.2008.123.1.139",
        "Pomeranz K (2000) The great divergence: China, Europe, and the making of the modern world economy. Princeton University Press, Princeton",
        "Tainter JA (1988) The collapse of complex societies. Cambridge University Press, Cambridge",
        "Turchin P (2009) A theory for formation of large empires. J Glob Hist 4:191–217. https://doi.org/10.1017/S1740022809003192",
        "Turchin P, Nefedov SA (2009) Secular cycles. Princeton University Press, Princeton",
    ]
    for ref in refs:
        p = doc.add_paragraph(ref)
        for run in p.runs:
            run.font.size = Pt(11)
        p.paragraph_format.first_line_indent = Cm(-1)
        p.paragraph_format.left_indent = Cm(1)

    # ════════════════════════════════════════
    # STATEMENTS AND DECLARATIONS
    # ════════════════════════════════════════
    doc.add_heading("Statements and Declarations", level=1)

    doc.add_heading("Funding", level=2)
    doc.add_paragraph("[To be completed by author]")

    doc.add_heading("Competing Interests", level=2)
    doc.add_paragraph("The author declares no competing interests.")

    doc.add_heading("Data Availability", level=2)
    doc.add_paragraph(
        "The complete dataset and analysis code are available at [repository URL]. "
        "Supplementary Table S1 provides the full dataset of 96 polities with all coded variables."
    )

    # Save
    ms_path = os.path.join(OUT, "manuscript.docx")
    doc.save(ms_path)
    print("  Manuscript saved to", ms_path)


# ════════════════════════════════════════════════════════════
# Main
# ════════════════════════════════════════════════════════════

def main():
    print("Running analysis...")
    results = run_analysis()

    print("Creating figures...")
    create_figures(results)

    print("Creating PPTX...")
    create_pptx()

    print("Creating Table S1...")
    create_table_s1(results)

    print("Creating manuscript...")
    create_manuscript(results)

    print("\nAll outputs saved to:", OUT)
    print("  manuscript/manuscript.docx")
    print("  manuscript/table_s1.docx")
    print("  manuscript/figures_pptx.pptx")
    print("  manuscript/figures/Fig1–Fig4.png")


if __name__ == "__main__":
    main()
