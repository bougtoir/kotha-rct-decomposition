#!/usr/bin/env python3
"""Build an editable PPTX with one slide per figure from 05_paper_cct.md."""
import argparse
import os
import re
import shutil
import tempfile
import zipfile
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def clean_math(latex):
    """Convert a small subset of LaTeX math to plain ASCII text for PPTX captions."""
    text = latex
    replacements = [
        ('\\text{', ''), ('}', ''),
        ('\\sim', '~'), ('\\cdot', '*'),
        ('\\sqrt', 'sqrt'), ('\\sum', 'sum'),
        ('\\prod', 'prod'), ('\\alpha', 'alpha'),
        ('\\beta', 'beta'), ('\\mu', 'mu'),
        ('\\tau', 'tau'), ('\\theta', 'theta'),
        ('\\sigma', 'sigma'), ('\\delta', 'delta'),
        ('\\rho', 'rho'), ('\\Phi', 'Phi'),
        ('\\phi', 'phi'), ('\\in', ' in '),
        ('\\leq', '<='), ('\\geq', '>='),
        ('\\neq', '!='), ('\\times', 'x'),
        ('\\to', '->'), ('\\approx', '~'),
        ('\\frac', ''), ('\\left', ''),
        ('\\right', ''), ('\\quad', '  '),
        ('\\hat', ''), ('\\log', 'log'),
        ('\\exp', 'exp'), ('\\mid', '|'),
        ('^2', '^2'),
    ]
    for old, new in replacements:
        text = text.replace(old, new)
    text = re.sub(r'_\{([^}]+)\}', r'_\1', text)
    text = re.sub(r'\^\{([^}]+)\}', r'^\1', text)
    text = text.replace('\\', '')
    return text


def _sanitize_pptx(path):
    """Replace typographic and multibyte characters in PPTX XML with ASCII."""
    tmp = path + '.tmp'
    try:
        with zipfile.ZipFile(path, 'r') as zin, zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as zout:
            for info in zin.infolist():
                data = zin.read(info.filename)
                if info.filename.endswith('.xml'):
                    text = data.decode('utf-8')
                    text = text.replace('\u2039', '').replace('\u203a', '')
                    text = text.replace('\u2018', "'").replace('\u2019', "'")
                    text = text.replace('\u201c', '"').replace('\u201d', '"')
                    text = text.replace('\u2013', '-').replace('\u2014', '-')
                    text = text.replace('\u00a0', ' ')
                    data = text.encode('utf-8')
                zout.writestr(info, data)
        shutil.move(tmp, path)
    except Exception:
        if os.path.exists(tmp):
            os.remove(tmp)
        raise


def build(figures_md_path, out_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    md = Path(figures_md_path).read_text()

    # Parse each **Fig. N** caption followed by ![...](path) anywhere in the document
    for cap_match in re.finditer(r'\*\*Fig\.\s*(\d+)\*\*\s*(.*?)\n\s*!\[.*?\]\((.*?)\)', md, re.S):
        fig_num = cap_match.group(1)
        caption = cap_match.group(2).strip()
        caption = re.sub(r'\$([^$]+)\$', lambda m: clean_math(m.group(1)), caption)
        img_path = cap_match.group(3).strip()

        blank = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(blank)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f'Figure {fig_num}'
        p.font.size = Pt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Image
        full_img = Path(figures_md_path).parent / img_path
        if full_img.exists():
            slide.shapes.add_picture(str(full_img), Inches(1.0), Inches(1.1), width=Inches(11.333))

        # Caption
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(1.0))
        tf2 = cap_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = f'Figure {fig_num}. {caption}'
        p2.font.size = Pt(14)
        p2.word_wrap = True
        p2.alignment = PP_ALIGN.LEFT

    prs.save(out_path)
    _sanitize_pptx(out_path)
    print(f'Saved {out_path}')


if __name__ == '__main__':
    parser = argparse.ArgumentParser()
    parser.add_argument('--md', default='05_paper_cct.md')
    parser.add_argument('--out', default='KOTHA_Framework_CCT_figures.pptx')
    args = parser.parse_args()
    build(args.md, args.out)
