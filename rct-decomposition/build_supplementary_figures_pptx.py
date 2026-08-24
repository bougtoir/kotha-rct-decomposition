#!/usr/bin/env python3
"""Build an editable PPTX for supplementary figures (S1a/b, S2-S4)."""
import argparse
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def clean_math(latex):
    """Convert a small subset of LaTeX math to plain Unicode text for PPTX captions."""
    text = latex
    replacements = [
        ('\\text{', ''), ('}', ''),
        ('\\sim', '~'), ('\\cdot', '\u00b7'),
        ('\\sqrt', '\u221a'), ('\\sum', '\u03a3'),
        ('\\prod', '\u220f'), ('\\alpha', '\u03b1'),
        ('\\beta', '\u03b2'), ('\\mu', '\u03bc'),
        ('\\tau', '\u03c4'), ('\\theta', '\u03b8'),
        ('\\sigma', '\u03c3'), ('\\delta', '\u03b4'),
        ('\\rho', '\u03c1'), ('\\Phi', '\u03a6'),
        ('\\phi', '\u03c6'), ('\\in', '\u2208'),
        ('\\leq', '\u2264'), ('\\geq', '\u2265'),
        ('\\neq', '\u2260'), ('\\times', '\u00d7'),
        ('\\to', '\u2192'), ('\\approx', '\u2248'),
        ('\\frac', ''), ('\\left', ''),
        ('\\right', ''), ('\\quad', '  '),
        ('\\hat', ''), ('\\log', 'log'),
        ('\\exp', 'exp'), ('\\mid', '|'),
        ('^2', '\u00b2'),
    ]
    for old, new in replacements:
        text = text.replace(old, new)
    text = re.sub(r'_\{([^}]+)\}', r'_\1', text)
    text = re.sub(r'\^\{([^}]+)\}', r'^\1', text)
    text = text.replace('\\', '')
    return text


# Trace plots are referenced in the Methods but are not given full figure blocks in the manuscript.
TRACE_CAPTIONS = {
    'figS1a_trace_mcmc_magnesium.png': (
        'MCMC trace plots for selected parameters in the Bayesian power-prior model for '
        'magnesium in acute myocardial infarction (discounting factor α = 0.3). Chains show no obvious non-stationarity.'
    ),
    'figS1b_trace_mcmc_statins.png': (
        'MCMC trace plots for selected parameters in the Bayesian power-prior model for '
        'statins in heart failure (discounting factor α = 0.3). Chains show no obvious non-stationarity.'
    ),
}


def _parse_cctc_figures(cctc_md_path):
    """Parse CCTC figure blocks for Figs 5/6/7 and map them to supplementary labels."""
    md = Path(cctc_md_path).read_text()
    mapping = {}
    label_map = {
        '5': 'S2',
        '6': 'S3',
        '7': 'S4',
    }
    pattern = r'\*\*Fig\.\s*(5|6|7)\*\*\s*(.*?)\n\s*!\[.*?\]\((.*?)\)'
    for m in re.finditer(pattern, md, re.S):
        fig_num = m.group(1)
        caption = m.group(2).strip()
        img_path = m.group(3).strip()
        caption = re.sub(r'\$([^$]+)\$', lambda x: clean_math(x.group(1)), caption)
        mapping[label_map[fig_num]] = {
            'caption': caption,
            'img': Path(cctc_md_path).parent / img_path,
        }
    return mapping


def build(cctc_md_path, out_path, base_dir=None):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    base_dir = Path(base_dir) if base_dir else Path(cctc_md_path).parent

    # S1a and S1b trace plots
    figures = []
    for img_name, caption in [
        ('figS1a_trace_mcmc_magnesium.png', 'Magnesium'),
        ('figS1b_trace_mcmc_statins.png', 'Statins'),
    ]:
        src = base_dir / 'validation' / 'figures' / img_name
        figures.append((f'S1{chr(ord("a") + len(figures))}', src, TRACE_CAPTIONS.get(img_name, '')))

    # S2-S4 from CCTC manuscript
    cctc_figs = _parse_cctc_figures(cctc_md_path)
    for label, info in [('S2', '5'), ('S3', '6'), ('S4', '7')]:
        if label in cctc_figs:
            figures.append((label, cctc_figs[label]['img'], cctc_figs[label]['caption']))

    for label, img_path, caption in figures:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f'Supplementary Figure {label}'
        p.font.size = Pt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        if img_path.exists():
            slide.shapes.add_picture(str(img_path), Inches(1.0), Inches(1.1), width=Inches(11.333))
        else:
            print(f'  WARNING: supplementary figure image not found: {img_path}')

        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(1.0))
        tf2 = cap_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = f'Supplementary Figure {label}. {caption}'
        p2.font.size = Pt(14)
        p2.word_wrap = True
        p2.alignment = PP_ALIGN.LEFT

    prs.save(out_path)
    print(f'Saved {out_path}')


if __name__ == '__main__':
    parser = argparse.ArgumentParser()
    parser.add_argument('--cctc-md', default='05_paper_cctc.md')
    parser.add_argument('--out', default='KOTHA_Framework_ClinicalTrials_supplementary_figures.pptx')
    parser.add_argument('--base-dir', default=None)
    args = parser.parse_args()
    build(args.cctc_md, args.out, args.base_dir)
