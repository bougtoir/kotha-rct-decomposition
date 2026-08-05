#!/usr/bin/env python3
"""Build an editable PPTX with one slide per figure from 05_paper_cct.md."""
import argparse
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def clean_math(latex):
    """Convert a small subset of LaTeX math to plain Unicode text for PPTX captions."""
    text = latex
    replacements = [
        ('\\text{', ''), ('}', ''),
        ('\\sim', '~'), ('\\cdot', '\u00b7'),
        ('\\sqrt', '\u221a'), ('\\sum', '\u03a3'),
        ('\\prod', '\u220f'), ('\\alpha', '\u03b1'),
        ('\\beta', '\u03b2'), ('\\mu', '\u03bc'),
        ('\\tau', '\u03c4'), ('\\theta', '\u03b8'),
        ('\\sigma', '\u03c3'), ('\\delta', '\u03b4'),
        ('\\rho', '\u03c1'), ('\\Phi', '\u03a6'),
        ('\\phi', '\u03c6'), ('\\in', '\u2208'),
        ('\\leq', '\u2264'), ('\\geq', '\u2265'),
        ('\\neq', '\u2260'), ('\\times', '\u00d7'),
        ('\\to', '\u2192'), ('\\approx', '\u2248'),
        ('\\frac', ''), ('\\left', ''),
        ('\\right', ''), ('\\quad', '  '),
        ('\\hat', ''), ('\\log', 'log'),
        ('\\exp', 'exp'), ('\\mid', '|'),
        ('^2', '\u00b2'),
    ]
    for old, new in replacements:
        text = text.replace(old, new)
    text = re.sub(r'_\{([^}]+)\}', r'_\1', text)
    text = re.sub(r'\^\{([^}]+)\}', r'^\1', text)
    text = text.replace('\\', '')
    return text


def build(figures_md_path, out_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    md = Path(figures_md_path).read_text()

    # Parse each **Fig. N** caption followed by ![...](path) anywhere in the document
    for cap_match in re.finditer(r'\*\*Fig\.\s*(\d+)\*\*\s*(.*?)\n\s*!\[.*?\]\((.*?)\)', md, re.S):
        fig_num = cap_match.group(1)
        caption = cap_match.group(2).strip()
        caption = re.sub(r'\$([^$]+)\$', lambda m: clean_math(m.group(1)), caption)
        img_path = cap_match.group(3).strip()

        blank = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(blank)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f'Figure {fig_num}'
        p.font.size = Pt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Image
        full_img = Path(figures_md_path).parent / img_path
        if full_img.exists():
            slide.shapes.add_picture(str(full_img), Inches(1.0), Inches(1.1), width=Inches(11.333))

        # Caption
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(1.0))
        tf2 = cap_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = f'Figure {fig_num}. {caption}'
        p2.font.size = Pt(14)
        p2.word_wrap = True
        p2.alignment = PP_ALIGN.LEFT

    prs.save(out_path)
    print(f'Saved {out_path}')


if __name__ == '__main__':
    parser = argparse.ArgumentParser()
    parser.add_argument('--md', default='05_paper_cct.md')
    parser.add_argument('--out', default='KOTHA_Framework_CCT_figures.pptx')
    args = parser.parse_args()
    build(args.md, args.out)
