"""Build final manuscript .docx (JA + EN) with inline figures/tables
and a single editable .pptx (EN figures) for per-slide editing.

Usage:  python build_docx_pptx.py
Outputs into ../manuscript/ :
  - manuscript_en.docx
  - manuscript_ja.docx
  - figures_en.pptx
  - table1_model_metrics.docx
  - table2_correspondence.docx
"""
from __future__ import annotations

import os
import re
import json
from dataclasses import dataclass

import pandas as pd
from docx import Document
from docx.shared import Pt, Inches, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_ALIGN_VERTICAL
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGB

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.abspath(os.path.join(HERE, ".."))
MS = os.path.join(ROOT, "manuscript")
FIG = os.path.join(ROOT, "figures")
TAB = os.path.join(ROOT, "tables")
os.makedirs(MS, exist_ok=True)

# ---------------------------------------------------------------------------
# Figure/table metadata

FIG_LIST = [
    ("fig1", "Fig. 1", "図1",
     "In-sample 1-year GDP growth RMSE across M0-M4 for 39 countries (lower is better).",
     "39カ国における単年GDP成長率RMSEの標本内比較（M0-M4、小さいほど良い）。",
     "fig1_m_ranking_{lang}.png"),
    ("fig2", "Fig. 2", "図2",
     "Out-of-sample MAPE on 2015-19 held-out window; M2 achieves 13% relative improvement vs M0.",
     "2015-19年のホールドアウト窓における標本外MAPE。M2はM0比で相対13%改善。",
     "fig2_oos_{lang}.png"),
    ("fig3", "Fig. 3", "図3",
     "PIM stock K_tang+betaK_I versus CWON PCA (within-country demeaned log).",
     "PIM資本ストックK_tang+βK_IとCWON PCAの軌跡比較（国内平均除去対数）。",
     "fig3_trajectories_{lang}.png"),
    ("fig4", "Fig. 4", "図4",
     "gamma_price sensitivity of PIM/CWON log-ratio; Japan gap closes around gamma=+0.02.",
     "PIM/CWON対数比のγ_price感度。日本の乖離はγ=+0.02で閉じる。",
     "fig4_gamma_price_{lang}.png"),
    ("fig5", "Fig. 5", "図5",
     "Conceptual diagram of the population-capital tempo correspondence.",
     "人口・資本テンポ対応関係の概念図。",
     "fig5_concept_{lang}.png"),
    ("fig6", "Fig. 6", "図6",
     "Relational PIM diagnostics: rho_2 across 39 countries under M0 vs M4.",
     "関係型PIM診断: M0とM4における39カ国のρ̂₂。",
     "fig6_rpim_{lang}.png"),
    ("fig7", "Fig. 7", "図7",
     "Delta-mu sensitivity: mu_hat under +/-20% depreciation perturbation.",
     "δ-μ感度分析: 減価償却率±20%変動下のμ̂。",
     "fig7_delta_sensitivity_{lang}.png"),
]


# ---------------------------------------------------------------------------
# docx helpers

def set_font(run, name="Times New Roman", size=11, bold=False, italic=False):
    run.font.name = name
    run.font.size = Pt(size)
    run.bold = bold
    run.italic = italic
    # Apply East Asian font for proper JA rendering
    rPr = run._element.get_or_add_rPr()
    rFonts = rPr.find(qn("w:rFonts"))
    if rFonts is None:
        rFonts = OxmlElement("w:rFonts")
        rPr.append(rFonts)
    rFonts.set(qn("w:ascii"), name)
    rFonts.set(qn("w:hAnsi"), name)
    rFonts.set(qn("w:eastAsia"), "MS Mincho")


def add_heading(doc, text, level, lang):
    h = doc.add_paragraph()
    h.paragraph_format.space_before = Pt(12)
    h.paragraph_format.space_after = Pt(6)
    run = h.add_run(text)
    set_font(run, size={1: 16, 2: 13, 3: 12}[level], bold=True)
    return h


def add_para(doc, text, lang, italic=False):
    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(6)
    p.paragraph_format.line_spacing = 1.5
    p.paragraph_format.first_line_indent = Pt(12)
    run = p.add_run(text)
    set_font(run, italic=italic)
    return p


def add_figure(doc, png_path, caption_prefix, caption_text):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(18)
    p.paragraph_format.space_after = Pt(6)
    run = p.add_run()
    run.add_picture(png_path, width=Inches(6.2))
    cap = doc.add_paragraph()
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    cap.paragraph_format.space_after = Pt(12)
    r1 = cap.add_run(f"{caption_prefix}. ")
    set_font(r1, size=10, bold=True)
    r2 = cap.add_run(caption_text)
    set_font(r2, size=10)


def add_dataframe_as_table(doc, df: pd.DataFrame, col_widths=None, font_size=10):
    tbl = doc.add_table(rows=1 + len(df), cols=len(df.columns))
    tbl.style = "Table Grid"
    # Header
    for j, col in enumerate(df.columns):
        c = tbl.rows[0].cells[j]
        c.text = ""
        p = c.paragraphs[0]
        r = p.add_run(str(col))
        set_font(r, size=font_size, bold=True)
    for i, row in enumerate(df.itertuples(index=False), start=1):
        for j, v in enumerate(row):
            c = tbl.rows[i].cells[j]
            c.text = ""
            p = c.paragraphs[0]
            r = p.add_run(str(v))
            set_font(r, size=font_size)
    if col_widths:
        for row in tbl.rows:
            for idx, w in enumerate(col_widths):
                row.cells[idx].width = w


def add_table_block(doc, title, df, caption, widths=None):
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(18)
    p.paragraph_format.space_after = Pt(6)
    r = p.add_run(title)
    set_font(r, size=11, bold=True)
    r2 = p.add_run(f"  {caption}")
    set_font(r2, size=11)
    add_dataframe_as_table(doc, df, col_widths=widths)


# ---------------------------------------------------------------------------
# Build docx manuscript from markdown, replacing figure/table placeholders

INSERT_FIG_RE = re.compile(r"^\*\*\[(Insert |)[Ff]ig(ure|\.)? ?(\d+)[^\]]*\]\*\*$")
INSERT_FIG_JA_RE = re.compile(r"^\*\*［図\s*(\d+)[^］]*］\*\*$")
INSERT_TAB_RE = re.compile(r"^\*\*\[(Insert |)[Tt]able ?(\d+)[^\]]*\]\*\*$")
INSERT_TAB_JA_RE = re.compile(r"^\*\*［表\s*(\d+)[^］]*］\*\*$")


def build_manuscript(lang: str):
    md_path = os.path.join(MS, f"manuscript_{lang}.md")
    with open(md_path, encoding="utf-8") as fh:
        lines = fh.readlines()

    # Load tables
    t1 = pd.read_csv(os.path.join(TAB, "table1_model_metrics.csv"))
    t2 = pd.read_csv(os.path.join(TAB, "table2_correspondence.csv"))
    t3_path = os.path.join(TAB, "table3_rpim.csv")
    t3 = pd.read_csv(t3_path) if os.path.exists(t3_path) else None

    # Figure caption lookup by index
    fig_cap = {}
    for key, en_prefix, ja_prefix, en_cap, ja_cap, pattern in FIG_LIST:
        idx = int(key[-1])
        fig_cap[idx] = {
            "en": (en_prefix, en_cap, pattern.format(lang="en")),
            "ja": (ja_prefix, ja_cap, pattern.format(lang="ja")),
        }

    t1_cap_en = "M0-M4: in-sample and out-of-sample performance across 39 countries. "\
                "Medians across countries; IQR in brackets."
    t1_cap_ja = "M0-M4: 39カ国の標本内・標本外パフォーマンス（国間中央値、IQRを括弧内）。"
    t2_cap_en = "Population-capital tempo correspondence."
    t2_cap_ja = "人口・資本テンポ対応関係。"
    t3_cap_en = "Relational PIM diagnostics: rho_2 summary under M0, M1, M2, M4."
    t3_cap_ja = "関係型PIM診断: M0, M1, M2, M4におけるρ̂₂の要約。"

    doc = Document()
    # Use sensible page margins
    for section in doc.sections:
        section.top_margin = Cm(2.5)
        section.bottom_margin = Cm(2.5)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)

    i = 0
    while i < len(lines):
        line = lines[i].rstrip("\n")
        stripped = line.strip()
        if not stripped:
            i += 1
            continue
        # Headings
        if stripped.startswith("# "):
            add_heading(doc, stripped[2:], 1, lang)
        elif stripped.startswith("## "):
            add_heading(doc, stripped[3:], 2, lang)
        elif stripped.startswith("### "):
            add_heading(doc, stripped[4:], 3, lang)
        elif stripped.startswith("---"):
            pass  # ignore horizontal rules
        else:
            # Figure inline placeholder (both lang patterns)
            m_fig = INSERT_FIG_RE.match(stripped) or INSERT_FIG_JA_RE.match(stripped)
            m_tab = INSERT_TAB_RE.match(stripped) or INSERT_TAB_JA_RE.match(stripped)
            if m_fig:
                try:
                    idx = int(m_fig.group(m_fig.lastindex))
                except Exception:
                    idx = None
                if idx in fig_cap:
                    prefix, cap, fname = fig_cap[idx][lang]
                    png = os.path.join(FIG, fname)
                    if os.path.exists(png):
                        add_figure(doc, png, prefix, cap)
            elif m_tab:
                try:
                    idx = int(m_tab.group(m_tab.lastindex))
                except Exception:
                    idx = None
                if idx == 1:
                    add_table_block(
                        doc,
                        "Table 1." if lang == "en" else "表1.",
                        t1,
                        t1_cap_en if lang == "en" else t1_cap_ja,
                    )
                elif idx == 2:
                    add_table_block(
                        doc,
                        "Table 2." if lang == "en" else "表2.",
                        t2,
                        t2_cap_en if lang == "en" else t2_cap_ja,
                        widths=[Inches(1.3), Inches(2.5), Inches(2.5)],
                    )
                elif idx == 3 and t3 is not None:
                    add_table_block(
                        doc,
                        "Table 3." if lang == "en" else "表3.",
                        t3,
                        t3_cap_en if lang == "en" else t3_cap_ja,
                    )
            else:
                # Drop markdown emphasis markers for clean rendering
                text = stripped
                # preserve bold sentence like "**Abstract** (146 words)."
                text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
                text = re.sub(r"\*(.+?)\*", r"\1", text)
                # section separator
                if text in ("Tables", "表", "References", "参考文献"):
                    add_heading(doc, text, 2, lang)
                else:
                    add_para(doc, text, lang)
        i += 1

    out = os.path.join(MS, f"manuscript_{lang}.docx")
    doc.save(out)
    print("wrote", out)


# ---------------------------------------------------------------------------
# Separate table .docx files

def build_standalone_tables():
    t1 = pd.read_csv(os.path.join(TAB, "table1_model_metrics.csv"))
    t2 = pd.read_csv(os.path.join(TAB, "table2_correspondence.csv"))
    t3_path = os.path.join(TAB, "table3_rpim.csv")
    t3 = pd.read_csv(t3_path) if os.path.exists(t3_path) else None
    table_list = [
        ("table1_model_metrics.docx", t1,
         "Table 1. M0-M4: in-sample and out-of-sample performance across 39 countries.",
         None),
        ("table2_correspondence.docx", t2,
         "Table 2. Population-capital tempo correspondence.",
         [Inches(1.3), Inches(2.5), Inches(2.5)]),
    ]
    if t3 is not None:
        table_list.append(
            ("table3_rpim.docx", t3,
             "Table 3. Relational PIM diagnostics: rho_2 summary under M0, M1, M2, M4.",
             None))
    for name, df, cap, widths in table_list:
        d = Document()
        add_heading(d, cap, 2, "en")
        add_dataframe_as_table(d, df, col_widths=widths)
        out = os.path.join(MS, name)
        d.save(out)
        print("wrote", out)


# ---------------------------------------------------------------------------
# Editable pptx (English figures only, one per slide)

def build_pptx():
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    blank = prs.slide_layouts[6]

    for key, en_prefix, _ja, en_cap, _jc, pattern in FIG_LIST:
        png = os.path.join(FIG, pattern.format(lang="en"))
        if not os.path.exists(png):
            print("skip missing", png)
            continue
        slide = prs.slides.add_slide(blank)
        # Title
        tb = slide.shapes.add_textbox(
            PptInches(0.4), PptInches(0.2),
            PptInches(12.5), PptInches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"{en_prefix}. {en_cap}"
        r.font.size = PptPt(20)
        r.font.bold = True
        # Image: fit to 12 x 5.5 area, centered
        slide.shapes.add_picture(
            png, PptInches(1.5), PptInches(1.1),
            width=PptInches(10.3), height=PptInches(5.6))
        # Caption bar at bottom
        cb = slide.shapes.add_textbox(
            PptInches(0.4), PptInches(6.85),
            PptInches(12.5), PptInches(0.5))
        cf = cb.text_frame
        cf.word_wrap = True
        pp = cf.paragraphs[0]
        rr = pp.add_run()
        rr.text = en_cap
        rr.font.size = PptPt(14)
        rr.font.italic = True
        rr.font.color.rgb = PptRGB(0x44, 0x44, 0x44)

    out = os.path.join(MS, "figures_en.pptx")
    prs.save(out)
    print("wrote", out)


def main():
    build_manuscript("en")
    build_manuscript("ja")
    build_standalone_tables()
    build_pptx()


if __name__ == "__main__":
    main()
