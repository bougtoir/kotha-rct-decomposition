"""
Generate editable .pptx with all figures (one per slide).
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIGURES_DIR = os.path.join(BASE_DIR, 'figures')


def create_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    figures = [
        ('fig1_verdict_distribution.png',
         'Figure 1: Verdict Distribution',
         'Distribution of verdicts across 52 canonical curves by domain.'),
        ('fig2_sensitivity_analysis.png',
         'Figure 2: Sensitivity Analysis',
         'F-test p-values before and after outlier removal (Cook\'s distance top 3).'),
        ('fig3_model_comparison.png',
         'Figure 3: Model Comparison (AIC/BIC)',
         'Model preference by AIC and BIC for linear, quadratic, and log models.'),
        ('fig4_loocv_comparison.png',
         'Figure 4: LOOCV Comparison',
         'Leave-one-out cross-validation RMSE: linear vs quadratic models.'),
        ('fig5_sample_size.png',
         'Figure 5: Sample Size vs Verdict',
         'Relationship between sample size and analysis verdict.'),
    ]

    for fname, title, caption in figures:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.7))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Image
        img_path = os.path.join(FIGURES_DIR, fname)
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.0),
                                     width=Inches(10.0))

        # Caption
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.CENTER

    output_path = os.path.join(BASE_DIR, 'figures_canonical_curves.pptx')
    prs.save(output_path)
    print(f"PPTX saved to: {output_path}")


if __name__ == "__main__":
    create_pptx()
