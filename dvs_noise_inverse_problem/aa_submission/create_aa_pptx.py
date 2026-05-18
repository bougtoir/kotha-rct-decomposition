#!/usr/bin/env python3
"""
Generate A&A figures PPTX (editable, 1 figure per slide, English).
A&A independent figure numbering: Fig. 1-7.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
PARENT_DIR = os.path.dirname(BASE_DIR)

# A&A figure mapping (independent from MLST Fig 1-4)
AA_FIGURES = [
    {
        'num': 1,
        'title': 'Fig. 1: PI-DC-DVS Pipeline Architecture',
        'file': os.path.join(PARENT_DIR, 'fig2_g3_pipeline_en.png'),
        'caption': 'System architecture of the PI-DC-DVS noise inverse problem pipeline. '
                   'Four stages: (1) noise forward model; (2) Bayesian inverse problem; '
                   '(3) residual generation; (4) calibration and verification.',
    },
    {
        'num': 2,
        'title': 'Fig. 2: Research Gap Map',
        'file': os.path.join(PARENT_DIR, 'fig1_gap_map_en.png'),
        'caption': 'Gap map showing four domains (A-D) and identified research gaps. '
                   'Cal-6 satellite trail calibration is highlighted as a novel contribution.',
    },
    {
        'num': 3,
        'title': 'Fig. 3: Systematic Evaluation — 3-Method Comparison',
        'file': os.path.join(PARENT_DIR, 'fig5_systematic_evaluation.png'),
        'caption': 'Systematic evaluation on 20 EBSSA recordings. Four-panel boxplot: '
                   'NRR, SPR, F1, AUC. Fano filter achieves best balance (AUC = 0.866).',
    },
    {
        'num': 4,
        'title': 'Fig. 4: A5 Noise Rate Simulation',
        'file': os.path.join(PARENT_DIR, 'fig6_a5_simulation.png'),
        'caption': 'A5-based noise rate simulation: noise rate, SNR improvement, '
                   'and factor map across temperature–illuminance space. Mean 5.4×, max 10.0×.',
    },
    {
        'num': 5,
        'title': 'Fig. 5: Per-Recording Noise Removal Comparison',
        'file': os.path.join(PARENT_DIR, 'fig7_per_recording_comparison.png'),
        'caption': 'Per-recording noise removal rate comparison. '
                   'Fano filter (blue) consistently outperforms temporal filter (orange).',
    },
    {
        'num': 6,
        'title': 'Fig. 6: Proof-of-Concept Demonstration',
        'file': os.path.join(PARENT_DIR, 'fig3_noise_inverse_demo.png'),
        'caption': 'Noise inverse problem demonstration on EBSSA data. '
                   '90.3% noise removal with satellite trajectory clearly visible.',
    },
    {
        'num': 7,
        'title': 'Fig. 7: SNR Improvement Analysis',
        'file': os.path.join(PARENT_DIR, 'fig4_sn_improvement.png'),
        'caption': 'SNR improvement analysis: Fano factor spatial map, temporal dynamics, '
                   'and per-pixel SNR distribution before and after noise subtraction.',
    },
]

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank_layout = prs.slide_layouts[6]  # blank

    for fig in AA_FIGURES:
        slide = prs.slides.add_slide(blank_layout)

        # Title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = fig['title']
        p.font.size = Pt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Image
        img_path = fig['file']
        if os.path.exists(img_path):
            # Scale to fit
            max_w = Inches(11.5)
            max_h = Inches(5.5)
            from PIL import Image
            with Image.open(img_path) as im:
                w, h = im.size
            aspect = w / h
            if aspect > (11.5 / 5.5):
                img_w = max_w
                img_h = int(max_w / aspect)
            else:
                img_h = max_h
                img_w = int(max_h * aspect)

            left = int((SLIDE_WIDTH - img_w) / 2)
            top = Inches(1.0)
            slide.shapes.add_picture(img_path, left, top, img_w, img_h)
        else:
            txBox2 = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(1))
            tf2 = txBox2.text_frame
            tf2.paragraphs[0].text = f"[Image not found: {os.path.basename(img_path)}]"
            tf2.paragraphs[0].font.size = Pt(14)
            tf2.paragraphs[0].font.color.rgb = None

        # Caption
        txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.8))
        tf3 = txBox3.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        p3.text = fig['caption']
        p3.font.size = Pt(12)
        p3.alignment = PP_ALIGN.CENTER

    out_path = os.path.join(BASE_DIR, 'dvs_noise_inverse_aa_figures.pptx')
    prs.save(out_path)
    print(f"Saved: {out_path}")


if __name__ == '__main__':
    main()
