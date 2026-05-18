#!/usr/bin/env python3
"""Generate English figures for DVS × Noise Inverse Problem review as PPTX."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch
import os


def create_gap_map_figure_en():
    """Create English version of the 4-domain, 5-gap map (G1-G5)."""
    fig, ax = plt.subplots(1, 1, figsize=(14, 9))

    domains = {
        'A': {'pos': (0.05, 0.72), 'label': 'A. DVS Noise\nPhysical Modeling\n(5 papers)', 'color': '#4ECDC4'},
        'B': {'pos': (0.37, 0.72), 'label': 'B. DVS Noise\nFiltering Methods\n(7 papers)', 'color': '#45B7D1'},
        'C': {'pos': (0.05, 0.12), 'label': 'C. DVS Astronomy\n& Space Apps\n(5 papers)', 'color': '#96CEB4'},
        'D': {'pos': (0.37, 0.12), 'label': 'D. Noise Inverse\nProblem (non-DVS)\n(7 papers)', 'color': '#FFEAA7'},
    }

    box_w, box_h = 0.25, 0.2
    for key, d in domains.items():
        x, y = d['pos']
        rect = FancyBboxPatch((x, y), box_w, box_h,
                              boxstyle="round,pad=0.02",
                              facecolor=d['color'], edgecolor='#333',
                              linewidth=2, alpha=0.85,
                              transform=ax.transAxes)
        ax.add_patch(rect)
        ax.text(x + box_w/2, y + box_h/2, d['label'],
                ha='center', va='center', fontsize=10, fontweight='bold',
                transform=ax.transAxes)

    gaps = [
        {'pos': (0.68, 0.80), 'label': 'G1\nBayesian Inverse\nFormulation', 'color': '#FF6B6B'},
        {'pos': (0.68, 0.62), 'label': 'G2\nAuxiliary Channel\nNoise Prediction', 'color': '#C44D58'},
        {'pos': (0.68, 0.44), 'label': 'G3 \u2605\nAstronomy-specific\nPipeline', 'color': '#E74C3C'},
        {'pos': (0.68, 0.26), 'label': 'G4 \u2605\nTemplate-free\nObject Detection', 'color': '#C0392B'},
        {'pos': (0.68, 0.08), 'label': 'G5\nSciDVS+Large Tel.\nIntegrated Demo', 'color': '#8E44AD'},
    ]

    for g in gaps:
        x, y = g['pos']
        ellipse = mpatches.Ellipse((x + 0.12, y + 0.05), 0.24, 0.13,
                                   facecolor=g['color'], edgecolor='white',
                                   linewidth=2, alpha=0.9,
                                   transform=ax.transAxes)
        ax.add_patch(ellipse)
        ax.text(x + 0.12, y + 0.05, g['label'],
                ha='center', va='center', fontsize=8, color='white',
                fontweight='bold', transform=ax.transAxes)

    arrow_style = dict(arrowstyle='->', color='#555', lw=1.5)
    ax.annotate('', xy=(0.68, 0.85), xytext=(0.30, 0.85),
                xycoords='axes fraction', arrowprops=arrow_style)
    ax.annotate('', xy=(0.68, 0.67), xytext=(0.62, 0.30),
                xycoords='axes fraction', arrowprops=arrow_style)
    ax.annotate('', xy=(0.80, 0.44+0.13), xytext=(0.80, 0.62),
                xycoords='axes fraction', arrowprops=dict(arrowstyle='->', color='#C0392B', lw=2))
    ax.annotate('', xy=(0.80, 0.26+0.13), xytext=(0.80, 0.44),
                xycoords='axes fraction', arrowprops=dict(arrowstyle='->', color='#C0392B', lw=2))
    ax.annotate('', xy=(0.80, 0.08+0.13), xytext=(0.80, 0.26),
                xycoords='axes fraction', arrowprops=dict(arrowstyle='->', color='#8E44AD', lw=2))

    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.set_aspect('equal')
    ax.axis('off')
    ax.set_title('DVS \u00d7 Noise Inverse Problem: Four Domains and Five Gaps (G1\u2013G5)',
                 fontsize=14, fontweight='bold', pad=20)

    plt.tight_layout()
    out = 'fig1_gap_map_en.png'
    fig.savefig(out, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f'Saved: {out}')
    return out


def create_pipeline_figure_en():
    """Create English version of the G3 pipeline."""
    fig, ax = plt.subplots(1, 1, figsize=(14, 10))
    ax.set_xlim(0, 14)
    ax.set_ylim(0, 10)
    ax.axis('off')

    ax.text(7, 9.7, 'G3 Astronomy-specific DVS Noise Inverse Pipeline + G4 Template-free Detection',
            ha='center', va='center', fontsize=14, fontweight='bold')

    inputs = [
        (1.5, 8.8, 'DVS Main Channel\ne(t,x,y,p)', '#4ECDC4'),
        (5.5, 8.8, 'Auxiliary Channels\nT(t), a(t), I(t)', '#45B7D1'),
        (9.5, 8.8, 'Physics Model\n(A5 Pixel Model)', '#96CEB4'),
    ]
    for x, y, label, color in inputs:
        rect = FancyBboxPatch((x-1.3, y-0.4), 2.6, 0.8,
                              boxstyle="round,pad=0.1",
                              facecolor=color, edgecolor='#333', linewidth=1.5)
        ax.add_patch(rect)
        ax.text(x, y, label, ha='center', va='center', fontsize=9)

    stages = [
        (7, 7.2, 'Stage 1: Noise Forward Model Construction\n\u03bb_noise(x,y,t) = F(\u03b8, T(t), bias, I_bg)', '#FFF3CD', 10, 0.8),
        (7, 5.8, 'Stage 2: Bayesian Inverse Problem Solving (G1)\n\u03b8\u0302 = argmin_\u03b8 D(e_obs, F(\u03b8)) \u2014 MLE / Variational / DeepClean-type NN (G2)', '#D1ECF1', 10, 0.8),
        (7, 4.4, 'Stage 3: Residual Event Stream Generation\ne_residual = e_obs \u2296 F(\u03b8\u0302) \u2014 Probabilistic Thinning / Rate Subtraction', '#D4EDDA', 10, 0.8),
        (7, 3.0, 'Stage 4 (G4): Template-free Object Detection\nResidual anomaly detection \u2192 shift-and-stack \u2192 Statistical test \u2192 Cataloging', '#F8D7DA', 10, 0.8),
        (7, 1.6, 'Stage 5: Physical Validation\nPSD test / Injection-Recovery / Known Object / Blind test', '#E2D5F1', 10, 0.8),
    ]

    for x, y, label, color, w, h in stages:
        rect = FancyBboxPatch((x - w/2, y - h/2), w, h,
                              boxstyle="round,pad=0.05",
                              facecolor=color, edgecolor='#333', linewidth=1.5)
        ax.add_patch(rect)
        ax.text(x, y, label, ha='center', va='center', fontsize=9)

    for y_from, y_to in [(8.4, 7.6), (6.8, 6.2), (5.4, 4.8), (4.0, 3.4), (2.6, 2.0)]:
        ax.annotate('', xy=(7, y_to), xytext=(7, y_from),
                    arrowprops=dict(arrowstyle='->', color='#333', lw=2))

    for x_from in [1.5, 5.5, 9.5]:
        ax.annotate('', xy=(7, 7.6), xytext=(x_from, 8.4),
                    arrowprops=dict(arrowstyle='->', color='#666', lw=1.5))

    plt.tight_layout()
    out = 'fig2_g3_pipeline_en.png'
    fig.savefig(out, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f'Saved: {out}')
    return out


def create_pptx(fig_files):
    """Create English PPTX with figures."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    titles = [
        'Fig. 1: DVS \u00d7 Noise Inverse Problem \u2014 Four Domains and Five Gaps (G1\u2013G5)',
        'Fig. 2: G3 Astronomy Pipeline + G4 Template-free Detection',
    ]
    captions = [
        'Prior work categorized into four domains (A: DVS noise modeling, B: DVS noise filtering, '
        'C: DVS astronomical applications, D: Noise inverse problem in non-DVS fields). '
        'Five gaps (G1\u2013G5) identified. G3 (astronomy pipeline) and G4 (template-free detection) are the highest-priority unexplored areas.',
        'Four-stage architecture transplanting the LIGO noise reconstruction pipeline (Vajente et al. 2020) to DVS astronomy (G3). '
        'Stage 4 integrates G4 template-free detection\u2014solving for noise and discovering objects in residuals without signal templates. '
        'Inputs: DVS main channel, auxiliary sensors, and physics-based pixel model (A5).',
    ]

    for fig_path, title, caption in zip(fig_files, titles, captions):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        from PIL import Image
        img = Image.open(fig_path)
        img_w, img_h = img.size
        max_w = Inches(11)
        max_h = Inches(5.5)
        scale = min(max_w / Emu(int(img_w * 914400 / 200)),
                    max_h / Emu(int(img_h * 914400 / 200)))
        final_w = int(img_w * 914400 / 200 * scale)
        final_h = int(img_h * 914400 / 200 * scale)
        left = (prs.slide_width - final_w) // 2
        top = Inches(1.0)
        slide.shapes.add_picture(fig_path, left, top, final_w, final_h)

        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.8))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = caption
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        p2.alignment = PP_ALIGN.CENTER

    out = 'dvs_noise_inverse_problem_figures_en.pptx'
    prs.save(out)
    print(f'Saved: {out}')


if __name__ == '__main__':
    os.chdir(os.path.dirname(os.path.abspath(__file__)))
    fig1 = create_gap_map_figure_en()
    fig2 = create_pipeline_figure_en()
    create_pptx([fig1, fig2])
