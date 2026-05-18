#!/usr/bin/env python3
"""Generate figures for DVS × Noise Inverse Problem review as PPTX."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch
import numpy as np
import os


def create_gap_map_figure():
    """Create a visual map of the 4 research domains and 5 gaps (G1-G5)."""
    fig, ax = plt.subplots(1, 1, figsize=(14, 9))

    # Domain boxes
    domains = {
        'A': {'pos': (0.05, 0.72), 'label': 'A. DVSノイズ\n物理モデリング\n(5件)', 'color': '#4ECDC4'},
        'B': {'pos': (0.37, 0.72), 'label': 'B. DVSノイズ\nフィルタリング\n(7件)', 'color': '#45B7D1'},
        'C': {'pos': (0.05, 0.12), 'label': 'C. DVS天文・\n宇宙応用\n(5件)', 'color': '#96CEB4'},
        'D': {'pos': (0.37, 0.12), 'label': 'D. ノイズ逆問題\n(非DVS)\n(7件)', 'color': '#FFEAA7'},
    }

    box_w, box_h = 0.25, 0.2
    for key, d in domains.items():
        x, y = d['pos']
        rect = FancyBboxPatch((x, y), box_w, box_h,
                              boxstyle="round,pad=0.02",
                              facecolor=d['color'], edgecolor='#333',
                              linewidth=2, alpha=0.85,
                              transform=ax.transAxes)
        ax.add_patch(rect)
        ax.text(x + box_w/2, y + box_h/2, d['label'],
                ha='center', va='center', fontsize=10, fontweight='bold',
                transform=ax.transAxes, fontfamily='Noto Sans CJK JP')

    # Gap labels — 5 gaps
    gaps = [
        {'pos': (0.68, 0.80), 'label': 'G1\nベイズ逆問題\n定式化', 'color': '#FF6B6B'},
        {'pos': (0.68, 0.62), 'label': 'G2\n補助チャンネル\nノイズ予測', 'color': '#C44D58'},
        {'pos': (0.68, 0.44), 'label': 'G3 ★\n天文特化\nパイプライン', 'color': '#E74C3C'},
        {'pos': (0.68, 0.26), 'label': 'G4 ★\nテンプレートフリー\n天体検出', 'color': '#C0392B'},
        {'pos': (0.68, 0.08), 'label': 'G5\nSciDVS+大口径\n統合実証', 'color': '#8E44AD'},
    ]

    for g in gaps:
        x, y = g['pos']
        ellipse = mpatches.Ellipse((x + 0.12, y + 0.05), 0.24, 0.13,
                                   facecolor=g['color'], edgecolor='white',
                                   linewidth=2, alpha=0.9,
                                   transform=ax.transAxes)
        ax.add_patch(ellipse)
        ax.text(x + 0.12, y + 0.05, g['label'],
                ha='center', va='center', fontsize=8, color='white',
                fontweight='bold', transform=ax.transAxes,
                fontfamily='Noto Sans CJK JP')

    # Arrows from domains to gaps
    arrow_style = dict(arrowstyle='->', color='#555', lw=1.5)
    # A→G1
    ax.annotate('', xy=(0.68, 0.85), xytext=(0.30, 0.85),
                xycoords='axes fraction', arrowprops=arrow_style)
    # D→G2
    ax.annotate('', xy=(0.68, 0.67), xytext=(0.62, 0.30),
                xycoords='axes fraction', arrowprops=arrow_style)
    # G1+G2→G3
    ax.annotate('', xy=(0.80, 0.44+0.13), xytext=(0.80, 0.62),
                xycoords='axes fraction', arrowprops=dict(arrowstyle='->', color='#C0392B', lw=2))
    # G3→G4
    ax.annotate('', xy=(0.80, 0.26+0.13), xytext=(0.80, 0.44),
                xycoords='axes fraction', arrowprops=dict(arrowstyle='->', color='#C0392B', lw=2))
    # G4→G5
    ax.annotate('', xy=(0.80, 0.08+0.13), xytext=(0.80, 0.26),
                xycoords='axes fraction', arrowprops=dict(arrowstyle='->', color='#8E44AD', lw=2))

    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.set_aspect('equal')
    ax.axis('off')
    ax.set_title('DVS × ノイズ逆問題: 4領域と5つの未開拓ギャップ（G1–G5）',
                 fontsize=14, fontweight='bold', pad=20,
                 fontfamily='Noto Sans CJK JP')

    plt.tight_layout()
    out = 'fig1_gap_map.png'
    fig.savefig(out, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f'Saved: {out}')
    return out


def create_pipeline_figure():
    """Create the G3 pipeline architecture figure."""
    fig, ax = plt.subplots(1, 1, figsize=(14, 10))
    ax.set_xlim(0, 14)
    ax.set_ylim(0, 10)
    ax.axis('off')

    # Title
    ax.text(7, 9.7, 'G3 天文特化DVSノイズ逆問題パイプライン + G4 テンプレートフリー検出',
            ha='center', va='center', fontsize=14, fontweight='bold',
            fontfamily='Noto Sans CJK JP')

    # Input boxes (top)
    inputs = [
        (1.5, 8.8, 'DVS主チャンネル\ne(t,x,y,p)', '#4ECDC4'),
        (5.5, 8.8, '補助チャンネル\nT(t), a(t), I(t)', '#45B7D1'),
        (9.5, 8.8, '物理モデル\n(A5 ピクセルモデル)', '#96CEB4'),
    ]
    for x, y, label, color in inputs:
        rect = FancyBboxPatch((x-1.3, y-0.4), 2.6, 0.8,
                              boxstyle="round,pad=0.1",
                              facecolor=color, edgecolor='#333', linewidth=1.5)
        ax.add_patch(rect)
        ax.text(x, y, label, ha='center', va='center', fontsize=9,
                fontfamily='Noto Sans CJK JP')

    # Stage boxes
    stages = [
        (7, 7.2, 'Stage 1: ノイズフォワードモデル構築\nλ_noise(x,y,t) = F(θ, T(t), bias, I_bg)', '#FFF3CD', 10, 0.8),
        (7, 5.8, 'Stage 2: ノイズ逆問題求解\nθ̂ = argmin_θ D(e_obs, F(θ)) — MLE / 変分推論 / DeepClean型NN', '#D1ECF1', 10, 0.8),
        (7, 4.4, 'Stage 3: 残差イベントストリーム生成\ne_residual = e_obs ⊖ F(θ̂) — 確率的薄化 / レート引き算', '#D4EDDA', 10, 0.8),
        (7, 3.0, 'Stage 4 (G4): テンプレートフリー天体検出\n残差統計逸脱検出 → shift-and-stack → 統計検定 → カタログ化', '#F8D7DA', 10, 0.8),
        (7, 1.6, 'Stage 5: 物理的検証\nPSDテスト / 注入・回収テスト / 既知天体テスト / ブラインドテスト', '#E2D5F1', 10, 0.8),
    ]

    for x, y, label, color, w, h in stages:
        rect = FancyBboxPatch((x - w/2, y - h/2), w, h,
                              boxstyle="round,pad=0.05",
                              facecolor=color, edgecolor='#333', linewidth=1.5)
        ax.add_patch(rect)
        ax.text(x, y, label, ha='center', va='center', fontsize=9,
                fontfamily='Noto Sans CJK JP')

    # Arrows between stages
    for y_from, y_to in [(8.4, 7.6), (6.8, 6.2), (5.4, 4.8), (4.0, 3.4), (2.6, 2.0)]:
        ax.annotate('', xy=(7, y_to), xytext=(7, y_from),
                    arrowprops=dict(arrowstyle='->', color='#333', lw=2))

    # Input arrows to stage 1
    for x_from in [1.5, 5.5, 9.5]:
        ax.annotate('', xy=(7, 7.6), xytext=(x_from, 8.4),
                    arrowprops=dict(arrowstyle='->', color='#666', lw=1.5,
                                   connectionstyle='arc3,rad=0'))

    plt.tight_layout()
    out = 'fig2_g3_pipeline.png'
    fig.savefig(out, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f'Saved: {out}')
    return out


def create_pptx(fig_files):
    """Create PPTX with figures on separate slides."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    titles = [
        'Fig. 1: DVS × ノイズ逆問題 — 4領域と5つの未開拓ギャップ（G1–G5）',
        'Fig. 2: G3天文パイプライン + G4テンプレートフリー検出',
    ]
    captions = [
        '先行研究を4領域（A: DVSノイズ物理モデリング、B: DVSノイズフィルタリング、C: DVS天文応用、D: ノイズ逆問題）に分類し、'
        '5つの未開拓ギャップ（G1–G5）を同定。G3（天文パイプライン）とG4（テンプレートフリー検出）が最も探索価値の高い領域。',
        'LIGOのノイズ再構成パイプライン（Vajente et al. 2020）をDVS天文観測に移植する4段階アーキテクチャ（G3）。'
        'Stage 4ではG4のテンプレートフリー検出——ノイズを解いて残差から未知天体を発見——を統合。',
    ]

    for fig_path, title, caption in zip(fig_files, titles, captions):
        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)

        # Title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Image
        from PIL import Image
        img = Image.open(fig_path)
        img_w, img_h = img.size
        max_w = Inches(11)
        max_h = Inches(5.5)
        scale = min(max_w / Emu(int(img_w * 914400 / 200)),
                    max_h / Emu(int(img_h * 914400 / 200)))
        final_w = int(img_w * 914400 / 200 * scale)
        final_h = int(img_h * 914400 / 200 * scale)
        left = (prs.slide_width - final_w) // 2
        top = Inches(1.0)
        slide.shapes.add_picture(fig_path, left, top, final_w, final_h)

        # Caption
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.8))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = caption
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        p2.alignment = PP_ALIGN.CENTER

    out = 'dvs_noise_inverse_problem_figures.pptx'
    prs.save(out)
    print(f'Saved: {out}')


if __name__ == '__main__':
    os.chdir(os.path.dirname(os.path.abspath(__file__)))
    fig1 = create_gap_map_figure()
    fig2 = create_pipeline_figure()
    create_pptx([fig1, fig2])
