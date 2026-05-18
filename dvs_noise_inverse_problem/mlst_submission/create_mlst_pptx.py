#!/usr/bin/env python3
"""Generate English PPTX with all figures for MLST submission (1 figure per slide)."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

BASE_DIR = os.path.dirname(os.path.abspath(__file__))

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

FIGURES = [
    {
        'file': 'fig1_gap_map_en.png',
        'title': 'Figure 1: Gap Map',
        'caption': (
            'Four surveyed domains (A: DVS noise physics, B: DVS denoising, '
            'C: DVS astronomical applications, D: noise inverse problem methods) '
            'and five identified research gaps (G1\u2013G5) at their intersections.'
        ),
    },
    {
        'file': 'fig2_g3_pipeline_en.png',
        'title': 'Figure 2: G3 Pipeline Architecture',
        'caption': (
            'System architecture of the G3 astronomy-specific DVS noise inverse '
            'problem pipeline. Four stages: (1) noise forward model construction, '
            '(2) Bayesian inverse problem solution, (3) residual event stream '
            'generation, (4) astronomical calibration and verification.'
        ),
    },
    {
        'file': 'fig4_sn_improvement.png',
        'title': 'Figure 3: SNR Improvement Analysis',
        'caption': (
            'Signal-to-noise ratio improvement analysis. '
            '(a) Fano factor spatial map showing noise-dominated vs. signal pixels; '
            '(b) Temporal dynamics of noise and signal event rates; '
            '(c) Per-pixel SNR distribution before and after noise subtraction.'
        ),
    },
    {
        'file': 'fig3_noise_inverse_demo.png',
        'title': 'Figure 4: Proof-of-Concept Demonstration',
        'caption': (
            'Noise inverse problem demonstration on EBSSA data. '
            '(a) Raw event accumulation (1,800,674 events); '
            '(b) Estimated noise rate map; '
            '(c) Per-event noise probability distribution; '
            '(d) Residual events after 90.3% noise removal (175,261 events).'
        ),
    },
]


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]  # blank

    for fig in FIGURES:
        slide = prs.slides.add_slide(blank_layout)

        # Title at top
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.7))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = fig['title']
        p.font.size = Pt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Image centered
        img_path = os.path.join(BASE_DIR, fig['file'])
        if os.path.exists(img_path):
            # Scale to fit
            max_w = Inches(11.0)
            max_h = Inches(5.2)
            from PIL import Image
            with Image.open(img_path) as im:
                iw, ih = im.size
            aspect = iw / ih
            if aspect > (11.0 / 5.2):
                w = max_w
                h = int(w / aspect)
            else:
                h = max_h
                w = int(h * aspect)
            left = (SLIDE_WIDTH - w) // 2
            top = Inches(1.0)
            slide.shapes.add_picture(img_path, left, top, width=w, height=h)

        # Caption at bottom
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(6.4), Inches(12.3), Inches(1.0))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = fig['caption']
        p2.font.size = Pt(12)
        p2.alignment = PP_ALIGN.CENTER

    out_path = os.path.join(BASE_DIR, 'dvs_noise_inverse_mlst_figures.pptx')
    prs.save(out_path)
    print(f"Saved: {out_path}")


if __name__ == '__main__':
    main()
