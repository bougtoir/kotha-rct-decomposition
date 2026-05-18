#!/usr/bin/env python3
"""Generate English figures for Clinical Noise Inverse Problem concept document as PPTX."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch
import numpy as np
import os


def create_conceptual_framework_figure():
    """Fig 1: Conceptual framework — Y_obs = S + N decomposition."""
    fig, ax = plt.subplots(1, 1, figsize=(14, 8))
    ax.set_xlim(0, 14)
    ax.set_ylim(0, 8)
    ax.axis('off')

    ax.text(7, 7.6, 'Clinical Noise Inverse Problem: Conceptual Framework',
            ha='center', va='center', fontsize=16, fontweight='bold')

    # Observed outcome box (top center)
    rect = FancyBboxPatch((4.5, 6.2), 5, 0.9,
                          boxstyle="round,pad=0.1",
                          facecolor='#D5E8D4', edgecolor='#82B366',
                          linewidth=2)
    ax.add_patch(rect)
    ax.text(7, 6.65, r'$Y_{obs}(i)$ = Observed Clinical Outcome',
            ha='center', va='center', fontsize=12, fontweight='bold')

    # Equals sign
    ax.text(7, 5.7, '=', ha='center', va='center', fontsize=24, fontweight='bold',
            color='#333')

    # Signal box (left)
    rect_s = FancyBboxPatch((1, 4.3), 4.5, 1.0,
                            boxstyle="round,pad=0.1",
                            facecolor='#DAE8FC', edgecolor='#6C8EBF',
                            linewidth=2)
    ax.add_patch(rect_s)
    ax.text(3.25, 5.0, 'S(i) = Signal', ha='center', va='center',
            fontsize=13, fontweight='bold', color='#1A5276')
    ax.text(3.25, 4.6, 'True Treatment Effect\n(Primary Endpoint)',
            ha='center', va='center', fontsize=10, color='#333')

    # Plus sign
    ax.text(7, 4.8, '+', ha='center', va='center', fontsize=24, fontweight='bold',
            color='#333')

    # Noise box (right)
    rect_n = FancyBboxPatch((8.5, 4.3), 4.5, 1.0,
                            boxstyle="round,pad=0.1",
                            facecolor='#FFF2CC', edgecolor='#D6B656',
                            linewidth=2)
    ax.add_patch(rect_n)
    ax.text(10.75, 5.0, 'N(i) = Noise', ha='center', va='center',
            fontsize=13, fontweight='bold', color='#7D6608')
    ax.text(10.75, 4.6, 'Confounders + Random Variation',
            ha='center', va='center', fontsize=10, color='#333')

    # Noise decomposition (bottom section)
    noise_components = [
        (1.5, 2.5, r'$N_{conf}$' + '\nMeasured\nConfounders',
         '#FADBD8', '#E74C3C', 'Age, sex, BMI,\ncomorbidities', 'High'),
        (4.75, 2.5, r'$N_{bio}$' + '\nBiological\nVariability',
         '#D5F5E3', '#27AE60', 'Genetic variation,\ncircadian rhythm', 'Medium'),
        (8.0, 2.5, r'$N_{meas}$' + '\nMeasurement\nError',
         '#D6EAF8', '#2980B9', 'Instrument precision,\ninter-observer var.', 'Medium'),
        (11.25, 2.5, r'$N_{rand}$' + '\nIrreducible\nRandomness',
         '#E8DAEF', '#8E44AD', 'Stochastic biological\nprocesses', 'Low'),
    ]

    for x, y, label, facecolor, edgecolor, examples, modelability in noise_components:
        rect = FancyBboxPatch((x - 1.3, y - 0.8), 2.6, 1.6,
                              boxstyle="round,pad=0.08",
                              facecolor=facecolor, edgecolor=edgecolor,
                              linewidth=1.5)
        ax.add_patch(rect)
        ax.text(x, y + 0.3, label, ha='center', va='center',
                fontsize=9, fontweight='bold')
        ax.text(x, y - 0.35, examples, ha='center', va='center',
                fontsize=7, color='#555', style='italic')
        ax.text(x, y - 0.65, f'Modelability: {modelability}',
                ha='center', va='center', fontsize=7, fontweight='bold',
                color=edgecolor)

    # Arrows from noise box to components
    for x in [1.5, 4.75, 8.0, 11.25]:
        ax.annotate('', xy=(x, 3.3), xytext=(10.75, 4.3),
                    arrowprops=dict(arrowstyle='->', color='#D6B656',
                                    lw=1.5, connectionstyle='arc3,rad=0'))

    # Key insight box
    rect_k = FancyBboxPatch((2, 0.2), 10, 0.7,
                            boxstyle="round,pad=0.1",
                            facecolor='#FCF3CF', edgecolor='#F39C12',
                            linewidth=2)
    ax.add_patch(rect_k)
    ax.text(7, 0.55, 'Key Insight: If we can model and reproduce N(i), '
            'we can subtract it to isolate S(i) with higher precision',
            ha='center', va='center', fontsize=11, fontweight='bold',
            color='#7D6608')

    plt.tight_layout()
    out = 'fig1_conceptual_framework.png'
    fig.savefig(out, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f'Saved: {out}')
    return out


def create_pipeline_figure():
    """Fig 2: Four-phase Clinical Noise Inverse Pipeline."""
    fig, ax = plt.subplots(1, 1, figsize=(14, 10))
    ax.set_xlim(0, 14)
    ax.set_ylim(0, 10)
    ax.axis('off')

    ax.text(7, 9.6, 'Four-Phase Clinical Noise Inverse Pipeline',
            ha='center', va='center', fontsize=16, fontweight='bold')

    # Input boxes (top row)
    inputs = [
        (2.5, 8.6, 'Clinical Outcome Data\n(Primary Endpoint)', '#D5E8D4', '#82B366'),
        (7, 8.6, 'Auxiliary Data\n(Covariates, Biomarkers,\nSerial Labs, Vitals)', '#DAE8FC', '#6C8EBF'),
        (11.5, 8.6, 'Domain Knowledge\n(Causal DAG,\nPhysiological Models)', '#E8DAEF', '#8E44AD'),
    ]
    for x, y, label, fc, ec in inputs:
        rect = FancyBboxPatch((x - 1.8, y - 0.5), 3.6, 1.0,
                              boxstyle="round,pad=0.08",
                              facecolor=fc, edgecolor=ec, linewidth=1.5)
        ax.add_patch(rect)
        ax.text(x, y, label, ha='center', va='center', fontsize=9, fontweight='bold')

    # Four phases
    phases = [
        (7, 6.8,
         'Phase 1: Noise Forward Model Construction',
         'Build F(θ; X) predicting noise from auxiliary data\n'
         'Methods: GLM / GAM / ML with DAG-guided variable selection',
         '#FFF3CD', '#F39C12'),
        (7, 5.2,
         'Phase 2: Bayesian Noise Parameter Estimation',
         'θ̂ = argmin_θ D(Y_obs, S_model + F(θ))\n'
         'Methods: MLE / MCMC / Variational Inference / Amortized NN',
         '#D1ECF1', '#17A2B8'),
        (7, 3.6,
         'Phase 3: Residual Generation (Clean Signal)',
         'Ŝ(i) = Y_obs(i) − F(θ̂; X_i)\n'
         'Methods: Direct subtraction / Probabilistic thinning / Doubly-robust',
         '#D4EDDA', '#28A745'),
        (7, 2.0,
         'Phase 4: Hypothesis-Free Discovery',
         'Examine residuals for unexpected patterns\n'
         'Methods: Anomaly detection / Subgroup discovery / Temporal analysis',
         '#F8D7DA', '#DC3545'),
    ]

    for x, y, title, desc, fc, ec in phases:
        rect = FancyBboxPatch((x - 5.5, y - 0.6), 11, 1.2,
                              boxstyle="round,pad=0.08",
                              facecolor=fc, edgecolor=ec, linewidth=2)
        ax.add_patch(rect)
        ax.text(x, y + 0.25, title, ha='center', va='center',
                fontsize=11, fontweight='bold', color='#333')
        ax.text(x, y - 0.25, desc, ha='center', va='center',
                fontsize=9, color='#555')

    # Arrows between phases
    for y_from, y_to in [(8.1, 7.4), (6.2, 5.8), (4.6, 4.2), (3.0, 2.6)]:
        ax.annotate('', xy=(7, y_to), xytext=(7, y_from),
                    arrowprops=dict(arrowstyle='->', color='#333', lw=2))

    # Input arrows to Phase 1
    for x_from in [2.5, 7, 11.5]:
        ax.annotate('', xy=(7, 7.4), xytext=(x_from, 8.1),
                    arrowprops=dict(arrowstyle='->', color='#666', lw=1.5))

    # Feedback arrow (Phase 3 → Phase 2 for iterative refinement)
    ax.annotate('', xy=(12.8, 5.8), xytext=(12.8, 3.0),
                arrowprops=dict(arrowstyle='->', color='#999',
                                lw=1.5, linestyle='dashed',
                                connectionstyle='arc3,rad=-0.3'))
    ax.text(13.5, 4.4, 'Iterative\nRefinement', ha='center', va='center',
            fontsize=8, color='#999', style='italic')

    # Output box
    rect_out = FancyBboxPatch((2.5, 0.3), 9, 0.7,
                              boxstyle="round,pad=0.1",
                              facecolor='#E8F8F5', edgecolor='#1ABC9C',
                              linewidth=2)
    ax.add_patch(rect_out)
    ax.text(7, 0.65, 'Output: Clean treatment effect estimate + discovered subgroup effects + uncertainty quantification',
            ha='center', va='center', fontsize=10, fontweight='bold', color='#0E6655')
    ax.annotate('', xy=(7, 1.0), xytext=(7, 1.4),
                arrowprops=dict(arrowstyle='->', color='#333', lw=2))

    plt.tight_layout()
    out = 'fig2_pipeline.png'
    fig.savefig(out, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f'Saved: {out}')
    return out


def create_application_comparison_figure():
    """Fig 3: CNIP application to RCT vs Retrospective studies."""
    fig, axes = plt.subplots(1, 2, figsize=(16, 9))

    for ax in axes:
        ax.set_xlim(0, 8)
        ax.set_ylim(0, 9)
        ax.axis('off')

    # Left panel: RCT
    ax = axes[0]
    ax.text(4, 8.6, 'RCT Setting', ha='center', va='center',
            fontsize=15, fontweight='bold', color='#1A5276')

    # Randomization box
    rect = FancyBboxPatch((0.5, 7.2), 7, 0.9,
                          boxstyle="round,pad=0.08",
                          facecolor='#AED6F1', edgecolor='#2980B9', linewidth=2)
    ax.add_patch(rect)
    ax.text(4, 7.65, 'Randomization → Confounders balanced on average',
            ha='center', va='center', fontsize=10, fontweight='bold')

    # But...
    ax.text(4, 6.7, 'BUT individual-level noise remains large',
            ha='center', va='center', fontsize=10, color='#E74C3C',
            fontweight='bold')

    # Auxiliary channels in RCT
    channels_rct = [
        'Baseline covariates (age, BMI, ASA-PS)',
        'Serial biomarkers (CRP, lactate, BNP)',
        'Concomitant medications',
        'Site / operator variables',
    ]
    ax.text(4, 6.1, 'Auxiliary Channels (underutilized in standard analysis):',
            ha='center', va='center', fontsize=9, fontweight='bold', color='#2C3E50')
    for i, ch in enumerate(channels_rct):
        rect = FancyBboxPatch((0.8, 5.4 - i * 0.55), 6.4, 0.45,
                              boxstyle="round,pad=0.05",
                              facecolor='#EBF5FB', edgecolor='#5DADE2',
                              linewidth=1)
        ax.add_patch(rect)
        ax.text(4, 5.6 - i * 0.55, ch, ha='center', va='center', fontsize=9)

    # Benefit box
    rect = FancyBboxPatch((0.5, 1.8), 7, 1.5,
                          boxstyle="round,pad=0.1",
                          facecolor='#D5F5E3', edgecolor='#27AE60', linewidth=2)
    ax.add_patch(rect)
    ax.text(4, 2.95, 'CNIP Benefit in RCTs', ha='center', va='center',
            fontsize=11, fontweight='bold', color='#1E8449')
    benefits_rct = [
        '• Noise model ρ²=0.4 → n_eff = 1.67× n_actual',
        '• Faster interim analyses',
        '• Hypothesis-free subgroup discovery',
    ]
    for i, b in enumerate(benefits_rct):
        ax.text(4, 2.45 - i * 0.35, b, ha='center', va='center',
                fontsize=9, color='#333')

    # Sample size illustration
    ax.text(4, 0.9, 'Example: 600 patients + CNIP ≈ 1000 conventional',
            ha='center', va='center', fontsize=10, fontweight='bold',
            color='#E67E22',
            bbox=dict(boxstyle='round,pad=0.3', facecolor='#FEF9E7',
                      edgecolor='#F39C12'))

    # Right panel: Retrospective
    ax = axes[1]
    ax.text(4, 8.6, 'Retrospective / Observational Setting', ha='center', va='center',
            fontsize=15, fontweight='bold', color='#7D3C98')

    # Challenge box
    rect = FancyBboxPatch((0.5, 7.2), 7, 0.9,
                          boxstyle="round,pad=0.08",
                          facecolor='#FADBD8', edgecolor='#E74C3C', linewidth=2)
    ax.add_patch(rect)
    ax.text(4, 7.65, 'No randomization → Confounders entangled with treatment',
            ha='center', va='center', fontsize=10, fontweight='bold')

    ax.text(4, 6.7, 'Rich EHR data available as auxiliary channels',
            ha='center', va='center', fontsize=10, color='#27AE60',
            fontweight='bold')

    # Auxiliary channels in retrospective
    channels_retro = [
        'Serial lab values (hourly–daily)',
        'Continuous vital signs (min–hourly)',
        'Medication records (timestamped)',
        'Clinical notes (NLP-extracted features)',
    ]
    ax.text(4, 6.1, 'Auxiliary Channels (EHR-derived):',
            ha='center', va='center', fontsize=9, fontweight='bold', color='#2C3E50')
    for i, ch in enumerate(channels_retro):
        rect = FancyBboxPatch((0.8, 5.4 - i * 0.55), 6.4, 0.45,
                              boxstyle="round,pad=0.05",
                              facecolor='#F5EEF8', edgecolor='#AF7AC5',
                              linewidth=1)
        ax.add_patch(rect)
        ax.text(4, 5.6 - i * 0.55, ch, ha='center', va='center', fontsize=9)

    # Benefit box
    rect = FancyBboxPatch((0.5, 1.8), 7, 1.5,
                          boxstyle="round,pad=0.1",
                          facecolor='#E8DAEF', edgecolor='#8E44AD', linewidth=2)
    ax.add_patch(rect)
    ax.text(4, 2.95, 'CNIP Benefit in Observational Studies', ha='center', va='center',
            fontsize=11, fontweight='bold', color='#6C3483')
    benefits_retro = [
        '• Explicit noise model complements propensity scores',
        '• Doubly-robust estimation (noise + outcome model)',
        '• Discovers hidden confounders via residual analysis',
    ]
    for i, b in enumerate(benefits_retro):
        ax.text(4, 2.45 - i * 0.35, b, ha='center', va='center',
                fontsize=9, color='#333')

    # Comparison note
    ax.text(4, 0.9, 'Combines with PS / IV / DiD for stronger inference',
            ha='center', va='center', fontsize=10, fontweight='bold',
            color='#8E44AD',
            bbox=dict(boxstyle='round,pad=0.3', facecolor='#F5EEF8',
                      edgecolor='#8E44AD'))

    fig.suptitle('CNIP Application: RCT vs Retrospective Studies',
                 fontsize=18, fontweight='bold', y=0.98)
    plt.tight_layout(rect=[0, 0, 1, 0.96])
    out = 'fig3_rct_vs_retrospective.png'
    fig.savefig(out, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f'Saved: {out}')
    return out


def create_sample_size_figure():
    """Fig 4: Sample size reduction as a function of noise model accuracy."""
    fig, axes = plt.subplots(1, 2, figsize=(14, 6))

    # Left panel: Effective sample size multiplier
    ax = axes[0]
    rho2 = np.linspace(0, 0.8, 100)
    n_multiplier = 1 / (1 - rho2)

    ax.plot(rho2, n_multiplier, 'b-', linewidth=2.5)
    ax.fill_between(rho2, 1, n_multiplier, alpha=0.15, color='blue')

    # Highlight key points
    key_points = [
        (0.2, 1/(1-0.2), 'ρ²=0.2\n1.25×'),
        (0.4, 1/(1-0.4), 'ρ²=0.4\n1.67×'),
        (0.6, 1/(1-0.6), 'ρ²=0.6\n2.50×'),
    ]
    for rho, mult, label in key_points:
        ax.plot(rho, mult, 'ro', markersize=10, zorder=5)
        ax.annotate(label, xy=(rho, mult), xytext=(rho + 0.05, mult + 0.15),
                    fontsize=10, fontweight='bold', color='#C0392B',
                    arrowprops=dict(arrowstyle='->', color='#C0392B', lw=1.5))

    ax.set_xlabel('Noise Model Accuracy (ρ²)', fontsize=12, fontweight='bold')
    ax.set_ylabel('Effective Sample Size Multiplier', fontsize=12, fontweight='bold')
    ax.set_title('Effective Sample Size Gain\nfrom Noise Modeling', fontsize=13, fontweight='bold')
    ax.set_xlim(0, 0.8)
    ax.set_ylim(1, 5.5)
    ax.grid(True, alpha=0.3)
    ax.axhline(y=1, color='gray', linestyle='--', alpha=0.5)

    # Right panel: Practical example — time to significance
    ax = axes[1]
    n_conventional = np.array([200, 400, 600, 800, 1000, 1200, 1500])
    power_conventional = 1 - np.exp(-n_conventional / 500)

    rho2_values = [0, 0.2, 0.4, 0.6]
    colors = ['#95A5A6', '#3498DB', '#E67E22', '#E74C3C']
    labels_r = ['Conventional (ρ²=0)', 'CNIP (ρ²=0.2)', 'CNIP (ρ²=0.4)', 'CNIP (ρ²=0.6)']

    for rho2_val, color, label in zip(rho2_values, colors, labels_r):
        effective_n = n_conventional / (1 - rho2_val) if rho2_val > 0 else n_conventional
        power = 1 - np.exp(-effective_n / 500)
        power = np.minimum(power, 0.999)
        ax.plot(n_conventional, power * 100, '-o', color=color, linewidth=2,
                markersize=5, label=label)

    ax.axhline(y=80, color='red', linestyle='--', alpha=0.5, linewidth=1.5)
    ax.text(1400, 81.5, '80% power', fontsize=9, color='red', fontweight='bold')

    ax.set_xlabel('Actual Sample Size (n)', fontsize=12, fontweight='bold')
    ax.set_ylabel('Statistical Power (%)', fontsize=12, fontweight='bold')
    ax.set_title('Power Curves: CNIP vs Conventional\n(Illustrative)', fontsize=13, fontweight='bold')
    ax.set_xlim(150, 1550)
    ax.set_ylim(0, 105)
    ax.legend(loc='lower right', fontsize=9, framealpha=0.9)
    ax.grid(True, alpha=0.3)

    plt.tight_layout()
    out = 'fig4_sample_size_reduction.png'
    fig.savefig(out, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f'Saved: {out}')
    return out


def create_pptx(fig_files):
    """Create English PPTX with all figures."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    titles = [
        'Fig. 1: Clinical Noise Inverse Problem — Conceptual Framework',
        'Fig. 2: Four-Phase Clinical Noise Inverse Pipeline',
        'Fig. 3: CNIP Application — RCT vs Retrospective Studies',
        'Fig. 4: Sample Size Reduction from Noise Modeling',
    ]
    captions = [
        'The observed clinical outcome Y_obs(i) is decomposed into signal S(i) '
        '(true treatment effect on the primary endpoint) and noise N(i) '
        '(confounders, biological variability, measurement error, and irreducible randomness). '
        'Each noise component has varying degrees of modelability using auxiliary clinical data.',

        'Four-phase pipeline for clinical noise inverse problem solving. '
        'Phase 1 constructs a forward model of noise from auxiliary data (covariates, biomarkers, serial measurements). '
        'Phase 2 estimates noise parameters via Bayesian inference. '
        'Phase 3 generates residuals by subtracting estimated noise. '
        'Phase 4 performs hypothesis-free discovery in the cleaned residuals.',

        'Left: In RCTs, randomization balances confounders on average, but individual-level noise remains. '
        'CNIP exploits underutilized auxiliary data to model this noise, effectively increasing sample size. '
        'Right: In retrospective studies, CNIP provides explicit noise modeling that complements '
        'propensity scores and enables doubly-robust estimation.',

        'Left: Effective sample size multiplier as a function of noise model accuracy (ρ²). '
        'A noise model explaining 40% of outcome variance yields an effective 1.67× sample size gain. '
        'Right: Illustrative power curves showing how CNIP enables earlier achievement of '
        '80% statistical power compared to conventional analysis.',
    ]

    from PIL import Image
    for fig_path, title, caption in zip(fig_files, titles, captions):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2),
                                         Inches(12.3), Inches(0.6))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Image
        img = Image.open(fig_path)
        img_w, img_h = img.size
        max_w = Inches(12)
        max_h = Inches(5.5)
        scale = min(max_w / Emu(int(img_w * 914400 / 200)),
                    max_h / Emu(int(img_h * 914400 / 200)))
        final_w = int(img_w * 914400 / 200 * scale)
        final_h = int(img_h * 914400 / 200 * scale)
        left = (prs.slide_width - final_w) // 2
        top = Inches(0.9)
        slide.shapes.add_picture(fig_path, left, top, final_w, final_h)

        # Caption
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(6.5),
                                          Inches(12.3), Inches(0.9))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = caption
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        p2.alignment = PP_ALIGN.CENTER

    out = 'clinical_noise_inverse_problem_figures_en.pptx'
    prs.save(out)
    print(f'Saved: {out}')
    return out


if __name__ == '__main__':
    os.chdir(os.path.dirname(os.path.abspath(__file__)))
    fig1 = create_conceptual_framework_figure()
    fig2 = create_pipeline_figure()
    fig3 = create_application_comparison_figure()
    fig4 = create_sample_size_figure()
    create_pptx([fig1, fig2, fig3, fig4])
