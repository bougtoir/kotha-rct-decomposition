#!/usr/bin/env python3
"""Generate Japanese figures for Clinical Noise Inverse Problem concept document as PPTX."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch
import matplotlib.font_manager as fm
import numpy as np
import os


def setup_japanese_font():
    """Configure matplotlib for Japanese text rendering."""
    jp_fonts = [f.name for f in fm.fontManager.ttflist
                if any(kw in f.name.lower() for kw in ['noto sans cjk', 'noto sans jp', 'ipa', 'takao', 'gothic'])]
    if jp_fonts:
        plt.rcParams['font.family'] = jp_fonts[0]
    else:
        import subprocess
        subprocess.run(['sudo', 'apt-get', 'install', '-y', 'fonts-noto-cjk'], capture_output=True)
        fm.fontManager.addfont('/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc')
        plt.rcParams['font.family'] = 'Noto Sans CJK JP'
    plt.rcParams['axes.unicode_minus'] = False


def create_conceptual_framework_figure_ja():
    """Fig 1: 概念フレームワーク — Y_obs = S + N 分解."""
    fig, ax = plt.subplots(1, 1, figsize=(14, 8))
    ax.set_xlim(0, 14)
    ax.set_ylim(0, 8)
    ax.axis('off')

    ax.text(7, 7.6, '臨床ノイズ逆問題: 概念フレームワーク',
            ha='center', va='center', fontsize=16, fontweight='bold')

    rect = FancyBboxPatch((4.5, 6.2), 5, 0.9,
                          boxstyle="round,pad=0.1",
                          facecolor='#D5E8D4', edgecolor='#82B366', linewidth=2)
    ax.add_patch(rect)
    ax.text(7, 6.65, r'$Y_{obs}(i)$ = 観察された臨床アウトカム',
            ha='center', va='center', fontsize=12, fontweight='bold')

    ax.text(7, 5.7, '=', ha='center', va='center', fontsize=24, fontweight='bold', color='#333')

    rect_s = FancyBboxPatch((1, 4.3), 4.5, 1.0,
                            boxstyle="round,pad=0.1",
                            facecolor='#DAE8FC', edgecolor='#6C8EBF', linewidth=2)
    ax.add_patch(rect_s)
    ax.text(3.25, 5.0, 'S(i) = シグナル', ha='center', va='center',
            fontsize=13, fontweight='bold', color='#1A5276')
    ax.text(3.25, 4.6, '真の治療効果\n（主エンドポイント）',
            ha='center', va='center', fontsize=10, color='#333')

    ax.text(7, 4.8, '+', ha='center', va='center', fontsize=24, fontweight='bold', color='#333')

    rect_n = FancyBboxPatch((8.5, 4.3), 4.5, 1.0,
                            boxstyle="round,pad=0.1",
                            facecolor='#FFF2CC', edgecolor='#D6B656', linewidth=2)
    ax.add_patch(rect_n)
    ax.text(10.75, 5.0, 'N(i) = ノイズ', ha='center', va='center',
            fontsize=13, fontweight='bold', color='#7D6608')
    ax.text(10.75, 4.6, '交絡因子 + 偶然変動',
            ha='center', va='center', fontsize=10, color='#333')

    noise_components = [
        (1.5, 2.5, r'$N_{conf}$' + '\n測定可能な\n交絡因子',
         '#FADBD8', '#E74C3C', '年齢, 性別, BMI,\n合併症', '高'),
        (4.75, 2.5, r'$N_{bio}$' + '\n生物学的\n変動',
         '#D5F5E3', '#27AE60', '遺伝的多型,\n概日リズム', '中'),
        (8.0, 2.5, r'$N_{meas}$' + '\n測定\n誤差',
         '#D6EAF8', '#2980B9', '機器精度,\n観察者間変動', '中'),
        (11.25, 2.5, r'$N_{rand}$' + '\n不可避の\nランダムネス',
         '#E8DAEF', '#8E44AD', '確率的生物学的\nプロセス', '低'),
    ]

    for x, y, label, facecolor, edgecolor, examples, modelability in noise_components:
        rect = FancyBboxPatch((x - 1.3, y - 0.8), 2.6, 1.6,
                              boxstyle="round,pad=0.08",
                              facecolor=facecolor, edgecolor=edgecolor, linewidth=1.5)
        ax.add_patch(rect)
        ax.text(x, y + 0.3, label, ha='center', va='center', fontsize=9, fontweight='bold')
        ax.text(x, y - 0.35, examples, ha='center', va='center',
                fontsize=7, color='#555', style='italic')
        ax.text(x, y - 0.65, f'モデル化可能性: {modelability}',
                ha='center', va='center', fontsize=7, fontweight='bold', color=edgecolor)

    for x in [1.5, 4.75, 8.0, 11.25]:
        ax.annotate('', xy=(x, 3.3), xytext=(10.75, 4.3),
                    arrowprops=dict(arrowstyle='->', color='#D6B656', lw=1.5,
                                    connectionstyle='arc3,rad=0'))

    rect_k = FancyBboxPatch((2, 0.2), 10, 0.7,
                            boxstyle="round,pad=0.1",
                            facecolor='#FCF3CF', edgecolor='#F39C12', linewidth=2)
    ax.add_patch(rect_k)
    ax.text(7, 0.55, '核心的洞察: N(i)をモデル化・再現できれば、'
            'それを差し引いてS(i)をより高精度に分離できる',
            ha='center', va='center', fontsize=11, fontweight='bold', color='#7D6608')

    plt.tight_layout()
    out = 'fig1_conceptual_framework_ja.png'
    fig.savefig(out, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f'Saved: {out}')
    return out


def create_pipeline_figure_ja():
    """Fig 2: 4フェーズ臨床ノイズ逆問題パイプライン."""
    fig, ax = plt.subplots(1, 1, figsize=(14, 10))
    ax.set_xlim(0, 14)
    ax.set_ylim(0, 10)
    ax.axis('off')

    ax.text(7, 9.6, '4フェーズ臨床ノイズ逆問題パイプライン',
            ha='center', va='center', fontsize=16, fontweight='bold')

    inputs = [
        (2.5, 8.6, '臨床アウトカムデータ\n（主エンドポイント）', '#D5E8D4', '#82B366'),
        (7, 8.6, '補助データ\n（共変量, バイオマーカー,\n経時的検査値, バイタル）', '#DAE8FC', '#6C8EBF'),
        (11.5, 8.6, 'ドメイン知識\n（因果DAG,\n生理学モデル）', '#E8DAEF', '#8E44AD'),
    ]
    for x, y, label, fc, ec in inputs:
        rect = FancyBboxPatch((x - 1.8, y - 0.5), 3.6, 1.0,
                              boxstyle="round,pad=0.08",
                              facecolor=fc, edgecolor=ec, linewidth=1.5)
        ax.add_patch(rect)
        ax.text(x, y, label, ha='center', va='center', fontsize=9, fontweight='bold')

    phases = [
        (7, 6.8,
         'Phase 1: ノイズ前方モデル構築',
         '補助データからノイズを予測するF(θ; X)を構築\n'
         '手法: GLM / GAM / 機械学習（DAGガイド付き変数選択）',
         '#FFF3CD', '#F39C12'),
        (7, 5.2,
         'Phase 2: ベイズノイズパラメータ推定',
         'θ̂ = argmin_θ D(Y_obs, S_model + F(θ))\n'
         '手法: MLE / MCMC / 変分推論 / 償却型ニューラルネットワーク',
         '#D1ECF1', '#17A2B8'),
        (7, 3.6,
         'Phase 3: 残差生成（クリーンシグナル抽出）',
         'Ŝ(i) = Y_obs(i) − F(θ̂; X_i)\n'
         '手法: 直接差引 / 確率的間引き / 二重頑健推定',
         '#D4EDDA', '#28A745'),
        (7, 2.0,
         'Phase 4: 仮説フリー発見',
         '残差中の予期しないパターンを探索\n'
         '手法: 異常検知 / サブグループ発見 / 時間的パターン分析',
         '#F8D7DA', '#DC3545'),
    ]

    for x, y, title, desc, fc, ec in phases:
        rect = FancyBboxPatch((x - 5.5, y - 0.6), 11, 1.2,
                              boxstyle="round,pad=0.08",
                              facecolor=fc, edgecolor=ec, linewidth=2)
        ax.add_patch(rect)
        ax.text(x, y + 0.25, title, ha='center', va='center',
                fontsize=11, fontweight='bold', color='#333')
        ax.text(x, y - 0.25, desc, ha='center', va='center', fontsize=9, color='#555')

    for y_from, y_to in [(8.1, 7.4), (6.2, 5.8), (4.6, 4.2), (3.0, 2.6)]:
        ax.annotate('', xy=(7, y_to), xytext=(7, y_from),
                    arrowprops=dict(arrowstyle='->', color='#333', lw=2))

    for x_from in [2.5, 7, 11.5]:
        ax.annotate('', xy=(7, 7.4), xytext=(x_from, 8.1),
                    arrowprops=dict(arrowstyle='->', color='#666', lw=1.5))

    ax.annotate('', xy=(12.8, 5.8), xytext=(12.8, 3.0),
                arrowprops=dict(arrowstyle='->', color='#999', lw=1.5, linestyle='dashed',
                                connectionstyle='arc3,rad=-0.3'))
    ax.text(13.5, 4.4, '反復\n精緻化', ha='center', va='center',
            fontsize=8, color='#999', style='italic')

    rect_out = FancyBboxPatch((2.5, 0.3), 9, 0.7,
                              boxstyle="round,pad=0.1",
                              facecolor='#E8F8F5', edgecolor='#1ABC9C', linewidth=2)
    ax.add_patch(rect_out)
    ax.text(7, 0.65, '出力: クリーンな治療効果推定 + 発見されたサブグループ効果 + 不確実性の定量化',
            ha='center', va='center', fontsize=10, fontweight='bold', color='#0E6655')
    ax.annotate('', xy=(7, 1.0), xytext=(7, 1.4),
                arrowprops=dict(arrowstyle='->', color='#333', lw=2))

    plt.tight_layout()
    out = 'fig2_pipeline_ja.png'
    fig.savefig(out, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f'Saved: {out}')
    return out


def create_application_comparison_figure_ja():
    """Fig 3: RCT vs 後ろ向き研究への適用比較."""
    fig, axes = plt.subplots(1, 2, figsize=(16, 9))

    for ax in axes:
        ax.set_xlim(0, 8)
        ax.set_ylim(0, 9)
        ax.axis('off')

    # 左パネル: RCT
    ax = axes[0]
    ax.text(4, 8.6, 'RCTセッティング', ha='center', va='center',
            fontsize=15, fontweight='bold', color='#1A5276')

    rect = FancyBboxPatch((0.5, 7.2), 7, 0.9,
                          boxstyle="round,pad=0.08",
                          facecolor='#AED6F1', edgecolor='#2980B9', linewidth=2)
    ax.add_patch(rect)
    ax.text(4, 7.65, 'ランダム化 → 交絡因子は平均的に均衡',
            ha='center', va='center', fontsize=10, fontweight='bold')

    ax.text(4, 6.7, 'しかし個人レベルのノイズは大きいまま',
            ha='center', va='center', fontsize=10, color='#E74C3C', fontweight='bold')

    channels_rct = [
        'ベースライン共変量（年齢, BMI, ASA-PS）',
        '経時的バイオマーカー（CRP, 乳酸, BNP）',
        '併用薬',
        '施設 / 術者変数',
    ]
    ax.text(4, 6.1, '補助チャンネル（標準解析では活用不足）:',
            ha='center', va='center', fontsize=9, fontweight='bold', color='#2C3E50')
    for i, ch in enumerate(channels_rct):
        rect = FancyBboxPatch((0.8, 5.4 - i * 0.55), 6.4, 0.45,
                              boxstyle="round,pad=0.05",
                              facecolor='#EBF5FB', edgecolor='#5DADE2', linewidth=1)
        ax.add_patch(rect)
        ax.text(4, 5.6 - i * 0.55, ch, ha='center', va='center', fontsize=9)

    rect = FancyBboxPatch((0.5, 1.8), 7, 1.5,
                          boxstyle="round,pad=0.1",
                          facecolor='#D5F5E3', edgecolor='#27AE60', linewidth=2)
    ax.add_patch(rect)
    ax.text(4, 2.95, 'RCTにおけるCNIPの利点', ha='center', va='center',
            fontsize=11, fontweight='bold', color='#1E8449')
    benefits_rct = [
        '• ノイズモデル ρ²=0.4 → n_eff = 1.67× n_actual',
        '• より早い中間解析',
        '• 仮説フリーのサブグループ発見',
    ]
    for i, b in enumerate(benefits_rct):
        ax.text(4, 2.45 - i * 0.35, b, ha='center', va='center', fontsize=9, color='#333')

    ax.text(4, 0.9, '例: 600例 + CNIP ≈ 従来法1000例相当',
            ha='center', va='center', fontsize=10, fontweight='bold', color='#E67E22',
            bbox=dict(boxstyle='round,pad=0.3', facecolor='#FEF9E7', edgecolor='#F39C12'))

    # 右パネル: 後ろ向き研究
    ax = axes[1]
    ax.text(4, 8.6, '後ろ向き / 観察研究セッティング', ha='center', va='center',
            fontsize=15, fontweight='bold', color='#7D3C98')

    rect = FancyBboxPatch((0.5, 7.2), 7, 0.9,
                          boxstyle="round,pad=0.08",
                          facecolor='#FADBD8', edgecolor='#E74C3C', linewidth=2)
    ax.add_patch(rect)
    ax.text(4, 7.65, 'ランダム化なし → 交絡因子が治療と絡み合う',
            ha='center', va='center', fontsize=10, fontweight='bold')

    ax.text(4, 6.7, '豊富なEHRデータが補助チャンネルとして利用可能',
            ha='center', va='center', fontsize=10, color='#27AE60', fontweight='bold')

    channels_retro = [
        '経時的検査値（時間〜日単位）',
        '連続バイタルサイン（分〜時間単位）',
        '投薬記録（タイムスタンプ付き）',
        '臨床記録（NLP抽出特徴量）',
    ]
    ax.text(4, 6.1, '補助チャンネル（EHR由来）:',
            ha='center', va='center', fontsize=9, fontweight='bold', color='#2C3E50')
    for i, ch in enumerate(channels_retro):
        rect = FancyBboxPatch((0.8, 5.4 - i * 0.55), 6.4, 0.45,
                              boxstyle="round,pad=0.05",
                              facecolor='#F5EEF8', edgecolor='#AF7AC5', linewidth=1)
        ax.add_patch(rect)
        ax.text(4, 5.6 - i * 0.55, ch, ha='center', va='center', fontsize=9)

    rect = FancyBboxPatch((0.5, 1.8), 7, 1.5,
                          boxstyle="round,pad=0.1",
                          facecolor='#E8DAEF', edgecolor='#8E44AD', linewidth=2)
    ax.add_patch(rect)
    ax.text(4, 2.95, '観察研究におけるCNIPの利点', ha='center', va='center',
            fontsize=11, fontweight='bold', color='#6C3483')
    benefits_retro = [
        '• 明示的ノイズモデルが傾向スコアを補完',
        '• 二重頑健推定（ノイズモデル＋アウトカムモデル）',
        '• 残差分析で隠れた交絡因子を発見',
    ]
    for i, b in enumerate(benefits_retro):
        ax.text(4, 2.45 - i * 0.35, b, ha='center', va='center', fontsize=9, color='#333')

    ax.text(4, 0.9, 'PS / IV / DiDと組み合わせてより強い推論',
            ha='center', va='center', fontsize=10, fontweight='bold', color='#8E44AD',
            bbox=dict(boxstyle='round,pad=0.3', facecolor='#F5EEF8', edgecolor='#8E44AD'))

    fig.suptitle('CNIP適用: RCT vs 後ろ向き研究',
                 fontsize=18, fontweight='bold', y=0.98)
    plt.tight_layout(rect=[0, 0, 1, 0.96])
    out = 'fig3_rct_vs_retrospective_ja.png'
    fig.savefig(out, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f'Saved: {out}')
    return out


def create_sample_size_figure_ja():
    """Fig 4: ノイズモデル精度に基づくサンプルサイズ縮小."""
    fig, axes = plt.subplots(1, 2, figsize=(14, 6))

    ax = axes[0]
    rho2 = np.linspace(0, 0.8, 100)
    n_multiplier = 1 / (1 - rho2)

    ax.plot(rho2, n_multiplier, 'b-', linewidth=2.5)
    ax.fill_between(rho2, 1, n_multiplier, alpha=0.15, color='blue')

    key_points = [
        (0.2, 1/(1-0.2), 'ρ²=0.2\n1.25倍'),
        (0.4, 1/(1-0.4), 'ρ²=0.4\n1.67倍'),
        (0.6, 1/(1-0.6), 'ρ²=0.6\n2.50倍'),
    ]
    for rho, mult, label in key_points:
        ax.plot(rho, mult, 'ro', markersize=10, zorder=5)
        ax.annotate(label, xy=(rho, mult), xytext=(rho + 0.05, mult + 0.15),
                    fontsize=10, fontweight='bold', color='#C0392B',
                    arrowprops=dict(arrowstyle='->', color='#C0392B', lw=1.5))

    ax.set_xlabel('ノイズモデル精度 (ρ²)', fontsize=12, fontweight='bold')
    ax.set_ylabel('有効サンプルサイズ倍率', fontsize=12, fontweight='bold')
    ax.set_title('ノイズモデリングによる\n有効サンプルサイズ増加', fontsize=13, fontweight='bold')
    ax.set_xlim(0, 0.8)
    ax.set_ylim(1, 5.5)
    ax.grid(True, alpha=0.3)
    ax.axhline(y=1, color='gray', linestyle='--', alpha=0.5)

    ax = axes[1]
    n_conventional = np.array([200, 400, 600, 800, 1000, 1200, 1500])
    rho2_values = [0, 0.2, 0.4, 0.6]
    colors = ['#95A5A6', '#3498DB', '#E67E22', '#E74C3C']
    labels_r = ['従来法 (ρ²=0)', 'CNIP (ρ²=0.2)', 'CNIP (ρ²=0.4)', 'CNIP (ρ²=0.6)']

    for rho2_val, color, label in zip(rho2_values, colors, labels_r):
        effective_n = n_conventional / (1 - rho2_val) if rho2_val > 0 else n_conventional
        power = 1 - np.exp(-effective_n / 500)
        power = np.minimum(power, 0.999)
        ax.plot(n_conventional, power * 100, '-o', color=color, linewidth=2,
                markersize=5, label=label)

    ax.axhline(y=80, color='red', linestyle='--', alpha=0.5, linewidth=1.5)
    ax.text(1400, 81.5, '80%検出力', fontsize=9, color='red', fontweight='bold')

    ax.set_xlabel('実際のサンプルサイズ (n)', fontsize=12, fontweight='bold')
    ax.set_ylabel('統計的検出力 (%)', fontsize=12, fontweight='bold')
    ax.set_title('検出力曲線: CNIP vs 従来法\n（概念図）', fontsize=13, fontweight='bold')
    ax.set_xlim(150, 1550)
    ax.set_ylim(0, 105)
    ax.legend(loc='lower right', fontsize=9, framealpha=0.9)
    ax.grid(True, alpha=0.3)

    plt.tight_layout()
    out = 'fig4_sample_size_reduction_ja.png'
    fig.savefig(out, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f'Saved: {out}')
    return out


def create_pptx_ja(fig_files):
    """Create Japanese PPTX with all figures."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    titles = [
        '図1: 臨床ノイズ逆問題 — 概念フレームワーク',
        '図2: 4フェーズ臨床ノイズ逆問題パイプライン',
        '図3: CNIP適用 — RCT vs 後ろ向き研究',
        '図4: ノイズモデリングによるサンプルサイズ縮小効果',
    ]
    captions = [
        '観察された臨床アウトカム Y_obs(i) をシグナル S(i)（主エンドポイントに対する真の治療効果）と'
        'ノイズ N(i)（交絡因子、生物学的変動、測定誤差、不可避のランダムネス）に分解する。'
        '各ノイズ成分は補助臨床データを用いたモデル化可能性が異なる。',

        '臨床ノイズ逆問題の4フェーズパイプライン。'
        'Phase 1: 補助データからノイズ前方モデルを構築。'
        'Phase 2: ベイズ推論によりノイズパラメータを推定。'
        'Phase 3: 推定ノイズを差し引いて残差を生成。'
        'Phase 4: クリーンな残差中で仮説フリーの発見を行う。',

        '左: RCTセッティング。ランダム化により交絡因子は平均的に均衡するが、個人レベルのノイズは残る。'
        'CNIPは標準解析で活用不足の補助データをモデル化し、有効サンプルサイズを増加させる。'
        '右: 後ろ向き/観察研究セッティング。CNIPは傾向スコアを補完する明示的ノイズモデリングと'
        '二重頑健推定を提供する。',

        '左: ノイズモデル精度（ρ²）に対する有効サンプルサイズ倍率。'
        'アウトカム分散の40%を説明するノイズモデル（ρ²=0.4）は有効1.67倍のサンプルサイズ増加をもたらす。'
        '右: CNIPにより従来法と比較して80%検出力をより早く達成できることを示す概念的な検出力曲線。',
    ]

    from PIL import Image
    for fig_path, title, caption in zip(fig_files, titles, captions):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        img = Image.open(fig_path)
        img_w, img_h = img.size
        max_w = Inches(12)
        max_h = Inches(5.5)
        scale = min(max_w / Emu(int(img_w * 914400 / 200)),
                    max_h / Emu(int(img_h * 914400 / 200)))
        final_w = int(img_w * 914400 / 200 * scale)
        final_h = int(img_h * 914400 / 200 * scale)
        left = (prs.slide_width - final_w) // 2
        top = Inches(0.9)
        slide.shapes.add_picture(fig_path, left, top, final_w, final_h)

        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.9))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = caption
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        p2.alignment = PP_ALIGN.CENTER

    out = 'clinical_noise_inverse_problem_figures_ja.pptx'
    prs.save(out)
    print(f'Saved: {out}')
    return out


if __name__ == '__main__':
    os.chdir(os.path.dirname(os.path.abspath(__file__)))
    setup_japanese_font()
    fig1 = create_conceptual_framework_figure_ja()
    fig2 = create_pipeline_figure_ja()
    fig3 = create_application_comparison_figure_ja()
    fig4 = create_sample_size_figure_ja()
    create_pptx_ja([fig1, fig2, fig3, fig4])
