#!/usr/bin/env python3
"""
Generate editable PowerPoint with all figures for Water Policy submission.
One figure per slide, widescreen format (13.333 x 7.5 inches).
Each slide: title at top, image centered, caption at bottom.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

OUTDIR = "/home/ubuntu/repos/wip/flood-prediction-review"
FIGDIR = os.path.join(OUTDIR, "figures")

prs = Presentation()

# Widescreen 13.333 x 7.5 inches
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Blank layout
blank_layout = prs.slide_layouts[6]  # Blank

figures = [
    {
        "filename": "fig1_energy_balance_en.png",
        "title": "Figure 1",
        "caption": (
            "Energy balance contour map showing generated power (MW) as a "
            "function of effective head and discharge. The red dashed line "
            "indicates the pumping power requirement (3.68 MW). The green "
            "marker shows the reference scenario (H = 100 m, Q = 50 m\u00b3/s)."
        ),
    },
    {
        "filename": "fig2_aquifer_storage_en.png",
        "title": "Figure 2",
        "caption": (
            "Aquifer storage capacity contours (10\u2076 m\u00b3) as a "
            "function of recharge area and water table drawdown, assuming a "
            "specific yield of 0.2. The dashed line indicates the 2018 flood "
            "volume (15.3 \u00d7 10\u2076 m\u00b3). The green marker shows "
            "the reference scenario."
        ),
    },
    {
        "filename": "fig3_hydrograph_en.png",
        "title": "Figure 3",
        "caption": (
            "Flood hydrograph comparison for the 2018 event under four "
            "scenarios: no measures (baseline), confluence relocation only, "
            "groundwater management only, and integrated measures. Upper "
            "panel shows the 72-hour rainfall hyetograph."
        ),
    },
    {
        "filename": "fig4_scenario_comparison_en.png",
        "title": "Figure 4",
        "caption": (
            "Comparison of peak discharge and total runoff volume across "
            "four scenarios. Error bars indicate \u00b110% sensitivity range."
        ),
    },
]

for fig in figures:
    slide = prs.slides.add_slide(blank_layout)

    # Title at top
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(12.333)
    height = Inches(0.6)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = fig["title"]
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.name = "Times New Roman"

    # Image centered
    fig_path = os.path.join(FIGDIR, fig["filename"])
    if os.path.exists(fig_path):
        # Read image to get aspect ratio
        from PIL import Image
        img = Image.open(fig_path)
        img_w, img_h = img.size
        aspect = img_w / img_h

        # Available area: 12.333 x 5.0 inches (leaving room for title/caption)
        max_w = Inches(11.0)
        max_h = Inches(4.8)

        if aspect > (11.0 / 4.8):
            # Width-constrained
            pic_w = max_w
            pic_h = int(max_w / aspect)
        else:
            # Height-constrained
            pic_h = max_h
            pic_w = int(max_h * aspect)

        # Center horizontally
        pic_left = int((Inches(13.333) - pic_w) / 2)
        pic_top = Inches(1.1)

        slide.shapes.add_picture(fig_path, pic_left, pic_top, pic_w, pic_h)
    else:
        # Placeholder text
        txBox2 = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(9), Inches(3))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = f"[Image: {fig['filename']}]"
        run2.font.size = Pt(18)
        run2.font.color.rgb = RGBColor(255, 0, 0)

    # Caption at bottom
    cap_left = Inches(1.0)
    cap_top = Inches(6.2)
    cap_width = Inches(11.333)
    cap_height = Inches(1.0)
    txBox3 = slide.shapes.add_textbox(cap_left, cap_top, cap_width, cap_height)
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = f"{fig['title']}. {fig['caption']}"
    run3.font.size = Pt(12)
    run3.font.name = "Times New Roman"

outpath = os.path.join(OUTDIR, "wp_figures_en.pptx")
prs.save(outpath)
print(f"Water Policy figures pptx saved: {outpath}")
