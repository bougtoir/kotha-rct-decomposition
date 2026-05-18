#!/usr/bin/env python3
"""
図表の編集可能pptxファイル生成（日本語版・英語版）
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN

FIGDIR = "/home/ubuntu/repos/wip/flood-prediction-review/figures"
OUTDIR = "/home/ubuntu/repos/wip/flood-prediction-review"

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

def create_pptx(lang, figures):
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for fig_path, title, caption in figures:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

        # Title at top
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Image centered
        if os.path.exists(fig_path):
            # Calculate proportional sizing
            max_w = Inches(11)
            max_h = Inches(5.5)
            pic = slide.shapes.add_picture(fig_path, Inches(1.2), Inches(1.0), max_w)
            # Scale to fit
            ratio = min(max_w / pic.width, max_h / pic.height)
            if ratio < 1:
                pic.width = int(pic.width * ratio)
                pic.height = int(pic.height * ratio)
            # Center horizontally
            pic.left = int((SLIDE_WIDTH - pic.width) / 2)

        # Caption at bottom
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.8))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = caption
        p2.font.size = Pt(14)
        p2.alignment = PP_ALIGN.CENTER

    return prs


# ── 日本語版 ──
figs_ja = [
    (f"{FIGDIR}/fig1_energy_balance_ja.png",
     "図-1 計画放水発電の出力と揚水所要出力の関係",
     "青破線は揚水所要出力3.7 MWの等出力線．破線より上の領域でエネルギー的に自立する．"),
    (f"{FIGDIR}/fig2_aquifer_storage_ja.png",
     "図-2 事前地下水位低下による帯水層貯留容量",
     "赤破線は2018年浸水量（15.3 x 10^6 m3）．バー上の数値は浸水量に対する割合．"),
    (f"{FIGDIR}/fig3_hydrograph_ja.png",
     "図-3 洪水シナリオ別ハイドログラフ比較（小田川流域）",
     "上段: 降雨強度，下段: 流出量．シナリオ1: 対策なし，2: 合流点付替え，3: 地下水管理型，4: 統合型．"),
    (f"{FIGDIR}/fig4_scenario_comparison_ja.png",
     "図-4 シナリオ別のピーク流量と総流出量の比較",
     "(a) ピーク流量の比較，(b) 総流出量の比較．数値は対策なしに対する削減率．"),
]

prs_ja = create_pptx("ja", figs_ja)
prs_ja.save(f"{OUTDIR}/figures_ja.pptx")
print(f"日本語版pptx保存: {OUTDIR}/figures_ja.pptx")

# ── 英語版 ──
figs_en = [
    (f"{FIGDIR}/fig1_energy_balance_en.png",
     "Fig. 1  Hydropower Output vs Pumping Requirement",
     "Blue dashed line: pumping power = 3.7 MW. Area above the line is energy-positive."),
    (f"{FIGDIR}/fig2_aquifer_storage_en.png",
     "Fig. 2  Aquifer Storage from Pre-Flood Drawdown",
     "Red dashed line: 2018 flood volume (15.3 x 10^6 m3). Percentages show fraction of flood volume absorbable."),
    (f"{FIGDIR}/fig3_hydrograph_en.png",
     "Fig. 3  Flood Hydrograph Comparison (Oda River Basin)",
     "Upper: rainfall intensity. Lower: discharge. Scenarios 1-4 as described in text."),
    (f"{FIGDIR}/fig4_scenario_comparison_en.png",
     "Fig. 4  Flood Control Effectiveness by Scenario",
     "(a) Peak discharge, (b) Total runoff volume. Percentages show reduction from no-action baseline."),
]

prs_en = create_pptx("en", figs_en)
prs_en.save(f"{OUTDIR}/figures_en.pptx")
print(f"英語版pptx保存: {OUTDIR}/figures_en.pptx")
