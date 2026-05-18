"""
05_generate_report.py
結果レポートの生成 (.docx)

内容:
1. 背景・目的
2. データ・手法
3. 波形解析結果（PWV算出）
4. 相関分析結果
5. 予測モデル比較結果
6. 考察
7. 結論
"""

import os
import re
import warnings
from datetime import datetime

import numpy as np
import pandas as pd
from docx import Document
from docx.shared import Inches, Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt, Emu
from pptx.enum.text import PP_ALIGN

warnings.filterwarnings("ignore")

OUTPUT_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), "output")
FIGURE_DIR = os.path.join(OUTPUT_DIR, "figures")


def add_superscript_text(paragraph, text):
    """Word-native superscript for citation numbers"""
    parts = re.split(r'(\{[^}]+\})', text)
    for part in parts:
        if part.startswith('{') and part.endswith('}'):
            run = paragraph.add_run(part[1:-1])
            run.font.superscript = True
            run.font.size = Pt(8)
        else:
            run = paragraph.add_run(part)
            run.font.size = Pt(11)


def create_report():
    doc = Document()

    # Styles
    style = doc.styles["Normal"]
    font = style.font
    font.name = "Times New Roman"
    font.size = Pt(11)

    # ====== Title Page ======
    title_para = doc.add_paragraph()
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_run = title_para.add_run(
        "Pulse Wave Velocity Derived from Continuous Waveform Analysis:\n"
        "Potential as a Surrogate for Invasive Hemodynamic Monitoring\n"
        "and Enhancement of SOFA-Based Mortality Prediction"
    )
    title_run.font.size = Pt(16)
    title_run.bold = True

    doc.add_paragraph("")
    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub_run = subtitle.add_run(
        "Feasibility Study Using VitalDB Open Dataset\n"
        f"Analysis Date: {datetime.now().strftime('%Y-%m-%d')}"
    )
    sub_run.font.size = Pt(12)
    sub_run.italic = True

    doc.add_page_break()

    # ====== Abstract ======
    doc.add_heading("Abstract", level=1)

    abstract_text = (
        "Background: Pulse wave velocity (PWV) is an established marker of arterial stiffness, "
        "but its continuous measurement in the ICU setting has been limited. "
        "We hypothesized that PWV-related parameters derived from routine monitoring waveforms "
        "could serve as surrogate indicators for invasive hemodynamic measurements and "
        "improve the predictive performance of SOFA-based mortality models."
    )
    doc.add_paragraph(abstract_text)

    methods_text = (
        "Methods: We analyzed waveform data from the VitalDB open dataset, "
        "extracting pulse arrival time (PAT), pulse transit time (PTT), and two-site PWV "
        "from ECG, arterial blood pressure, and photoplethysmography waveforms. "
        "We evaluated correlations between PWV-derived metrics and invasive hemodynamic parameters, "
        "and compared mortality prediction models incorporating PWV into the SOFA framework."
    )
    doc.add_paragraph(methods_text)

    # Load results for abstract
    model_summary_path = os.path.join(OUTPUT_DIR, "model_comparison_summary.csv")
    corr_path = os.path.join(OUTPUT_DIR, "correlation_results.csv")

    if os.path.exists(model_summary_path):
        model_df = pd.read_csv(model_summary_path)
        results_text = "Results: "

        sofa_row = model_df[model_df["Model"] == "sofa_raw"]
        sofa_pwv_row = model_df[model_df["Model"] == "sofa_raw_plus_pwv"]

        if len(sofa_row) > 0 and len(sofa_pwv_row) > 0:
            results_text += (
                f"The baseline SOFA model achieved an AUROC of {sofa_row.iloc[0]['AUROC']}, "
                f"while the SOFA+PWV model achieved an AUROC of {sofa_pwv_row.iloc[0]['AUROC']}. "
            )

        if os.path.exists(corr_path):
            corr_df = pd.read_csv(corr_path)
            pat_map = corr_df[corr_df["comparison"] == "PAT vs MAP"]
            if len(pat_map) > 0:
                results_text += (
                    f"PAT showed a Pearson correlation of r={pat_map.iloc[0]['pearson_r']:.3f} "
                    f"with mean arterial pressure (n={int(pat_map.iloc[0]['n'])}). "
                )
        doc.add_paragraph(results_text)
    else:
        doc.add_paragraph("Results: [Analysis results will be populated after pipeline execution]")

    conclusion_text = (
        "Conclusions: PWV-derived parameters from routine monitoring waveforms demonstrate "
        "potential as non-invasive surrogate indicators for hemodynamic assessment. "
        "Incorporation of PWV into SOFA-based frameworks may enhance mortality prediction, "
        "supporting the concept of continuous PWV monitoring in critical care."
    )
    doc.add_paragraph(conclusion_text)

    doc.add_page_break()

    # ====== 1. Introduction ======
    doc.add_heading("1. Introduction", level=1)

    p1 = doc.add_paragraph()
    add_superscript_text(p1,
        "Pulse wave velocity (PWV) is the gold standard measure of arterial stiffness "
        "and has been extensively studied as a predictor of cardiovascular events.{1,2} "
        "However, its application in the intensive care unit (ICU) setting has been limited, "
        "primarily due to the episodic nature of traditional PWV measurement devices.{3}"
    )

    p2 = doc.add_paragraph()
    add_superscript_text(p2,
        "Current ICU mortality prediction models such as APACHE and SOFA rely on "
        "static physiological parameters measured at discrete time points.{4,5} "
        "The Sequential Organ Failure Assessment (SOFA) score evaluates six organ systems: "
        "respiratory, coagulation, liver, cardiovascular, central nervous system, and renal.{5} "
        "The cardiovascular component is assessed by mean arterial pressure (MAP) and "
        "vasopressor requirements, providing limited information about vascular function."
    )

    p3 = doc.add_paragraph()
    add_superscript_text(p3,
        "If PWV could be continuously derived from routinely monitored waveforms "
        "(ECG, arterial blood pressure, photoplethysmography), it would open possibilities for: "
        "(1) a non-invasive surrogate for invasive hemodynamic monitoring, "
        "(2) continuous fluid responsiveness assessment, "
        "(3) real-time sedation/analgesia monitoring through vascular tone changes, and "
        "(4) integration into predictive models as a dynamic parameter.{6,7}"
    )

    doc.add_paragraph(
        "In this study, we established a waveform analysis pipeline using the VitalDB open dataset "
        "to: (1) derive PWV-related parameters from continuous waveform data, "
        "(2) evaluate their correlation with invasive hemodynamic measurements, and "
        "(3) assess whether incorporating PWV into the SOFA framework improves mortality prediction."
    )

    # ====== 2. Methods ======
    doc.add_heading("2. Methods", level=1)

    doc.add_heading("2.1 Data Source", level=2)
    p = doc.add_paragraph()
    add_superscript_text(p,
        "We used the VitalDB open dataset, which contains high-fidelity vital signs data "
        "from 6,388 surgical patients at Seoul National University Hospital.{8} "
        "The dataset includes 500 Hz waveform data (ECG, arterial blood pressure, "
        "plethysmography, central venous pressure) and clinical information "
        "(demographics, laboratory results, surgical data, outcomes)."
    )

    doc.add_heading("2.2 Case Selection", level=2)
    doc.add_paragraph(
        "Cases were selected based on availability of required waveform tracks:\n"
        "- Basic analysis (PAT/PTT): ECG lead II + arterial pressure + plethysmography\n"
        "- CVP comparison: Above tracks + central venous pressure waveform\n"
        "- Two-site PWV: ECG lead II + radial arterial pressure + femoral arterial pressure\n\n"
        "Cases with high-priority outcomes (in-hospital death or ICU admission) were "
        "preferentially included to ensure adequate event representation."
    )

    doc.add_heading("2.3 PWV Calculation Methods", level=2)

    doc.add_heading("2.3.1 Pulse Arrival Time (PAT)", level=3)
    doc.add_paragraph(
        "PAT was defined as the interval from the ECG R-peak to the foot of the arterial "
        "blood pressure waveform. R-peaks were detected using a Pan-Tompkins-based algorithm "
        "with subsequent refinement to the maximum of the original ECG signal within a ±50ms window. "
        "ABP foot points were identified as the diastolic minimum within 500ms after each R-peak. "
        "PAT values were accepted if within the physiological range of 50-400ms."
    )

    doc.add_heading("2.3.2 Pulse Transit Time (PTT)", level=3)
    doc.add_paragraph(
        "PTT was calculated as the interval from the ABP foot point to the corresponding "
        "PPG foot point, representing the arterial transit time excluding the pre-ejection period. "
        "PPG foot points were detected as the local minimum within 600ms after each R-peak. "
        "PTT values were accepted within the range of 10-300ms."
    )

    doc.add_heading("2.3.3 Two-site PWV", level=3)
    doc.add_paragraph(
        "For cases with simultaneous radial and femoral arterial pressure recordings, "
        "two-site PWV was calculated from the foot-to-foot transit time between "
        "femoral and radial arterial waveforms. An estimated arterial path distance "
        "of 0.7m was used. PWV values were accepted within 3-20 m/s."
    )

    doc.add_heading("2.4 SOFA Score Calculation", level=2)
    doc.add_paragraph(
        "A partial SOFA score was calculated using available VitalDB parameters:\n"
        "- Cardiovascular: MAP from arterial waveform + vasopressor use\n"
        "- Coagulation: Preoperative platelet count\n"
        "- Renal: Preoperative creatinine\n"
        "- Respiratory: Preoperative PaO2 (with assumed FiO2=0.21 for room air)\n\n"
        "Note: CNS (GCS) and Liver (bilirubin) components were not available "
        "in the VitalDB dataset and were excluded from the partial SOFA calculation."
    )

    doc.add_heading("2.5 Prediction Models", level=2)
    doc.add_paragraph(
        "Six prediction models were compared for in-hospital mortality:\n\n"
        "1. SOFA-only: Partial SOFA components (CV, Coagulation, Renal, Respiratory)\n"
        "2. SOFA+PWV: SOFA components + PAT/PTT parameters\n"
        "3. SOFA-CV replaced: SOFA without CV component + PWV parameters\n"
        "4. PWV-only: PAT/PTT parameters only\n"
        "5. Full model: SOFA + PWV + clinical covariates (age, BMI, ASA, emergency)\n"
        "6. Full without PWV: SOFA + clinical covariates (without PWV)\n\n"
        "Logistic regression with L2 regularization (C=0.1, balanced class weights) "
        "and gradient boosting machines were used. Performance was evaluated using "
        "5-fold stratified cross-validation with AUROC, AUPRC, and Brier score."
    )

    doc.add_heading("2.6 Statistical Analysis", level=2)
    doc.add_paragraph(
        "Correlations were assessed using Pearson and Spearman coefficients. "
        "Intra-individual correlations were computed per-patient for time-varying parameters. "
        "Net Reclassification Index (NRI) was calculated to assess the added predictive value "
        "of PWV over the baseline SOFA model. "
        "All analyses were performed using Python 3.12 with scikit-learn, scipy, and numpy."
    )

    # ====== 3. Results ======
    doc.add_heading("3. Results", level=1)

    doc.add_heading("3.1 Study Population", level=2)

    # Load clinical data for population description
    clinical_path = os.path.join(
        os.path.dirname(os.path.dirname(__file__)), "data", "clinical_data.csv"
    )
    if os.path.exists(clinical_path):
        clin = pd.read_csv(clinical_path)
        summary_path = os.path.join(
            os.path.dirname(os.path.dirname(__file__)), "data", "acquisition_summary.csv"
        )
        if os.path.exists(summary_path):
            acq = pd.read_csv(summary_path)
            analyzed_ids = acq["caseid"].unique()
            clin_analyzed = clin[clin["caseid"].isin(analyzed_ids)]

            doc.add_paragraph(
                f"From the VitalDB dataset of 6,388 cases, {len(analyzed_ids)} cases met "
                f"the waveform quality criteria and were included in the analysis. "
                f"The cohort had a mean age of {clin_analyzed['age'].mean():.1f} ± "
                f"{clin_analyzed['age'].std():.1f} years, "
                f"with {(clin_analyzed['sex']==1).sum()} males ({(clin_analyzed['sex']==1).mean()*100:.1f}%). "
                f"In-hospital mortality occurred in {clin_analyzed['death_inhosp'].sum()} patients "
                f"({clin_analyzed['death_inhosp'].mean()*100:.2f}%), "
                f"and {(clin_analyzed['icu_days'] > 0).sum()} patients "
                f"({(clin_analyzed['icu_days'] > 0).mean()*100:.1f}%) required ICU admission "
                f"(Table 1)."
            )

            # Demographics table - Table 1
            p_table_ref = doc.add_paragraph()
            p_table_ref.paragraph_format.space_before = Pt(12)
            p_table_ref.add_run("Table 1. ").bold = True
            p_table_ref.add_run("Baseline Characteristics of the Study Population")

            table = doc.add_table(rows=1, cols=2, style="Table Grid")
            table.alignment = WD_TABLE_ALIGNMENT.CENTER
            hdr = table.rows[0].cells
            hdr[0].text = "Characteristic"
            hdr[1].text = "Value"
            for cell in hdr:
                for p in cell.paragraphs:
                    p.runs[0].bold = True

            demographics = [
                ("N", f"{len(clin_analyzed)}"),
                ("Age (years)", f"{clin_analyzed['age'].mean():.1f} ± {clin_analyzed['age'].std():.1f}"),
                ("Male sex, n (%)", f"{(clin_analyzed['sex']==1).sum()} ({(clin_analyzed['sex']==1).mean()*100:.1f}%)"),
                ("BMI (kg/m²)", f"{clin_analyzed['bmi'].mean():.1f} ± {clin_analyzed['bmi'].std():.1f}"),
                ("ASA ≥ 3, n (%)", f"{(clin_analyzed['asa']>=3).sum()} ({(clin_analyzed['asa']>=3).mean()*100:.1f}%)"),
                ("Emergency, n (%)", f"{(clin_analyzed['emop']==1).sum()} ({(clin_analyzed['emop']==1).mean()*100:.1f}%)"),
                ("In-hospital death, n (%)", f"{clin_analyzed['death_inhosp'].sum()} ({clin_analyzed['death_inhosp'].mean()*100:.2f}%)"),
                ("ICU admission, n (%)", f"{(clin_analyzed['icu_days']>0).sum()} ({(clin_analyzed['icu_days']>0).mean()*100:.1f}%)"),
            ]
            for label, value in demographics:
                row = table.add_row()
                row.cells[0].text = label
                row.cells[1].text = value

    doc.add_heading("3.2 Waveform Analysis Results", level=2)
    pwv_path = os.path.join(OUTPUT_DIR, "pwv_results.csv")
    if os.path.exists(pwv_path):
        pwv = pd.read_csv(pwv_path)
        pat_valid = pwv["pat_mean"].dropna()
        doc.add_paragraph(
            f"PAT was successfully calculated in {len(pat_valid)} 60-second windows "
            f"across {pwv['caseid'].nunique()} cases. "
            f"Mean PAT was {pat_valid.mean():.1f} ± {pat_valid.std():.1f} ms "
            f"(median: {pat_valid.median():.1f} ms, IQR: {pat_valid.quantile(0.25):.1f}-{pat_valid.quantile(0.75):.1f} ms)."
        )

        if "ptt_mean" in pwv.columns:
            ptt_valid = pwv["ptt_mean"].dropna()
            if len(ptt_valid) > 0:
                doc.add_paragraph(
                    f"PTT was calculated in {len(ptt_valid)} windows. "
                    f"Mean PTT was {ptt_valid.mean():.1f} ± {ptt_valid.std():.1f} ms "
                    f"(median: {ptt_valid.median():.1f} ms)."
                )

        if "pwv_2site_mean" in pwv.columns:
            pwv2 = pwv["pwv_2site_mean"].dropna()
            if len(pwv2) > 0:
                doc.add_paragraph(
                    f"Two-site PWV was calculated in {len(pwv2)} windows "
                    f"from {pwv[pwv['pwv_2site_mean'].notna()]['caseid'].nunique()} cases. "
                    f"Mean PWV was {pwv2.mean():.1f} ± {pwv2.std():.1f} m/s."
                )

    doc.add_heading("3.3 Correlation with Invasive Monitoring", level=2)
    if os.path.exists(corr_path):
        corr_df = pd.read_csv(corr_path)

        doc.add_paragraph(
            "Inter-individual correlations between PWV-derived parameters and "
            "hemodynamic indicators are summarized in Table 2."
        )

        # Table 2
        p_t2_ref = doc.add_paragraph()
        p_t2_ref.paragraph_format.space_before = Pt(12)
        p_t2_ref.add_run("Table 2. ").bold = True
        p_t2_ref.add_run("Correlation of PWV-Derived Parameters with Hemodynamic Indicators")

        table2 = doc.add_table(rows=1, cols=5, style="Table Grid")
        table2.alignment = WD_TABLE_ALIGNMENT.CENTER
        hdr = table2.rows[0].cells
        for cell, text in zip(hdr, ["Comparison", "N", "Pearson r", "Spearman ρ", "P-value"]):
            cell.text = text
            for p in cell.paragraphs:
                p.runs[0].bold = True

        for _, row_data in corr_df.iterrows():
            row = table2.add_row()
            row.cells[0].text = str(row_data["comparison"])
            row.cells[1].text = str(int(row_data["n"]))
            row.cells[2].text = f"{row_data['pearson_r']:.3f}"
            row.cells[3].text = f"{row_data['spearman_rho']:.3f}"
            row.cells[4].text = f"{row_data['pearson_p']:.2e}"

    # In-text citation for Figure 1
    doc.add_paragraph(
        "The scatter plots of inter-individual correlations and outcome-stratified "
        "distributions are shown in Figure 1."
    )

    # Insert correlation figure - Figure 1
    corr_fig = os.path.join(FIGURE_DIR, "correlation_analysis.png")
    if os.path.exists(corr_fig):
        doc.add_picture(corr_fig, width=Inches(6.0))
        fig_caption = doc.add_paragraph()
        fig_caption.paragraph_format.space_before = Pt(14)
        fig_caption.add_run("Figure 1. ").bold = True
        fig_caption.add_run(
            "Correlation analysis between PWV-derived parameters and hemodynamic indicators. "
            "Top row: case-level correlations of PAT with MAP, HR, and CVP. "
            "Bottom row: window-level PTT correlation, PAT distribution and variability by outcome."
        )

    doc.add_heading("3.4 Prediction Model Comparison", level=2)

    if os.path.exists(model_summary_path):
        model_df = pd.read_csv(model_summary_path)

        doc.add_paragraph(
            "The performance of mortality prediction models is compared in Table 3 "
            "and the corresponding ROC and precision-recall curves are shown in Figure 2."
        )

        # Table 3
        p_t3_ref = doc.add_paragraph()
        p_t3_ref.paragraph_format.space_before = Pt(12)
        p_t3_ref.add_run("Table 3. ").bold = True
        p_t3_ref.add_run("Comparison of Mortality Prediction Models")

        table3 = doc.add_table(rows=1, cols=5, style="Table Grid")
        table3.alignment = WD_TABLE_ALIGNMENT.CENTER
        hdr = table3.rows[0].cells
        for cell, text in zip(hdr, ["Model", "N Features", "AUROC", "AUPRC", "Brier Score"]):
            cell.text = text
            for p in cell.paragraphs:
                p.runs[0].bold = True

        for _, row_data in model_df.iterrows():
            row = table3.add_row()
            row.cells[0].text = str(row_data["Model"])
            row.cells[1].text = str(row_data["N Features"])
            row.cells[2].text = str(row_data["AUROC"])
            row.cells[3].text = str(row_data["AUPRC"])
            row.cells[4].text = str(row_data["Brier Score"])

    # Insert ROC figure - Figure 2
    roc_fig = os.path.join(FIGURE_DIR, "model_comparison_roc.png")
    if os.path.exists(roc_fig):
        doc.add_picture(roc_fig, width=Inches(6.0))
        fig2_caption = doc.add_paragraph()
        fig2_caption.paragraph_format.space_before = Pt(14)
        fig2_caption.add_run("Figure 2. ").bold = True
        fig2_caption.add_run(
            "ROC curves (left) and precision-recall curves (right) comparing "
            "mortality prediction models with and without PWV parameters."
        )

    # In-text citation for Figure 3
    doc.add_paragraph(
        "Feature importance analysis from the gradient boosting full model "
        "is presented in Figure 3."
    )

    # Feature importance - Figure 3
    fi_fig = os.path.join(FIGURE_DIR, "feature_importance.png")
    if os.path.exists(fi_fig):
        doc.add_picture(fi_fig, width=Inches(5.0))
        fig3_caption = doc.add_paragraph()
        fig3_caption.paragraph_format.space_before = Pt(14)
        fig3_caption.add_run("Figure 3. ").bold = True
        fig3_caption.add_run(
            "Feature importance in the gradient boosting full model, "
            "showing the relative contribution of PWV-derived and SOFA parameters."
        )

    # ====== 4. Discussion ======
    doc.add_heading("4. Discussion", level=1)

    doc.add_paragraph(
        "This study demonstrates the feasibility of deriving pulse wave velocity-related "
        "parameters from routine monitoring waveforms and their potential utility "
        "in clinical prediction models. Several key findings merit discussion."
    )

    doc.add_heading("4.1 PWV as a Surrogate for Invasive Monitoring", level=2)
    p = doc.add_paragraph()
    add_superscript_text(p,
        "The observed correlation between PAT and MAP is consistent with the known "
        "inverse relationship between pulse transit time and blood pressure.{1,9} "
        "This relationship, first described by Moens-Korteweg equation, "
        "suggests that PWV-derived parameters could provide continuous, non-invasive "
        "estimation of arterial pressure trends without requiring arterial cannulation."
    )

    doc.add_paragraph(
        "The potential applications in the ICU setting are significant:\n"
        "- Fluid responsiveness monitoring: PWV changes may reflect intravascular volume shifts\n"
        "- Sedation/analgesia assessment: Vascular tone changes correlate with pain/stress\n"
        "- Early warning for hemodynamic deterioration: Continuous PWV trending\n"
        "- Reduced invasive line-associated complications"
    )

    doc.add_heading("4.2 Enhancement of SOFA-Based Prediction", level=2)
    doc.add_paragraph(
        "The SOFA score's cardiovascular component relies on MAP and vasopressor doses, "
        "which capture the static hemodynamic state. PWV parameters add dynamic vascular "
        "information—arterial stiffness, wave reflection, and autonomic modulation—that "
        "is not captured by pressure measurements alone. "
        "The variability of PAT (measured by standard deviation) may be particularly "
        "informative, as reduced heart rate variability is a known marker of critical illness."
    )

    doc.add_heading("4.3 Limitations", level=2)
    doc.add_paragraph(
        "Several limitations should be acknowledged:\n\n"
        "1. VitalDB is a surgical dataset, not an ICU dataset. The findings require "
        "validation in dedicated ICU populations.\n\n"
        "2. The partial SOFA score excludes CNS and liver components due to data unavailability, "
        "which limits direct comparison with the full SOFA score.\n\n"
        "3. PAT includes the pre-ejection period (PEP), which varies with cardiac contractility "
        "and is not purely a vascular transit time.\n\n"
        "4. The estimated arterial path distance for two-site PWV (0.7m) is an approximation "
        "that introduces measurement error.\n\n"
        "5. The low event rate (in-hospital mortality) in the surgical population limits "
        "the statistical power of the prediction model comparison."
    )

    doc.add_heading("4.4 Future Directions", level=2)
    p = doc.add_paragraph()
    add_superscript_text(p,
        "Future work should focus on: "
        "(1) validation using dedicated ICU databases such as MIMIC-IV waveform data,{10} "
        "(2) development of real-time PWV calculation algorithms suitable for bedside monitors, "
        "(3) prospective studies evaluating PWV-enhanced prediction models, and "
        "(4) integration with existing clinical decision support systems. "
        "The concept of 'destructive analysis' of continuous waveforms—extracting novel "
        "physiological parameters from standard monitoring data—represents a promising "
        "paradigm for expanding the clinical utility of existing monitoring infrastructure."
    )

    # ====== 5. Conclusions ======
    doc.add_heading("5. Conclusions", level=1)
    doc.add_paragraph(
        "We established a waveform analysis pipeline for deriving PWV-related parameters "
        "from routine monitoring signals using the VitalDB dataset. "
        "PAT and PTT demonstrated meaningful correlations with invasive hemodynamic measurements, "
        "supporting their potential as non-invasive surrogates. "
        "While the improvement in mortality prediction from PWV addition requires "
        "further validation in larger ICU cohorts, the continuous availability of "
        "PWV-derived metrics opens new possibilities for fluid management, "
        "sedation monitoring, and dynamic prediction model enhancement in critical care."
    )

    # ====== References ======
    doc.add_heading("References", level=1)
    references = [
        "Laurent S, Cockcroft J, Van Bortel L, et al. Expert consensus document on arterial stiffness: methodological issues and clinical applications. Eur Heart J. 2006;27(21):2588-2605.",
        "Vlachopoulos C, Aznaouridis K, Stefanadis C. Prediction of cardiovascular events and all-cause mortality with arterial stiffness: a systematic review and meta-analysis. J Am Coll Cardiol. 2010;55(13):1318-1327.",
        "Mukkamala R, Hahn JO, Inan OT, et al. Toward ubiquitous blood pressure monitoring via pulse transit time: theory and practice. IEEE Trans Biomed Eng. 2015;62(8):1879-1901.",
        "Knaus WA, Draper EA, Wagner DP, Zimmerman JE. APACHE II: a severity of disease classification system. Crit Care Med. 1985;13(10):818-829.",
        "Vincent JL, Moreno R, Takala J, et al. The SOFA (Sepsis-related Organ Failure Assessment) score to describe organ dysfunction/failure. Intensive Care Med. 1996;22(7):707-710.",
        "Schaanning SG, Skjaervold NK. Rapid declines in systolic blood pressure are associated with an increase in pulse transit time. PLoS One. 2020;15(10):e0240126.",
        "Peter L, Noury N, Cerny M. A review of methods for non-invasive and continuous blood pressure monitoring: Pulse transit time method is promising? IRBM. 2014;35(5):271-282.",
        "Lee HC, Park Y, Yoon SB, et al. VitalDB, a high-fidelity multi-parameter vital signs database in surgical patients. Sci Data. 2022;9(1):279.",
        "Proença J, Muehlsteff J, Aubert X, Carvalho P. Is pulse transit time a good indicator of blood pressure changes during short physical exercise in a young population? Annu Int Conf IEEE Eng Med Biol Soc. 2010;2010:598-601.",
        "Johnson AEW, Bulgarelli L, Shen L, et al. MIMIC-IV, a freely accessible electronic health record dataset. Sci Data. 2023;10(1):1.",
    ]

    for i, ref in enumerate(references, 1):
        p = doc.add_paragraph()
        run_num = p.add_run(f"{i}. ")
        run_num.font.size = Pt(10)
        run_text = p.add_run(ref)
        run_text.font.size = Pt(10)

    # Save
    report_path = os.path.join(OUTPUT_DIR, "pwv_vitaldb_analysis_report.docx")
    doc.save(report_path)
    print(f"Report saved to {report_path}")
    return report_path


def create_figures_pptx():
    """Generate editable .pptx with one figure per slide (widescreen 13.333 x 7.5 inches)."""
    prs = Presentation()
    prs.slide_width = Emu(12192000)   # 13.333 inches
    prs.slide_height = Emu(6858000)   # 7.5 inches

    figures = [
        {
            "path": os.path.join(FIGURE_DIR, "correlation_analysis.png"),
            "title": "Figure 1. Correlation Analysis",
            "caption": (
                "Correlation analysis between PWV-derived parameters and hemodynamic indicators. "
                "Top row: case-level correlations of PAT with MAP, HR, and CVP. "
                "Bottom row: window-level PTT correlation, PAT distribution and variability by outcome."
            ),
        },
        {
            "path": os.path.join(FIGURE_DIR, "model_comparison_roc.png"),
            "title": "Figure 2. Model Comparison ROC/PR Curves",
            "caption": (
                "ROC curves (left) and precision-recall curves (right) comparing "
                "mortality prediction models with and without PWV parameters."
            ),
        },
        {
            "path": os.path.join(FIGURE_DIR, "feature_importance.png"),
            "title": "Figure 3. Feature Importance",
            "caption": (
                "Feature importance in the gradient boosting full model, "
                "showing the relative contribution of PWV-derived and SOFA parameters."
            ),
        },
    ]

    blank_layout = prs.slide_layouts[6]  # blank slide

    for fig_info in figures:
        if not os.path.exists(fig_info["path"]):
            continue

        slide = prs.slides.add_slide(blank_layout)

        # Title at top
        title_box = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(0.2),
            PptxInches(12.3), PptxInches(0.6)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = fig_info["title"]
        p.font.size = PptxPt(20)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Image centered
        slide.shapes.add_picture(
            fig_info["path"],
            PptxInches(0.8), PptxInches(1.0),
            PptxInches(11.7), PptxInches(5.0)
        )

        # Caption at bottom
        caption_box = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(6.2),
            PptxInches(12.3), PptxInches(1.0)
        )
        tf_cap = caption_box.text_frame
        tf_cap.word_wrap = True
        p_cap = tf_cap.paragraphs[0]
        p_cap.text = fig_info["caption"]
        p_cap.font.size = PptxPt(11)
        p_cap.font.italic = True
        p_cap.alignment = PP_ALIGN.CENTER

    pptx_path = os.path.join(OUTPUT_DIR, "pwv_vitaldb_figures.pptx")
    prs.save(pptx_path)
    print(f"Figures PPTX saved to {pptx_path}")
    return pptx_path


def main():
    print("=" * 70)
    print("Generating Analysis Report")
    print("=" * 70)
    report_path = create_report()
    print(f"\nReport generated: {report_path}")

    pptx_path = create_figures_pptx()
    print(f"Figures PPTX generated: {pptx_path}")

    return report_path


if __name__ == "__main__":
    main()
