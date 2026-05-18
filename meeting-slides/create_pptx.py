#!/usr/bin/env python3
"""入職後初面談スライド (16枚) を PowerPoint (.pptx) で生成するスクリプト
   テキストはシェイプ内のtext_frameに直接書き込み、z-order問題を回避する"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ===== Colors =====
PRIMARY = RGBColor(0x1A, 0x36, 0x5D)
ACCENT = RGBColor(0x2B, 0x6C, 0xB0)
BG = RGBColor(0xF7, 0xFA, 0xFC)
TEXT_COLOR = RGBColor(0x1A, 0x20, 0x2C)
TEXT_LIGHT = RGBColor(0x4A, 0x55, 0x68)
GOLD = RGBColor(0xB7, 0x79, 0x1F)
GOLD_LIGHT = RGBColor(0xFE, 0xFC, 0xBF)
GREEN = RGBColor(0x27, 0x67, 0x49)
GREEN_LIGHT = RGBColor(0xC6, 0xF6, 0xD5)
RED = RGBColor(0x9B, 0x2C, 0x2C)
RED_LIGHT = RGBColor(0xFE, 0xD7, 0xD7)
PURPLE = RGBColor(0x55, 0x3C, 0x9A)
PURPLE_LIGHT = RGBColor(0xE9, 0xD8, 0xFD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG = RGBColor(0x1A, 0x36, 0x5D)
TEAL = RGBColor(0x2C, 0x7A, 0x7B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
FONT = "Yu Gothic"


def add_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _set_run(p, text, size, color, bold=False):
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = FONT
    return run


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_COLOR, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    _set_run(p, text, font_size, color, bold)
    return txBox


def add_rich_tf(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def add_bullet(tf, text, font_size=17, color=TEXT_COLOR, bold=False, sub_text=None):
    p = tf.add_paragraph()
    p.space_before = Pt(4)
    p.space_after = Pt(3)
    _set_run(p, text, font_size, color, bold)
    if sub_text:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(1)
        _set_run(p2, f"    {sub_text}", 13, TEXT_LIGHT)


def add_numbered_item(tf, num, text, num_color=ACCENT, font_size=16, sub_text=None):
    p = tf.add_paragraph()
    p.space_before = Pt(4)
    _set_run(p, f"{num}  ", font_size - 1, num_color, bold=True)
    _set_run(p, text, font_size, TEXT_COLOR)
    if sub_text:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(1)
        _set_run(p2, f"      {sub_text}", 12, TEXT_LIGHT)


def add_heading(slide, text):
    add_text_box(slide, Inches(0.7), Inches(0.35), Inches(11.5), Inches(0.6),
                 text, font_size=32, color=PRIMARY, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(0.95),
                                  Inches(11.5), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def add_shape_with_text(slide, left, top, width, height, fill_color,
                        texts, border_color=None, anchor=MSO_ANCHOR.TOP,
                        shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    """Create a shape with text directly inside it (no separate text boxes).
    texts: list of (text, font_size, color, bold, alignment) tuples"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    for i, (txt, sz, clr, bld, align) in enumerate(texts):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(2)
        _set_run(p, txt, sz, clr, bld)
    return shape


def add_highlight_box(slide, left, top, width, height, label, text,
                      label_color=ACCENT, bg_color=RGBColor(0xEB, 0xF8, 0xFF),
                      border_color=ACCENT):
    """Highlight box: colored left border + label + body text, all inside one shape."""
    texts = []
    if label:
        texts.append((label, 10, label_color, True, PP_ALIGN.LEFT))
    texts.append((text, 14, TEXT_COLOR, False, PP_ALIGN.LEFT))
    shape = add_shape_with_text(slide, left, top, width, height, bg_color, texts,
                                border_color=border_color)
    shape.line.color.rgb = border_color
    shape.line.width = Pt(2)
    return shape


def add_flow_box(slide, left, top, title, subtitle, border_color=ACCENT):
    w, h = Inches(2.5), Inches(1.0)
    texts = [
        (title, 13, PRIMARY, True, PP_ALIGN.CENTER),
        (subtitle, 10, TEXT_LIGHT, False, PP_ALIGN.CENTER),
    ]
    add_shape_with_text(slide, left, top, w, h, WHITE, texts, border_color=border_color)


def add_flow_arrow(slide, left, top):
    add_text_box(slide, left, top, Inches(0.5), Inches(1.0),
                 "→", font_size=24, color=ACCENT, alignment=PP_ALIGN.CENTER)


def add_tag(slide, left, top, text, color=ACCENT):
    w = Inches(0.18 + len(text) * 0.17)
    texts = [(text, 10, WHITE, True, PP_ALIGN.CENTER)]
    add_shape_with_text(slide, left, top, w, Inches(0.3), color, texts)
    return w


def add_page_num(slide, num, total):
    add_text_box(slide, Inches(0.4), Inches(7.0), Inches(1), Inches(0.3),
                 f"{num} / {total}", font_size=10, color=TEXT_LIGHT)


def add_col_box(slide, left, top, width, height, title, items, border_color=ACCENT):
    """Column box with title + bullet items, all inside one shape."""
    texts = [(title, 16, PRIMARY, True, PP_ALIGN.LEFT)]
    for item in items:
        texts.append((f"▸ {item}", 13, TEXT_COLOR, False, PP_ALIGN.LEFT))
    shape = add_shape_with_text(slide, left, top, width, height, WHITE, texts,
                                border_color=border_color)
    return shape


TOTAL = 16


# ===== 1. Title =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, DARK_BG)
add_text_box(s, Inches(1), Inches(2.0), Inches(11.3), Inches(1.2),
             "入職後初面談", font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.3), Inches(2.3), Pt(2))
div.fill.solid(); div.fill.fore_color.rgb = WHITE; div.line.fill.background()
add_text_box(s, Inches(1), Inches(3.6), Inches(11.3), Inches(0.5),
             "キャリア・研究ビジョン・産学連携のご共有", font_size=20,
             color=RGBColor(0xCC, 0xCC, 0xCC), alignment=PP_ALIGN.CENTER)
add_text_box(s, Inches(1), Inches(4.6), Inches(11.3), Inches(0.4),
             "2026年5月15日", font_size=15,
             color=RGBColor(0x99, 0x99, 0x99), alignment=PP_ALIGN.CENTER)
add_page_num(s, 1, TOTAL)


# ===== 2. Agenda =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG); add_heading(s, "本日のアジェンダ")
tf = add_rich_tf(s, Inches(1), Inches(1.4), Inches(11), Inches(5.5))
tf.paragraphs[0].clear()
for item in [
    "キャリア概要 — 医学部→麻酔科→医工連携→公衆衛生→起業→DS",
    "麻酔科臨床期のハイライト",
    "医工連携・公衆衛生・国際連携",
    "データサイエンティストとしての現在",
    "研究ビジョン — データ上流への遡行",
    "Devin × 伊藤忠 × 滋賀大学",
    "産学連携 — シフト最適化・学術支援・論文ダッシュボード",
    "法人補助金の大学還流モデル",
]:
    add_bullet(tf, f"◆  {item}", font_size=20)
add_page_num(s, 2, TOTAL)


# ===== 3. Career Overview =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG); add_heading(s, "キャリア概要")
phases = [
    ("Phase 1", "医学部\n麻酔科臨床", "1-1〜1-5", ACCENT),
    ("Phase 2", "医工連携", "2-1〜2-3", TEAL),
    ("Phase 3", "公衆衛生\n学位取得", "3-1〜3-3", GREEN),
    ("Phase 4", "起業\n国際連携", "4-1〜4-6", GOLD),
    ("Phase 5", "データ\nサイエンティスト", "5-1〜5-5", PURPLE),
]
x = Inches(0.8)
for num, title, sub, color in phases:
    w = Inches(2.2)
    texts = [
        (num, 9, RGBColor(0xDD, 0xDD, 0xDD), True, PP_ALIGN.CENTER),
        (title, 14, WHITE, True, PP_ALIGN.CENTER),
        (sub, 10, RGBColor(0xDD, 0xDD, 0xDD), False, PP_ALIGN.CENTER),
    ]
    add_shape_with_text(s, x, Inches(1.5), w, Inches(1.3), color, texts)
    x += w + Inches(0.25)
add_highlight_box(s, Inches(0.8), Inches(3.2), Inches(5.5), Inches(1.0),
                  "専門", "「好奇心」", ACCENT)
add_highlight_box(s, Inches(6.7), Inches(3.2), Inches(5.8), Inches(1.0),
                  "座右の銘",
                  "「大隠朝市」\n「未解決の課題は他の分野では90%解決されている」",
                  GOLD, GOLD_LIGHT, GOLD)
add_highlight_box(s, Inches(0.8), Inches(4.5), Inches(11.7), Inches(0.9),
                  "研究者としての目標",
                  "蝋人形になること — 後世に残る研究業績を確立し、不朽の存在として認知される",
                  PURPLE, PURPLE_LIGHT, PURPLE)
add_page_num(s, 3, TOTAL)


# ===== 4. Phase 1: Anesthesiology =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG); add_heading(s, "Phase 1：麻酔科臨床期")
tf = add_rich_tf(s, Inches(0.8), Inches(1.3), Inches(8.5), Inches(5.5))
tf.paragraphs[0].clear()
add_numbered_item(tf, "1-1", "医学部卒業・麻酔科入局")
add_numbered_item(tf, "1-2", "東京都立墨東病院（三次救命救急）麻酔科に勤務",
                  sub_text="のべ200〜300名の医師の中で「他の医師に手に負えない症例は知識も技術も大西を呼ぶ」ルール")
add_numbered_item(tf, "1-3", "酒井哲郎先生の著書で「東京の市中病院にいる天才麻酔科医」として紹介")
add_numbered_item(tf, "1-4", "天野篤先生（上皇の執刀医）からお見合いを紹介されるほど気に入られる")
add_numbered_item(tf, "1-5", "麻酔科でしたいことは全て終えて「卒業」")
book_path = os.path.join(SCRIPT_DIR, "book_cover.png")
if os.path.exists(book_path):
    s.shapes.add_picture(book_path, Inches(9.8), Inches(1.3), height=Inches(5.0))
add_page_num(s, 4, TOTAL)


# ===== 5. Phase 1 cont: Book quote =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG); add_heading(s, "Phase 1：「天才麻酔科医」の紹介")
page_path = os.path.join(SCRIPT_DIR, "book_page.jpg")
if os.path.exists(page_path):
    s.shapes.add_picture(page_path, Inches(0.8), Inches(1.3), height=Inches(5.2))
add_highlight_box(s, Inches(7.0), Inches(1.5), Inches(5.8), Inches(2.2),
                  "書籍より引用",
                  "「例えば東京の市中病院で働く天才麻酔科医はおるのだが…」\n\n"
                  "— 酒井哲郎先生\n"
                  "『麻酔科診療にみる医学留学へのパスポート』",
                  ACCENT)
add_highlight_box(s, Inches(7.0), Inches(4.0), Inches(5.8), Inches(1.5),
                  "天野篤先生との関係",
                  "天野篤先生（上皇の心臓バイパス手術執刀医）から\n"
                  "お見合いを紹介されるほどの信頼関係",
                  GOLD, GOLD_LIGHT, GOLD)
add_page_num(s, 5, TOTAL)


# ===== 6. Phase 2: Medical Engineering =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG); add_heading(s, "Phase 2：医工連携期")
tf = add_rich_tf(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(1.6))
tf.paragraphs[0].clear()
add_numbered_item(tf, "2-1", "墨東病院在籍中、市中病院と医学部を持たない工学部での医工連携に着手", num_color=TEAL)
add_numbered_item(tf, "2-2", "全身麻酔の機序研究 — 京都工業繊維大学（医学部なし）の論文が確信に最も近い", num_color=TEAL)
add_numbered_item(tf, "2-3", "生成AI前のプログラミングコンペ受賞歴 → コーディング能力の客観的証明", num_color=TEAL)
add_highlight_box(s, Inches(0.8), Inches(3.2), Inches(11.5), Inches(1.2),
                  "麻酔科の特異性",
                  "全身麻酔の効く仕組みには十分な科学的根拠がまだない。\n"
                  "医学部を持たない京都工業繊維大学の研究が最も核心に迫っている — 分野横断的アプローチの重要性",
                  TEAL, RGBColor(0xE6, 0xFF, 0xFA), TEAL)
add_col_box(s, Inches(0.8), Inches(4.7), Inches(5.5), Inches(2.3),
            "臨床面の特徴",
            ["リアルタイム生体モニターデータ", "PK/PDモデリングとの親和性", "手術スケジュールの複雑な最適化"],
            border_color=RGBColor(0xE2, 0xE8, 0xF0))
add_col_box(s, Inches(6.7), Inches(4.7), Inches(5.8), Inches(2.3),
            "研究面の特徴",
            ["数理モデル・シミュレーション需要大", "多変量時系列データの宝庫", "シフト制勤務 → 最適化問題と直結"],
            border_color=RGBColor(0xE2, 0xE8, 0xF0))
add_page_num(s, 6, TOTAL)


# ===== 7. Phase 3: Public Health =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG); add_heading(s, "Phase 3：公衆衛生・学位取得期")
tf = add_rich_tf(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(2.5))
tf.paragraphs[0].clear()
add_numbered_item(tf, "3-1", "滋賀医科大学で感染症論文により学位（博士号）取得", num_color=GREEN,
                  sub_text="生成AIのない時代に執筆。ChatGPTにコード改善を依頼したところデュアルユースを疑われ拒否された")
add_numbered_item(tf, "3-2", "学位の縁で滋賀医科大学 救急集中治療講座 非常勤助教に就任", num_color=GREEN)
add_numbered_item(tf, "3-3", "医療の (not so) big data を使った研究支援を推進したい", num_color=GREEN,
                  sub_text="ドメイン知識を使って滋賀大学と滋賀医科大学の間の通訳を果たす")
add_highlight_box(s, Inches(0.8), Inches(4.0), Inches(11.5), Inches(1.2),
                  "通訳者としての役割",
                  "データサイエンス（滋賀大学）× 臨床医学（滋賀医科大学）\n"
                  "両大学の間でドメイン知識による橋渡しができる稀有なポジション",
                  GREEN, GREEN_LIGHT, GREEN)
add_highlight_box(s, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.2),
                  "倫理と研究加速",
                  "研究において倫理は再現性のパーツであるに過ぎないのに過大視されている。\n"
                  "オープンデータを活用して研究を加速させるほうが生産的",
                  ACCENT)
add_page_num(s, 7, TOTAL)


# ===== 8. Phase 4: Venture =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG); add_heading(s, "Phase 4：起業・国際連携期")
tf = add_rich_tf(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(3.8))
tf.paragraphs[0].clear()
add_numbered_item(tf, "4-1", "スリランカとのSATREPS計画 — 3年越しでスリランカ側が書類作成できず頓挫", num_color=GOLD,
                  sub_text="現地JICAにも連絡を絶たれた。国内では文科省・前スリランカ大使・SATREPS発案者と知己を得た")
add_numbered_item(tf, "4-2", "成育医療研究センターとの共同研究 — まさかの成育側の倫理審査が未了というサプライズ", num_color=GOLD)
add_numbered_item(tf, "4-3", "シフト作成の最適化（数値最適化）— 導入実績あり", num_color=GOLD)
add_numbered_item(tf, "4-4", "学術支援サービス — 導入実績あり", num_color=GOLD)
add_numbered_item(tf, "4-5", "大学発ベンチャー法人補助金の大学還流モデル構想", num_color=GOLD)
add_numbered_item(tf, "4-6", "産学連携 — シフト最適化・学術支援の展開", num_color=GOLD)
add_highlight_box(s, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.2),
                  "教訓",
                  "国際連携は相手側の事務能力に大きく依存する。国内でのネットワーク構築は成果として残る。\n"
                  "倫理審査の未了は共同研究における予測困難なリスク。",
                  RED, RED_LIGHT, RED)
add_page_num(s, 8, TOTAL)


# ===== 9. Phase 5: Data Scientist =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG); add_heading(s, "Phase 5：データサイエンティスト期（現在）")
tf = add_rich_tf(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.5))
tf.paragraphs[0].clear()
add_numbered_item(tf, "5-1", "滋賀大学にデータサイエンティストとして入職", num_color=PURPLE)
add_numbered_item(tf, "5-2", "hikone — 独自の因果推論手法論文を執筆、JMLRに投稿中", num_color=PURPLE)
add_numbered_item(tf, "5-3", "論文を100本書いてもいいですか？", num_color=PURPLE)
add_numbered_item(tf, "5-4", "論文管理ダッシュボードを自作 → センターのプロジェクト管理にも活用", num_color=PURPLE)
add_numbered_item(tf, "5-5", "Devin × 伊藤忠 × 滋賀大学 — AI支援の大学導入提案", num_color=PURPLE)
add_page_num(s, 9, TOTAL)


# ===== 10. Research Vision =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG); add_heading(s, "研究ビジョン — データの上流へ")
fy = Inches(1.6)
add_flow_box(s, Inches(0.8), fy, "医学", "エビデンスを使う立場", ACCENT)
add_flow_arrow(s, Inches(3.3), fy)
add_flow_box(s, Inches(3.8), fy, "公衆衛生", "手法を使う立場", GREEN)
add_flow_arrow(s, Inches(6.3), fy)
add_flow_box(s, Inches(6.8), fy, "DS・統計学", "手法をつくる立場", PURPLE)
add_flow_arrow(s, Inches(9.3), fy)
add_shape_with_text(s, Inches(9.8), fy, Inches(2.5), Inches(1.0), GOLD_LIGHT,
                    [("hikone", 13, PRIMARY, True, PP_ALIGN.CENTER),
                     ("因果推論手法\nJMLR投稿中", 10, TEXT_LIGHT, False, PP_ALIGN.CENTER)],
                    border_color=GOLD)
add_highlight_box(s, Inches(0.8), Inches(3.0), Inches(5.5), Inches(1.0),
                  "データの下流＝医学の強み",
                  "医学はデータの下流だが本人のプラクティスは変わる\n→ 必ず社会実装される")
add_highlight_box(s, Inches(6.7), Inches(3.0), Inches(5.8), Inches(1.0),
                  "エビデンスの階梯への異論",
                  "エビデンスの階梯は誤解されている。\nケースリポートはRCTに劣らない",
                  GOLD, GOLD_LIGHT, GOLD)
add_highlight_box(s, Inches(0.8), Inches(4.3), Inches(11.7), Inches(1.0),
                  "戦略",
                  "データの上流に遡ってキャリアを積み、手法をつくる側へ。\n"
                  "同時に臨床経験に基づく社会実装力を武器に、研究と実践の両輪を回す。",
                  PURPLE, PURPLE_LIGHT, PURPLE)
add_page_num(s, 10, TOTAL)


# ===== 11. Coding Ability =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG); add_heading(s, "コーディング能力の証明")
add_highlight_box(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.0),
                  "課題認識", "生成AI時代においてはコード能力の客観的評価が困難に")
add_highlight_box(s, Inches(0.8), Inches(2.8), Inches(11.5), Inches(1.2),
                  "強み",
                  "生成AI登場以前のプログラミングコンペティションにおける受賞歴あり\n"
                  "→ コーディング能力、およびコードレビュー能力の客観的裏付け",
                  GREEN, GREEN_LIGHT, GREEN)
x = Inches(0.8)
for text, color in [("数値最適化", ACCENT), ("統計モデリング", ACCENT),
                    ("コードレビュー", ACCENT), ("AI前受賞実績", GREEN),
                    ("因果推論", PURPLE), ("論文ダッシュボード", GOLD)]:
    w = add_tag(s, x, Inches(4.3), text, color)
    x += w + Inches(0.1)
add_page_num(s, 11, TOTAL)


# ===== 12. Devin × Itochu =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG); add_heading(s, "Devin × 伊藤忠 × 滋賀大学")
fy = Inches(1.8)
add_flow_box(s, Inches(1.2), fy, "Cognition AI", "Devin 開発元", ACCENT)
add_flow_arrow(s, Inches(3.7), fy)
add_flow_box(s, Inches(4.2), fy, "伊藤忠テクノ\nソリューションズ", "日本代理店", GOLD)
add_flow_arrow(s, Inches(6.7), fy)
add_flow_box(s, Inches(7.2), fy, "滋賀大学", "伊藤忠と縁あり", GREEN)
add_highlight_box(s, Inches(0.8), Inches(3.2), Inches(11.5), Inches(1.0),
                  "提案",
                  "伊藤忠が日本におけるDevin代理店となった今、滋賀大学でのDevin導入を検討できないか？",
                  GOLD, GOLD_LIGHT, GOLD)
tf = add_rich_tf(s, Inches(1.0), Inches(4.5), Inches(11), Inches(2))
tf.paragraphs[0].clear()
add_bullet(tf, "◆ AI駆動型ソフトウェアエンジニアリングによる研究加速",
           sub_text="コーディング、データ分析、論文整備の自動化")
add_bullet(tf, "◆ 大学全体のDX推進の先行事例になり得る")
add_page_num(s, 12, TOTAL)


# ===== 13. Subsidy Model =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG); add_heading(s, "法人補助金の大学還流モデル")
fy = Inches(1.7)
add_flow_box(s, Inches(1.2), fy, "大学発ベンチャー", "技術・知見の事業化", GREEN)
add_flow_arrow(s, Inches(3.7), fy)
add_flow_box(s, Inches(4.2), fy, "法人補助金取得", "公的資金の獲得", GOLD)
add_flow_arrow(s, Inches(6.7), fy)
add_flow_box(s, Inches(7.2), fy, "大学への還流", "研究基盤の強化", ACCENT)
add_highlight_box(s, Inches(0.8), Inches(3.1), Inches(11.5), Inches(1.0),
                  "ビジョン",
                  "大学発ベンチャーが獲得した法人補助金を大学の研究基盤に還流させるサステナブルなモデルを構築",
                  GOLD, GOLD_LIGHT, GOLD)
tf = add_rich_tf(s, Inches(1.0), Inches(4.4), Inches(11), Inches(2))
tf.paragraphs[0].clear()
add_bullet(tf, "◆ 産学連携の成果を大学運営に直接貢献させる仕組み")
add_bullet(tf, "◆ ベンチャー → 補助金 → 大学 の好循環サイクル")
add_page_num(s, 13, TOTAL)


# ===== 14. Industry-Academia =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG); add_heading(s, "産学連携の展開")
add_col_box(s, Inches(0.8), Inches(1.3), Inches(3.8), Inches(4.0),
            "シフト作成の最適化",
            ["手法：数値最適化", "実績：導入済 ✓", "公平性・連続勤務制限", "コスト削減 + 満足度両立"],
            border_color=ACCENT)
add_col_box(s, Inches(4.9), Inches(1.3), Inches(3.8), Inches(4.0),
            "学術支援",
            ["実績：導入済 ✓", "論文執筆・データ解析", "統計コンサルティング", "研究デザイン助言"],
            border_color=GREEN)
add_col_box(s, Inches(9.0), Inches(1.3), Inches(3.8), Inches(4.0),
            "論文ダッシュボード",
            ["論文管理用に自作", "センターのプロジェクト管理に活用", "目標：論文100本"],
            border_color=PURPLE)
# Tags under columns
for (x, tags) in [
    (Inches(0.8), [("最適化", ACCENT), ("OR", ACCENT)]),
    (Inches(4.9), [("統計", GREEN), ("論文支援", GREEN)]),
    (Inches(9.0), [("管理ツール", PURPLE)]),
]:
    tx = x
    for ttext, tcolor in tags:
        w = add_tag(s, tx, Inches(5.5), ttext, tcolor)
        tx += w + Inches(0.08)
add_highlight_box(s, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.8),
                  "目標",
                  "導入実績のあるシフト最適化・学術支援を基盤に学内外へ展開。論文ダッシュボードでセンター業務を効率化。",
                  GREEN, GREEN_LIGHT, GREEN)
add_page_num(s, 14, TOTAL)


# ===== 15. Research Policy =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG); add_heading(s, "研究上の方針")
add_highlight_box(s, Inches(0.8), Inches(1.8), Inches(11.5), Inches(1.0),
                  "留意事項",
                  "特定の研究者（佐藤俊哉氏）との共同研究は行わない方針",
                  RED, RED_LIGHT, RED)
tf = add_rich_tf(s, Inches(1.0), Inches(3.2), Inches(11), Inches(3))
tf.paragraphs[0].clear()
add_bullet(tf, "◆ 独自の研究ラインを確立・維持する", font_size=18,
           sub_text="研究の独立性と方向性を自ら決定")
add_bullet(tf, "◆ 共同研究は研究ビジョンが合致する相手と進める", font_size=18,
           sub_text="相互補完的で建設的なパートナーシップを重視")
add_page_num(s, 15, TOTAL)


# ===== 16. Summary =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, DARK_BG)
add_text_box(s, Inches(1), Inches(0.8), Inches(11.3), Inches(0.8),
             "まとめ", font_size=42, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(1.7), Inches(2.3), Pt(2))
div.fill.solid(); div.fill.fore_color.rgb = WHITE; div.line.fill.background()
tf = add_rich_tf(s, Inches(2.5), Inches(2.2), Inches(8.3), Inches(4.5))
tf.paragraphs[0].clear()
for item in [
    "医学部→麻酔科→医工連携→公衆衛生→起業→DSのキャリア遷移",
    "墨東病院「天才麻酔科医」— 麻酔科を卒業して新たな挑戦へ",
    "データの上流に遡り手法をつくる側へ（hikone → JMLR）",
    "滋賀大学 × 滋賀医科大学のドメイン知識通訳",
    "Devin導入・法人補助金還流・シフト最適化・学術支援",
    "論文100本を目標に、蝋人形を目指す",
]:
    add_bullet(tf, f"◆  {item}", font_size=18, color=RGBColor(0xDD, 0xDD, 0xDD))
add_text_box(s, Inches(1), Inches(6.4), Inches(11.3), Inches(0.5),
             "ご清聴ありがとうございました", font_size=16,
             color=RGBColor(0x99, 0x99, 0x99), alignment=PP_ALIGN.CENTER)
add_page_num(s, 16, TOTAL)


# ===== Save =====
output_path = os.path.join(SCRIPT_DIR, "meeting_slides.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
