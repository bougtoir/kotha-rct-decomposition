#!/usr/bin/env python3
"""入職後初面談スライド (16枚) — テキストのみ版（図形なし）"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

PRIMARY = RGBColor(0x1A, 0x36, 0x5D)
ACCENT = RGBColor(0x2B, 0x6C, 0xB0)
TEXT_COLOR = RGBColor(0x1A, 0x20, 0x2C)
TEXT_LIGHT = RGBColor(0x4A, 0x55, 0x68)
GOLD = RGBColor(0xB7, 0x79, 0x1F)
GREEN = RGBColor(0x27, 0x67, 0x49)
RED = RGBColor(0x9B, 0x2C, 0x2C)
PURPLE = RGBColor(0x55, 0x3C, 0x9A)
TEAL = RGBColor(0x2C, 0x7A, 0x7B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF7, 0xFA, 0xFC)
DARK_BG = RGBColor(0x1A, 0x36, 0x5D)

FONT = "Yu Gothic"
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _run(p, text, size, color, bold=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = FONT
    return r


def add_tf(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def heading(slide, text):
    tf = add_tf(slide, Inches(0.7), Inches(0.3), Inches(11.5), Inches(0.7))
    p = tf.paragraphs[0]
    _run(p, text, 32, PRIMARY, bold=True)
    p.space_after = Pt(8)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(0.95),
                                  Inches(11.5), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def bullet_list(slide, items, left=0.8, top=1.3, width=11.5, height=5.5, font_size=18, color=TEXT_COLOR):
    tf = add_tf(slide, Inches(left), Inches(top), Inches(width), Inches(height))
    tf.paragraphs[0].clear()
    for item in items:
        p = tf.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(3)
        _run(p, f"◆  {item}", font_size, color)
    return tf


def numbered_list(slide, items, left=0.8, top=1.3, width=11.5, height=5.5):
    """items: list of (num, text, num_color, sub_text_or_None)"""
    tf = add_tf(slide, Inches(left), Inches(top), Inches(width), Inches(height))
    tf.paragraphs[0].clear()
    for num, text, num_color, sub in items:
        p = tf.add_paragraph()
        p.space_before = Pt(5)
        _run(p, f"{num}  ", 15, num_color, bold=True)
        _run(p, text, 16, TEXT_COLOR)
        if sub:
            p2 = tf.add_paragraph()
            p2.space_before = Pt(1)
            _run(p2, f"      {sub}", 12, TEXT_LIGHT)
    return tf


def label_body(slide, label, body, label_color=ACCENT, top=None, left=0.8, width=11.5):
    """Label line + body text, no shapes."""
    tf = add_tf(slide, Inches(left), Inches(top), Inches(width), Inches(1.5))
    p = tf.paragraphs[0]
    _run(p, f"【{label}】", 12, label_color, bold=True)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(2)
    _run(p2, body, 14, TEXT_COLOR)
    return tf


def page_num(slide, num, total=16):
    tf = add_tf(slide, Inches(0.4), Inches(7.0), Inches(1), Inches(0.3))
    _run(tf.paragraphs[0], f"{num} / {total}", 10, TEXT_LIGHT)


# ===== 1. Title =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, DARK_BG)
tf = add_tf(s, Inches(1), Inches(2.0), Inches(11.3), Inches(1.2))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
_run(p, "入職後初面談", 52, WHITE, bold=True)
tf2 = add_tf(s, Inches(1), Inches(3.5), Inches(11.3), Inches(0.5))
p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
_run(p2, "キャリア・研究ビジョン・産学連携のご共有", 20, RGBColor(0xCC, 0xCC, 0xCC))
tf3 = add_tf(s, Inches(1), Inches(4.3), Inches(11.3), Inches(0.4))
p3 = tf3.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
_run(p3, "2026年5月15日", 15, RGBColor(0x99, 0x99, 0x99))
page_num(s, 1)

# ===== 2. Agenda =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG); heading(s, "本日のアジェンダ")
bullet_list(s, [
    "キャリア概要 — 医学部→麻酔科→医工連携→公衆衛生→起業→DS",
    "麻酔科臨床期のハイライト",
    "医工連携・公衆衛生・国際連携",
    "データサイエンティストとしての現在",
    "研究ビジョン — データ上流への遡行",
    "Devin × 伊藤忠 × 滋賀大学",
    "産学連携 — シフト最適化・学術支援・論文ダッシュボード",
    "法人補助金の大学還流モデル",
], font_size=20)
page_num(s, 2)

# ===== 3. Career Overview =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG); heading(s, "キャリア概要")
tf = add_tf(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(1.2))
p = tf.paragraphs[0]
phases = [
    ("Phase 1: 医学部・麻酔科臨床 (1-1〜1-5)", ACCENT),
    ("Phase 2: 医工連携 (2-1〜2-3)", TEAL),
    ("Phase 3: 公衆衛生・学位取得 (3-1〜3-3)", GREEN),
    ("Phase 4: 起業・国際連携 (4-1〜4-6)", GOLD),
    ("Phase 5: データサイエンティスト (5-1〜5-5)", PURPLE),
]
for i, (txt, clr) in enumerate(phases):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.space_before = Pt(4)
    _run(p, f"■ {txt}", 16, clr, bold=True)

label_body(s, "専門", "「好奇心」", ACCENT, top=3.8)
label_body(s, "座右の銘", "「大隠朝市」／「どこかの分野で未解決の課題は他の分野では90%解決されている」", GOLD, top=4.6)
label_body(s, "研究者としての目標", "蝋人形になること — 後世に残る研究業績を確立し、不朽の存在として認知される", PURPLE, top=5.4)
page_num(s, 3)

# ===== 4. Phase 1: Anesthesiology =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG); heading(s, "Phase 1：麻酔科臨床期")
numbered_list(s, [
    ("1-1", "医学部卒業・麻酔科入局", ACCENT, None),
    ("1-2", "東京都立墨東病院（三次救命救急）麻酔科に勤務", ACCENT,
     "のべ200〜300名の医師の中で「他の医師に手に負えない症例は知識も技術も大西を呼ぶ」ルール"),
    ("1-3", "酒井哲郎先生の著書で「東京の市中病院にいる天才麻酔科医」として紹介", ACCENT, None),
    ("1-4", "天野篤先生（上皇の執刀医）からお見合いを紹介されるほど気に入られる", ACCENT, None),
    ("1-5", "麻酔科でしたいことは全て終えて「卒業」", ACCENT, None),
], width=8.5)
book_path = os.path.join(SCRIPT_DIR, "book_cover.png")
if os.path.exists(book_path):
    s.shapes.add_picture(book_path, Inches(9.8), Inches(1.3), height=Inches(5.0))
page_num(s, 4)

# ===== 5. Phase 1 cont: Book =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG); heading(s, "Phase 1：「天才麻酔科医」の紹介")
page_path = os.path.join(SCRIPT_DIR, "book_page.jpg")
if os.path.exists(page_path):
    s.shapes.add_picture(page_path, Inches(0.8), Inches(1.3), height=Inches(5.2))
label_body(s, "書籍より引用",
           "「例えば東京の市中病院で働く天才麻酔科医はおるのだが…」\n"
           "— 酒井哲郎先生『麻酔科診療にみる医学留学へのパスポート』",
           ACCENT, top=1.5, left=7.0, width=5.5)
label_body(s, "天野篤先生との関係",
           "天野篤先生（上皇の心臓バイパス手術執刀医）から\nお見合いを紹介されるほどの信頼関係",
           GOLD, top=3.8, left=7.0, width=5.5)
page_num(s, 5)

# ===== 6. Phase 2: Medical Engineering =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG); heading(s, "Phase 2：医工連携期")
numbered_list(s, [
    ("2-1", "墨東病院在籍中、市中病院と医学部を持たない工学部での医工連携に着手", TEAL, None),
    ("2-2", "全身麻酔の機序研究 — 京都工業繊維大学（医学部なし）の論文が確信に最も近い", TEAL, None),
    ("2-3", "生成AI前のプログラミングコンペ受賞歴 → コーディング能力の客観的証明", TEAL, None),
], height=2.0)
label_body(s, "麻酔科の特異性",
           "全身麻酔の効く仕組みには十分な科学的根拠がまだない。\n"
           "医学部を持たない京都工業繊維大学の研究が最も核心に迫っている — 分野横断的アプローチの重要性",
           TEAL, top=3.3)
tf = add_tf(s, Inches(0.8), Inches(4.8), Inches(5.5), Inches(2.0))
p = tf.paragraphs[0]
_run(p, "＜臨床面の特徴＞", 16, PRIMARY, bold=True)
for item in ["リアルタイム生体モニターデータ", "PK/PDモデリングとの親和性", "手術スケジュールの複雑な最適化"]:
    p2 = tf.add_paragraph(); p2.space_before = Pt(2)
    _run(p2, f"  ▸ {item}", 14, TEXT_COLOR)
tf2 = add_tf(s, Inches(6.7), Inches(4.8), Inches(5.8), Inches(2.0))
p = tf2.paragraphs[0]
_run(p, "＜研究面の特徴＞", 16, PRIMARY, bold=True)
for item in ["数理モデル・シミュレーション需要大", "多変量時系列データの宝庫", "シフト制勤務 → 最適化問題と直結"]:
    p2 = tf2.add_paragraph(); p2.space_before = Pt(2)
    _run(p2, f"  ▸ {item}", 14, TEXT_COLOR)
page_num(s, 6)

# ===== 7. Phase 3: Public Health =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG); heading(s, "Phase 3：公衆衛生・学位取得期")
numbered_list(s, [
    ("3-1", "滋賀医科大学で感染症論文により学位（博士号）取得", GREEN,
     "生成AIのない時代に執筆。ChatGPTにコード改善を依頼したところデュアルユースを疑われ拒否された"),
    ("3-2", "学位の縁で滋賀医科大学 救急集中治療講座 非常勤助教に就任", GREEN, None),
    ("3-3", "医療の (not so) big data を使った研究支援を推進したい", GREEN,
     "ドメイン知識を使って滋賀大学と滋賀医科大学の間の通訳を果たす"),
], height=2.8)
label_body(s, "通訳者としての役割",
           "データサイエンス（滋賀大学）× 臨床医学（滋賀医科大学）\n"
           "両大学の間でドメイン知識による橋渡しができる稀有なポジション",
           GREEN, top=4.3)
label_body(s, "倫理と研究加速",
           "研究において倫理は再現性のパーツであるに過ぎないのに過大視されている。\n"
           "オープンデータを活用して研究を加速させるほうが生産的",
           ACCENT, top=5.5)
page_num(s, 7)

# ===== 8. Phase 4: Venture =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG); heading(s, "Phase 4：起業・国際連携期")
numbered_list(s, [
    ("4-1", "スリランカとのSATREPS計画 — 3年越しでスリランカ側が書類作成できず頓挫", GOLD,
     "現地JICAにも連絡を絶たれた。国内では文科省・前スリランカ大使・SATREPS発案者と知己を得た"),
    ("4-2", "成育医療研究センターとの共同研究 — まさかの成育側の倫理審査が未了というサプライズ", GOLD, None),
    ("4-3", "シフト作成の最適化（数値最適化）— 導入実績あり", GOLD, None),
    ("4-4", "学術支援サービス — 導入実績あり", GOLD, None),
    ("4-5", "大学発ベンチャー法人補助金の大学還流モデル構想", GOLD, None),
    ("4-6", "産学連携 — シフト最適化・学術支援の展開", GOLD, None),
], height=4.0)
label_body(s, "教訓",
           "国際連携は相手側の事務能力に大きく依存する。国内でのネットワーク構築は成果として残る。\n"
           "倫理審査の未了は共同研究における予測困難なリスク。",
           RED, top=5.8)
page_num(s, 8)

# ===== 9. Phase 5: Data Scientist =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG); heading(s, "Phase 5：データサイエンティスト期（現在）")
numbered_list(s, [
    ("5-1", "滋賀大学にデータサイエンティストとして入職", PURPLE, None),
    ("5-2", "hikone — 独自の因果推論手法論文を執筆、JMLRに投稿中", PURPLE, None),
    ("5-3", "論文を100本書いてもいいですか？", PURPLE, None),
    ("5-4", "論文管理ダッシュボードを自作 → センターのプロジェクト管理にも活用", PURPLE, None),
    ("5-5", "Devin × 伊藤忠 × 滋賀大学 — AI支援の大学導入提案", PURPLE, None),
])
page_num(s, 9)

# ===== 10. Research Vision =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG); heading(s, "研究ビジョン — データの上流へ")
tf = add_tf(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(1.0))
p = tf.paragraphs[0]
_run(p, "医学（エビデンスを使う）", 16, ACCENT, bold=True)
_run(p, "  →  ", 16, TEXT_LIGHT)
_run(p, "公衆衛生（手法を使う）", 16, GREEN, bold=True)
_run(p, "  →  ", 16, TEXT_LIGHT)
_run(p, "DS・統計学（手法をつくる）", 16, PURPLE, bold=True)
_run(p, "  →  ", 16, TEXT_LIGHT)
_run(p, "hikone（因果推論手法 / JMLR投稿中）", 16, GOLD, bold=True)
label_body(s, "データの下流＝医学の強み",
           "医学はデータの下流だが本人のプラクティスは変わる → 必ず社会実装される",
           ACCENT, top=2.8)
label_body(s, "エビデンスの階梯への異論",
           "エビデンスの階梯は誤解されている。ケースリポートはRCTに劣らない",
           GOLD, top=3.8)
label_body(s, "戦略",
           "データの上流に遡ってキャリアを積み、手法をつくる側へ。\n"
           "同時に臨床経験に基づく社会実装力を武器に、研究と実践の両輪を回す。",
           PURPLE, top=4.8)
page_num(s, 10)

# ===== 11. Coding Ability =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG); heading(s, "コーディング能力の証明")
label_body(s, "課題認識", "生成AI時代においてはコード能力の客観的評価が困難に", ACCENT, top=1.5)
label_body(s, "強み",
           "生成AI登場以前のプログラミングコンペティションにおける受賞歴あり\n"
           "→ コーディング能力、およびコードレビュー能力の客観的裏付け",
           GREEN, top=2.8)
tf = add_tf(s, Inches(0.8), Inches(4.3), Inches(11.5), Inches(0.5))
p = tf.paragraphs[0]
tags = ["数値最適化", "統計モデリング", "コードレビュー", "AI前受賞実績", "因果推論", "論文ダッシュボード"]
_run(p, "  /  ".join(tags), 14, ACCENT, bold=True)
page_num(s, 11)

# ===== 12. Devin × Itochu =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG); heading(s, "Devin × 伊藤忠 × 滋賀大学")
tf = add_tf(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(1.0))
p = tf.paragraphs[0]
_run(p, "Cognition AI（Devin開発元）", 16, ACCENT, bold=True)
_run(p, "  →  ", 16, TEXT_LIGHT)
_run(p, "伊藤忠テクノソリューションズ（日本代理店）", 16, GOLD, bold=True)
_run(p, "  →  ", 16, TEXT_LIGHT)
_run(p, "滋賀大学（伊藤忠と縁あり）", 16, GREEN, bold=True)
label_body(s, "提案",
           "伊藤忠が日本におけるDevin代理店となった今、滋賀大学でのDevin導入を検討できないか？",
           GOLD, top=2.8)
bullet_list(s, [
    "AI駆動型ソフトウェアエンジニアリングによる研究加速（コーディング、データ分析、論文整備の自動化）",
    "大学全体のDX推進の先行事例になり得る",
], top=4.0)
page_num(s, 12)

# ===== 13. Subsidy Model =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG); heading(s, "法人補助金の大学還流モデル")
tf = add_tf(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(1.0))
p = tf.paragraphs[0]
_run(p, "大学発ベンチャー（技術の事業化）", 16, GREEN, bold=True)
_run(p, "  →  ", 16, TEXT_LIGHT)
_run(p, "法人補助金取得（公的資金）", 16, GOLD, bold=True)
_run(p, "  →  ", 16, TEXT_LIGHT)
_run(p, "大学への還流（研究基盤の強化）", 16, ACCENT, bold=True)
label_body(s, "ビジョン",
           "大学発ベンチャーが獲得した法人補助金を大学の研究基盤に還流させるサステナブルなモデルを構築",
           GOLD, top=2.8)
bullet_list(s, [
    "産学連携の成果を大学運営に直接貢献させる仕組み",
    "ベンチャー → 補助金 → 大学 の好循環サイクル",
], top=4.0)
page_num(s, 13)

# ===== 14. Industry-Academia =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG); heading(s, "産学連携の展開")
tf1 = add_tf(s, Inches(0.8), Inches(1.3), Inches(3.8), Inches(4.0))
p = tf1.paragraphs[0]
_run(p, "＜シフト作成の最適化＞", 16, ACCENT, bold=True)
for item in ["手法：数値最適化", "実績：導入済 ✓", "公平性・連続勤務制限", "コスト削減 + 満足度両立"]:
    p2 = tf1.add_paragraph(); p2.space_before = Pt(3)
    _run(p2, f"  ▸ {item}", 14, TEXT_COLOR)

tf2 = add_tf(s, Inches(4.9), Inches(1.3), Inches(3.8), Inches(4.0))
p = tf2.paragraphs[0]
_run(p, "＜学術支援＞", 16, GREEN, bold=True)
for item in ["実績：導入済 ✓", "論文執筆・データ解析", "統計コンサルティング", "研究デザイン助言"]:
    p2 = tf2.add_paragraph(); p2.space_before = Pt(3)
    _run(p2, f"  ▸ {item}", 14, TEXT_COLOR)

tf3 = add_tf(s, Inches(9.0), Inches(1.3), Inches(3.8), Inches(4.0))
p = tf3.paragraphs[0]
_run(p, "＜論文ダッシュボード＞", 16, PURPLE, bold=True)
for item in ["論文管理用に自作", "センターのプロジェクト管理に活用", "目標：論文100本"]:
    p2 = tf3.add_paragraph(); p2.space_before = Pt(3)
    _run(p2, f"  ▸ {item}", 14, TEXT_COLOR)

label_body(s, "目標",
           "導入実績のあるシフト最適化・学術支援を基盤に学内外へ展開。論文ダッシュボードでセンター業務を効率化。",
           GREEN, top=5.8)
page_num(s, 14)

# ===== 15. Research Policy =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG); heading(s, "研究上の方針")
label_body(s, "留意事項",
           "特定の研究者（佐藤俊哉氏）との共同研究は行わない方針",
           RED, top=1.8)
bullet_list(s, [
    "独自の研究ラインを確立・維持する（研究の独立性と方向性を自ら決定）",
    "共同研究は研究ビジョンが合致する相手と進める（相互補完的で建設的なパートナーシップを重視）",
], top=3.2)
page_num(s, 15)

# ===== 16. Summary =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, DARK_BG)
tf = add_tf(s, Inches(1), Inches(0.8), Inches(11.3), Inches(0.8))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
_run(p, "まとめ", 42, WHITE, bold=True)
bullet_list(s, [
    "医学部→麻酔科→医工連携→公衆衛生→起業→DSのキャリア遷移",
    "墨東病院「天才麻酔科医」— 麻酔科を卒業して新たな挑戦へ",
    "データの上流に遡り手法をつくる側へ（hikone → JMLR）",
    "滋賀大学 × 滋賀医科大学のドメイン知識通訳",
    "Devin導入・法人補助金還流・シフト最適化・学術支援",
    "論文100本を目標に、蝋人形を目指す",
], top=2.2, font_size=18, color=RGBColor(0xDD, 0xDD, 0xDD))
tf_end = add_tf(s, Inches(1), Inches(6.4), Inches(11.3), Inches(0.5))
p = tf_end.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
_run(p, "ご清聴ありがとうございました", 16, RGBColor(0x99, 0x99, 0x99))
page_num(s, 16)


# ===== Save =====
output_path = os.path.join(SCRIPT_DIR, "meeting_slides_textonly.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
