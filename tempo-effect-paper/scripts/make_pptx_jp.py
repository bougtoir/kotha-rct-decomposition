# -*- coding: utf-8 -*-
"""Create PPTX figure file — Japanese. One figure per slide.
Code-output figures: embedded as images.
Conceptual/flow diagrams: editable shapes."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

FIG_DIR = '/home/ubuntu/figures'
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def add_image_slide(prs, img_path, title, caption):
    slide = prs.slides.add_slide(BLANK)
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1B, 0x1B, 0x1B)
    p.alignment = PP_ALIGN.CENTER
    if os.path.exists(img_path):
        from PIL import Image
        im = Image.open(img_path)
        w, h = im.size
        aspect = w / h
        max_w, max_h = 11.5, 5.5
        if max_w / aspect <= max_h:
            disp_w, disp_h = max_w, max_w / aspect
        else:
            disp_w, disp_h = max_h * aspect, max_h
        left = (13.333 - disp_w) / 2
        slide.shapes.add_picture(img_path, Inches(left), Inches(1.0), Inches(disp_w), Inches(disp_h))
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.7))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = caption
    p2.font.size = Pt(11)
    p2.font.italic = True
    p2.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
    p2.alignment = PP_ALIGN.CENTER
    return slide

def add_box(slide, left, top, width, height, text, font_size=12, fill_rgb=(0x41,0x82,0xC4),
            font_rgb=(0xFF,0xFF,0xFF), bold=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    shp = slide.shapes.add_shape(shape, Inches(left), Inches(top), Inches(width), Inches(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(*fill_rgb)
    shp.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    shp.line.width = Pt(1)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*font_rgb)
    return shp

def add_arrow(slide, x1, y1, x2, y2, color=(0x33,0x33,0x33)):
    from pptx.oxml.ns import qn
    connector = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = RGBColor(*color)
    connector.line.width = Pt(2)
    ln = connector.line._ln
    tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail)
    return connector

def add_text_label(slide, left, top, width, height, text, font_size=11, bold=False, italic=False,
                   align=PP_ALIGN.CENTER, color=(0x22,0x22,0x22)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = RGBColor(*color)
    p.alignment = align
    return txBox

# ========================
# Slide 1: fig1
# ========================
add_image_slide(prs, f'{FIG_DIR}/fig1_showcase.png',
    "\u56f31. \u30e2\u30c7\u30eb vs \u89b3\u6e2c\u4eba\u53e3\uff1a\u4ee3\u88686\u30ab\u56fd\uff081970\u20132023\uff09",
    "\u52d5\u7684\u30e2\u30c7\u30eb\uff08\u9752\u7834\u7dda\uff09\u306f10\u5e74\u3054\u3068\u306b\u30d1\u30e9\u30e1\u30fc\u30bf\u66f4\u65b0\u3002\u9759\u7684\u30e2\u30c7\u30eb\uff08\u8d64\u70b9\u7dda\uff09\u306f1970\u5e74\u30d1\u30e9\u30e1\u30fc\u30bf\u56fa\u5b9a\u3002"
    "\u9ed2\u5b9f\u7dda\uff1dUN WPP 2024\u63a8\u8a08\u5024\u3002\u5bfe\u8c61\u56fd\uff1a\u65e5\u672c\u3001\u4e2d\u56fd\u3001\u7c73\u56fd\u3001\u97d3\u56fd\u3001\u30c9\u30a4\u30c4\u3001DRC\u3002")

# ========================
# Slide 2: fig2
# ========================
add_image_slide(prs, f'{FIG_DIR}/fig2_all_countries.png',
    "\u56f32. \u516840\u30ab\u56fd\u306e\u30e2\u30c7\u30eb\u691c\u8a3c",
    "\u5404\u30d1\u30cd\u30eb\u53f3\u4e0a\u306b\u52d5\u7684\u30e2\u30c7\u30ebMAPE\u3092\u8868\u793a\u3002\u30a2\u30eb\u30d5\u30a1\u30d9\u30c3\u30c8\u9806\u3002"
    "40\u30ab\u56fd\u4e2d30\u30ab\u56fd\u3067MAPE<10%\u300120\u30ab\u56fd\u3067<5%\u3002")

# ========================
# Slide 3: fig3
# ========================
add_image_slide(prs, f'{FIG_DIR}/fig3_heatmap.png',
    "\u56f33. \u9759\u7684\u30e2\u30c7\u30ebMAPE\uff08%\uff09\uff1a\u56fd\u00d7\u57fa\u6e96\u5e74",
    "\u7dd1\uff1d\u826f\u597d\u306a\u9069\u5408\u3001\u8d64\uff1d\u4e0d\u826f\u3002\u30b9\u30b1\u30fc\u30eb\u4e0a\u965030%\u3002"
    "\u57fa\u6e96\u5e74\u304c\u65b0\u3057\u3044\u307b\u3069\u9069\u5408\u304c\u6539\u5584\u3002")

# ========================
# Slide 4: fig4
# ========================
add_image_slide(prs, f'{FIG_DIR}/fig4_comparison.png',
    "\u56f34. \u9759\u7684 vs \u52d5\u7684\u30e2\u30c7\u30eb\u6bd4\u8f03",
    "\u5de6\uff1a\u56fd\u5225MAPE\u3002\u53f3\uff1a\u6700\u7d42\u4eba\u53e3\u6bd4\u7387\uff08\u30e2\u30c7\u30eb/\u89b3\u6e2c\u30012023\u5e74\uff09\u3002"
    "\u52d5\u7684\u30e2\u30c7\u30eb\u306fMAPE\u4e2d\u592e\u5024\u3067\u7d042\u500d\u306e\u6539\u5584\u3002")

# ========================
# Slide 5: fig5
# ========================
add_image_slide(prs, f'{FIG_DIR}/fig5_bias.png',
    "\u56f35. \u30e2\u30c7\u30eb\u30d0\u30a4\u30a2\u30b9\u5206\u6790\uff08\u57fa\u6e96\u5e742000\uff09",
    "(A) vs TFR\u3001(B) vs \u5e73\u5747\u5bff\u547d\u3001(C) vs MAC\u3002"
    "\u4f53\u7cfb\u7684\u95a2\u4fc2\u306f\u89b3\u5bdf\u3055\u308c\u305a\u3001\u30e2\u30c7\u30eb\u306e\u9811\u5065\u6027\u3092\u793a\u3059\u3002")

# ========================
# Slide 6: Editable Model Flow Diagram (JP)
# ========================
slide6 = prs.slides.add_slide(BLANK)
add_text_label(slide6, 0.5, 0.2, 12.3, 0.6,
    "\u56f36. \u5185\u751f\u66f4\u65b0\u30e2\u30c7\u30eb\uff1a\u69cb\u9020\u3068\u30c7\u30fc\u30bf\u30d5\u30ed\u30fc", font_size=24, bold=True)
add_text_label(slide6, 0.5, 0.75, 12.3, 0.4,
    "\uff08\u7de8\u96c6\u53ef\u80fd\u306a\u56f3\u2014\u2014\u5168\u3066\u306e\u56f3\u5f62\u3092\u79fb\u52d5\u30fb\u30ea\u30b5\u30a4\u30ba\u30fb\u8272\u5909\u66f4\u53ef\u80fd\uff09",
    font_size=12, italic=True, color=(0x88,0x88,0x88))

# Input parameters
add_text_label(slide6, 0.3, 1.25, 3.0, 0.3, "\u5165\u529b\u30d1\u30e9\u30e1\u30fc\u30bf", font_size=13, bold=True, color=(0x88,0x44,0x00))
add_box(slide6, 0.5, 1.6, 2.2, 0.7, "TFR\n\uff08\u5408\u8a08\u7279\u6b8a\u51fa\u751f\u7387\uff09", 11, (0xE8,0x6C,0x00))
add_box(slide6, 3.0, 1.6, 2.2, 0.7, "e\u2080\n\uff08\u5e73\u5747\u5bff\u547d\uff09", 11, (0xE8,0x6C,0x00))
add_box(slide6, 5.5, 1.6, 2.2, 0.7, "MAC\n\uff08\u5e73\u5747\u51fa\u7523\u5e74\u9f62\uff09", 11, (0xE8,0x6C,0x00))
add_box(slide6, 8.0, 1.6, 2.2, 0.7, "\u03c3\n\uff08\u51fa\u751f\u30b9\u30b1\u30b8\u30e5\u30fc\u30eb\u5e45\uff09", 11, (0xE8,0x6C,0x00))

add_arrow(slide6, 1.6, 2.35, 1.6, 2.9)
add_arrow(slide6, 4.1, 2.35, 4.1, 2.9)
add_arrow(slide6, 6.6, 2.35, 6.6, 2.9)
add_arrow(slide6, 9.1, 2.35, 9.1, 2.9)

# Model components
add_text_label(slide6, 0.3, 2.75, 3.0, 0.3, "\u30e2\u30c7\u30eb\u69cb\u6210\u8981\u7d20", font_size=13, bold=True, color=(0x00,0x55,0x88))
add_box(slide6, 0.3, 3.05, 2.6, 0.85, "\u6b63\u898f\u51fa\u751f\u30b9\u30b1\u30b8\u30e5\u30fc\u30eb\nASFR ~ N(MAC, \u03c3\u00b2)", 10, (0x41,0x82,0xC4))
add_box(slide6, 3.3, 3.05, 2.3, 0.85, "Gompertz\u751f\u5b58\u95a2\u6570\nS(x) = exp[-(a/b)\n\u00d7(exp(bx)-1)]", 10, (0x41,0x82,0xC4))
add_box(slide6, 6.0, 3.05, 2.6, 0.85, "\u5e74\u9f62\u5225\u51fa\u751f\u7387\nf(x) = TFR \u00d7 \u03c6(x)", 10, (0x41,0x82,0xC4))
add_box(slide6, 9.0, 3.05, 2.0, 0.85, "\u30ad\u30e3\u30ea\u30d6\u30ec\u30fc\u30b7\u30e7\u30f3\na: e\u2080 = \u222bS(x)dx\nb = 0.085 (\u56fa\u5b9a)", 9, (0x41,0x82,0xC4))

add_arrow(slide6, 1.6, 3.95, 5.5, 4.6)
add_arrow(slide6, 4.45, 3.95, 5.5, 4.6)
add_arrow(slide6, 7.3, 3.95, 7.0, 4.6)
add_arrow(slide6, 10.0, 3.95, 7.5, 4.6)

# Core engine
add_text_label(slide6, 0.3, 4.35, 3.0, 0.3, "\u30b3\u30a2\u30a8\u30f3\u30b8\u30f3", font_size=13, bold=True, color=(0x00,0x66,0x33))
add_box(slide6, 4.0, 4.55, 5.3, 0.9,
    "\u5185\u751f\u66f4\u65b0: P(t+1)\n"
    "\u51fa\u751f = \u03a3 P_x(t) \u00d7 f \u00d7 ASFR(x)  |  \u52a0\u9f62: P_{x+1}(t+1) = P_x(t) \u00d7 s(x)",
    11, (0x2E,0x7D,0x32))

add_arrow(slide6, 6.65, 5.5, 6.65, 5.95)

# Output
add_text_label(slide6, 0.3, 5.75, 3.0, 0.3, "\u51fa\u529b", font_size=13, bold=True, color=(0x88,0x00,0x44))
add_box(slide6, 3.7, 5.95, 5.9, 0.7,
    "\u540c\u6642\u5728\u751f\u4eba\u53e3 (SLP)\nP(t) = [P\u2080(t), P\u2081(t), ..., P\u2081\u2080\u2080(t)]",
    11, (0xC4,0x41,0x82))

# Dynamic update
add_box(slide6, 10.5, 4.55, 2.3, 0.9,
    "\u52d5\u7684\u66f4\u65b0\n\uff0810\u5e74\u3054\u3068\uff09:\nTFR, e\u2080, MAC \u2190 WPP",
    10, (0x99,0x66,0xCC))
add_arrow(slide6, 10.5, 5.0, 9.35, 5.0)

add_box(slide6, 10.5, 6.0, 2.3, 0.7,
    "\u56fd\u9023WPP 2024\n\uff08\u691c\u8a3c\u30c7\u30fc\u30bf\uff09", 10, (0x88,0x88,0x88))

# ========================
# Slide 7: Tempo Effect Conceptual (JP)
# ========================
slide7 = prs.slides.add_slide(BLANK)
add_text_label(slide7, 0.5, 0.2, 12.3, 0.6,
    "\u56f37. \u30c6\u30f3\u30dd\u52b9\u679c\u304c\u540c\u6642\u5728\u751f\u4eba\u53e3\u306b\u4e0e\u3048\u308b\u5f71\u97ff", font_size=24, bold=True)
add_text_label(slide7, 0.5, 0.75, 12.3, 0.4,
    "\uff08\u7de8\u96c6\u53ef\u80fd\u306a\u56f3\u2014\u2014\u5168\u3066\u306e\u56f3\u5f62\u3092\u79fb\u52d5\u30fb\u30ea\u30b5\u30a4\u30ba\u30fb\u8272\u5909\u66f4\u53ef\u80fd\uff09",
    font_size=12, italic=True, color=(0x88,0x88,0x88))

# Left: MAC=25
add_text_label(slide7, 0.5, 1.2, 5.5, 0.4, "\u30b7\u30ca\u30ea\u30aaA\uff1aMAC = 25\u6b73", font_size=18, bold=True,
               color=(0x00,0x66,0xCC))
gen_colors = [(0x1B,0x7A,0x2F), (0x41,0x82,0xC4), (0xE8,0x6C,0x00), (0xC4,0x41,0x82)]
gen_labels_25 = ["\u7b2c1\u4e16\u4ee3\uff0875\u6b73\uff09", "\u7b2c2\u4e16\u4ee3\uff0850\u6b73\uff09",
                 "\u7b2c3\u4e16\u4ee3\uff0825\u6b73\uff09", "\u7b2c4\u4e16\u4ee3\uff080\u6b73\uff09"]
bar_top = 1.7
bar_h = 0.55
for i, label in enumerate(gen_labels_25):
    y = bar_top + i * 0.7
    add_box(slide7, 0.7, y, 4.5, bar_h, label, 12, gen_colors[i], bold=False)

add_text_label(slide7, 0.7, bar_top + 3.0, 4.5, 0.5,
    "\u2192 4\u4e16\u4ee3\u304c\u540c\u6642\u306b\u751f\u5b58\n\u2192 \u3088\u308a\u9ad8\u3044SLP",
    font_size=14, bold=True, color=(0x00,0x66,0xCC))

# Right: MAC=30
add_text_label(slide7, 7.0, 1.2, 5.5, 0.4, "\u30b7\u30ca\u30ea\u30aaB\uff1aMAC = 30\u6b73", font_size=18, bold=True,
               color=(0xCC,0x33,0x00))
gen_labels_30 = ["\u7b2c1\u4e16\u4ee3\uff0890\u6b73\uff09", "\u7b2c2\u4e16\u4ee3\uff0860\u6b73\uff09",
                 "\u7b2c3\u4e16\u4ee3\uff0830\u6b73\uff09", "\u7b2c4\u4e16\u4ee3\uff080\u6b73\uff09\u2014\u90e8\u5206\u7684*"]
gen_colors_30 = [(0x1B,0x7A,0x2F), (0x41,0x82,0xC4), (0xE8,0x6C,0x00), (0xCC,0xCC,0xCC)]
for i, label in enumerate(gen_labels_30):
    y = bar_top + i * 0.7
    w = 4.5 if i < 3 else 3.0
    fr = (0xFF,0xFF,0xFF) if i < 3 else (0x66,0x66,0x66)
    add_box(slide7, 7.2, y, w, bar_h, label, 12, gen_colors_30[i], bold=False, font_rgb=fr)

add_text_label(slide7, 7.2, bar_top + 3.0, 5.0, 0.5,
    "\u2192 \u7d043.3\u4e16\u4ee3\u304c\u540c\u6642\u306b\u751f\u5b58\n\u2192 SLP\u304c\u7d041/6\u6e1b\u5c11",
    font_size=14, bold=True, color=(0xCC,0x33,0x00))

add_text_label(slide7, 0.7, bar_top + 3.7, 4.5, 0.3,
    "*\u4e21\u30b7\u30ca\u30ea\u30aa\u3067TFR\u306f\u540c\u4e00", font_size=11, italic=True, color=(0x66,0x66,0x66))

# Central arrow
add_box(slide7, 5.5, 2.8, 1.3, 0.8, "MAC\n+5\u5e74", 12, (0xFF,0xFF,0xFF),
        font_rgb=(0x33,0x33,0x33), bold=True, shape=MSO_SHAPE.RIGHT_ARROW)

# Key insight
add_box(slide7, 1.5, 5.8, 10.3, 1.2,
    "\u6838\u5fc3\u7684\u77e5\u898b\uff1aMAC\u306e5\u5e74\u4e0a\u6607\u306fTFR\u3068\u306f\u72ec\u7acb\u306bSLP\u3092\u7d041/6\u6e1b\u5c11\u3055\u305b\u308b\u3002\n"
    "\u3053\u308c\u306f\u4eba\u53e3\u8ee2\u63db\u306e\u300c\u901f\u5ea6\u300d\u3092\u5236\u5fa1\u3057\u3001\n"
    "\u5236\u5ea6\u7684\u9069\u5fdc\u306b\u5229\u7528\u53ef\u80fd\u306a\u6642\u9593\u3092\u5727\u7e2e\u3059\u308b\u3002",
    13, (0x2E,0x2E,0x5E), bold=True)

# ========================
# Slide 8: Policy Framework (JP)
# ========================
slide8 = prs.slides.add_slide(BLANK)
add_text_label(slide8, 0.5, 0.2, 12.3, 0.6,
    "\u56f38. \u653f\u7b56\u30d5\u30ec\u30fc\u30e0\u30ef\u30fc\u30af\uff1a\u30c6\u30f3\u30dd\u611f\u5fdc\u578b\u4eba\u53e3\u7ba1\u7406", font_size=24, bold=True)
add_text_label(slide8, 0.5, 0.75, 12.3, 0.4,
    "\uff08\u7de8\u96c6\u53ef\u80fd\u306a\u56f3\u2014\u2014\u5168\u3066\u306e\u56f3\u5f62\u3092\u79fb\u52d5\u30fb\u30ea\u30b5\u30a4\u30ba\u30fb\u8272\u5909\u66f4\u53ef\u80fd\uff09",
    font_size=12, italic=True, color=(0x88,0x88,0x88))

# Current approach
add_text_label(slide8, 0.5, 1.2, 6.0, 0.35, "\u73fe\u884c\u30a2\u30d7\u30ed\u30fc\u30c1\uff08\u30ab\u30f3\u30bf\u30e0\u306e\u307f\uff09", font_size=16, bold=True,
               color=(0xCC,0x33,0x00))
add_box(slide8, 0.5, 1.6, 2.5, 0.8, "\u5c11\u5b50\u5316\u5bfe\u7b56", 12, (0xCC,0x55,0x55))
add_arrow(slide8, 3.05, 2.0, 3.7, 2.0)
add_box(slide8, 3.7, 1.6, 2.5, 0.8, "TFR\u5f15\u304d\u4e0a\u3052\n\uff08\u2191\u51fa\u751f\u6570\uff09", 12, (0xCC,0x55,0x55))
add_arrow(slide8, 6.25, 2.0, 6.9, 2.0)
add_box(slide8, 6.9, 1.6, 2.8, 0.8, "\u4eba\u53e3\u898f\u6a21\u306e\n\u5897\u52a0", 12, (0xCC,0x55,0x55))
add_arrow(slide8, 9.75, 2.0, 10.4, 2.0)
add_box(slide8, 10.4, 1.6, 2.5, 0.8, "\u5f71\u97ff\u3092\u904e\u5927\u8a55\u4fa1\n\uff08\u30c6\u30f3\u30dd\u3092\u7121\u8996\uff09", 10, (0xFF,0xCC,0xCC),
        font_rgb=(0x88,0x00,0x00))

# Proposed approach
add_text_label(slide8, 0.5, 3.0, 6.0, 0.35, "\u63d0\u6848\u30a2\u30d7\u30ed\u30fc\u30c1\uff08\u30ab\u30f3\u30bf\u30e0\uff0b\u30c6\u30f3\u30dd\uff09", font_size=16, bold=True,
               color=(0x00,0x66,0x33))

add_box(slide8, 0.5, 3.5, 2.5, 0.8, "\u30ab\u30f3\u30bf\u30e0\uff1a\nTFR\u5f15\u304d\u4e0a\u3052", 12, (0x2E,0x7D,0x32))
add_arrow(slide8, 3.05, 3.9, 5.2, 4.8)

add_box(slide8, 0.5, 4.6, 2.5, 0.8, "\u30c6\u30f3\u30dd\uff1a\nMAC/AFB\u5f15\u304d\u4e0b\u3052", 12, (0x41,0x82,0xC4))
add_arrow(slide8, 3.05, 5.0, 5.2, 5.0)

# Tempo policies
add_box(slide8, 3.3, 5.7, 2.0, 0.65, "\u624b\u9802\u306a\n\u4f4f\u5b85\u652f\u63f4", 10, (0x66,0xAA,0xDD))
add_box(slide8, 5.5, 5.7, 2.0, 0.65, "\u666e\u904d\u7684\n\u4fdd\u80b2\u30a4\u30f3\u30d5\u30e9", 10, (0x66,0xAA,0xDD))
add_box(slide8, 7.7, 5.7, 2.3, 0.65, "\u6559\u80b2\u8ab2\u7a0b\u306e\n\u518d\u69cb\u7bc9", 10, (0x66,0xAA,0xDD))

# Combined effect
add_box(slide8, 5.0, 4.5, 3.0, 0.9,
    "\u8907\u5408\u52b9\u679c\uff1a\n\u2191 \u4e16\u4ee3\u306e\u91cd\u306a\u308a\n\u2193 \u6e1b\u5c11\u306e\u30b9\u30d4\u30fc\u30c9", 11, (0x2E,0x7D,0x32))
add_arrow(slide8, 8.05, 4.95, 8.8, 4.95)

# Outcome
add_box(slide8, 8.7, 3.8, 3.8, 1.5,
    "\u6210\u679c\uff1a\n\u2022 \u4eba\u53e3\u8ee2\u63db\u306e\u30da\u30fc\u30b9\u3092\n  \u7ba1\u7406\n\u2022 \u5236\u5ea6\u7684\u9069\u5fdc\u306e\u305f\u3081\u306e\n  \u6642\u9593\u78ba\u4fdd\n\u2022 \u6301\u7d9a\u53ef\u80fd\u306aSLP\u7d4c\u8def",
    11, (0x00,0x44,0x22), bold=False)

add_arrow(slide8, 4.3, 5.7, 5.5, 5.45)
add_arrow(slide8, 6.5, 5.7, 6.5, 5.45)
add_arrow(slide8, 8.3, 5.7, 7.5, 5.45)

prs.save('/home/ubuntu/Figures_JP.pptx')
print("OK: Figures_JP.pptx")
