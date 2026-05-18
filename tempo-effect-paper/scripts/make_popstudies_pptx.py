"""Create editable PPTX with all figures for Population Studies submission."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
FIG_DIR = os.path.join(os.path.dirname(SCRIPT_DIR), 'figures')
OUT_DIR = os.path.join(os.path.dirname(SCRIPT_DIR), 'manuscripts')
os.makedirs(OUT_DIR, exist_ok=True)

# Widescreen dimensions
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


def add_figure_slide(prs, fig_path, title, caption, img_width=10.0):
    """Add one slide with title at top, figure centred, caption at bottom."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Title text box
    from pptx.util import Emu
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.6))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Image (centred)
    if os.path.exists(fig_path):
        from PIL import Image
        img = Image.open(fig_path)
        aspect = img.height / img.width
        w = Inches(img_width)
        h = Inches(img_width * aspect)
        # Constrain height
        max_h = Inches(5.2)
        if h > max_h:
            h = max_h
            w = Inches(float(max_h) / aspect / 914400)  # recalculate
            w = Inches(5.2 / aspect)
        left = (SLIDE_WIDTH - w) // 2
        top = Inches(1.0)
        slide.shapes.add_picture(fig_path, left, top, w, h)
    else:
        # Placeholder text if figure missing
        tb = slide.shapes.add_textbox(
            Inches(2), Inches(3), Inches(9), Inches(1))
        tb.text_frame.paragraphs[0].text = f"[Figure not found: {fig_path}]"

    # Caption text box at bottom
    cap_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.8))
    tf2 = cap_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = caption
    p2.font.size = Pt(11)
    p2.font.italic = True
    p2.alignment = PP_ALIGN.CENTER


# Figures from the manuscript
figures = [
    {
        'path': os.path.join(FIG_DIR, 'fig1_showcase.png'),
        'title': 'Figure 1: Five model variants versus observed population, '
                 'six representative countries, 1970\u20132023',
        'caption': 'Observed = black solid (UN WPP 2024); Tempo-responsive = blue; '
                   'Tempo-invariant = orange dashed; Tempo-adjusted (TFR*) = green '
                   'dash-dot; Fixed-parameter (1970) = red dotted.',
    },
    {
        'path': os.path.join(FIG_DIR, 'fig2_all_countries.png'),
        'title': 'Figure 2: Five-model validation across all 40 countries',
        'caption': 'Observed = black; Responsive = blue; Invariant = orange; '
                   'Adjusted (TFR*) = green; Fixed = red. MAPE shown in '
                   'upper-right corner. Countries sorted alphabetically.',
    },
    {
        'path': os.path.join(FIG_DIR, 'fig3_heatmap.png'),
        'title': 'Figure 3: Fixed-parameter model MAPE (%) by country and base year',
        'caption': 'Greener cells indicate better fit; redder cells indicate '
                   'poorer fit. Scale capped at 30%.',
    },
    {
        'path': os.path.join(FIG_DIR, 'fig4_comparison.png'),
        'title': 'Figure 4: Four-model comparison \u2014 MAPE and final population ratio',
        'caption': 'Left: MAPE by country for fixed (red), invariant (orange), '
                   'responsive (blue), adjusted TFR* (green). Right: final ratio '
                   '(model/observed, 2023). TFR* achieves best overall fit.',
    },
    {
        'path': os.path.join(FIG_DIR, 'fig5_bias.png'),
        'title': 'Figure 5: Model bias analysis using base year 2000',
        'caption': '(a) Fit versus TFR; (b) fit versus life expectancy; '
                   '(c) bias versus MAC. No systematic relationship observed.',
    },
]

for fig in figures:
    add_figure_slide(prs, fig['path'], fig['title'], fig['caption'])

outpath = os.path.join(OUT_DIR, 'PopStudies_Figures_EN.pptx')
prs.save(outpath)
print(f'OK: {outpath}')
