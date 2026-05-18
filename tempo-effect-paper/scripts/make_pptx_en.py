# -*- coding: utf-8 -*-
"""Create PPTX figure file — English. One figure per slide.
Code-output figures: embedded as images.
Conceptual/flow diagrams: editable shapes."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

FIG_DIR = '/home/ubuntu/figures'
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]  # blank layout

def add_image_slide(prs, img_path, title, caption):
    """Add a slide with a code-output figure (as image) + title + caption."""
    slide = prs.slides.add_slide(BLANK)
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1B, 0x1B, 0x1B)
    p.alignment = PP_ALIGN.CENTER
    # Image
    if os.path.exists(img_path):
        from PIL import Image
        im = Image.open(img_path)
        w, h = im.size
        aspect = w / h
        max_w, max_h = 11.5, 5.5
        if max_w / aspect <= max_h:
            disp_w, disp_h = max_w, max_w / aspect
        else:
            disp_w, disp_h = max_h * aspect, max_h
        left = (13.333 - disp_w) / 2
        slide.shapes.add_picture(img_path, Inches(left), Inches(1.0), Inches(disp_w), Inches(disp_h))
    # Caption
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.7))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = caption
    p2.font.size = Pt(11)
    p2.font.italic = True
    p2.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
    p2.alignment = PP_ALIGN.CENTER
    return slide

def add_box(slide, left, top, width, height, text, font_size=12, fill_rgb=(0x41, 0x82, 0xC4),
            font_rgb=(0xFF, 0xFF, 0xFF), bold=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    """Add an editable shape box."""
    shp = slide.shapes.add_shape(shape, Inches(left), Inches(top), Inches(width), Inches(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(*fill_rgb)
    shp.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    shp.line.width = Pt(1)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*font_rgb)
    return shp

def add_arrow(slide, x1, y1, x2, y2, color=(0x33, 0x33, 0x33)):
    """Add a connector arrow."""
    from pptx.oxml.ns import qn
    connector = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))  # 1 = straight
    connector.line.color.rgb = RGBColor(*color)
    connector.line.width = Pt(2)
    # Add arrowhead
    ln = connector.line._ln
    tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail)
    return connector

def add_text_label(slide, left, top, width, height, text, font_size=11, bold=False, italic=False,
                   align=PP_ALIGN.CENTER, color=(0x22, 0x22, 0x22)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = RGBColor(*color)
    p.alignment = align
    return txBox

# ========================
# Slide 1: fig1 showcase
# ========================
add_image_slide(prs, f'{FIG_DIR}/fig1_showcase.png',
    "Figure 1. Model vs Observed Population: Six Representative Countries (1970\u20132023)",
    "Dynamic model (blue dashed) updates parameters every 10 years; static model (red dotted) uses "
    "fixed 1970 parameters. Black line = UN WPP 2024 estimates. Countries: Japan, China, USA, Korea, Germany, DRC.")

# ========================
# Slide 2: fig2 all countries
# ========================
add_image_slide(prs, f'{FIG_DIR}/fig2_all_countries.png',
    "Figure 2. Model Validation Across All 40 Countries",
    "Dynamic model MAPE shown in upper-right corner of each panel. Countries sorted alphabetically. "
    "30 of 40 countries achieve MAPE < 10%; 20 countries < 5%.")

# ========================
# Slide 3: fig3 heatmap
# ========================
add_image_slide(prs, f'{FIG_DIR}/fig3_heatmap.png',
    "Figure 3. Static Model MAPE (%) by Country and Base Year",
    "Greener cells indicate better fit; redder cells indicate poorer fit. Scale capped at 30%. "
    "Fit improves with more recent base years as expected.")

# ========================
# Slide 4: fig4 comparison
# ========================
add_image_slide(prs, f'{FIG_DIR}/fig4_comparison.png',
    "Figure 4. Static vs Dynamic Model Comparison",
    "Left: MAPE by country. Right: final population ratio (model/observed in 2023). "
    "Dynamic model consistently outperforms static variant with ~2x improvement in median MAPE.")

# ========================
# Slide 5: fig5 bias
# ========================
add_image_slide(prs, f'{FIG_DIR}/fig5_bias.png',
    "Figure 5. Model Bias Analysis (Base Year 2000)",
    "(A) Fit vs TFR; (B) Fit vs life expectancy; (C) Bias vs MAC. "
    "No systematic relationship observed, suggesting robustness across demographic contexts.")

# ========================
# Slide 6: Editable Model Flow Diagram
# ========================
slide6 = prs.slides.add_slide(BLANK)
add_text_label(slide6, 0.5, 0.2, 12.3, 0.6,
    "Figure 6. Endogenous Renewal Model: Structure and Data Flow", font_size=24, bold=True)
add_text_label(slide6, 0.5, 0.75, 12.3, 0.4,
    "(Editable diagram \u2014 all shapes can be moved, resized, and recolored)",
    font_size=12, italic=True, color=(0x88, 0x88, 0x88))

# Input parameters row
add_text_label(slide6, 0.3, 1.25, 3.0, 0.3, "INPUT PARAMETERS", font_size=13, bold=True, color=(0x88,0x44,0x00))
add_box(slide6, 0.5, 1.6, 2.2, 0.7, "TFR\n(Total Fertility Rate)", 11, (0xE8,0x6C,0x00))
add_box(slide6, 3.0, 1.6, 2.2, 0.7, "e\u2080\n(Life Expectancy)", 11, (0xE8,0x6C,0x00))
add_box(slide6, 5.5, 1.6, 2.2, 0.7, "MAC\n(Mean Age at\nChildbearing)", 10, (0xE8,0x6C,0x00))
add_box(slide6, 8.0, 1.6, 2.2, 0.7, "\u03c3\n(Fertility Schedule\nWidth)", 10, (0xE8,0x6C,0x00))

# Arrows from inputs to model components
add_arrow(slide6, 1.6, 2.35, 1.6, 2.9)
add_arrow(slide6, 4.1, 2.35, 4.1, 2.9)
add_arrow(slide6, 6.6, 2.35, 6.6, 2.9)
add_arrow(slide6, 9.1, 2.35, 9.1, 2.9)

# Model components row
add_text_label(slide6, 0.3, 2.75, 3.0, 0.3, "MODEL COMPONENTS", font_size=13, bold=True, color=(0x00,0x55,0x88))
add_box(slide6, 0.3, 3.05, 2.6, 0.85, "Normal Fertility\nSchedule\nASFR ~ N(MAC, \u03c3\u00b2)", 10, (0x41,0x82,0xC4))
add_box(slide6, 3.3, 3.05, 2.3, 0.85, "Gompertz Survival\nS(x) = exp[-(a/b)\n\u00d7(exp(bx)-1)]", 10, (0x41,0x82,0xC4))
add_box(slide6, 6.0, 3.05, 2.6, 0.85, "Age-Specific\nFertility Rate\nf(x) = TFR \u00d7 \u03c6(x)", 10, (0x41,0x82,0xC4))
add_box(slide6, 9.0, 3.05, 2.0, 0.85, "Calibration\na: e\u2080 = \u222bS(x)dx\nb = 0.085 (fixed)", 9, (0x41,0x82,0xC4))

# Arrows to core engine
add_arrow(slide6, 1.6, 3.95, 5.5, 4.6)
add_arrow(slide6, 4.45, 3.95, 5.5, 4.6)
add_arrow(slide6, 7.3, 3.95, 7.0, 4.6)
add_arrow(slide6, 10.0, 3.95, 7.5, 4.6)

# Core engine
add_text_label(slide6, 0.3, 4.35, 3.0, 0.3, "CORE ENGINE", font_size=13, bold=True, color=(0x00,0x66,0x33))
add_box(slide6, 4.0, 4.55, 5.3, 0.9,
    "Endogenous Renewal: P(t+1)\n"
    "Births = \u03a3 P_x(t) \u00d7 f \u00d7 ASFR(x)  |  Aging: P_{x+1}(t+1) = P_x(t) \u00d7 s(x)",
    11, (0x2E,0x7D,0x32))

# Arrow to output
add_arrow(slide6, 6.65, 5.5, 6.65, 5.95)

# Output
add_text_label(slide6, 0.3, 5.75, 3.0, 0.3, "OUTPUT", font_size=13, bold=True, color=(0x88,0x00,0x44))
add_box(slide6, 3.7, 5.95, 5.9, 0.7,
    "Simultaneously Living Population (SLP)\nP(t) = [P\u2080(t), P\u2081(t), ..., P\u2081\u2080\u2080(t)]",
    11, (0xC4,0x41,0x82))

# Dynamic update loop
add_box(slide6, 10.5, 4.55, 2.3, 0.9,
    "Dynamic Update\n(every 10 years):\nTFR, e\u2080, MAC \u2190 WPP",
    10, (0x99,0x66,0xCC))
add_arrow(slide6, 10.5, 5.0, 9.35, 5.0)

# Data source
add_box(slide6, 10.5, 6.0, 2.3, 0.7,
    "UN WPP 2024\n(Validation Data)",
    10, (0x88, 0x88, 0x88))

# ========================
# Slide 7: Editable Tempo Effect Conceptual Diagram
# ========================
slide7 = prs.slides.add_slide(BLANK)
add_text_label(slide7, 0.5, 0.2, 12.3, 0.6,
    "Figure 7. The Tempo Effect on Simultaneously Living Population", font_size=24, bold=True)
add_text_label(slide7, 0.5, 0.75, 12.3, 0.4,
    "(Editable diagram \u2014 all shapes can be moved, resized, and recolored)",
    font_size=12, italic=True, color=(0x88, 0x88, 0x88))

# Left panel: MAC = 25
add_text_label(slide7, 0.5, 1.2, 5.5, 0.4, "Scenario A: MAC = 25 years", font_size=18, bold=True,
               color=(0x00,0x66,0xCC))
# Generation bars
gen_colors = [(0x1B,0x7A,0x2F), (0x41,0x82,0xC4), (0xE8,0x6C,0x00), (0xC4,0x41,0x82)]
gen_labels_25 = [
    ("Gen 1 (age 75)", 0, 75), ("Gen 2 (age 50)", 0, 50),
    ("Gen 3 (age 25)", 0, 25), ("Gen 4 (age 0)", 0, 0)
]
bar_top = 1.7
bar_h = 0.55
for i, (label, _, age) in enumerate(gen_labels_25):
    y = bar_top + i * 0.7
    add_box(slide7, 0.7, y, 4.5, bar_h, label, 12, gen_colors[i], bold=False)

add_text_label(slide7, 0.7, bar_top + 3.0, 4.5, 0.5,
    "\u2192 4 generations simultaneously alive\n\u2192 Higher SLP",
    font_size=14, bold=True, color=(0x00,0x66,0xCC))

# Right panel: MAC = 30
add_text_label(slide7, 7.0, 1.2, 5.5, 0.4, "Scenario B: MAC = 30 years", font_size=18, bold=True,
               color=(0xCC,0x33,0x00))
gen_labels_30 = [
    ("Gen 1 (age 90)", 0, 90), ("Gen 2 (age 60)", 0, 60),
    ("Gen 3 (age 30)", 0, 30), ("Gen 4 (age 0) \u2014 partial*", 0, 0)
]
gen_colors_30 = [(0x1B,0x7A,0x2F), (0x41,0x82,0xC4), (0xE8,0x6C,0x00), (0xCC,0xCC,0xCC)]
for i, (label, _, age) in enumerate(gen_labels_30):
    y = bar_top + i * 0.7
    w = 4.5 if i < 3 else 3.0  # partial bar for 4th gen (fewer survive to 90)
    add_box(slide7, 7.2, y, w, bar_h, label, 12, gen_colors_30[i],
            bold=False, font_rgb=(0xFF,0xFF,0xFF) if i < 3 else (0x66,0x66,0x66))

add_text_label(slide7, 7.2, bar_top + 3.0, 5.0, 0.5,
    "\u2192 ~3.3 generations simultaneously alive\n\u2192 Lower SLP (\u22481/6 reduction)",
    font_size=14, bold=True, color=(0xCC,0x33,0x00))

# Bottom: Key message
add_text_label(slide7, 0.7, bar_top + 3.7, 4.5, 0.3,
    "*Same TFR in both scenarios", font_size=11, italic=True, color=(0x66,0x66,0x66))

# Central comparison arrow
add_box(slide7, 5.5, 2.8, 1.3, 0.8, "+5 years\nMAC shift", 12, (0xFF,0xFF,0xFF),
        font_rgb=(0x33,0x33,0x33), bold=True, shape=MSO_SHAPE.RIGHT_ARROW)

# Bottom key insight box
add_box(slide7, 1.5, 5.8, 10.3, 1.2,
    "KEY INSIGHT: A 5-year increase in MAC reduces SLP by ~1/6,\n"
    "independent of TFR. This controls the SPEED of demographic transition,\n"
    "compressing the time available for institutional adaptation.",
    13, (0x2E,0x2E,0x5E), bold=True)

# ========================
# Slide 8: Editable Policy Implications Flow
# ========================
slide8 = prs.slides.add_slide(BLANK)
add_text_label(slide8, 0.5, 0.2, 12.3, 0.6,
    "Figure 8. Policy Framework: Tempo-Sensitive Population Management", font_size=24, bold=True)
add_text_label(slide8, 0.5, 0.75, 12.3, 0.4,
    "(Editable diagram \u2014 all shapes can be moved, resized, and recolored)",
    font_size=12, italic=True, color=(0x88, 0x88, 0x88))

# Top: Current approach
add_text_label(slide8, 0.5, 1.2, 6.0, 0.35, "CURRENT APPROACH (Quantum-only)", font_size=16, bold=True,
               color=(0xCC,0x33,0x00))
add_box(slide8, 0.5, 1.6, 2.5, 0.8, "Pronatalist\nPolicies", 12, (0xCC,0x55,0x55))
add_arrow(slide8, 3.05, 2.0, 3.7, 2.0)
add_box(slide8, 3.7, 1.6, 2.5, 0.8, "Raise TFR\n(\u2191 Births)", 12, (0xCC,0x55,0x55))
add_arrow(slide8, 6.25, 2.0, 6.9, 2.0)
add_box(slide8, 6.9, 1.6, 2.8, 0.8, "Increase\nPopulation Size", 12, (0xCC,0x55,0x55))
add_arrow(slide8, 9.75, 2.0, 10.4, 2.0)
add_box(slide8, 10.4, 1.6, 2.5, 0.8, "Overestimates\nimpact\n(ignores tempo)", 10, (0xFF,0xCC,0xCC),
        font_rgb=(0x88,0x00,0x00))

# Bottom: Proposed approach
add_text_label(slide8, 0.5, 3.0, 6.0, 0.35, "PROPOSED APPROACH (Quantum + Tempo)", font_size=16, bold=True,
               color=(0x00,0x66,0x33))

# Quantum path
add_box(slide8, 0.5, 3.5, 2.5, 0.8, "Quantum:\nRaise TFR", 12, (0x2E,0x7D,0x32))
add_arrow(slide8, 3.05, 3.9, 5.2, 4.8)

# Tempo path
add_box(slide8, 0.5, 4.6, 2.5, 0.8, "Tempo:\nLower MAC/AFB", 12, (0x41,0x82,0xC4))
add_arrow(slide8, 3.05, 5.0, 5.2, 5.0)

# Tempo policy details
add_box(slide8, 3.3, 5.7, 2.0, 0.65, "Affordable\nHousing", 10, (0x66,0xAA,0xDD))
add_box(slide8, 5.5, 5.7, 2.0, 0.65, "Universal\nChildcare", 10, (0x66,0xAA,0xDD))
add_box(slide8, 7.7, 5.7, 2.3, 0.65, "Educational\nPathway Reform", 10, (0x66,0xAA,0xDD))

# Combined effect
add_box(slide8, 5.0, 4.5, 3.0, 0.9,
    "Combined Effect:\n\u2191 Generational Overlap\n\u2193 Speed of Decline", 11, (0x2E,0x7D,0x32))
add_arrow(slide8, 8.05, 4.95, 8.8, 4.95)

# Outcome
add_box(slide8, 8.7, 3.8, 3.8, 1.5,
    "OUTCOME:\n\u2022 Managed pace of\n  demographic transition\n\u2022 Time for institutional\n  adaptation\n\u2022 Sustainable SLP path",
    11, (0x00,0x44,0x22), bold=False)

# Arrow from tempo policies to combined
add_arrow(slide8, 4.3, 5.7, 5.5, 5.45)
add_arrow(slide8, 6.5, 5.7, 6.5, 5.45)
add_arrow(slide8, 8.3, 5.7, 7.5, 5.45)

prs.save('/home/ubuntu/Figures_EN.pptx')
print("OK: Figures_EN.pptx")
