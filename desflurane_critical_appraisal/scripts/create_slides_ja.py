#!/usr/bin/env python3
"""
デスフルラン文献批判的吟味 - 発表用スライド（pptx）
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

OUTPUT_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "output")
os.makedirs(OUTPUT_DIR, exist_ok=True)

# Widescreen dimensions
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def add_slide(prs, title_text, content_items=None, subtitle=None, layout_idx=1):
    """Add a slide with title and bullet content."""
    slide_layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    if slide.shapes.title:
        title = slide.shapes.title
        title.text = title_text
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(28)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0x1A, 0x47, 0x7A)

    # Body
    if content_items and len(slide.placeholders) > 1:
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        for i, item in enumerate(content_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
            p.space_after = Pt(6)
            if item.startswith('→'):
                p.font.color.rgb = RGBColor(0xCC, 0x00, 0x00)
                p.font.bold = True

    return slide


def add_title_slide(prs, title, subtitle):
    """Add a title-only slide."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    if slide.shapes.title:
        slide.shapes.title.text = title
        for p in slide.shapes.title.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(36)
                r.font.bold = True
                r.font.color.rgb = RGBColor(0x1A, 0x47, 0x7A)

    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
        for p in slide.placeholders[1].text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(20)

    return slide


def create_slides():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # === Slide 1: Title ===
    add_title_slide(prs,
        'デスフルラン優位性を主張する文献の批判的吟味',
        '― セボフルランに対する真の優位性は存在するか ―\n'
        '臨床麻酔学会発表用'
    )

    # === Slide 2: Background ===
    add_slide(prs, '背景', [
        'EU規則 (EU) 2024/573: デスフルラン原則使用禁止（2026/1/1〜）',
        'GWP: デスフルラン 2,540 vs セボフルラン 130（約20倍）',
        '例外: 「医学的根拠により必要な場合」のみ',
        'NICEレビュー（2024）: デスフルランの臨床的優位性を示す質の高いエビデンスなし',
        '',
        '→ 本発表の目的: 「デスフルランを使用すべき症例」は存在するか？',
    ])

    # === Slide 3: Method ===
    add_slide(prs, '方法', [
        'デスフルラン優位を主張する10のシチュエーションを特定',
        '各シチュエーションの代表的論文を選定',
        '方法論的問題点を批判的に検証',
        '批判レター（Letter to the Editor）と著者回答を調査',
        '「引用の非対称性」を検討',
    ])

    # === Slide 4: Scenario 1 - Obesity ===
    add_slide(prs, '1. 肥満患者における覚醒の速さ', [
        '主張: 低い血液/ガス分配係数（0.42 vs 0.65）で覚醒が早い',
        'La Colla 2007 (BJA): 非再呼吸回路（FGF 8L/min）、手術30分',
        '  → 臨床的に使用される低流量回路とは条件が異なる',
        'Kaur 2013 (JOACP): FGF未報告、覚醒時のプロトコル不明',
        '  → デスフルラン群のみハイフローの可能性を排除できない',
        '反証: Arain 2005 — BIS-guided titrationで有意差なし',
        'エディトリアル: Eger & Shafer 2005 — 覚醒数分の差は臨床的に無意味',
        '',
        '→ 覚醒の速さ ≠ 退院の早さ。FGF条件で結果が操作される',
    ])

    # === Slide 5: FI-FA gap ===
    add_slide(prs, '肥満患者: FI-FA差の問題', [
        'デスフルラン代謝率 0.02% → 定常状態ではFI ≈ FA のはず',
        '臨床現場ではFI-FA差が持続的に観察される',
        '原因: 大量の脂肪組織への分配が継続',
        '脂肪/血液分配係数: デスフルラン 27 vs セボフルラン 48',
        '  → 血液/ガス分配係数ほどの差はない',
        '',
        '→ 「肥満でも速い覚醒」の理論的根拠が矛盾',
    ])

    # === Slide 6: Scenario 2 - Elderly ===
    add_slide(prs, '2. 高齢者における認知機能保存', [
        'Tachibana et al. 2015 (J Anesth): n=42, MMSE 術前後24h',
        '問題点:',
        '  ・MMSE変化量わずか1点（測定誤差±2-3点の範囲内）',
        '  ・p=0.048（ぎりぎりの有意水準）',
        '  ・認知機能低下例の除外基準が不明確',
        '  ・低スコア症例を除外すると有意差消失',
        '  ・BISモニタリングなし、等MAC濃度でない',
        '',
        '→ 著者自身が批判レターへの回答で事実上の譲歩',
    ])

    # === Slide 7: Citation asymmetry ===
    add_slide(prs, '引用の非対称性（Tachibana 2015の例）', [
        '原著（Tachibana 2015）: 31回引用',
        '批判レター（Ding et al. 2015）: 2回引用',
        '著者回答/譲歩（Tachibana 2017）: ほぼ引用ゼロ',
        '',
        '著者回答の要点:',
        '  「MMSEの群間差については言及していない」',
        '  「両群のMMSEスコアは同等であった」',
        '  「認知機能の結果は慎重に解釈すべき」',
        '',
        '→ 企業・後続研究は原著の「デスフルラン有利」のみ引用',
    ])

    # === Slide 8: Scenarios 3-4 ===
    add_slide(prs, '3. 長時間手術 / 4. 脳外科', [
        '【長時間手術】',
        '理論的にはウォッシュアウトが速い',
        'しかし退院時間、PACU滞在時間に差なし（メタアナリシス）',
        'BIS-guided titrationなしでは交絡因子',
        '',
        '【脳外科】',
        'Dube 2015: 抜管時間差あり、GCS回復・転帰に差なし',
        '2025年メタアナリシス: 臨床的転帰に差なし',
        'NICEレビュー: 脳外科でもデスフルランの優位性なし',
        'デスフルランの気道刺激 → 術後バッキング・ICP上昇リスク',
    ])

    # === Slide 9: Scenarios 5-6 ===
    add_slide(prs, '5. 日帰り手術 / 6. 小児', [
        '【日帰り手術】',
        'White 2009: 早期覚醒は速いが退院時間に差なし',
        'デスフルラン群で咳嗽増加（67% vs 33%, p<0.01）',
        'ボトルネックは覚醒ではなくPONV・疼痛管理',
        '',
        '【小児】',
        'Lim 2016 メタアナリシス: EA発生率25% vs 25%で同等',
        'デスフルランは気道刺激が強くマスク導入不可',
        '導入はセボフルラン→維持のみデスフルランは非現実的',
    ])

    # === Slide 10: Scenarios 7-8 ===
    add_slide(prs, '7. 心臓手術 / 8. 呼吸器疾患', [
        '【心臓手術】',
        '心筋保護効果は揮発性麻酔薬のクラス効果',
        'セボフルランの方がエビデンス豊富',
        'De Hert 2009 CABG 1年死亡率: D 6.9% vs S 3.3% vs TIVA 12.3%',
        '→ デスフルランはセボフルランの約2倍の死亡率',
        '',
        '【呼吸器疾患】',
        'デスフルラン: TRPA1活性化 → 咳嗽・喉頭痙攣・気管支攣縮',
        'セボフルラン: 気管支拡張作用あり',
        '→ 呼吸器疾患患者ではデスフルランはむしろ「劣位」',
    ])

    # === Slide 11: Scenarios 9-10 ===
    add_slide(prs, '9. 肝腎障害 / 10. 筋疾患', [
        '【肝腎障害】',
        'Compound A腎毒性: ヒトで未確認（ラットのみ）',
        '20年以上の臨床使用で問題なし',
        '理論的メリットはあるが臨床的差はなし',
        '',
        '【筋疾患】',
        'MHトリガーは全揮発性麻酔薬共通',
        'エビデンスは症例報告レベルのみ',
        'TIVA推奨例が多い',
    ])

    # === Slide 12: Summary Table ===
    add_slide(prs, '総括: 10シチュエーション全てで優位性なし', [
        '1. 肥満 → FGF条件操作、退院時間差なし',
        '2. 高齢者POCD → MMSE 1点差、著者自身が譲歩',
        '3. 長時間手術 → 退院時間・転帰に差なし',
        '4. 脳外科 → 転帰に差なし、気道刺激リスク',
        '5. 日帰り手術 → 退院時間差なし、咳嗽増加',
        '6. 小児 → EA同率、マスク導入不可',
        '7. 心臓手術 → クラス効果、CABG死亡率D6.9% vs S3.3%',
        '8. 呼吸器疾患 → むしろ劣位（気道刺激）',
        '9. 肝腎障害 → 理論的のみ、臨床差なし',
        '10. 筋疾患 → 症例報告レベル、差なし',
    ])

    # === Slide 13: Patterns of bias ===
    add_slide(prs, '産業バイアスのパターン', [
        '1. 原著のみ引用、批判レターと著者回答を無視',
        '2. 主要エンドポイントの事前指定なし（6/10論文）',
        '3. 多重比較補正の系統的欠如 → ぎりぎりのp値',
        '4. 抜管時間を主要EPとし退院・転帰を軽視',
        '5. FGF未報告またはハイフロー条件で差を人工的に拡大',
        '6. BISモニタリングなしで麻酔深度が不均衡',
        '7. 少数例、検出力不足、利益相反の開示不十分',
        '',
        '→ GWP 2,540のデスフルランを正当化するエビデンスは存在しない',
    ])

    # === Slide 14: Conclusion ===
    add_slide(prs, '結論', [
        'デスフルランを使用すべき症例は存在しない',
        '',
        '10シチュエーション全てにおいて:',
        '  セボフルランに対する臨床的優位性を示す',
        '  質の高いエビデンスは見出されなかった',
        '',
        'EU規則 (EU) 2024/573による使用禁止は',
        '  環境的にも臨床的にも妥当である',
        '',
        '→ 「医学的根拠による例外」を主張する根拠はない',
    ])

    # === Slide 15: COI ===
    add_slide(prs, '利益相反開示', [
        '開示すべき利益相反はありません',
    ])

    output_path = os.path.join(OUTPUT_DIR, 'desflurane_critical_appraisal_ja.pptx')
    prs.save(output_path)
    print(f"Slides saved to: {output_path}")
    return output_path


if __name__ == '__main__':
    create_slides()
